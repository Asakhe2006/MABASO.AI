import asyncio
import base64
import binascii
from datetime import datetime, timedelta, timezone
from email.message import EmailMessage
import html
import hashlib
import hmac
from http.cookiejar import MozillaCookieJar
from io import BytesIO
import json
import logging
import mimetypes
import os
import re
import shutil
import smtplib
import sqlite3
import secrets
import ssl
import subprocess
import tempfile
import unicodedata
import zipfile
from pathlib import Path, PurePosixPath
from typing import Any
from urllib.parse import parse_qs, quote, urlparse
from uuid import uuid4
from xml.etree import ElementTree as ET

from fastapi import Depends, FastAPI, File, Form, Header, HTTPException, Request, UploadFile
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import Response
from openai import APIStatusError, InternalServerError, OpenAI
from pydantic import BaseModel
import requests
try:
    from curl_cffi import requests as curl_requests
except ImportError:
    curl_requests = None

try:
    import yt_dlp
except ImportError:
    yt_dlp = None

try:
    from youtube_transcript_api import YouTubeTranscriptApi
    try:
        from youtube_transcript_api.proxies import GenericProxyConfig, WebshareProxyConfig
    except ImportError:
        GenericProxyConfig = None
        WebshareProxyConfig = None
except ImportError:
    YouTubeTranscriptApi = None
    GenericProxyConfig = None
    WebshareProxyConfig = None

try:
    from reportlab.lib import colors
    from reportlab.lib.enums import TA_LEFT
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
    from reportlab.platypus import Paragraph, SimpleDocTemplate, Spacer, Table, TableStyle
except ImportError:
    A4 = None

try:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
    from pptx.util import Inches, Pt
except ImportError:
    Presentation = None
    RGBColor = None
    MSO_SHAPE = None
    MSO_ANCHOR = None
    PP_ALIGN = None
    Inches = None
    Pt = None

try:
    from pypdf import PdfReader
except ImportError:
    PdfReader = None

try:
    from google.auth.transport.requests import Request as GoogleRequest
    from google.oauth2 import id_token as google_id_token
except ImportError:
    GoogleRequest = None
    google_id_token = None

try:
    import jwt
    from jwt import PyJWKClient
except ImportError:
    jwt = None
    PyJWKClient = None


logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

app = FastAPI()
client = OpenAI()

TRANSCRIPTION_MODEL = os.getenv("TRANSCRIPTION_MODEL", "gpt-4o-transcribe")
FALLBACK_TRANSCRIPTION_MODEL = os.getenv("FALLBACK_TRANSCRIPTION_MODEL", "whisper-1")
STUDY_GUIDE_MODEL = os.getenv("STUDY_GUIDE_MODEL", "gpt-4.1-mini")
VISION_MODEL = os.getenv("VISION_MODEL", "gpt-4.1-mini")
STUDY_CHAT_MODEL = os.getenv("STUDY_CHAT_MODEL", STUDY_GUIDE_MODEL)
ASSET_GENERATION_MODEL = os.getenv("ASSET_GENERATION_MODEL", STUDY_GUIDE_MODEL)
PODCAST_SCRIPT_MODEL = os.getenv("PODCAST_SCRIPT_MODEL", STUDY_GUIDE_MODEL)
PODCAST_TTS_MODEL = os.getenv("PODCAST_TTS_MODEL", "gpt-4o-mini-tts")
TEACHER_SCRIPT_MODEL = os.getenv("TEACHER_SCRIPT_MODEL", STUDY_GUIDE_MODEL)
PRESENTATION_MODEL = os.getenv("PRESENTATION_MODEL", STUDY_GUIDE_MODEL)
GOOGLE_CLIENT_ID = os.getenv("GOOGLE_CLIENT_ID", "")
APPLE_CLIENT_ID = os.getenv("APPLE_CLIENT_ID", "").strip()
APPLE_TEAM_ID = os.getenv("APPLE_TEAM_ID", "").strip()
APPLE_KEY_ID = os.getenv("APPLE_KEY_ID", "").strip()
APPLE_REDIRECT_URI = os.getenv("APPLE_REDIRECT_URI", "").strip()
APPLE_PRIVATE_KEY = os.getenv("APPLE_PRIVATE_KEY", "").strip()
MAX_COMPLETION_TOKENS = int(os.getenv("MAX_COMPLETION_TOKENS", "8000"))
MAX_FILE_SIZE_BYTES = int(os.getenv("MAX_FILE_SIZE_BYTES", str(600 * 1024 * 1024)))
OPENAI_AUDIO_LIMIT_BYTES = int(os.getenv("OPENAI_AUDIO_LIMIT_BYTES", str(25 * 1024 * 1024)))
CHUNK_DURATION_SECONDS = int(os.getenv("CHUNK_DURATION_SECONDS", "300"))
CHUNK_OVERLAP_SECONDS = int(os.getenv("CHUNK_OVERLAP_SECONDS", "20"))
TRANSCRIPTION_AUDIO_BITRATE = os.getenv("TRANSCRIPTION_AUDIO_BITRATE", "48k")
TRANSCRIPTION_AUDIO_SAMPLE_RATE = os.getenv("TRANSCRIPTION_AUDIO_SAMPLE_RATE", "16000")
MAX_STUDY_GUIDE_INPUT_CHARS = int(os.getenv("MAX_STUDY_GUIDE_INPUT_CHARS", "30000"))
MAX_TRANSCRIPT_STUDY_GUIDE_INPUT_CHARS = int(os.getenv("MAX_TRANSCRIPT_STUDY_GUIDE_INPUT_CHARS", "45000"))
STUDY_GUIDE_REQUEST_TIMEOUT = float(os.getenv("STUDY_GUIDE_REQUEST_TIMEOUT", "90"))
VISION_REQUEST_TIMEOUT = float(os.getenv("VISION_REQUEST_TIMEOUT", "45"))
TRANSCRIPTION_REQUEST_TIMEOUT = float(os.getenv("TRANSCRIPTION_REQUEST_TIMEOUT", "1200"))
TRANSCRIPTION_JOB_TIMEOUT = float(
    os.getenv("TRANSCRIPTION_JOB_TIMEOUT", str(max(int(TRANSCRIPTION_REQUEST_TIMEOUT * 4), 3600)))
)
VIDEO_DOWNLOAD_TIMEOUT = float(os.getenv("VIDEO_DOWNLOAD_TIMEOUT", "1200"))
TRANSCRIPTION_RETRIES = int(os.getenv("TRANSCRIPTION_RETRIES", "2"))
MAX_IMAGE_UPLOAD_BYTES = int(os.getenv("MAX_IMAGE_UPLOAD_BYTES", str(15 * 1024 * 1024)))
MAX_SLIDE_UPLOAD_BYTES = int(os.getenv("MAX_SLIDE_UPLOAD_BYTES", str(30 * 1024 * 1024)))
MAX_CHAT_CONTEXT_CHARS = int(os.getenv("MAX_CHAT_CONTEXT_CHARS", "36000"))
MAX_PODCAST_CONTEXT_CHARS = int(os.getenv("MAX_PODCAST_CONTEXT_CHARS", "42000"))
LOGIN_CODE_TTL_MINUTES = int(os.getenv("LOGIN_CODE_TTL_MINUTES", "10"))
REGISTRATION_TOKEN_TTL_MINUTES = int(os.getenv("REGISTRATION_TOKEN_TTL_MINUTES", str(LOGIN_CODE_TTL_MINUTES)))
SESSION_TTL_MINUTES = int(os.getenv("SESSION_TTL_MINUTES", "90"))
SESSION_REFRESH_WINDOW_MINUTES = int(os.getenv("SESSION_REFRESH_WINDOW_MINUTES", "20"))
SUPPORT_EMAIL = os.getenv("SUPPORT_EMAIL", "mabasoasakhe@gmail.com").strip()
MAX_HISTORY_ITEMS = int(os.getenv("MAX_HISTORY_ITEMS", "24"))
ADMIN_DASHBOARD_AUDIT_LOG_LIMIT = int(os.getenv("ADMIN_DASHBOARD_AUDIT_LOG_LIMIT", "8000"))
ADMIN_DASHBOARD_HISTORY_LIMIT = int(os.getenv("ADMIN_DASHBOARD_HISTORY_LIMIT", "1200"))
MIN_PASSWORD_LENGTH = int(os.getenv("MIN_PASSWORD_LENGTH", "8"))
PASSWORD_HASH_ITERATIONS = int(os.getenv("PASSWORD_HASH_ITERATIONS", "200000"))
ADMIN_EMAILS_RAW = os.getenv("ADMIN_EMAILS", os.getenv("ADMIN_EMAIL", "")).strip()
ADMIN_LOGIN_MAX_ATTEMPTS = int(os.getenv("ADMIN_LOGIN_MAX_ATTEMPTS", "3"))
ADMIN_LOGIN_LOCKOUT_MINUTES = int(os.getenv("ADMIN_LOGIN_LOCKOUT_MINUTES", "15"))
PODCAST_REQUEST_TIMEOUT = float(os.getenv("PODCAST_REQUEST_TIMEOUT", "180"))
PODCAST_TTS_TIMEOUT = float(os.getenv("PODCAST_TTS_TIMEOUT", "120"))
TEACHER_REQUEST_TIMEOUT = float(os.getenv("TEACHER_REQUEST_TIMEOUT", "120"))
PRESENTATION_REQUEST_TIMEOUT = float(os.getenv("PRESENTATION_REQUEST_TIMEOUT", "150"))
STUDY_IMAGE_QUERY_TIMEOUT = float(os.getenv("STUDY_IMAGE_QUERY_TIMEOUT", "45"))
STUDY_IMAGE_SEARCH_TIMEOUT = float(os.getenv("STUDY_IMAGE_SEARCH_TIMEOUT", "12"))
MAX_STUDY_IMAGES = int(os.getenv("MAX_STUDY_IMAGES", "6"))
UPLOAD_DIR = Path(tempfile.gettempdir()) / "lecture-ai-project"
UPLOAD_DIR.mkdir(parents=True, exist_ok=True)
PODCAST_OUTPUT_DIR = UPLOAD_DIR / "podcasts"
PODCAST_OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
PRESENTATION_OUTPUT_DIR = UPLOAD_DIR / "presentations"
PRESENTATION_OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
DB_PATH = Path(__file__).with_name("mabaso_ai.db")
APP_SECRET = os.getenv("APP_SECRET", os.getenv("OPENAI_API_KEY", "mabaso-dev-secret"))
SESSION_TOKEN_PREFIX = "mabaso.v1"
WIKIMEDIA_API_URL = "https://commons.wikimedia.org/w/api.php"
APPLE_IDENTITY_ISSUER = "https://appleid.apple.com"
APPLE_JWKS_URL = f"{APPLE_IDENTITY_ISSUER}/auth/keys"
APPLE_TOKEN_VALIDATION_URL = f"{APPLE_IDENTITY_ISSUER}/auth/token"
APPLE_AUTH_HTTP_TIMEOUT = float(os.getenv("APPLE_AUTH_HTTP_TIMEOUT", "15"))
YOUTUBE_LANGUAGE_PREFERENCES = ("en", "en-US", "en-GB")
YTDLP_YOUTUBE_PLAYER_CLIENTS = ("ios", "android_vr", "web_safari", "mweb", "tv_simply")
YOUTUBE_USER_AGENT = os.getenv(
    "YOUTUBE_USER_AGENT",
    (
        "Mozilla/5.0 (Windows NT 10.0; Win64; x64) "
        "AppleWebKit/537.36 (KHTML, like Gecko) "
        "Chrome/135.0.0.0 Safari/537.36"
    ),
).strip()
YOUTUBE_PROXY_HTTP_URL = os.getenv("YOUTUBE_PROXY_HTTP_URL", "").strip()
YOUTUBE_PROXY_HTTPS_URL = os.getenv("YOUTUBE_PROXY_HTTPS_URL", "").strip()
YOUTUBE_WEBSHARE_PROXY_USERNAME = os.getenv("YOUTUBE_WEBSHARE_PROXY_USERNAME", "").strip()
YOUTUBE_WEBSHARE_PROXY_PASSWORD = os.getenv("YOUTUBE_WEBSHARE_PROXY_PASSWORD", "").strip()
YOUTUBE_WEBSHARE_PROXY_HOST = os.getenv("YOUTUBE_WEBSHARE_PROXY_HOST", "p.webshare.io").strip()
YOUTUBE_WEBSHARE_PROXY_PORT = os.getenv("YOUTUBE_WEBSHARE_PROXY_PORT", "80").strip()
YTDLP_IMPERSONATE_TARGET = os.getenv("YTDLP_IMPERSONATE_TARGET", "chrome").strip()
YOUTUBE_WEBSHARE_PROXY_LOCATIONS = tuple(
    location.strip()
    for location in os.getenv("YOUTUBE_WEBSHARE_PROXY_LOCATIONS", "").split(",")
    if location.strip()
)
YOUTUBE_COOKIES_FILE = os.getenv("YOUTUBE_COOKIES_FILE", "").strip()
YOUTUBE_COOKIES_TXT = os.getenv("YOUTUBE_COOKIES_TXT", "").strip()

jobs: dict[str, dict] = {}
apple_jwk_client = PyJWKClient(APPLE_JWKS_URL) if PyJWKClient is not None else None

LEGACY_STUDY_GUIDE_PROMPT = """
You are an expert academic assistant for university students.

Convert the available lecture material into concise, high-value study notes.
Keep the response practical and faster to generate than a long textbook rewrite.

Return clean Markdown with these sections:
- LECTURE TITLE
- SHORT SUMMARY
- KEY CONCEPTS
- IMPORTANT DEFINITIONS
- IMPORTANT FORMULAS
- WORKED EXAMPLES
- STEP-BY-STEP EXPLANATIONS
- ADVANTAGES AND DISADVANTAGES
- COMMON MISTAKES TO AVOID
- QUICK REVISION PLAN
- VISUAL AIDS
- REAL-WORLD EXAMPLES
- PRACTICE QUESTIONS AND ANSWERS
- FLASHCARDS
- EXAM TIPS

Rules:
- Use clear section headings in uppercase and make each heading bold using Markdown.
- Leave a blank line between every section.
- Keep the notes compact and easy to revise.
- Write like a strong school or university revision handout: high-level first, then detailed support.
- Use bullet points where useful, especially for summaries, concepts, definitions, exam tips, and revision steps.
- Prefer short readable bullets over long dense paragraphs.
- Do not simply paraphrase the transcript line by line. Reorganize the material into teachable notes.
- In LECTURE TITLE, write one clean topic line only so the student can see the lecture topic immediately at the top.
- If formulas appear, rewrite them in readable human style.
- Never use LaTeX syntax or math delimiters such as \\, \\, $$, \\frac, \\int, \\mathcal, or \\begin.
- Do not write exponents with caret notation like s^2, t^n, or e^(-at) in the final answer.
- Write exponents and indices in textbook style using proper symbols where possible, for example s², tⁿ, e⁻ᵃᵗ, ∫₀∞.
- Write formulas the way a lecturer would write them on a board using plain readable text.
- In IMPORTANT FORMULAS, prefer a study-sheet layout using short lines in this style:
  t -> 1 / s²
  tⁿ -> n! / s⁽ⁿ ⁺ ¹⁾
  e⁻ᵃᵗ -> 1 / (s + a)
  eᵃᵗ -> 1 / (s - a)
- If the lecture covers transform pairs, conversion rules, or standard results, list them as readable two-column style mappings.
- Example style:
  F(s) = ∫₀∞ e⁻ˢᵗ f(t) dt
  u(t - a) = 0 for t < a, and 1 for t >= a
- Put a blank line before and after each formula block so it is easy to read.
- In WORKED EXAMPLES, include at least one step-by-step solved example when the lecture contains a problem, calculation, derivation, or sum.
- In STEP-BY-STEP EXPLANATIONS, explain what the student should notice, why it matters, and how to avoid confusing it with similar ideas.
- In ADVANTAGES AND DISADVANTAGES, give practical study-focused pros, limits, or caution points that help a student know when the idea is useful and where it becomes confusing.
- In COMMON MISTAKES TO AVOID, list short warnings about misunderstandings, skipped steps, wrong formula use, or revision traps.
- In QUICK REVISION PLAN, give a short sequence a student can follow before a class test or exam.
- In VISUAL AIDS, include simple ASCII diagrams, neat comparison tables, flow layouts, or graph sketches when they help explain the topic.
- If the lecture covers concrete physical things, types, components, specimens, machines, instruments, valves, organs, or structures, name the important visual subtypes clearly and describe the visible features students should recognise.
- In VISUAL AIDS, explicitly mention concrete objects or subtype names that would benefit from a real photo reference.
- Only include a bar graph, line graph, axis sketch, or trend diagram when the lecture discusses data, change over time, or relationships between variables. Do not invent fake numerical data.
- Use simple text layouts that students can read easily in plain Markdown.
- In PRACTICE QUESTIONS AND ANSWERS, include 4 to 8 short exam-style revision questions with brief model answers.
- Keep PRACTICE QUESTIONS AND ANSWERS clearly separate from FLASHCARDS.
- In FLASHCARDS, use this exact style for every card:
  Q: ...

  A: ...
- Leave a blank line between Q and A, and a blank line between flashcards so they are easy to read.
- If lecture notes or lecture slides are provided together with the transcript, use all sources together. Prefer explicit formulas, worked examples, definitions, and likely assessment points from the slides or notes when they improve clarity.
- Prefer well-structured notes, slides, and past-paper references over messy transcript wording when they explain the topic more clearly.
- Ignore clearly corrupted OCR, random fragments, duplicate scraps, broken symbols, or unrelated extraction noise instead of mixing them into the final notes.
- Do not let a noisy transcript dominate cleaner note or slide material.
- If past question papers are provided, use them as an assessment reference. Infer recurring topics, phrasing style, and mark patterns from them, but do not copy their questions verbatim.
- If the transcript is missing but lecture notes, lecture slides, or past question papers are provided, still generate the study guide from those sources and mention when the source material is limited.
- If the lecture skips steps, fill in only the most important missing steps.
- If a topic is unclear, say "Not clearly covered in transcript".
- Do not force photo-oriented commentary unless the topic clearly involves physical things students should recognise by sight.
- Do not include YouTube links or long export suggestions.
"""

STUDY_GUIDE_PROMPT = """
You are an advanced academic study-guide generation engine.

Your task is NOT to simply summarize.

Your task is to:
1. analyze the academic content,
2. detect the subject and topic structure,
3. identify what is educationally important,
4. generate a clean university-level study guide.

CORE BEHAVIOR RULES
- Do NOT use the same sections every time.
- Dynamically adapt the structure to the content.
- Only include sections that are relevant to the topic.
- Avoid filler information.
- Avoid unnecessary advantages/disadvantages unless central to the material.
- Avoid repeating concepts.
- Avoid overly long explanations.
- Prioritize clarity, organization, and revision efficiency.

The notes should feel like:
- a premium edtech platform,
- a university study guide,
- concise but intelligent,
- visually organized,
- optimized for exam revision.

STEP 1 - ANALYZE CONTENT
Before generating notes, identify:
- academic subject
- topic type
- educational intent
- major concepts
- relationships between ideas
- exam-relevant material
- definitions and terminology
- theories, processes, mechanisms, or frameworks if present

Then decide the best note structure for this specific content.

STEP 2 - CHOOSE STRUCTURE DYNAMICALLY
Use only sections that naturally fit the material.

Examples:
- Theory content: Key Concepts, Main Assumptions, Important Thinkers, Examples, Criticisms, Applications
- Science content: Definitions, Mechanisms, Functions, Processes, Stages, Comparisons
- History content: Causes, Timeline, Major Events, Consequences, Historical Significance
- Law content: Legal Principles, Elements, Case Applications, Exceptions
- Medical content: Symptoms, Causes, Diagnosis, Treatment, Risk Factors
- Technical content: Components, Architecture, Workflow, Inputs/Outputs, Advantages only if relevant

Do not force sections unnecessarily.

OUTPUT STYLE
- Return clean Markdown.
- Start with one H1 line only: # Topic Title
- After the title, include a short overview section using either ## SHORT SUMMARY or ## Introduction / Overview.
- Use clear headings and subheadings.
- Mix bullets with short explanations.
- Keep paragraphs short.
- Use spacing generously.
- Make notes easy to skim.
- Keep the tone academic, simple, and clear.
- Important concepts should get deeper explanation.
- Minor concepts should stay concise.
- Definitions should be short and precise.
- Complex ideas should be simplified without losing meaning.
- Examples should appear only when genuinely useful.
- Add Key Questions only when educationally valuable.

MANDATORY COMPATIBILITY RULES
- Because the app builds formulas, worked-example, flashcard, quiz, teacher-mode, and presentation assets from the guide, keep these exact headings whenever the content supports them:
  - ## IMPORTANT FORMULAS
  - ## WORKED EXAMPLES
  - ## STEP-BY-STEP EXPLANATIONS
  - ## PRACTICE QUESTIONS AND ANSWERS
  - ## FLASHCARDS
- If the topic includes formulas, always include ## IMPORTANT FORMULAS.
- If the topic includes calculations, derivations, worked procedures, or problem-solving, always include ## WORKED EXAMPLES.
- Always include ## STEP-BY-STEP EXPLANATIONS with 3 to 6 clear sequenced steps or step-labeled bullets that teach the learner how to move through the method, process, argument, or reasoning. This section is mandatory and is never optional.
- Always include ## PRACTICE QUESTIONS AND ANSWERS with 4 to 8 short exam-style questions and brief model answers.
- Always include ## FLASHCARDS and use this exact card format for every card:
  Q: ...

  A: ...
- Leave a blank line between Q and A, and a blank line between flashcards.

FORMATTING AND DEPTH RULES
- Do not simply paraphrase the transcript line by line. Reorganize the material into teachable notes.
- Keep the notes concise but not shallow.
- Use bullets where they improve revision speed.
- Do not include generic motivational text.
- Do not add filler conclusions unless they genuinely help.
- Do not add unnecessary summaries after every section.
- If multiple distinct topics appear in the source material, separate them clearly instead of blending them together.
- If one source belongs to a different topic, isolate it under its own heading.

FORMULA RULES
- If formulas appear, rewrite them in readable human style.
- Never use LaTeX syntax or math delimiters such as \\, $$, \\frac, \\int, \\mathcal, or \\begin.
- Do not use caret notation like s^2 or t^n in the final answer.
- Write formulas the way a lecturer would write them on a board using plain readable text.
- Prefer short readable mappings when listing standard transform pairs, rules, or conversions.
- Put a blank line before and after each formula block.

VISUAL LEARNING RULES
- After difficult concepts, add [Suggested Visual: ...] only when it genuinely improves understanding.
- Good examples:
  - [Suggested Visual: Flowchart of cellular respiration]
  - [Suggested Visual: Diagram comparing Functionalism vs Conflict Theory]
- Only suggest visuals that improve understanding.
- Make each suggested visual explicit enough for the app to render it. Include the main stages, compared sides, plotted signals, axes, or labels inside the suggestion itself.
- Prefer render-friendly phrasing such as:
  - [Suggested Visual: Flowchart - Input -> Transform -> Output]
  - [Suggested Visual: Comparison table - Continuous vs Discrete | formula, operation, output]
  - [Suggested Visual: Plot - x(t), h(t), y(t)]
- If the lecture covers concrete physical things such as organs, instruments, valves, structures, machines, or components, mention the visual subtypes students should recognize.
- Only include charts, graphs, axes, or trend sketches when the lecture discusses data or variable relationships. Do not invent fake numerical data.

SOURCE PRIORITY RULES
- If lecture notes, slides, and transcript are all provided, use all of them together.
- Prefer clearer lecturer notes, slides, formulas, definitions, and worked examples over messy transcript wording.
- Prefer well-structured notes, slides, and past-paper references over noisy OCR or broken transcript fragments.
- Ignore corrupted OCR, duplicate scraps, unrelated fragments, and broken symbols.
- Do not let a noisy transcript dominate cleaner source material.
- If past question papers are provided, use them to infer recurring themes, command words, and assessment style, but do not copy their questions verbatim.
- If the transcript is missing, still generate the guide from the other supplied sources and mention when the source material is limited.
- If the lecture skips steps, fill in only the most important missing steps.
- If a topic is unclear, say "Not clearly covered in the supplied material."

QUALITY RULES
- The output should feel intelligent, adaptive, academically useful, modern, polished, and revision-friendly.
- Think like an expert university tutor creating revision notes for students before exams.
- Write naturally like a highly organized top university student.
- Use more comparison tables, process flows, and visual-learning suggestions when concepts are complex.
- Focus on conceptual understanding and intuitive explanations, not just memorization.
- Do not include YouTube links or long export suggestions.
"""

GUIDE_SECTION_HEADINGS = [
    "LECTURE TITLE",
    "SHORT SUMMARY",
    "KEY CONCEPTS",
    "IMPORTANT DEFINITIONS",
    "IMPORTANT FORMULAS",
    "WORKED EXAMPLES",
    "STEP-BY-STEP EXPLANATIONS",
    "ADVANTAGES AND DISADVANTAGES",
    "COMMON MISTAKES TO AVOID",
    "QUICK REVISION PLAN",
    "VISUAL AIDS",
    "REAL-WORLD EXAMPLES",
    "PRACTICE QUESTIONS AND ANSWERS",
    "FLASHCARDS",
    "EXAM TIPS",
]

GUIDE_SECTION_ALIASES = {
    "lecture title": "LECTURE TITLE",
    "topic title": "LECTURE TITLE",
    "title": "LECTURE TITLE",
    "short summary": "SHORT SUMMARY",
    "summary": "SHORT SUMMARY",
    "overview": "SHORT SUMMARY",
    "brief overview": "SHORT SUMMARY",
    "introduction": "SHORT SUMMARY",
    "introduction / overview": "SHORT SUMMARY",
    "introduction/overview": "SHORT SUMMARY",
    "key concepts": "KEY CONCEPTS",
    "main concepts": "KEY CONCEPTS",
    "important definitions": "IMPORTANT DEFINITIONS",
    "definitions": "IMPORTANT DEFINITIONS",
    "important formulas": "IMPORTANT FORMULAS",
    "formulas": "IMPORTANT FORMULAS",
    "worked examples": "WORKED EXAMPLES",
    "examples": "WORKED EXAMPLES",
    "step-by-step explanations": "STEP-BY-STEP EXPLANATIONS",
    "step by step explanations": "STEP-BY-STEP EXPLANATIONS",
    "step-by-step explanation": "STEP-BY-STEP EXPLANATIONS",
    "step by step explanation": "STEP-BY-STEP EXPLANATIONS",
    "advantages and disadvantages": "ADVANTAGES AND DISADVANTAGES",
    "common mistakes to avoid": "COMMON MISTAKES TO AVOID",
    "quick revision plan": "QUICK REVISION PLAN",
    "visual aids": "VISUAL AIDS",
    "suggested visuals": "VISUAL AIDS",
    "suggested visual": "VISUAL AIDS",
    "visual learning suggestions": "VISUAL AIDS",
    "real-world examples": "REAL-WORLD EXAMPLES",
    "real world examples": "REAL-WORLD EXAMPLES",
    "practice questions and answers": "PRACTICE QUESTIONS AND ANSWERS",
    "practice questions": "PRACTICE QUESTIONS AND ANSWERS",
    "flashcards": "FLASHCARDS",
    "exam tips": "EXAM TIPS",
    "exam-focused takeaways": "EXAM TIPS",
    "exam focused takeaways": "EXAM TIPS",
    "quick recap": "EXAM TIPS",
}


class StudyGuideRequest(BaseModel):
    transcript: str
    lecture_notes: str = ""
    lecture_slides: str = ""
    past_question_papers: str = ""
    language: str = "English"
    reference_images: list[str] = []


class PodcastGenerationRequest(BaseModel):
    transcript: str = ""
    summary: str = ""
    lecture_notes: str = ""
    lecture_slides: str = ""
    past_question_papers: str = ""
    speaker_count: int = 2
    target_minutes: int = 10
    language: str = "English"


class PresentationGenerationRequest(BaseModel):
    transcript: str = ""
    summary: str = ""
    lecture_notes: str = ""
    lecture_slides: str = ""
    past_question_papers: str = ""
    design_id: str = "emerald-scholar"
    language: str = "English"
    reference_images: list[str] = []


class QuizGenerationRequest(BaseModel):
    transcript: str = ""
    summary: str = ""
    lecture_notes: str = ""
    lecture_slides: str = ""
    past_question_papers: str = ""
    language: str = "English"


class TeacherLessonRequest(BaseModel):
    transcript: str = ""
    summary: str = ""
    lecture_notes: str = ""
    lecture_slides: str = ""
    past_question_papers: str = ""
    language: str = "English"


class VideoUrlTranscriptionRequest(BaseModel):
    video_url: str


class RequestCodeRequest(BaseModel):
    email: str


class VerifyCodeRequest(BaseModel):
    email: str
    code: str


class EmailPasswordAuthRequest(BaseModel):
    email: str
    password: str
    mode: str = "login"


class EmailPasswordVerifyRequest(BaseModel):
    email: str
    password: str
    code: str
    mode: str = "login"


class EmailPasswordRegistrationCompleteRequest(BaseModel):
    email: str
    registration_token: str
    password: str


class GoogleAuthRequest(BaseModel):
    credential: str


class AppleAuthRequest(BaseModel):
    authorization_code: str = ""
    id_token: str = ""
    nonce: str = ""
    state: str = ""
    user: dict[str, Any] | None = None
    redirect_uri: str = ""


class SupportMessageRequest(BaseModel):
    message: str
    page: str = ""
    client_request_id: str = ""


class HistorySyncRequest(BaseModel):
    items: list[dict[str, Any]] = []


class ChatTurn(BaseModel):
    role: str
    content: str


class StudyChatRequest(BaseModel):
    question: str
    transcript: str = ""
    summary: str = ""
    lecture_notes: str = ""
    lecture_slides: str = ""
    past_question_papers: str = ""
    history: list[ChatTurn] = []
    reference_images: list[str] = []
    language: str = "English"


class SessionModeRequest(BaseModel):
    mode: str = "user"


class PdfSection(BaseModel):
    title: str
    content: str = ""


class PdfExportRequest(BaseModel):
    title: str
    sections: list[PdfSection]


class CollaborationRoomCreateRequest(BaseModel):
    title: str
    transcript: str = ""
    summary: str = ""
    formula: str = ""
    example: str = ""
    lecture_notes: str = ""
    lecture_slides: str = ""
    shared_notes: str = ""
    flashcards: list[dict[str, str]] = []
    quiz_questions: list[dict[str, Any]] = []
    invited_emails: list[str] = []
    active_tab: str = "guide"
    test_visibility: str = "private"


class CollaborationMessageRequest(BaseModel):
    content: str


class CollaborationSharedNotesRequest(BaseModel):
    shared_notes: str = ""


class CollaborationActiveTabRequest(BaseModel):
    active_tab: str


class CollaborationTestVisibilityRequest(BaseModel):
    test_visibility: str


class CollaborationQuizAnswerRequest(BaseModel):
    question_number: str
    answer_text: str = ""


class AdminUserStatusRequest(BaseModel):
    status: str = "active"


app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)


def utc_now() -> datetime:
    return datetime.now(timezone.utc)


def iso_in_future(*, minutes: int = 0, days: int = 0) -> str:
    return (utc_now() + timedelta(minutes=minutes, days=days)).isoformat()


def get_db_connection() -> sqlite3.Connection:
    connection = sqlite3.connect(DB_PATH)
    connection.row_factory = sqlite3.Row
    return connection


def init_db():
    with get_db_connection() as connection:
        connection.execute(
            """
            CREATE TABLE IF NOT EXISTS users (
                email TEXT PRIMARY KEY,
                created_at TEXT NOT NULL,
                verified_at TEXT
            )
            """
        )
        connection.execute(
            """
            CREATE TABLE IF NOT EXISTS login_codes (
                email TEXT PRIMARY KEY,
                code_hash TEXT NOT NULL,
                expires_at TEXT NOT NULL,
                created_at TEXT NOT NULL
            )
            """
        )
        connection.execute(
            """
            CREATE TABLE IF NOT EXISTS collaboration_rooms (
                id TEXT PRIMARY KEY,
                owner_email TEXT NOT NULL,
                title TEXT NOT NULL,
                transcript TEXT NOT NULL,
                summary TEXT NOT NULL,
                formula TEXT NOT NULL,
                example TEXT NOT NULL,
                lecture_notes TEXT NOT NULL,
                lecture_slides TEXT NOT NULL,
                shared_notes TEXT NOT NULL,
                flashcards_json TEXT NOT NULL,
                quiz_questions_json TEXT NOT NULL,
                active_tab TEXT NOT NULL,
                test_visibility TEXT NOT NULL,
                created_at TEXT NOT NULL,
                updated_at TEXT NOT NULL
            )
            """
        )
        connection.execute(
            """
            CREATE TABLE IF NOT EXISTS collaboration_room_members (
                room_id TEXT NOT NULL,
                email TEXT NOT NULL,
                role TEXT NOT NULL,
                created_at TEXT NOT NULL,
                PRIMARY KEY (room_id, email)
            )
            """
        )
        connection.execute(
            """
            CREATE TABLE IF NOT EXISTS collaboration_room_messages (
                id TEXT PRIMARY KEY,
                room_id TEXT NOT NULL,
                author_email TEXT NOT NULL,
                content TEXT NOT NULL,
                created_at TEXT NOT NULL
            )
            """
        )
        connection.execute(
            """
            CREATE TABLE IF NOT EXISTS collaboration_room_answers (
                room_id TEXT NOT NULL,
                question_number TEXT NOT NULL,
                author_email TEXT NOT NULL,
                answer_text TEXT NOT NULL,
                updated_at TEXT NOT NULL,
                PRIMARY KEY (room_id, question_number, author_email)
            )
            """
        )
        connection.execute(
            """
            CREATE TABLE IF NOT EXISTS sessions (
                token_hash TEXT PRIMARY KEY,
                email TEXT NOT NULL,
                expires_at TEXT NOT NULL,
                created_at TEXT NOT NULL
            )
            """
        )
        connection.execute(
            """
            CREATE TABLE IF NOT EXISTS revoked_sessions (
                token_hash TEXT PRIMARY KEY,
                expires_at TEXT NOT NULL,
                created_at TEXT NOT NULL
            )
            """
        )
        connection.execute(
            """
            CREATE TABLE IF NOT EXISTS study_history_items (
                email TEXT NOT NULL,
                id TEXT NOT NULL,
                payload_json TEXT NOT NULL,
                created_at TEXT NOT NULL,
                updated_at TEXT NOT NULL,
                PRIMARY KEY (email, id)
            )
            """
        )
        connection.execute(
            """
            CREATE INDEX IF NOT EXISTS idx_study_history_items_email_updated_at
            ON study_history_items (email, updated_at DESC)
            """
        )
        connection.execute(
            """
            CREATE TABLE IF NOT EXISTS email_password_credentials (
                email TEXT PRIMARY KEY,
                password_hash TEXT NOT NULL,
                password_salt TEXT NOT NULL,
                created_at TEXT NOT NULL,
                updated_at TEXT NOT NULL
            )
            """
        )
        connection.execute(
            """
            CREATE TABLE IF NOT EXISTS pending_email_password_registrations (
                email TEXT PRIMARY KEY,
                token_hash TEXT NOT NULL,
                expires_at TEXT NOT NULL,
                created_at TEXT NOT NULL
            )
            """
        )
        connection.execute(
            """
            CREATE TABLE IF NOT EXISTS audit_logs (
                id TEXT PRIMARY KEY,
                email TEXT NOT NULL,
                ip_address TEXT NOT NULL,
                user_agent TEXT NOT NULL,
                action TEXT NOT NULL,
                resource_type TEXT NOT NULL,
                resource_name TEXT NOT NULL,
                duration_ms INTEGER NOT NULL,
                status TEXT NOT NULL,
                metadata_json TEXT NOT NULL,
                created_at TEXT NOT NULL
            )
            """
        )
        connection.execute(
            """
            CREATE INDEX IF NOT EXISTS idx_audit_logs_created_at
            ON audit_logs (created_at DESC)
            """
        )
        connection.execute(
            """
            CREATE INDEX IF NOT EXISTS idx_audit_logs_email_created_at
            ON audit_logs (email, created_at DESC)
            """
        )
        connection.execute(
            """
            CREATE TABLE IF NOT EXISTS admin_login_attempts (
                email TEXT NOT NULL,
                ip_address TEXT NOT NULL,
                failure_count INTEGER NOT NULL,
                last_failed_at TEXT NOT NULL,
                locked_until TEXT NOT NULL,
                PRIMARY KEY (email, ip_address)
            )
            """
        )
        connection.execute(
            """
            CREATE TABLE IF NOT EXISTS user_account_states (
                email TEXT PRIMARY KEY,
                status TEXT NOT NULL,
                updated_at TEXT NOT NULL,
                updated_by TEXT NOT NULL
            )
            """
        )
        connection.execute(
            """
            CREATE TABLE IF NOT EXISTS support_messages (
                id TEXT PRIMARY KEY,
                client_request_id TEXT NOT NULL UNIQUE,
                email TEXT NOT NULL,
                page TEXT NOT NULL,
                message TEXT NOT NULL,
                email_delivery_status TEXT NOT NULL,
                email_error TEXT NOT NULL,
                created_at TEXT NOT NULL,
                updated_at TEXT NOT NULL
            )
            """
        )
        connection.execute(
            """
            CREATE INDEX IF NOT EXISTS idx_support_messages_created_at
            ON support_messages (created_at DESC)
            """
        )


def hash_value(value: str) -> str:
    payload = f"{APP_SECRET}:{value}"
    return hashlib.sha256(payload.encode("utf-8")).hexdigest()


def encode_token_component(value: str) -> str:
    return base64.urlsafe_b64encode(value.encode("utf-8")).decode("ascii").rstrip("=")


def decode_token_component(value: str) -> str:
    padding = "=" * (-len(value) % 4)
    return base64.urlsafe_b64decode(f"{value}{padding}".encode("ascii")).decode("utf-8")


def is_signed_session_token(token: str) -> bool:
    return bool(token and token.startswith(f"{SESSION_TOKEN_PREFIX}."))


def build_signed_session_token(
    email: str,
    expires_at: datetime | None = None,
    session_mode: str = "user",
) -> str:
    resolved_mode = normalize_session_mode(session_mode, email)
    expiry = expires_at or (utc_now() + timedelta(minutes=SESSION_TTL_MINUTES))
    payload = json.dumps(
        {
            "email": email,
            "exp": int(expiry.timestamp()),
            "iat": int(utc_now().timestamp()),
            "nonce": uuid4().hex,
            "mode": resolved_mode,
        },
        separators=(",", ":"),
        sort_keys=True,
    )
    encoded_payload = encode_token_component(payload)
    signature = hmac.new(
        APP_SECRET.encode("utf-8"),
        encoded_payload.encode("ascii"),
        hashlib.sha256,
    ).hexdigest()
    return f"{SESSION_TOKEN_PREFIX}.{encoded_payload}.{signature}"


def decode_signed_session_token(token: str) -> dict[str, Any] | None:
    if not is_signed_session_token(token):
        return None

    parts = token.split(".")
    if len(parts) != 4:
        return None

    _, version, encoded_payload, signature = parts
    if version != "v1":
        return None

    expected_signature = hmac.new(
        APP_SECRET.encode("utf-8"),
        encoded_payload.encode("ascii"),
        hashlib.sha256,
    ).hexdigest()
    if not hmac.compare_digest(expected_signature, signature):
        return None

    try:
        payload = json.loads(decode_token_component(encoded_payload))
    except (json.JSONDecodeError, UnicodeDecodeError, binascii.Error, ValueError):
        return None

    email = normalize_email(payload.get("email", ""))
    expires_at = int(payload.get("exp", 0) or 0)
    if not email or not expires_at:
        return None
    return {
        "email": email,
        "exp": expires_at,
        "mode": normalize_session_mode(str(payload.get("mode") or "user"), email),
    }


def normalize_email(email: str) -> str:
    return (email or "").strip().lower()


def validate_email_address(email: str) -> str:
    normalized = normalize_email(email)
    if not normalized:
        raise HTTPException(status_code=400, detail="Email address is required.")
    if not re.fullmatch(r"[^@\s]+@[^@\s]+\.[^@\s]+", normalized):
        raise HTTPException(status_code=400, detail="Enter a valid email address.")
    return normalized


SUPPORTED_OUTPUT_LANGUAGES: dict[str, str] = {
    "english": "English",
    "en": "English",
    "isizulu": "isiZulu",
    "zulu": "isiZulu",
    "zu": "isiZulu",
    "afrikaans": "Afrikaans",
    "af": "Afrikaans",
    "isixhosa": "isiXhosa",
    "xhosa": "isiXhosa",
    "xh": "isiXhosa",
    "sesotho": "Sesotho",
    "st": "Sesotho",
    "setswana": "Setswana",
    "tswana": "Setswana",
    "tn": "Setswana",
    "sepedi": "Sepedi",
    "northern sotho": "Sepedi",
    "nso": "Sepedi",
    "portuguese": "Portuguese",
    "pt": "Portuguese",
    "french": "French",
    "fr": "French",
}


def normalize_output_language(value: str) -> str:
    cleaned = re.sub(r"\s+", " ", (value or "").strip()).lower()
    if not cleaned:
        return "English"
    return SUPPORTED_OUTPUT_LANGUAGES.get(cleaned, cleaned.title())


def get_admin_email_set() -> set[str]:
    return {
        normalize_email(item)
        for item in re.split(r"[\s,;]+", ADMIN_EMAILS_RAW)
        if normalize_email(item)
    }


def is_admin_email(email: str) -> bool:
    return normalize_email(email) in get_admin_email_set()


def get_available_auth_modes(email: str) -> list[str]:
    return ["user", "admin"] if is_admin_email(email) else ["user"]


def normalize_session_mode(mode: str, email: str) -> str:
    requested_mode = (mode or "user").strip().lower() or "user"
    if requested_mode not in {"user", "admin"}:
        raise HTTPException(status_code=400, detail="Session mode must be either user or admin.")
    if requested_mode == "admin" and not is_admin_email(email):
        raise HTTPException(status_code=403, detail="Admin access is not available for this account.")
    return requested_mode


def normalize_user_account_status(value: str) -> str:
    status = (value or "active").strip().lower() or "active"
    if status not in {"active", "suspended"}:
        raise HTTPException(status_code=400, detail="User status must be either active or suspended.")
    return status


def get_user_account_status(email: str) -> str:
    normalized_email = normalize_email(email)
    if not normalized_email:
        return "active"
    with get_db_connection() as connection:
        row = connection.execute(
            "SELECT status FROM user_account_states WHERE email = ?",
            (normalized_email,),
        ).fetchone()
    return normalize_user_account_status(row["status"]) if row else "active"


def ensure_user_account_is_active(email: str):
    if get_user_account_status(email) != "active":
        raise HTTPException(status_code=403, detail="This account has been suspended. Contact the administrator.")


def set_user_account_status(email: str, status: str, updated_by: str):
    normalized_email = validate_email_address(email)
    normalized_status = normalize_user_account_status(status)
    with get_db_connection() as connection:
        connection.execute(
            "INSERT OR IGNORE INTO users (email, created_at) VALUES (?, ?)",
            (normalized_email, utc_now().isoformat()),
        )
        connection.execute(
            """
            INSERT INTO user_account_states (email, status, updated_at, updated_by)
            VALUES (?, ?, ?, ?)
            ON CONFLICT(email) DO UPDATE SET
                status = excluded.status,
                updated_at = excluded.updated_at,
                updated_by = excluded.updated_by
            """,
            (normalized_email, normalized_status, utc_now().isoformat(), normalize_email(updated_by)),
        )


def revoke_all_sessions_for_user(email: str):
    normalized_email = normalize_email(email)
    if not normalized_email:
        return
    with get_db_connection() as connection:
        rows = connection.execute(
            "SELECT token_hash, expires_at FROM sessions WHERE email = ?",
            (normalized_email,),
        ).fetchall()
        for row in rows:
            connection.execute(
                """
                INSERT OR REPLACE INTO revoked_sessions (token_hash, expires_at, created_at)
                VALUES (?, ?, ?)
                """,
                (row["token_hash"], row["expires_at"], utc_now().isoformat()),
            )
        connection.execute("DELETE FROM sessions WHERE email = ?", (normalized_email,))


def get_client_ip(request: Request | None) -> str:
    if request is None:
        return ""

    forwarded = (request.headers.get("x-forwarded-for") or "").split(",")[0].strip()
    if forwarded:
        return forwarded[:120]

    real_ip = (request.headers.get("x-real-ip") or "").strip()
    if real_ip:
        return real_ip[:120]

    client_host = getattr(request.client, "host", "") or ""
    return client_host[:120]


def record_audit_log(
    *,
    action: str,
    status: str = "success",
    email: str = "",
    request: Request | None = None,
    resource_type: str = "",
    resource_name: str = "",
    duration_ms: int = 0,
    metadata: dict[str, Any] | None = None,
):
    normalized_email = normalize_email(email)
    safe_status = (status or "success").strip().lower() or "success"
    payload = metadata or {}
    with get_db_connection() as connection:
        connection.execute(
            """
            INSERT INTO audit_logs (
                id, email, ip_address, user_agent, action, resource_type,
                resource_name, duration_ms, status, metadata_json, created_at
            )
            VALUES (?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?)
            """,
            (
                uuid4().hex,
                normalized_email,
                get_client_ip(request),
                (request.headers.get("user-agent") or "")[:400] if request is not None else "",
                (action or "").strip()[:120],
                (resource_type or "").strip()[:80],
                (resource_name or "").strip()[:200],
                max(0, int(duration_ms or 0)),
                safe_status[:40],
                json.dumps(payload, ensure_ascii=False),
                utc_now().isoformat(),
            ),
        )


def get_admin_login_attempt_state(email: str, ip_address: str) -> dict[str, Any]:
    normalized_email = normalize_email(email)
    with get_db_connection() as connection:
        row = connection.execute(
            """
            SELECT failure_count, last_failed_at, locked_until
            FROM admin_login_attempts
            WHERE email = ? AND ip_address = ?
            """,
            (normalized_email, ip_address),
        ).fetchone()

    if not row:
        return {"failure_count": 0, "locked_until": None}

    locked_until_raw = (row["locked_until"] or "").strip()
    locked_until = datetime.fromisoformat(locked_until_raw) if locked_until_raw else None
    if locked_until and locked_until <= utc_now():
        clear_admin_login_attempts(normalized_email, ip_address)
        return {"failure_count": 0, "locked_until": None}

    return {
        "failure_count": int(row["failure_count"] or 0),
        "locked_until": locked_until,
    }


def clear_admin_login_attempts(email: str, ip_address: str):
    with get_db_connection() as connection:
        connection.execute(
            "DELETE FROM admin_login_attempts WHERE email = ? AND ip_address = ?",
            (normalize_email(email), ip_address),
        )


def ensure_admin_login_allowed(email: str, ip_address: str):
    if not is_admin_email(email):
        return

    state = get_admin_login_attempt_state(email, ip_address)
    locked_until = state.get("locked_until")
    if locked_until and locked_until > utc_now():
        readable_time = locked_until.astimezone(timezone.utc).strftime("%Y-%m-%d %H:%M UTC")
        raise HTTPException(
            status_code=429,
            detail=(
                "Admin login is temporarily locked after too many failed attempts. "
                f"Try again after {readable_time}."
            ),
        )


def record_admin_login_failure(email: str, ip_address: str) -> dict[str, Any]:
    current_state = get_admin_login_attempt_state(email, ip_address)
    failure_count = int(current_state.get("failure_count") or 0) + 1
    locked_until = ""
    if failure_count >= ADMIN_LOGIN_MAX_ATTEMPTS:
        locked_until = iso_in_future(minutes=ADMIN_LOGIN_LOCKOUT_MINUTES)

    with get_db_connection() as connection:
        connection.execute(
            """
            INSERT INTO admin_login_attempts (email, ip_address, failure_count, last_failed_at, locked_until)
            VALUES (?, ?, ?, ?, ?)
            ON CONFLICT(email, ip_address) DO UPDATE SET
                failure_count = excluded.failure_count,
                last_failed_at = excluded.last_failed_at,
                locked_until = excluded.locked_until
            """,
            (normalize_email(email), ip_address, failure_count, utc_now().isoformat(), locked_until),
        )

    return {
        "failure_count": failure_count,
        "locked_until": datetime.fromisoformat(locked_until) if locked_until else None,
    }


def verify_google_auth_is_configured():
    if not GOOGLE_CLIENT_ID:
        raise HTTPException(
            status_code=500,
            detail="Google login is not configured on the server yet. Missing GOOGLE_CLIENT_ID.",
        )
    if google_id_token is None or GoogleRequest is None:
        raise HTTPException(
            status_code=500,
            detail="Google login support is not installed on the server yet. Install backend requirements, including google-auth and requests, then redeploy.",
        )


def verify_apple_auth_is_configured():
    if not APPLE_CLIENT_ID:
        raise HTTPException(
            status_code=500,
            detail="Apple login is not configured on the server yet. Missing APPLE_CLIENT_ID.",
        )
    if jwt is None or apple_jwk_client is None:
        raise HTTPException(
            status_code=500,
            detail="Apple login support is not installed on the server yet. Install backend requirements, including PyJWT[crypto], then redeploy.",
        )


def normalize_apple_private_key() -> str:
    return APPLE_PRIVATE_KEY.replace("\\n", "\n").strip()


def can_exchange_apple_authorization_code() -> bool:
    return all([APPLE_TEAM_ID, APPLE_KEY_ID, normalize_apple_private_key()])


def mark_user_verified(email: str):
    now_iso = utc_now().isoformat()
    with get_db_connection() as connection:
        connection.execute(
            "INSERT OR IGNORE INTO users (email, created_at) VALUES (?, ?)",
            (email, now_iso),
        )
        connection.execute(
            "UPDATE users SET verified_at = COALESCE(verified_at, ?) WHERE email = ?",
            (now_iso, email),
        )


def normalize_email_password_auth_mode(mode: str) -> str:
    normalized = compact_text(mode, "login").lower()
    if normalized not in {"login", "register", "reset"}:
        raise HTTPException(status_code=400, detail="Authentication mode must be login, register, or reset.")
    return normalized


def validate_password_value(password: str) -> str:
    value = password or ""
    if len(value) < MIN_PASSWORD_LENGTH:
        raise HTTPException(
            status_code=400,
            detail=f"Use a password with at least {MIN_PASSWORD_LENGTH} characters.",
        )
    if len(value) > 256:
        raise HTTPException(status_code=400, detail="Use a shorter password.")
    return value


def hash_password_value(password: str, salt: str) -> str:
    derived = hashlib.pbkdf2_hmac(
        "sha256",
        password.encode("utf-8"),
        salt.encode("utf-8"),
        PASSWORD_HASH_ITERATIONS,
    )
    return base64.urlsafe_b64encode(derived).decode("ascii")


def get_password_credential(email: str) -> sqlite3.Row | None:
    with get_db_connection() as connection:
        return connection.execute(
            """
            SELECT email, password_hash, password_salt, created_at, updated_at
            FROM email_password_credentials
            WHERE email = ?
            """,
            (email,),
        ).fetchone()


def has_password_credential(email: str) -> bool:
    return get_password_credential(email) is not None


def store_password_credential(email: str, password: str):
    validated_password = validate_password_value(password)
    now_iso = utc_now().isoformat()
    salt = secrets.token_hex(16)
    password_hash = hash_password_value(validated_password, salt)
    with get_db_connection() as connection:
        connection.execute(
            "INSERT OR IGNORE INTO users (email, created_at) VALUES (?, ?)",
            (email, now_iso),
        )
        connection.execute(
            """
            INSERT INTO email_password_credentials (email, password_hash, password_salt, created_at, updated_at)
            VALUES (?, ?, ?, ?, ?)
            ON CONFLICT(email) DO UPDATE SET
                password_hash = excluded.password_hash,
                password_salt = excluded.password_salt,
                updated_at = excluded.updated_at
            """,
            (email, password_hash, salt, now_iso, now_iso),
        )


def verify_password_credential(email: str, password: str):
    validated_password = validate_password_value(password)
    credential = get_password_credential(email)
    if not credential:
        raise HTTPException(status_code=400, detail="Email or password is incorrect.")

    expected_hash = credential["password_hash"]
    actual_hash = hash_password_value(validated_password, credential["password_salt"])
    if not hmac.compare_digest(expected_hash, actual_hash):
        raise HTTPException(status_code=400, detail="Email or password is incorrect.")


def build_apple_client_secret() -> str:
    if not can_exchange_apple_authorization_code():
        raise HTTPException(
            status_code=500,
            detail="Apple code exchange is not configured on the server yet. Add APPLE_TEAM_ID, APPLE_KEY_ID, and APPLE_PRIVATE_KEY.",
        )

    now_timestamp = int(utc_now().timestamp())
    claims = {
        "iss": APPLE_TEAM_ID,
        "iat": now_timestamp,
        "exp": now_timestamp + 60 * 60 * 24 * 180,
        "aud": APPLE_IDENTITY_ISSUER,
        "sub": APPLE_CLIENT_ID,
    }
    return jwt.encode(
        claims,
        normalize_apple_private_key(),
        algorithm="ES256",
        headers={"kid": APPLE_KEY_ID},
    )


def exchange_apple_authorization_code(authorization_code: str, redirect_uri: str) -> dict[str, Any]:
    verify_apple_auth_is_configured()
    if not authorization_code:
        raise HTTPException(status_code=400, detail="Apple authorization code is required.")
    if not can_exchange_apple_authorization_code():
        return {}

    resolved_redirect_uri = (redirect_uri or APPLE_REDIRECT_URI).strip()
    payload = {
        "client_id": APPLE_CLIENT_ID,
        "client_secret": build_apple_client_secret(),
        "code": authorization_code,
        "grant_type": "authorization_code",
    }
    if resolved_redirect_uri:
        payload["redirect_uri"] = resolved_redirect_uri

    try:
        response = requests.post(
            APPLE_TOKEN_VALIDATION_URL,
            data=payload,
            timeout=APPLE_AUTH_HTTP_TIMEOUT,
        )
        data = response.json()
    except requests.RequestException as exc:
        raise HTTPException(status_code=502, detail="Apple sign-in could not be completed right now.") from exc
    except ValueError as exc:
        raise HTTPException(status_code=502, detail="Apple sign-in returned an unreadable response.") from exc

    if not response.ok:
        detail = compact_text(data.get("error_description") or data.get("error"), "Apple sign-in could not be verified.")
        raise HTTPException(status_code=401, detail=detail)

    return data


def verify_apple_identity_token(identity_token: str, nonce: str = "") -> dict[str, Any]:
    verify_apple_auth_is_configured()
    if not identity_token:
        raise HTTPException(status_code=400, detail="Apple identity token is required.")

    try:
        signing_key = apple_jwk_client.get_signing_key_from_jwt(identity_token)
        claims = jwt.decode(
            identity_token,
            signing_key.key,
            algorithms=["RS256"],
            audience=APPLE_CLIENT_ID,
            issuer=APPLE_IDENTITY_ISSUER,
        )
    except Exception as exc:
        raise HTTPException(status_code=401, detail="Apple sign-in could not be verified.") from exc

    provided_nonce = compact_text(nonce)
    token_nonce = compact_text(claims.get("nonce"))
    if provided_nonce and token_nonce and provided_nonce != token_nonce:
        raise HTTPException(status_code=401, detail="Apple sign-in nonce verification failed.")

    email = validate_email_address(claims.get("email", ""))
    email_verified = claims.get("email_verified")
    if str(email_verified).lower() not in {"true", "1"} and email_verified is not True:
        raise HTTPException(status_code=401, detail="Apple account email is not verified.")

    claims["email"] = email
    return claims


SMTP_ENV_ALIASES = {
    "host": ("SMTP_HOST", "MAIL_HOST", "EMAIL_HOST", "SMTP_SERVER"),
    "port": ("SMTP_PORT", "MAIL_PORT", "EMAIL_PORT"),
    "username": ("SMTP_USERNAME", "SMTP_USER", "MAIL_USERNAME", "MAIL_USER", "EMAIL_USERNAME"),
    "password": ("SMTP_PASSWORD", "SMTP_PASS", "MAIL_PASSWORD", "MAIL_PASS", "EMAIL_PASSWORD"),
    "from_email": ("SMTP_FROM_EMAIL", "FROM_EMAIL", "MAIL_FROM", "MAIL_FROM_EMAIL", "EMAIL_FROM"),
    "use_ssl": ("SMTP_USE_SSL", "MAIL_USE_SSL", "EMAIL_USE_SSL"),
    "use_tls": ("SMTP_USE_TLS", "MAIL_USE_TLS", "EMAIL_USE_TLS"),
    "timeout_seconds": ("SMTP_TIMEOUT_SECONDS", "MAIL_TIMEOUT_SECONDS", "EMAIL_TIMEOUT_SECONDS"),
}


def get_first_configured_env(*names: str, default: str = "") -> str:
    for name in names:
        value = os.getenv(name)
        if value is None:
            continue
        cleaned = str(value).strip()
        if cleaned:
            return cleaned
    return default


def format_env_aliases(*names: str) -> str:
    return " or ".join(names)


def verify_smtp_is_configured():
    required = {
        "SMTP host": SMTP_ENV_ALIASES["host"],
        "sender email": SMTP_ENV_ALIASES["from_email"],
    }
    missing = [
        f"{label} ({format_env_aliases(*names)})"
        for label, names in required.items()
        if not get_first_configured_env(*names)
    ]
    if missing:
        raise HTTPException(
            status_code=500,
            detail=(
                "Email login is not configured on the server yet. "
                f"Missing environment variables: {', '.join(missing)}."
            ),
        )


def parse_smtp_port(value: str) -> int:
    try:
        port = int((value or "").strip() or "587")
    except ValueError as exc:
        raise HTTPException(status_code=500, detail="SMTP_PORT must be a valid port number.") from exc
    if port <= 0 or port > 65535:
        raise HTTPException(status_code=500, detail="SMTP_PORT must be between 1 and 65535.")
    return port


def parse_smtp_host(value: str, default_port: int) -> tuple[str, int, bool | None]:
    raw_value = compact_text(value)
    parsed = urlparse(raw_value if "://" in raw_value else f"//{raw_value}")
    host = compact_text(parsed.hostname or parsed.netloc or raw_value)
    if not host:
        raise HTTPException(status_code=500, detail="SMTP_HOST is empty or invalid.")
    try:
        port = parsed.port or default_port
    except ValueError as exc:
        raise HTTPException(status_code=500, detail="SMTP_HOST contains an invalid port.") from exc

    scheme = compact_text(parsed.scheme).lower()
    use_ssl_hint = True if scheme == "smtps" else False if scheme == "smtp" else None
    return host, port, use_ssl_hint


def parse_smtp_timeout_seconds(value: str) -> float:
    try:
        timeout_seconds = float((value or "").strip() or "15")
    except ValueError as exc:
        raise HTTPException(status_code=500, detail="SMTP_TIMEOUT_SECONDS must be a number.") from exc
    return min(max(timeout_seconds, 5.0), 60.0)


def get_smtp_settings() -> dict[str, Any]:
    verify_smtp_is_configured()
    configured_port = parse_smtp_port(get_first_configured_env(*SMTP_ENV_ALIASES["port"], default="587"))
    host, port, use_ssl_hint = parse_smtp_host(get_first_configured_env(*SMTP_ENV_ALIASES["host"]), configured_port)
    from_email = compact_text(get_first_configured_env(*SMTP_ENV_ALIASES["from_email"]))
    password = compact_text(get_first_configured_env(*SMTP_ENV_ALIASES["password"]))
    username = compact_text(get_first_configured_env(*SMTP_ENV_ALIASES["username"]))
    if not username and password and from_email:
        username = from_email

    smtp_use_ssl_value = compact_text(get_first_configured_env(*SMTP_ENV_ALIASES["use_ssl"])).lower()
    if smtp_use_ssl_value:
        use_ssl = smtp_use_ssl_value in {"1", "true", "yes", "on"}
    elif use_ssl_hint is not None:
        use_ssl = use_ssl_hint
    else:
        use_ssl = port == 465

    smtp_use_tls_value = compact_text(get_first_configured_env(*SMTP_ENV_ALIASES["use_tls"])).lower()
    if use_ssl:
        use_tls = False
    elif smtp_use_tls_value:
        use_tls = smtp_use_tls_value in {"1", "true", "yes", "on"}
    else:
        use_tls = port in {25, 587, 2525}

    return {
        "host": host,
        "port": port,
        "username": username,
        "password": password,
        "from_email": from_email,
        "use_tls": use_tls,
        "use_ssl": use_ssl,
        "timeout_seconds": parse_smtp_timeout_seconds(get_first_configured_env(*SMTP_ENV_ALIASES["timeout_seconds"], default="15")),
    }


def build_smtp_connection_plans(smtp_settings: dict[str, Any]) -> list[dict[str, Any]]:
    shared = {
        "host": smtp_settings["host"],
        "username": smtp_settings["username"],
        "password": smtp_settings["password"],
        "from_email": smtp_settings["from_email"],
        "timeout_seconds": smtp_settings["timeout_seconds"],
    }
    candidates = [
        {**shared, "port": smtp_settings["port"], "use_ssl": smtp_settings["use_ssl"], "use_tls": smtp_settings["use_tls"]},
        {**shared, "port": 465, "use_ssl": True, "use_tls": False},
        {**shared, "port": 587, "use_ssl": False, "use_tls": True},
    ]
    plans: list[dict[str, Any]] = []
    seen: set[tuple[str, int, bool, bool]] = set()
    for candidate in candidates:
        plan_key = (
            candidate["host"],
            int(candidate["port"]),
            bool(candidate["use_ssl"]),
            bool(candidate["use_tls"]),
        )
        if plan_key in seen:
            continue
        seen.add(plan_key)
        plans.append(candidate)
    return plans


def send_smtp_message(message: EmailMessage):
    smtp_settings = get_smtp_settings()
    connection_plans = build_smtp_connection_plans(smtp_settings)
    last_exc: Exception | None = None

    for attempt_index, plan in enumerate(connection_plans, start=1):
        try:
            smtp_factory = smtplib.SMTP_SSL if plan["use_ssl"] else smtplib.SMTP
            with smtp_factory(plan["host"], plan["port"], timeout=plan["timeout_seconds"]) as server:
                try:
                    server.ehlo()
                except smtplib.SMTPException:
                    pass
                if plan["use_tls"]:
                    server.starttls(context=ssl.create_default_context())
                    try:
                        server.ehlo()
                    except smtplib.SMTPException:
                        pass
                if plan["username"]:
                    server.login(plan["username"], plan["password"])
                server.send_message(message)
            return
        except smtplib.SMTPAuthenticationError as exc:
            logger.exception("SMTP authentication failed for host %s", plan["host"])
            raise HTTPException(
                status_code=502,
                detail=(
                    "SMTP login failed. Recheck SMTP_USERNAME, SMTP_PASSWORD, "
                    "and your Gmail app password."
                ),
            ) from exc
        except (smtplib.SMTPConnectError, smtplib.SMTPServerDisconnected, smtplib.SMTPNotSupportedError, OSError) as exc:
            last_exc = exc
            logger.warning(
                "SMTP delivery attempt %s failed for %s:%s (ssl=%s tls=%s): %s",
                attempt_index,
                plan["host"],
                plan["port"],
                plan["use_ssl"],
                plan["use_tls"],
                exc,
            )
            continue
        except smtplib.SMTPException as exc:
            logger.exception("SMTP error while sending email")
            raise HTTPException(
                status_code=502,
                detail=f"SMTP error while sending email ({exc.__class__.__name__}).",
            ) from exc

    if isinstance(last_exc, smtplib.SMTPConnectError):
        logger.exception("SMTP connection failed for host %s", smtp_settings["host"])
        raise HTTPException(
            status_code=502,
            detail="The backend could not connect to the SMTP server. Recheck SMTP_HOST and SMTP_PORT.",
        ) from last_exc
    if isinstance(last_exc, smtplib.SMTPNotSupportedError):
        logger.exception("SMTP TLS mode was rejected for host %s", smtp_settings["host"])
        raise HTTPException(
            status_code=502,
            detail="The SMTP server rejected the current TLS mode. Recheck SMTP_USE_TLS and SMTP_USE_SSL.",
        ) from last_exc
    if isinstance(last_exc, OSError):
        logger.exception("SMTP network error while sending email")
        raise HTTPException(
            status_code=502,
            detail="The backend could not reach the SMTP server. Recheck SMTP_HOST, SMTP_PORT, SMTP_USE_TLS, and SMTP_USE_SSL.",
        ) from last_exc
    if last_exc is not None:
        logger.exception("SMTP error while sending email")
        raise HTTPException(
            status_code=502,
            detail=f"SMTP error while sending email ({last_exc.__class__.__name__}).",
        ) from last_exc

    raise HTTPException(status_code=502, detail="SMTP delivery failed for an unknown reason.")


def send_verification_email(email: str, code: str):
    smtp_settings = get_smtp_settings()

    message = EmailMessage()
    message["Subject"] = "Your MABASO.AI verification code"
    message["From"] = smtp_settings["from_email"]
    message["To"] = email
    message.set_content(
        (
            "Your MABASO.AI verification code is:\n\n"
            f"{code}\n\n"
            f"This code expires in {LOGIN_CODE_TTL_MINUTES} minutes."
        )
    )
    send_smtp_message(message)


def send_support_email(reply_email: str, message_text: str, page: str = ""):
    cleaned_message = compact_text(message_text)
    if not cleaned_message:
        raise HTTPException(status_code=400, detail="Support message cannot be empty.")

    smtp_settings = get_smtp_settings()
    support_email = SUPPORT_EMAIL or smtp_settings["from_email"]
    message = EmailMessage()
    message["Subject"] = f"MABASO support request from {reply_email}"
    message["From"] = smtp_settings["from_email"]
    message["To"] = support_email
    message["Reply-To"] = reply_email
    message.set_content(
        (
            "A new support request was sent from MABASO.\n\n"
            f"From: {reply_email}\n"
            f"Page: {page or 'unknown'}\n"
            f"Sent at (UTC): {utc_now().isoformat()}\n\n"
            "Message:\n"
            f"{message_text.strip()}\n"
        )
    )
    send_smtp_message(message)


def normalize_client_request_id(value: Any) -> str:
    cleaned = re.sub(r"[^A-Za-z0-9._:-]+", "-", compact_text(value)).strip("-._:")
    return cleaned[:120] or uuid4().hex


def serialize_support_message_row(row: sqlite3.Row | None) -> dict[str, Any]:
    if row is None:
        return {}
    return {
        "id": row["id"],
        "client_request_id": row["client_request_id"],
        "email": row["email"],
        "page": row["page"],
        "message": row["message"],
        "email_delivery_status": row["email_delivery_status"],
        "email_error": row["email_error"],
        "created_at": row["created_at"],
        "updated_at": row["updated_at"],
    }


def create_or_load_support_message(
    reply_email: str,
    message_text: str,
    page: str = "",
    client_request_id: str = "",
) -> dict[str, Any]:
    normalized_email = validate_email_address(reply_email)
    cleaned_message = compact_text(message_text)
    if not cleaned_message:
        raise HTTPException(status_code=400, detail="Support message cannot be empty.")

    normalized_request_id = normalize_client_request_id(client_request_id)
    cleaned_page = compact_text(page, "unknown-page")[:120]

    with get_db_connection() as connection:
        existing_row = connection.execute(
            "SELECT * FROM support_messages WHERE client_request_id = ?",
            (normalized_request_id,),
        ).fetchone()
        if existing_row:
            return serialize_support_message_row(existing_row)

        created_at = utc_now().isoformat()
        message_id = uuid4().hex
        connection.execute(
            """
            INSERT INTO support_messages (
                id, client_request_id, email, page, message, email_delivery_status,
                email_error, created_at, updated_at
            )
            VALUES (?, ?, ?, ?, ?, ?, ?, ?, ?)
            """,
            (
                message_id,
                normalized_request_id,
                normalized_email,
                cleaned_page,
                cleaned_message,
                "queued",
                "",
                created_at,
                created_at,
            ),
        )
        created_row = connection.execute(
            "SELECT * FROM support_messages WHERE id = ?",
            (message_id,),
        ).fetchone()
    return serialize_support_message_row(created_row)


def update_support_message_delivery(message_id: str, status: str, email_error: str = ""):
    with get_db_connection() as connection:
        connection.execute(
            """
            UPDATE support_messages
            SET email_delivery_status = ?, email_error = ?, updated_at = ?
            WHERE id = ?
            """,
            (compact_text(status, "queued")[:40], compact_text(email_error)[:600], utc_now().isoformat(), message_id),
        )


async def deliver_support_message_email(
    message_id: str,
    reply_email: str,
    message_text: str,
    page: str,
) -> dict[str, str]:
    try:
        await asyncio.to_thread(send_support_email, reply_email, message_text, page)
    except HTTPException as exc:
        logger.warning("Support email delivery failed for %s: %s", message_id, exc.detail)
        await asyncio.to_thread(update_support_message_delivery, message_id, "email_failed", str(exc.detail))
        return {"status": "email_failed", "error": str(exc.detail)}
    except Exception as exc:
        logger.exception("Unexpected support email delivery failure for %s", message_id)
        await asyncio.to_thread(update_support_message_delivery, message_id, "email_failed", str(exc))
        return {"status": "email_failed", "error": str(exc)}
    else:
        await asyncio.to_thread(update_support_message_delivery, message_id, "sent", "")
        return {"status": "sent", "error": ""}


def create_login_code(email: str) -> str:
    code = f"{secrets.randbelow(1_000_000):06d}"
    now_iso = utc_now().isoformat()
    expiry_iso = iso_in_future(minutes=LOGIN_CODE_TTL_MINUTES)

    with get_db_connection() as connection:
        connection.execute(
            "INSERT OR IGNORE INTO users (email, created_at) VALUES (?, ?)",
            (email, now_iso),
        )
        connection.execute(
            """
            INSERT INTO login_codes (email, code_hash, expires_at, created_at)
            VALUES (?, ?, ?, ?)
            ON CONFLICT(email) DO UPDATE SET
                code_hash = excluded.code_hash,
                expires_at = excluded.expires_at,
                created_at = excluded.created_at
            """,
            (email, hash_value(code), expiry_iso, now_iso),
        )

    return code


def consume_login_code(email: str, code: str):
    with get_db_connection() as connection:
        row = connection.execute(
            "SELECT code_hash, expires_at FROM login_codes WHERE email = ?",
            (email,),
        ).fetchone()

    if not row:
        raise HTTPException(status_code=400, detail="No verification code was requested for this email yet.")

    if datetime.fromisoformat(row["expires_at"]) <= utc_now():
        raise HTTPException(status_code=400, detail="That verification code has expired. Request a new code.")

    if not hmac.compare_digest(row["code_hash"], hash_value(code)):
        raise HTTPException(status_code=400, detail="The verification code is incorrect.")

    with get_db_connection() as connection:
        connection.execute("DELETE FROM login_codes WHERE email = ?", (email,))


def create_session(email: str, session_mode: str = "user") -> str:
    ensure_user_account_is_active(email)
    raw_token = build_signed_session_token(email, session_mode=session_mode)
    token_payload = decode_signed_session_token(raw_token) or {}
    expires_at = datetime.fromtimestamp(
        int(token_payload.get("exp", int((utc_now() + timedelta(minutes=SESSION_TTL_MINUTES)).timestamp()))),
        tz=timezone.utc,
    ).isoformat()
    with get_db_connection() as connection:
        existing_sessions = connection.execute(
            "SELECT token_hash, expires_at FROM sessions WHERE email = ?",
            (email,),
        ).fetchall()
        for session_row in existing_sessions:
            connection.execute(
                """
                INSERT OR REPLACE INTO revoked_sessions (token_hash, expires_at, created_at)
                VALUES (?, ?, ?)
                """,
                (session_row["token_hash"], session_row["expires_at"], utc_now().isoformat()),
            )
        connection.execute(
            "DELETE FROM sessions WHERE email = ?",
            (email,),
        )
        connection.execute(
            """
            INSERT INTO sessions (token_hash, email, expires_at, created_at)
            VALUES (?, ?, ?, ?)
            """,
            (hash_value(raw_token), email, expires_at, utc_now().isoformat()),
        )
        connection.execute(
            "UPDATE users SET verified_at = COALESCE(verified_at, ?) WHERE email = ?",
            (utc_now().isoformat(), email),
        )
    return raw_token


def revoke_session(token: str):
    with get_db_connection() as connection:
        connection.execute("DELETE FROM sessions WHERE token_hash = ?", (hash_value(token),))
        payload = decode_signed_session_token(token)
        if payload:
            expires_at = datetime.fromtimestamp(int(payload["exp"]), tz=timezone.utc).isoformat()
            connection.execute(
                """
                INSERT OR REPLACE INTO revoked_sessions (token_hash, expires_at, created_at)
                VALUES (?, ?, ?)
                """,
                (hash_value(token), expires_at, utc_now().isoformat()),
            )


def refresh_session_expiry(token_hash: str):
    with get_db_connection() as connection:
        connection.execute(
            "UPDATE sessions SET expires_at = ? WHERE token_hash = ?",
            (iso_in_future(minutes=SESSION_TTL_MINUTES), token_hash),
        )


def should_refresh_session_token(token: str) -> bool:
    if not is_signed_session_token(token):
        return True

    payload = decode_signed_session_token(token)
    if not payload:
        return False

    expires_at = datetime.fromtimestamp(int(payload.get("exp", 0) or 0), tz=timezone.utc)
    return expires_at <= utc_now() + timedelta(minutes=SESSION_REFRESH_WINDOW_MINUTES)


def parse_history_datetime(value: str | None, fallback: datetime | None = None) -> datetime:
    text = (value or "").strip()
    if not text:
        return fallback or utc_now()
    try:
        parsed = datetime.fromisoformat(text.replace("Z", "+00:00"))
    except ValueError:
        return fallback or utc_now()
    if parsed.tzinfo is None:
        return parsed.replace(tzinfo=timezone.utc)
    return parsed.astimezone(timezone.utc)


def normalize_history_item_payload(raw_item: Any) -> dict[str, Any]:
    if not isinstance(raw_item, dict):
        raise HTTPException(status_code=400, detail="Each history item must be an object.")

    try:
        item = json.loads(json.dumps(raw_item))
    except (TypeError, ValueError) as exc:
        raise HTTPException(status_code=400, detail="History items must be JSON serializable.") from exc

    now = utc_now()
    item_id = compact_text(item.get("id"))
    if not item_id:
        item_id = uuid4().hex

    created_at = parse_history_datetime(compact_text(item.get("createdAt")), now)
    updated_at = parse_history_datetime(compact_text(item.get("updatedAt")), created_at)

    item["id"] = item_id
    item["createdAt"] = created_at.isoformat()
    item["updatedAt"] = updated_at.isoformat()
    return item


def sort_history_items(items: list[dict[str, Any]]) -> list[dict[str, Any]]:
    return sorted(
        items,
        key=lambda item: (
            parse_history_datetime(compact_text(item.get("updatedAt")), utc_now()),
            parse_history_datetime(compact_text(item.get("createdAt")), utc_now()),
        ),
        reverse=True,
    )


def get_history_items_for_user(email: str) -> list[dict[str, Any]]:
    with get_db_connection() as connection:
        rows = connection.execute(
            """
            SELECT payload_json
            FROM study_history_items
            WHERE email = ?
            ORDER BY updated_at DESC
            LIMIT ?
            """,
            (email, MAX_HISTORY_ITEMS),
        ).fetchall()

    items: list[dict[str, Any]] = []
    for row in rows:
        try:
            items.append(normalize_history_item_payload(json.loads(row["payload_json"])))
        except (json.JSONDecodeError, HTTPException):
            continue
    return sort_history_items(items)[:MAX_HISTORY_ITEMS]


def replace_history_items_for_user(email: str, items: list[dict[str, Any]]) -> list[dict[str, Any]]:
    normalized_items = sort_history_items([normalize_history_item_payload(item) for item in items])[:MAX_HISTORY_ITEMS]

    with get_db_connection() as connection:
        connection.execute("DELETE FROM study_history_items WHERE email = ?", (email,))
        for item in normalized_items:
            created_at = compact_text(item.get("createdAt"), utc_now().isoformat())
            updated_at = compact_text(item.get("updatedAt"), created_at)
            connection.execute(
                """
                INSERT INTO study_history_items (email, id, payload_json, created_at, updated_at)
                VALUES (?, ?, ?, ?, ?)
                """,
                (email, item["id"], json.dumps(item), created_at, updated_at),
            )

    return normalized_items


def is_session_revoked(token_hash: str) -> bool:
    with get_db_connection() as connection:
        connection.execute(
            "DELETE FROM revoked_sessions WHERE expires_at <= ?",
            (utc_now().isoformat(),),
        )
        row = connection.execute(
            "SELECT 1 FROM revoked_sessions WHERE token_hash = ?",
            (token_hash,),
        ).fetchone()
    return bool(row)


def get_legacy_session_email(token: str) -> str | None:
    token_hash = hash_value(token)
    with get_db_connection() as connection:
        row = connection.execute(
            "SELECT email, expires_at FROM sessions WHERE token_hash = ?",
            (token_hash,),
        ).fetchone()

    if not row:
        return None

    if datetime.fromisoformat(row["expires_at"]) <= utc_now():
        with get_db_connection() as connection:
            connection.execute("DELETE FROM sessions WHERE token_hash = ?", (token_hash,))
        return None

    refresh_session_expiry(token_hash)
    return normalize_email(row["email"])


def get_session_context(token: str) -> dict[str, Any] | None:
    if not token:
        return None

    token_hash = hash_value(token)
    if is_signed_session_token(token):
        payload = decode_signed_session_token(token)
        if not payload or is_session_revoked(token_hash):
            return None

        expires_at = datetime.fromtimestamp(int(payload["exp"]), tz=timezone.utc)
        if expires_at <= utc_now():
            return None

        email = normalize_email(payload["email"])
        ensure_user_account_is_active(email)
        with get_db_connection() as connection:
            connection.execute(
                "DELETE FROM sessions WHERE expires_at <= ?",
                (utc_now().isoformat(),),
            )
            matching_row = connection.execute(
                "SELECT email, expires_at FROM sessions WHERE token_hash = ?",
                (token_hash,),
            ).fetchone()
            if matching_row:
                if normalize_email(matching_row["email"]) != email:
                    return None
                return {
                    "email": email,
                    "mode": normalize_session_mode(str(payload.get("mode") or "user"), email),
                    "available_modes": get_available_auth_modes(email),
                }

            superseded_row = connection.execute(
                "SELECT token_hash FROM sessions WHERE email = ? AND expires_at > ? LIMIT 1",
                (email, utc_now().isoformat()),
            ).fetchone()

        if superseded_row:
            return None
        return {
            "email": email,
            "mode": normalize_session_mode(str(payload.get("mode") or "user"), email),
            "available_modes": get_available_auth_modes(email),
        }

    legacy_email = get_legacy_session_email(token)
    if not legacy_email:
        return None
    ensure_user_account_is_active(legacy_email)
    return {
        "email": legacy_email,
        "mode": "user",
        "available_modes": get_available_auth_modes(legacy_email),
    }


def get_session_email(token: str) -> str | None:
    context = get_session_context(token)
    return context["email"] if context else None


def get_authorization_token(authorization: str | None) -> str:
    if not authorization or not authorization.lower().startswith("bearer "):
        raise HTTPException(status_code=401, detail="Authentication is required.")
    return authorization.split(" ", 1)[1].strip()


def require_authenticated_user(authorization: str | None = Header(None)) -> str:
    token = get_authorization_token(authorization)
    context = get_session_context(token)
    if not context:
        raise HTTPException(status_code=401, detail="Your session is invalid or has expired.")
    return context["email"]


def require_admin_user(authorization: str | None = Header(None)) -> str:
    token = get_authorization_token(authorization)
    context = get_session_context(token)
    if not context:
        raise HTTPException(status_code=401, detail="Your session is invalid or has expired.")
    if context["mode"] != "admin" or not is_admin_email(context["email"]):
        raise HTTPException(status_code=403, detail="Admin access is required.")
    return context["email"]


def verify_login_code(email: str, code: str) -> str:
    consume_login_code(email, code)
    return create_session(email)


def login_with_email_password(email: str, password: str) -> str:
    verify_password_credential(email, password)
    return create_session(email)


def register_with_email_password(email: str, password: str) -> str:
    ensure_user_account_is_active(email)
    validated_password = validate_password_value(password)
    if has_password_credential(email):
        raise HTTPException(status_code=400, detail="An account with this email already exists. Sign in instead.")
    store_password_credential(email, validated_password)
    mark_user_verified(email)
    return create_session(email)


def request_email_password_registration_code(email: str) -> str:
    ensure_user_account_is_active(email)
    if has_password_credential(email):
        raise HTTPException(status_code=400, detail="An account with this email already exists. Sign in instead.")
    with get_db_connection() as connection:
        connection.execute("DELETE FROM pending_email_password_registrations WHERE email = ?", (email,))
    return create_login_code(email)


def create_pending_registration_token(email: str) -> str:
    token = secrets.token_urlsafe(32)
    now_iso = utc_now().isoformat()
    expiry_iso = iso_in_future(minutes=REGISTRATION_TOKEN_TTL_MINUTES)

    with get_db_connection() as connection:
        connection.execute(
            """
            INSERT INTO pending_email_password_registrations (email, token_hash, expires_at, created_at)
            VALUES (?, ?, ?, ?)
            ON CONFLICT(email) DO UPDATE SET
                token_hash = excluded.token_hash,
                expires_at = excluded.expires_at,
                created_at = excluded.created_at
            """,
            (email, hash_value(token), expiry_iso, now_iso),
        )

    return token


def consume_pending_registration_token(email: str, token: str):
    with get_db_connection() as connection:
        row = connection.execute(
            "SELECT token_hash, expires_at FROM pending_email_password_registrations WHERE email = ?",
            (email,),
        ).fetchone()

    if not row:
        raise HTTPException(status_code=400, detail="Verify your email first before creating a password.")

    if datetime.fromisoformat(row["expires_at"]) <= utc_now():
        raise HTTPException(status_code=400, detail="Your password setup link expired. Verify your email again.")

    if not hmac.compare_digest(row["token_hash"], hash_value(token)):
        raise HTTPException(status_code=400, detail="Your password setup step is no longer valid. Verify your email again.")

    with get_db_connection() as connection:
        connection.execute("DELETE FROM pending_email_password_registrations WHERE email = ?", (email,))


def verify_email_password_registration_code(email: str, code: str) -> str:
    ensure_user_account_is_active(email)
    if has_password_credential(email):
        raise HTTPException(status_code=400, detail="An account with this email already exists. Sign in instead.")
    consume_login_code(email, code)
    if has_password_credential(email):
        raise HTTPException(status_code=400, detail="An account with this email already exists. Sign in instead.")
    return create_pending_registration_token(email)


def complete_email_password_registration(email: str, registration_token: str, password: str) -> str:
    ensure_user_account_is_active(email)
    validated_password = validate_password_value(password)
    if has_password_credential(email):
        raise HTTPException(status_code=400, detail="An account with this email already exists. Sign in instead.")
    consume_pending_registration_token(email, registration_token.strip())
    if has_password_credential(email):
        raise HTTPException(status_code=400, detail="An account with this email already exists. Sign in instead.")
    store_password_credential(email, validated_password)
    mark_user_verified(email)
    return create_session(email)


def request_email_password_code(email: str, password: str, mode: str) -> str:
    ensure_user_account_is_active(email)
    resolved_mode = normalize_email_password_auth_mode(mode)
    validated_password = validate_password_value(password)
    if resolved_mode == "register":
        if has_password_credential(email):
            raise HTTPException(status_code=400, detail="An account with this email already exists. Sign in instead.")
    elif resolved_mode == "login":
        verify_password_credential(email, validated_password)
    else:
        if not has_password_credential(email):
            raise HTTPException(status_code=400, detail="No account with this email exists yet.")
    return create_login_code(email)


def verify_email_password_auth(email: str, password: str, code: str, mode: str) -> str:
    ensure_user_account_is_active(email)
    resolved_mode = normalize_email_password_auth_mode(mode)
    validated_password = validate_password_value(password)
    consume_login_code(email, code)
    if resolved_mode == "register":
        if has_password_credential(email):
            raise HTTPException(status_code=400, detail="An account with this email already exists. Sign in instead.")
        store_password_credential(email, validated_password)
        mark_user_verified(email)
    elif resolved_mode == "login":
        verify_password_credential(email, validated_password)
    else:
        if not has_password_credential(email):
            raise HTTPException(status_code=400, detail="No account with this email exists yet.")
        store_password_credential(email, validated_password)
        mark_user_verified(email)
    return create_session(email)


def create_session_from_google_credential(credential: str) -> tuple[str, str]:
    verify_google_auth_is_configured()

    try:
        token_info = google_id_token.verify_oauth2_token(
            credential,
            GoogleRequest(),
            GOOGLE_CLIENT_ID,
        )
    except Exception as exc:
        raise HTTPException(status_code=401, detail="Google sign-in could not be verified.") from exc

    email = validate_email_address(token_info.get("email", ""))
    if not token_info.get("email_verified"):
        raise HTTPException(status_code=401, detail="Google account email is not verified.")

    ensure_user_account_is_active(email)
    mark_user_verified(email)
    return create_session(email), email


def create_session_from_apple_auth(payload: AppleAuthRequest) -> tuple[str, str]:
    verify_apple_auth_is_configured()

    authorization_code = compact_text(payload.authorization_code)
    redirect_uri = compact_text(payload.redirect_uri)
    identity_token = compact_text(payload.id_token)

    if authorization_code and can_exchange_apple_authorization_code():
        exchanged = exchange_apple_authorization_code(authorization_code, redirect_uri)
        identity_token = compact_text(exchanged.get("id_token"), identity_token)

    claims = verify_apple_identity_token(identity_token, nonce=payload.nonce)
    email = claims["email"]
    ensure_user_account_is_active(email)
    mark_user_verified(email)
    return create_session(email), email


def build_auth_response(email: str, token: str) -> dict[str, Any]:
    available_modes = get_available_auth_modes(email)
    payload = decode_signed_session_token(token) or {}
    session_mode = normalize_session_mode(str(payload.get("mode") or "user"), email)
    return {
        "token": token,
        "email": email,
        "session_mode": session_mode,
        "available_modes": available_modes,
        "is_admin": "admin" in available_modes,
    }


def build_pdf_document(title: str, sections: list[PdfSection]) -> bytes:
    if A4 is None:
        raise HTTPException(
            status_code=500,
            detail="PDF export is not configured on the server yet. Install reportlab and redeploy.",
        )

    pdf_symbol_replacements = {
        "≥": ">=",
        "≤": "<=",
        "≠": "!=",
        "≈": "approx",
        "→": "->",
        "←": "<-",
        "↔": "<->",
        "−": "-",
        "–": "-",
        "—": "-",
        "∞": "infinity",
        "∫": "Integral",
        "∑": "Sum",
        "∏": "Product",
        "√": "sqrt",
        "Δ": "Delta",
        "δ": "delta",
        "τ": "tau",
        "θ": "theta",
        "λ": "lambda",
        "μ": "mu",
        "σ": "sigma",
        "π": "pi",
        "ω": "omega",
        "Ω": "Omega",
        "α": "alpha",
        "β": "beta",
        "γ": "gamma",
        "ε": "epsilon",
        "η": "eta",
        "ρ": "rho",
        "φ": "phi",
        "∂": "partial",
        "∈": "in",
        "∉": "not in",
        "∪": "union",
        "∩": "intersection",
        "·": " * ",
        "×": " x ",
        "÷": " / ",
        "°": " degrees",
        "â‰¥": ">=",
        "â‰¤": "<=",
        "â‰ ": "!=",
        "â†’": "->",
    }
    superscript_translation = str.maketrans({
        "⁰": "0",
        "¹": "1",
        "²": "2",
        "³": "3",
        "⁴": "4",
        "⁵": "5",
        "⁶": "6",
        "⁷": "7",
        "⁸": "8",
        "⁹": "9",
        "⁺": "+",
        "⁻": "-",
        "⁼": "=",
        "⁽": "(",
        "⁾": ")",
        "ᵃ": "a",
        "ᵇ": "b",
        "ᶜ": "c",
        "ᵈ": "d",
        "ᵉ": "e",
        "ᶠ": "f",
        "ᵍ": "g",
        "ʰ": "h",
        "ᶦ": "i",
        "ʲ": "j",
        "ᵏ": "k",
        "ˡ": "l",
        "ᵐ": "m",
        "ⁿ": "n",
        "ᵒ": "o",
        "ᵖ": "p",
        "ʳ": "r",
        "ˢ": "s",
        "ᵗ": "t",
        "ᵘ": "u",
        "ᵛ": "v",
        "ʷ": "w",
        "ˣ": "x",
        "ʸ": "y",
        "ᶻ": "z",
    })
    subscript_translation = str.maketrans({
        "₀": "0",
        "₁": "1",
        "₂": "2",
        "₃": "3",
        "₄": "4",
        "₅": "5",
        "₆": "6",
        "₇": "7",
        "₈": "8",
        "₉": "9",
        "₊": "+",
        "₋": "-",
        "₌": "=",
        "₍": "(",
        "₎": ")",
        "ₐ": "a",
        "ₑ": "e",
        "ₕ": "h",
        "ᵢ": "i",
        "ⱼ": "j",
        "ₖ": "k",
        "ₗ": "l",
        "ₘ": "m",
        "ₙ": "n",
        "ₒ": "o",
        "ₚ": "p",
        "ᵣ": "r",
        "ₛ": "s",
        "ₜ": "t",
        "ᵤ": "u",
        "ᵥ": "v",
        "ₓ": "x",
    })

    def normalize_pdf_export_text(value: str) -> str:
        cleaned = compact_text(value)
        if not cleaned:
            return ""
        cleaned = cleaned.translate(superscript_translation)
        cleaned = cleaned.translate(subscript_translation)
        for old, new in pdf_symbol_replacements.items():
            cleaned = cleaned.replace(old, new)
        cleaned = re.sub(r"^```(?:\w+)?\s*", "", cleaned, flags=re.MULTILINE)
        cleaned = cleaned.replace("```", "")
        cleaned = cleaned.replace("\t", "    ").replace("\xa0", " ")
        cleaned = unicodedata.normalize("NFKD", cleaned).encode("ascii", "ignore").decode("ascii")
        cleaned = re.sub(r"[ ]+\n", "\n", cleaned)
        cleaned = re.sub(r"\n{3,}", "\n\n", cleaned)
        return cleaned.strip()

    def build_pdf_markup(value: str) -> str:
        safe = html.escape(normalize_pdf_export_text(value), quote=False)
        safe = re.sub(r"\*\*(.+?)\*\*", r"<b>\1</b>", safe)
        safe = re.sub(r"__(.+?)__", r"<b>\1</b>", safe)
        safe = safe.replace("`", "")
        return safe.replace("\n", "<br/>")

    def is_markdown_table_separator(line: str) -> bool:
        stripped = line.strip()
        return bool(stripped) and set(stripped) <= {"|", "-", ":", " "}

    def build_pdf_table(table_lines: list[str]):
        raw_rows = [
            [normalize_pdf_export_text(cell) for cell in line.strip().strip("|").split("|")]
            for line in table_lines
            if line.strip().startswith("|")
        ]
        rows = [row for row in raw_rows if row and any(cell.strip("-: ") for cell in row)]
        if not rows:
            return None
        column_count = max(len(row) for row in rows)
        normalized_rows = [
            row + [""] * (column_count - len(row))
            for row in rows
        ]
        paragraph_rows = [
            [Paragraph(build_pdf_markup(cell), body_style) for cell in row]
            for row in normalized_rows
        ]
        table = Table(paragraph_rows, hAlign="LEFT")
        table.setStyle(
            TableStyle(
                [
                    ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#eff6ff")),
                    ("TEXTCOLOR", (0, 0), (-1, 0), colors.HexColor("#0f172a")),
                    ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
                    ("GRID", (0, 0), (-1, -1), 0.6, colors.HexColor("#dbeafe")),
                    ("BOX", (0, 0), (-1, -1), 0.8, colors.HexColor("#cbd5e1")),
                    ("VALIGN", (0, 0), (-1, -1), "TOP"),
                    ("LEFTPADDING", (0, 0), (-1, -1), 8),
                    ("RIGHTPADDING", (0, 0), (-1, -1), 8),
                    ("TOPPADDING", (0, 0), (-1, -1), 6),
                    ("BOTTOMPADDING", (0, 0), (-1, -1), 6),
                    ("BACKGROUND", (0, 1), (-1, -1), colors.white),
                ]
            )
        )
        return table

    styles = getSampleStyleSheet()
    title_style = styles["Heading1"]
    title_style.textColor = colors.HexColor("#0f172a")
    title_style.spaceAfter = 12
    title_style.fontName = "Helvetica-Bold"
    heading_style = styles["Heading2"]
    heading_style.textColor = colors.HexColor("#0f172a")
    heading_style.spaceBefore = 12
    heading_style.spaceAfter = 8
    heading_style.fontName = "Helvetica-Bold"
    heading_style.fontSize = 16
    subheading_style = ParagraphStyle(
        "MabasoSubheading",
        parent=styles["Heading3"],
        fontName="Helvetica-Bold",
        fontSize=13,
        leading=16,
        textColor=colors.HexColor("#2563eb"),
        spaceBefore=10,
        spaceAfter=6,
    )
    minor_heading_style = ParagraphStyle(
        "MabasoMinorHeading",
        parent=styles["Heading4"],
        fontName="Helvetica-Bold",
        fontSize=11,
        leading=14,
        textColor=colors.HexColor("#16a34a"),
        spaceBefore=8,
        spaceAfter=4,
    )
    body_style = ParagraphStyle(
        "MabasoBody",
        parent=styles["BodyText"],
        fontName="Helvetica",
        fontSize=10.5,
        leading=15.5,
        textColor=colors.HexColor("#1f2937"),
        alignment=TA_LEFT,
        spaceAfter=8,
    )
    bullet_style = ParagraphStyle(
        "MabasoBullet",
        parent=body_style,
        leftIndent=14,
        firstLineIndent=-10,
        spaceAfter=5,
    )

    buffer = BytesIO()
    document = SimpleDocTemplate(
        buffer,
        pagesize=A4,
        leftMargin=36,
        rightMargin=36,
        topMargin=36,
        bottomMargin=36,
    )

    story: list = [Paragraph(title or "MABASO Study Pack", title_style), Spacer(1, 8)]

    def flush_paragraph_lines(paragraph_lines: list[str]) -> None:
        text = "\n".join(paragraph_lines).strip()
        if not text:
            return
        story.append(Paragraph(build_pdf_markup(text), body_style))
        story.append(Spacer(1, 6))

    def append_structured_pdf_content(value: str) -> None:
        cleaned = (value or "").replace("\r\n", "\n").strip()
        if not cleaned:
            return
        lines = cleaned.splitlines()
        paragraph_lines: list[str] = []
        index = 0

        while index < len(lines):
            line = lines[index].rstrip()
            stripped = line.strip()
            next_line = lines[index + 1].rstrip() if index + 1 < len(lines) else ""

            if not stripped:
                flush_paragraph_lines(paragraph_lines)
                paragraph_lines = []
                index += 1
                continue

            if stripped.startswith("|") and next_line.strip().startswith("|") and is_markdown_table_separator(next_line):
                flush_paragraph_lines(paragraph_lines)
                paragraph_lines = []
                table_lines = [line, next_line]
                index += 2
                while index < len(lines) and lines[index].strip().startswith("|"):
                    table_lines.append(lines[index].rstrip())
                    index += 1
                table = build_pdf_table(table_lines)
                if table is not None:
                    story.append(table)
                    story.append(Spacer(1, 10))
                continue

            markdown_heading_match = re.match(r"^(#{1,6})\s+(.+?)\s*$", stripped)
            if markdown_heading_match:
                flush_paragraph_lines(paragraph_lines)
                paragraph_lines = []
                level = len(markdown_heading_match.group(1))
                heading_text = markdown_heading_match.group(2).strip().strip("*").strip()
                style = heading_style if level <= 2 else subheading_style if level == 3 else minor_heading_style
                story.append(Paragraph(build_pdf_markup(heading_text), style))
                story.append(Spacer(1, 4))
                index += 1
                continue

            bold_heading_match = re.match(r"^\*\*(.+?)\*\*\s*:?\s*$", stripped)
            if bold_heading_match:
                flush_paragraph_lines(paragraph_lines)
                paragraph_lines = []
                story.append(Paragraph(build_pdf_markup(bold_heading_match.group(1).strip()), minor_heading_style))
                story.append(Spacer(1, 4))
                index += 1
                continue

            if re.match(r"^[-*_]{3,}\s*$", stripped):
                flush_paragraph_lines(paragraph_lines)
                paragraph_lines = []
                story.append(Spacer(1, 6))
                index += 1
                continue

            bullet_match = re.match(r"^([-*])\s+(.*)$", stripped)
            numbered_match = re.match(r"^(\d+\.)\s+(.*)$", stripped)
            if bullet_match or numbered_match:
                flush_paragraph_lines(paragraph_lines)
                paragraph_lines = []
                prefix = bullet_match.group(1) if bullet_match else numbered_match.group(1)
                text = bullet_match.group(2) if bullet_match else numbered_match.group(2)
                story.append(Paragraph(build_pdf_markup(f"{prefix} {text}"), bullet_style))
                story.append(Spacer(1, 2))
                index += 1
                continue

            paragraph_lines.append(normalize_pdf_export_text(stripped))
            index += 1

        flush_paragraph_lines(paragraph_lines)

    for section in sections:
        if not section.content.strip():
            continue
        story.append(Paragraph(section.title, heading_style))
        append_structured_pdf_content(section.content)
        story.append(Spacer(1, 6))

    document.build(story)
    return buffer.getvalue()


init_db()


def create_job(job_type: str, owner_email: str = "") -> str:
    job_id = uuid4().hex
    jobs[job_id] = {
        "job_id": job_id,
        "job_type": job_type,
        "owner_email": owner_email,
        "created_at": utc_now().isoformat(),
        "status": "queued",
        "stage": "Waiting",
        "progress": 0,
        "error": None,
        "transcript": "",
        "summary": "",
        "formula": "",
        "worked_example": "",
        "flashcards": [],
        "quiz_questions": [],
        "used_fallback": False,
        "podcast_title": "",
        "podcast_overview": "",
        "podcast_script": "",
        "podcast_segments": [],
        "teacher_title": "",
        "teacher_overview": "",
        "teacher_segments": [],
        "presentation_title": "",
        "presentation_subtitle": "",
        "presentation_design_id": "",
        "presentation_template_name": "",
        "presentation_slides": [],
        "study_images": [],
        "_podcast_audio_files": [],
        "_podcast_download_file": "",
        "_presentation_download_file": "",
        "_output_language": "English",
    }
    return job_id


def update_job(job_id: str, **fields):
    job = jobs.get(job_id)
    if not job:
        return
    job.update(fields)


def serialize_job(job: dict[str, Any]) -> dict[str, Any]:
    return {
        key: value
        for key, value in job.items()
        if not str(key).startswith("_")
    }


def build_chat_messages(payload: StudyChatRequest) -> list[dict[str, object]]:
    section_limit = max(2000, MAX_CHAT_CONTEXT_CHARS // 4)
    output_language = normalize_output_language(payload.language)

    def trimmed_block(label: str, value: str, limit: int = section_limit) -> str:
        cleaned = (value or "").strip()
        if not cleaned:
            return ""
        if len(cleaned) > limit:
            shortened = cleaned[:limit].rsplit(" ", 1)[0].strip() or cleaned[:limit].strip()
            cleaned = f"{shortened}\n\n[Trimmed for chat context]"
        return f"{label}\n{cleaned}"

    context_parts = [
        trimmed_block("STUDY GUIDE", payload.summary, section_limit),
        trimmed_block("LECTURE NOTES", payload.lecture_notes, section_limit),
        trimmed_block("LECTURE SLIDES", payload.lecture_slides, section_limit),
        trimmed_block("PAST QUESTION PAPERS", payload.past_question_papers, section_limit),
        trimmed_block("LECTURE TRANSCRIPT", payload.transcript, max(4000, MAX_CHAT_CONTEXT_CHARS // 2)),
    ]
    context_text = "\n\n".join(part for part in context_parts if part).strip()
    if len(context_text) > MAX_CHAT_CONTEXT_CHARS:
        context_text = context_text[:MAX_CHAT_CONTEXT_CHARS].rsplit(" ", 1)[0].strip()

    messages: list[dict[str, object]] = [
        {
            "role": "system",
            "content": (
                "You are MABASO.AI, a lecture study assistant. "
                "Answer only from the provided lecture context. "
                "If the material does not support an answer, say that it was not clearly covered. "
                "Be helpful, concise, and use bullets when they make the answer easier to study. "
                "After every answer, end with exactly one short follow-up question that is tailored to the exact concept, "
                "formula, worked example, or confusion the student just asked about. "
                "The follow-up should feel like a real tutor guiding the next step, not a generic closing line. "
                "For calculations or derivations, prefer offering to show the next step or the full step-by-step method. "
                "Put that follow-up question in its own final paragraph. "
                f"Reply in {output_language}."
            ),
        }
    ]

    if context_text:
        messages.append({"role": "system", "content": f"Lecture context:\n\n{context_text}"})

    for turn in payload.history[-6:]:
        role = "assistant" if turn.role == "assistant" else "user"
        content = (turn.content or "").strip()
        if content:
            messages.append({"role": role, "content": content[:1200]})

    question_text = payload.question.strip()
    reference_images = [image for image in payload.reference_images[:4] if (image or "").strip()]
    if reference_images:
        user_content: list[dict[str, object]] = [
            {
                "type": "text",
                "text": f"{question_text}\n\nUse the attached reference image(s) only if they help answer the question from the lecture context.",
            }
        ]
        for image in reference_images:
            user_content.append({"type": "image_url", "image_url": {"url": image}})
        messages.append({"role": "user", "content": user_content})
    else:
        messages.append({"role": "user", "content": question_text})
    return messages


def ensure_study_chat_follow_up(answer: str, question: str) -> str:
    cleaned = (answer or "").strip()
    if not cleaned:
        return "I could not form a clear answer from the lecture context.\n\nWould you like me to narrow it down and walk through that exact part step by step?"

    paragraphs = [part.strip() for part in re.split(r"\n\s*\n", cleaned) if part.strip()]
    if paragraphs:
        last_paragraph = paragraphs[-1]
        if last_paragraph.endswith("?") or "?" in last_paragraph:
            return cleaned

    focus = compact_text(question).replace('"', "'")
    if len(focus) > 140:
        focus = (focus[:137].rsplit(" ", 1)[0].strip() or focus[:137].strip()) + "..."

    if focus:
        follow_up = (
            f'Would you like me to unpack "{focus}" step by step and show how each part connects to the lecture method?'
        )
    else:
        follow_up = "Would you like me to unpack that step by step and show how each part connects to the lecture method?"

    return f"{cleaned}\n\n{follow_up}"


def sanitize_download_filename(value: str) -> str:
    cleaned = re.sub(r"[^A-Za-z0-9._-]+", "-", (value or "mabaso-study-pack").strip()).strip("-._")
    return cleaned[:80] or "mabaso-study-pack"


def dump_json(value: object) -> str:
    return json.dumps(value, ensure_ascii=False)


def load_json_list(value: str) -> list:
    if not value:
        return []
    try:
        parsed = json.loads(value)
    except json.JSONDecodeError:
        return []
    return parsed if isinstance(parsed, list) else []


def parse_json_object(value: str) -> dict[str, Any]:
    cleaned = (value or "").strip()
    if not cleaned:
        return {}
    if cleaned.startswith("```"):
        cleaned = re.sub(r"^```(?:json)?\s*", "", cleaned)
        cleaned = re.sub(r"\s*```$", "", cleaned)
    try:
        parsed = json.loads(cleaned)
    except json.JSONDecodeError:
        start = cleaned.find("{")
        end = cleaned.rfind("}")
        if start < 0 or end <= start:
            return {}
        try:
            parsed = json.loads(cleaned[start : end + 1])
        except json.JSONDecodeError:
            return {}
    return parsed if isinstance(parsed, dict) else {}


def parse_json_list(value: str) -> list[Any]:
    cleaned = (value or "").strip()
    if not cleaned:
        return []
    if cleaned.startswith("```"):
        cleaned = re.sub(r"^```(?:json)?\s*", "", cleaned)
        cleaned = re.sub(r"\s*```$", "", cleaned)
    try:
        parsed = json.loads(cleaned)
    except json.JSONDecodeError:
        start = cleaned.find("[")
        end = cleaned.rfind("]")
        if start < 0 or end <= start:
            return []
        try:
            parsed = json.loads(cleaned[start : end + 1])
        except json.JSONDecodeError:
            return []
    return parsed if isinstance(parsed, list) else []


def compact_text(value: Any, fallback: str = "") -> str:
    if value is None:
        return fallback
    text = str(value).replace("\r\n", "\n").strip()
    return text or fallback


def trimmed_context_block(label: str, value: str, limit: int) -> str:
    cleaned = (value or "").strip()
    if not cleaned:
        return ""
    if len(cleaned) > limit:
        cleaned = f"{cleaned[:limit].rstrip()}\n\nNOTE: This source was shortened for the structured asset step."
    return f"{label}\n{cleaned}"


def normalize_test_visibility(value: str) -> str:
    normalized = (value or "").strip().lower()
    if normalized not in {"shared", "private"}:
        raise HTTPException(status_code=400, detail="Test visibility must be either 'shared' or 'private'.")
    return normalized


def sanitize_collaboration_tab(value: str) -> str:
    normalized = (value or "guide").strip().lower()
    allowed_tabs = {"guide", "transcript", "formulas", "examples", "flashcards", "quiz", "chat"}
    return normalized if normalized in allowed_tabs else "guide"


def normalize_invited_emails(emails: list[str], owner_email: str) -> list[str]:
    cleaned: list[str] = []
    seen = {owner_email}
    for email in emails or []:
        raw_value = (email or "").strip()
        if not raw_value:
            continue
        normalized = validate_email_address(raw_value)
        if normalized in seen:
            continue
        seen.add(normalized)
        cleaned.append(normalized)
    return cleaned


def touch_collaboration_room(room_id: str):
    with get_db_connection() as connection:
        connection.execute(
            "UPDATE collaboration_rooms SET updated_at = ? WHERE id = ?",
            (utc_now().isoformat(), room_id),
        )


def get_collaboration_room_members(room_id: str) -> list[dict[str, str]]:
    with get_db_connection() as connection:
        rows = connection.execute(
            """
            SELECT email, role, created_at
            FROM collaboration_room_members
            WHERE room_id = ?
            ORDER BY CASE role WHEN 'owner' THEN 0 ELSE 1 END, email ASC
            """,
            (room_id,),
        ).fetchall()

    return [
        {"email": row["email"], "role": row["role"], "created_at": row["created_at"]}
        for row in rows
    ]


def get_collaboration_room_messages(room_id: str, limit: int = 80) -> list[dict[str, str]]:
    with get_db_connection() as connection:
        rows = connection.execute(
            """
            SELECT id, author_email, content, created_at
            FROM collaboration_room_messages
            WHERE room_id = ?
            ORDER BY created_at DESC
            LIMIT ?
            """,
            (room_id, limit),
        ).fetchall()

    ordered_rows = list(reversed(rows))
    return [
        {
            "id": row["id"],
            "author_email": row["author_email"],
            "content": row["content"],
            "created_at": row["created_at"],
        }
        for row in ordered_rows
    ]


def get_collaboration_room_answers(room_id: str, current_user: str, test_visibility: str) -> list[dict[str, str]]:
    query = """
        SELECT question_number, author_email, answer_text, updated_at
        FROM collaboration_room_answers
        WHERE room_id = ?
    """
    params: list[str] = [room_id]
    if test_visibility != "shared":
        query += " AND author_email = ?"
        params.append(current_user)
    query += " ORDER BY question_number ASC, author_email ASC"

    with get_db_connection() as connection:
        rows = connection.execute(query, params).fetchall()

    return [
        {
            "question_number": row["question_number"],
            "author_email": row["author_email"],
            "answer_text": row["answer_text"],
            "updated_at": row["updated_at"],
        }
        for row in rows
    ]


def get_accessible_collaboration_room(room_id: str, current_user: str) -> sqlite3.Row:
    with get_db_connection() as connection:
        row = connection.execute(
            """
            SELECT DISTINCT r.*
            FROM collaboration_rooms r
            LEFT JOIN collaboration_room_members m
                ON m.room_id = r.id
            WHERE r.id = ?
              AND (r.owner_email = ? OR m.email = ?)
            """,
            (room_id, current_user, current_user),
        ).fetchone()

    if not row:
        raise HTTPException(status_code=404, detail="Collaboration room not found.")
    return row


def serialize_collaboration_room(room_row: sqlite3.Row, current_user: str) -> dict:
    room_id = room_row["id"]
    test_visibility = normalize_test_visibility(room_row["test_visibility"])
    members = get_collaboration_room_members(room_id)

    return {
        "id": room_id,
        "title": room_row["title"],
        "owner_email": room_row["owner_email"],
        "transcript": room_row["transcript"],
        "summary": room_row["summary"],
        "formula": room_row["formula"],
        "example": room_row["example"],
        "lecture_notes": room_row["lecture_notes"],
        "lecture_slides": room_row["lecture_slides"],
        "shared_notes": room_row["shared_notes"],
        "flashcards": load_json_list(room_row["flashcards_json"]),
        "quiz_questions": load_json_list(room_row["quiz_questions_json"]),
        "active_tab": sanitize_collaboration_tab(room_row["active_tab"]),
        "test_visibility": test_visibility,
        "created_at": room_row["created_at"],
        "updated_at": room_row["updated_at"],
        "members": members,
        "messages": get_collaboration_room_messages(room_id),
        "quiz_answers": get_collaboration_room_answers(room_id, current_user, test_visibility),
        "is_owner": room_row["owner_email"] == current_user,
    }


def serialize_collaboration_room_card(room_row: sqlite3.Row) -> dict:
    members = get_collaboration_room_members(room_row["id"])
    return {
        "id": room_row["id"],
        "title": room_row["title"],
        "owner_email": room_row["owner_email"],
        "created_at": room_row["created_at"],
        "updated_at": room_row["updated_at"],
        "active_tab": sanitize_collaboration_tab(room_row["active_tab"]),
        "test_visibility": normalize_test_visibility(room_row["test_visibility"]),
        "member_count": len(members),
        "members": members,
    }


def load_audit_logs(limit: int = 400, *, days: int = 30) -> list[dict[str, Any]]:
    cutoff = (utc_now() - timedelta(days=days)).isoformat()
    with get_db_connection() as connection:
        rows = connection.execute(
            """
            SELECT *
            FROM audit_logs
            WHERE created_at >= ?
            ORDER BY created_at DESC
            LIMIT ?
            """,
            (cutoff, limit),
        ).fetchall()

    logs: list[dict[str, Any]] = []
    for row in rows:
        try:
            metadata = json.loads(row["metadata_json"] or "{}")
        except json.JSONDecodeError:
            metadata = {}
        logs.append(
            {
                "id": row["id"],
                "email": normalize_email(row["email"]),
                "ip_address": row["ip_address"],
                "user_agent": row["user_agent"],
                "action": row["action"],
                "resource_type": row["resource_type"],
                "resource_name": row["resource_name"],
                "duration_ms": int(row["duration_ms"] or 0),
                "status": row["status"],
                "metadata": metadata if isinstance(metadata, dict) else {},
                "created_at": row["created_at"],
            }
        )
    return logs


def load_history_items_with_owners(limit: int = 240) -> list[dict[str, Any]]:
    with get_db_connection() as connection:
        rows = connection.execute(
            """
            SELECT email, payload_json, created_at, updated_at
            FROM study_history_items
            ORDER BY updated_at DESC
            LIMIT ?
            """,
            (limit,),
        ).fetchall()

    items: list[dict[str, Any]] = []
    for row in rows:
        try:
            payload = normalize_history_item_payload(json.loads(row["payload_json"]))
        except (json.JSONDecodeError, HTTPException):
            continue
        payload["owner_email"] = normalize_email(row["email"])
        items.append(payload)
    return items


def classify_audit_feature(log: dict[str, Any]) -> str:
    action = (log.get("action") or "").lower()
    resource_type = (log.get("resource_type") or "").lower()
    if "study_guide" in action or resource_type == "study_guide":
        return "Study Guides"
    if "quiz" in action or resource_type == "quiz":
        return "Tests"
    if "presentation" in action or resource_type == "presentation":
        return "Presentations"
    if "podcast" in action or resource_type == "podcast":
        return "Podcasts"
    if "transcription" in action or "upload" in action or resource_type == "transcription":
        return "Lecture Capture"
    if "chat" in action:
        return "Study Chat"
    if "collaboration" in action:
        return "Collaboration"
    if "auth" in action:
        return "Authentication"
    return "Other"


def compute_retention_rate(logs: list[dict[str, Any]], days_after_signup: int) -> float:
    activity_by_email: dict[str, list[datetime]] = {}
    for log in logs:
        email = normalize_email(log.get("email", ""))
        if not email:
            continue
        activity_by_email.setdefault(email, []).append(parse_history_datetime(log.get("created_at")))

    if not activity_by_email:
        return 0.0

    retained = 0
    for timestamps in activity_by_email.values():
        ordered = sorted(timestamps)
        first_day = ordered[0].date()
        target_day = first_day + timedelta(days=days_after_signup)
        if any(timestamp.date() >= target_day for timestamp in ordered[1:]):
            retained += 1
    return round((retained / len(activity_by_email)) * 100, 1)


def estimate_history_item_storage_bytes(item: dict[str, Any]) -> dict[str, int]:
    lecture_sources_payload = "\n".join(
        [
            compact_text(item.get("transcript")),
            compact_text(item.get("lectureNotes")),
            compact_text(item.get("lectureSlides")),
            compact_text(item.get("pastQuestionPapers")),
        ]
    )
    generated_payload = json.dumps(
        {
            "summary": item.get("summary", ""),
            "formula": item.get("formula", ""),
            "example": item.get("example", ""),
            "flashcards": item.get("flashcards", []),
            "quizQuestions": item.get("quizQuestions", []),
            "presentationData": item.get("presentationData", {}),
            "podcastData": item.get("podcastData", {}),
        },
        ensure_ascii=False,
    )
    other_payload = json.dumps(
        {
            "studyImages": item.get("studyImages", []),
            "chatMessages": item.get("chatMessages", []),
            "metadata": {
                "title": item.get("title", ""),
                "fileName": item.get("fileName", ""),
                "createdAt": item.get("createdAt", ""),
                "updatedAt": item.get("updatedAt", ""),
            },
        },
        ensure_ascii=False,
    )

    lecture_sources_bytes = len(lecture_sources_payload.encode("utf-8"))
    generated_content_bytes = len(generated_payload.encode("utf-8"))
    other_data_bytes = len(other_payload.encode("utf-8"))
    return {
        "lecture_sources_bytes": lecture_sources_bytes,
        "generated_content_bytes": generated_content_bytes,
        "other_data_bytes": other_data_bytes,
        "total_bytes": lecture_sources_bytes + generated_content_bytes + other_data_bytes,
    }


def compute_storage_usage_breakdown(history_items: list[dict[str, Any]]) -> dict[str, int]:
    lecture_sources_bytes = 0
    generated_content_bytes = 0
    other_data_bytes = 0

    for item in history_items:
        item_sizes = estimate_history_item_storage_bytes(item)
        lecture_sources_bytes += item_sizes["lecture_sources_bytes"]
        generated_content_bytes += item_sizes["generated_content_bytes"]
        other_data_bytes += item_sizes["other_data_bytes"]

    return {
        "lecture_sources_bytes": lecture_sources_bytes,
        "generated_content_bytes": generated_content_bytes,
        "other_data_bytes": other_data_bytes,
        "total_bytes": lecture_sources_bytes + generated_content_bytes + other_data_bytes,
    }


def compute_session_analytics(
    logs: list[dict[str, Any]],
    session_rows: list[sqlite3.Row],
    last_login_by_email: dict[str, dict[str, str]],
    now: datetime,
) -> dict[str, Any]:
    active_sessions_by_email: dict[str, list[datetime]] = {}
    for row in session_rows:
        email = normalize_email(row["email"])
        if not email:
            continue
        expires_at = parse_history_datetime(row["expires_at"], now)
        active_sessions_by_email.setdefault(email, []).append(expires_at)

    session_gap = timedelta(minutes=30)
    session_logs_by_email: dict[str, list[dict[str, Any]]] = {}
    for log in sorted(logs, key=lambda entry: parse_history_datetime(entry["created_at"], now)):
        email = normalize_email(log.get("email", ""))
        if not email:
            continue
        session_logs_by_email.setdefault(email, []).append(log)

    completed_sessions: list[dict[str, Any]] = []
    sessions_by_email: dict[str, list[dict[str, Any]]] = {}

    for email, email_logs in session_logs_by_email.items():
        current_logs: list[dict[str, Any]] = []
        previous_timestamp: datetime | None = None

        for log in email_logs:
            timestamp = parse_history_datetime(log["created_at"], now)
            if previous_timestamp and timestamp - previous_timestamp > session_gap and current_logs:
                session_start = parse_history_datetime(current_logs[0]["created_at"], now)
                session_end = parse_history_datetime(current_logs[-1]["created_at"], now)
                session_record = {
                    "email": email,
                    "started_at": session_start.isoformat(),
                    "ended_at": session_end.isoformat(),
                    "duration_seconds": max(0, int((session_end - session_start).total_seconds())),
                    "actions": len(current_logs),
                    "is_bounce": len(current_logs) <= 1,
                }
                completed_sessions.append(session_record)
                sessions_by_email.setdefault(email, []).append(session_record)
                current_logs = []

            current_logs.append(log)
            previous_timestamp = timestamp

        if current_logs:
            session_start = parse_history_datetime(current_logs[0]["created_at"], now)
            session_end = parse_history_datetime(current_logs[-1]["created_at"], now)
            session_record = {
                "email": email,
                "started_at": session_start.isoformat(),
                "ended_at": session_end.isoformat(),
                "duration_seconds": max(0, int((session_end - session_start).total_seconds())),
                "actions": len(current_logs),
                "is_bounce": len(current_logs) <= 1,
            }
            completed_sessions.append(session_record)
            sessions_by_email.setdefault(email, []).append(session_record)

    total_sessions = len(completed_sessions)
    total_duration_seconds = sum(item["duration_seconds"] for item in completed_sessions)
    bounce_sessions = sum(1 for item in completed_sessions if item["is_bounce"])
    avg_duration_seconds = round(total_duration_seconds / total_sessions, 1) if total_sessions else 0.0
    bounce_rate_percent = round((bounce_sessions / total_sessions) * 100, 1) if total_sessions else 0.0

    session_timeline: list[dict[str, Any]] = []
    for day_offset in range(6, -1, -1):
        day_start = (now - timedelta(days=day_offset)).replace(hour=0, minute=0, second=0, microsecond=0)
        day_end = day_start + timedelta(days=1)
        matching_sessions = [
            session
            for session in completed_sessions
            if day_start <= parse_history_datetime(session["started_at"], now) < day_end
        ]
        day_total = len(matching_sessions)
        day_bounce = sum(1 for session in matching_sessions if session["is_bounce"])
        day_avg = round(sum(session["duration_seconds"] for session in matching_sessions) / day_total, 1) if day_total else 0.0
        session_timeline.append(
            {
                "date": day_start.date().isoformat(),
                "sessions": day_total,
                "avg_duration_seconds": day_avg,
                "bounce_rate_percent": round((day_bounce / day_total) * 100, 1) if day_total else 0.0,
            }
        )

    expiring_soon_cutoff = now + timedelta(minutes=20)
    expiring_soon = []
    session_table = []
    all_emails = set(active_sessions_by_email) | set(last_login_by_email) | set(sessions_by_email)
    for email in sorted(all_emails):
        active_expiries = sorted(active_sessions_by_email.get(email, []))
        user_sessions = sessions_by_email.get(email, [])
        next_timeout = active_expiries[0] if active_expiries else None
        latest_timeout = active_expiries[-1] if active_expiries else None
        user_avg_duration = round(sum(item["duration_seconds"] for item in user_sessions) / len(user_sessions), 1) if user_sessions else 0.0
        user_bounce_rate = round((sum(1 for item in user_sessions if item["is_bounce"]) / len(user_sessions)) * 100, 1) if user_sessions else 0.0
        total_actions = sum(item["actions"] for item in user_sessions)
        session_table.append(
            {
                "email": email,
                "active_sessions": len(active_expiries),
                "last_login_at": last_login_by_email.get(email, {}).get("created_at", ""),
                "next_timeout_at": next_timeout.isoformat() if next_timeout else "",
                "latest_timeout_at": latest_timeout.isoformat() if latest_timeout else "",
                "avg_session_duration_seconds": user_avg_duration,
                "bounce_rate_percent": user_bounce_rate,
                "total_sessions_30d": len(user_sessions),
                "total_actions": total_actions,
            }
        )
        if next_timeout and next_timeout <= expiring_soon_cutoff:
            expiring_soon.append(
                {
                    "email": email,
                    "expires_at": next_timeout.isoformat(),
                    "minutes_left": max(0, int((next_timeout - now).total_seconds() // 60)),
                    "active_sessions": len(active_expiries),
                }
            )

    session_table.sort(
        key=lambda item: (
            item["active_sessions"],
            item["total_sessions_30d"],
            item["total_actions"],
        ),
        reverse=True,
    )
    expiring_soon.sort(key=lambda item: item["expires_at"])

    return {
        "totals": {
            "active_sessions": len(session_rows),
            "tracked_sessions_30d": total_sessions,
            "avg_session_duration_seconds": avg_duration_seconds,
            "bounce_rate_percent": bounce_rate_percent,
            "expiring_soon_count": len(expiring_soon),
        },
        "timeline": session_timeline,
        "table_full": session_table,
        "table": session_table[:40],
        "expiring_soon": expiring_soon[:12],
    }


def build_admin_dashboard_snapshot() -> dict[str, Any]:
    now = utc_now()
    logs = load_audit_logs(limit=ADMIN_DASHBOARD_AUDIT_LOG_LIMIT, days=35)
    history_items = load_history_items_with_owners(limit=ADMIN_DASHBOARD_HISTORY_LIMIT)

    with get_db_connection() as connection:
        user_rows = connection.execute(
            "SELECT email, created_at, verified_at FROM users ORDER BY created_at DESC"
        ).fetchall()
        session_rows = connection.execute(
            "SELECT email, expires_at, created_at FROM sessions WHERE expires_at > ?",
            (now.isoformat(),),
        ).fetchall()
        account_state_rows = connection.execute(
            "SELECT email, status, updated_at, updated_by FROM user_account_states"
        ).fetchall()
        admin_attempt_rows = connection.execute(
            "SELECT email, ip_address, failure_count, last_failed_at, locked_until FROM admin_login_attempts"
        ).fetchall()

    account_state_by_email = {
        normalize_email(row["email"]): {
            "status": normalize_user_account_status(row["status"]),
            "updated_at": row["updated_at"],
            "updated_by": row["updated_by"],
        }
        for row in account_state_rows
    }
    sessions_by_email: dict[str, int] = {}
    for row in session_rows:
        email = normalize_email(row["email"])
        if not email:
            continue
        sessions_by_email[email] = sessions_by_email.get(email, 0) + 1

    last_login_by_email: dict[str, dict[str, str]] = {}
    uploads_by_email: dict[str, int] = {}
    generations_by_email: dict[str, int] = {}
    failed_auth_by_email: dict[str, int] = {}
    transcriptions_by_email: dict[str, int] = {}
    unique_ips_by_email: dict[str, set[str]] = {}
    feature_counts: dict[str, int] = {}
    saved_items_by_email: dict[str, int] = {}
    tests_by_email: dict[str, int] = {}
    storage_bytes_by_email: dict[str, int] = {}

    for item in history_items:
        email = normalize_email(item.get("owner_email", ""))
        if not email:
            continue
        saved_items_by_email[email] = saved_items_by_email.get(email, 0) + 1
        if item.get("quizQuestions"):
            tests_by_email[email] = tests_by_email.get(email, 0) + 1
        storage_bytes_by_email[email] = (
            storage_bytes_by_email.get(email, 0) + estimate_history_item_storage_bytes(item)["total_bytes"]
        )

    successful_login_actions = {
        "auth.email_password.login",
        "auth.email_password.verify_code",
        "auth.google.login",
        "auth.apple.login",
        "auth.code.verify",
    }
    upload_actions = {"lecture.upload.request", "lecture.video_link.request"}
    generation_actions = {
        "study_guide.request",
        "study_guide.completed",
        "quiz.request",
        "quiz.completed",
        "presentation.request",
        "presentation.completed",
        "podcast.request",
        "podcast.completed",
    }

    for log in logs:
        email = normalize_email(log["email"])
        action = (log["action"] or "").lower()
        if email:
            unique_ips_by_email.setdefault(email, set()).add(log["ip_address"])
        if action in successful_login_actions and log["status"] == "success" and email and email not in last_login_by_email:
            last_login_by_email[email] = {
                "created_at": log["created_at"],
                "ip_address": log["ip_address"],
            }
        if action in upload_actions and email:
            uploads_by_email[email] = uploads_by_email.get(email, 0) + 1
        if action in generation_actions and email:
            generations_by_email[email] = generations_by_email.get(email, 0) + 1
        if action == "transcription.completed" and log["status"] == "success" and email:
            transcriptions_by_email[email] = transcriptions_by_email.get(email, 0) + 1
        if action.startswith("auth.") and log["status"] == "failed" and email:
            failed_auth_by_email[email] = failed_auth_by_email.get(email, 0) + 1
        feature = classify_audit_feature(log)
        feature_counts[feature] = feature_counts.get(feature, 0) + 1

    session_analytics = compute_session_analytics(logs, session_rows, last_login_by_email, now)
    session_table_by_email = {
        normalize_email(item.get("email", "")): item
        for item in session_analytics.get("table_full", [])
        if normalize_email(item.get("email", ""))
    }
    storage_usage_breakdown = compute_storage_usage_breakdown(history_items)

    user_records: list[dict[str, Any]] = []
    for row in user_rows:
        email = normalize_email(row["email"])
        status = account_state_by_email.get(email, {}).get("status", "active")
        failed_attempts = failed_auth_by_email.get(email, 0)
        ip_count = len(unique_ips_by_email.get(email, set()))
        risk_score = "high" if failed_attempts >= 3 or ip_count >= 4 else "medium" if failed_attempts or ip_count >= 2 else "low"
        last_login = last_login_by_email.get(email, {})
        session_profile = session_table_by_email.get(email, {})
        user_records.append(
            {
                "email": email,
                "role": "admin" if is_admin_email(email) else "user",
                "status": status,
                "created_at": row["created_at"],
                "last_login_at": last_login.get("created_at", ""),
                "last_login_ip": last_login.get("ip_address", ""),
                "sessions_count": sessions_by_email.get(email, 0),
                "total_uploads": uploads_by_email.get(email, 0),
                "total_generations": generations_by_email.get(email, 0),
                "lectures_transcribed": transcriptions_by_email.get(email, 0),
                "study_materials": saved_items_by_email.get(email, 0),
                "tests_generated": tests_by_email.get(email, 0),
                "avg_session_duration_seconds": session_profile.get("avg_session_duration_seconds", 0),
                "avg_session_duration": session_profile.get("avg_session_duration_seconds", 0),
                "next_timeout_at": session_profile.get("next_timeout_at", ""),
                "latest_timeout_at": session_profile.get("latest_timeout_at", ""),
                "total_actions_30d": session_profile.get("total_actions", 0),
                "storage_bytes": storage_bytes_by_email.get(email, 0),
                "risk_score": risk_score,
            }
        )

    total_users = len(user_records)
    new_users_7d = sum(
        1 for row in user_rows if parse_history_datetime(row["created_at"], now) >= now - timedelta(days=7)
    )
    active_1h = {
        log["email"]
        for log in logs
        if log["email"] and parse_history_datetime(log["created_at"], now) >= now - timedelta(hours=1)
    }
    active_24h = {
        log["email"]
        for log in logs
        if log["email"] and parse_history_datetime(log["created_at"], now) >= now - timedelta(hours=24)
    }
    active_7d = {
        log["email"]
        for log in logs
        if log["email"] and parse_history_datetime(log["created_at"], now) >= now - timedelta(days=7)
    }

    lectures_today = sum(
        1
        for log in logs
        if log["action"] in upload_actions and parse_history_datetime(log["created_at"], now) >= now - timedelta(days=1)
    )
    lectures_week = sum(
        1
        for log in logs
        if log["action"] in upload_actions and parse_history_datetime(log["created_at"], now) >= now - timedelta(days=7)
    )
    transcriptions_total = sum(
        1 for log in logs if log["action"] == "transcription.completed" and log["status"] == "success"
    )
    transcriptions_week = sum(
        1
        for log in logs
        if log["action"] == "transcription.completed"
        and log["status"] == "success"
        and parse_history_datetime(log["created_at"], now) >= now - timedelta(days=7)
    )

    guide_count = len(history_items)
    study_material_count = sum(
        1
        for item in history_items
        if compact_text(item.get("summary"))
        or compact_text(item.get("formula"))
        or compact_text(item.get("example"))
        or item.get("flashcards")
        or item.get("quizQuestions")
        or item.get("presentationData")
        or item.get("podcastData")
    )
    test_count = sum(1 for item in history_items if item.get("quizQuestions"))
    processing_durations = [
        log["duration_ms"]
        for log in logs
        if log["duration_ms"] > 0 and log["status"] == "success" and log["action"].endswith(".completed")
    ]
    avg_processing_time = round(sum(processing_durations) / len(processing_durations), 1) if processing_durations else 0
    recent_logs_24h = [
        log for log in logs if parse_history_datetime(log["created_at"], now) >= now - timedelta(days=1)
    ]
    failed_logs_24h = [log for log in recent_logs_24h if log["status"] != "success"]
    error_rate = round((len(failed_logs_24h) / len(recent_logs_24h)) * 100, 1) if recent_logs_24h else 0.0
    current_jobs = list(jobs.values())
    processing_jobs = [job for job in current_jobs if job.get("status") in {"queued", "processing"}]
    transcription_jobs = [
        job for job in current_jobs if job.get("job_type") in {"transcription", "video_transcription"}
    ]
    transcription_queue = {
        "in_queue": sum(1 for job in transcription_jobs if job.get("status") == "queued"),
        "processing": sum(1 for job in transcription_jobs if job.get("status") == "processing"),
        "completed_7d": sum(
            1
            for log in logs
            if log["action"] == "transcription.completed"
            and log["status"] == "success"
            and parse_history_datetime(log["created_at"], now) >= now - timedelta(days=7)
        ),
        "failed_7d": sum(
            1
            for log in logs
            if log["action"] == "transcription.completed"
            and log["status"] != "success"
            and parse_history_datetime(log["created_at"], now) >= now - timedelta(days=7)
        ),
    }
    system_state = "green" if error_rate < 10 and len(processing_jobs) <= 6 else "yellow" if error_rate < 25 else "red"

    real_time_activity: list[dict[str, Any]] = []
    for bucket_index in range(12):
        bucket_end = now - timedelta(minutes=(11 - bucket_index) * 5)
        bucket_start = bucket_end - timedelta(minutes=5)
        real_time_activity.append(
            {
                "label": bucket_end.strftime("%H:%M"),
                "count": sum(
                    1
                    for log in logs
                    if bucket_start <= parse_history_datetime(log["created_at"], now) < bucket_end
                ),
            }
        )

    daily_activity: list[dict[str, Any]] = []
    for day_offset in range(29, -1, -1):
        day_start = (now - timedelta(days=day_offset)).replace(hour=0, minute=0, second=0, microsecond=0)
        day_end = day_start + timedelta(days=1)
        active_emails = {
            log["email"]
            for log in logs
            if log["email"] and day_start <= parse_history_datetime(log["created_at"], now) < day_end
        }
        new_user_count = sum(
            1
            for row in user_rows
            if day_start <= parse_history_datetime(row["created_at"], now) < day_end
        )
        daily_activity.append(
            {
                "date": day_start.date().isoformat(),
                "active_users": len(active_emails),
                "new_users": new_user_count,
            }
        )

    upload_count = sum(1 for log in logs if log["action"] in upload_actions)
    processed_count = sum(
        1 for log in logs if log["action"] == "transcription.completed" and log["status"] == "success"
    )
    generated_count = sum(
        1 for log in logs if log["action"] == "study_guide.completed" and log["status"] == "success"
    )

    activity_logs = [
        {
            "timestamp": log["created_at"],
            "user": log["email"] or "anonymous",
            "ip_address": log["ip_address"],
            "device": log["user_agent"],
            "action": log["action"],
            "resource": log["resource_name"] or log["resource_type"],
            "duration_ms": log["duration_ms"],
            "status": log["status"],
        }
        for log in logs[:120]
    ]

    all_user_emails = {
        normalize_email(row["email"])
        for row in user_rows
        if normalize_email(row["email"])
    } | set(saved_items_by_email) | set(last_login_by_email) | set(transcriptions_by_email)
    top_users_by_usage = sorted(
        (
            {
                "email": email,
                "lectures": max(uploads_by_email.get(email, 0), transcriptions_by_email.get(email, 0)),
                "materials": saved_items_by_email.get(email, 0),
                "tests": tests_by_email.get(email, 0),
                "sessions": sessions_by_email.get(email, 0),
                "total_actions": session_table_by_email.get(email, {}).get("total_actions", 0),
                "storage_bytes": storage_bytes_by_email.get(email, 0),
            }
            for email in all_user_emails
        ),
        key=lambda item: (
            item["lectures"] + item["materials"] + item["tests"],
            item["sessions"],
            item["total_actions"],
        ),
        reverse=True,
    )[:10]

    content_items = []
    for item in history_items[:120]:
        item_size = estimate_history_item_storage_bytes(item)["total_bytes"]
        has_output = bool(
            compact_text(item.get("summary"))
            or compact_text(item.get("formula"))
            or compact_text(item.get("example"))
            or item.get("flashcards")
            or item.get("quizQuestions")
            or item.get("presentationData")
            or item.get("podcastData")
        )
        content_items.append(
            {
                "file_name": item.get("fileName") or item.get("title") or "Saved lecture",
                "owner_email": item.get("owner_email", ""),
                "title": item.get("title") or "Saved lecture",
                "upload_date": item.get("updatedAt") or item.get("createdAt") or "",
                "processing_status": "done",
                "output_generated": "Y" if has_output else "N",
                "size_bytes": item_size,
                "size_label": "--",
                "duration_label": "--",
            }
        )

    failed_jobs = [
        {
            "timestamp": log["created_at"],
            "email": log["email"],
            "action": log["action"],
            "message": str(
                log["metadata"].get("error")
                or log["metadata"].get("reason")
                or log["resource_name"]
                or "Request failed."
            ),
        }
        for log in logs
        if log["status"] != "success"
    ][:40]

    session_heatmap = [
        {
            "hour": f"{hour:02d}:00",
            "actions": sum(
                1 for log in logs if parse_history_datetime(log["created_at"], now).hour == hour
            ),
        }
        for hour in range(24)
    ]

    suspicious_activity = []
    for email, failure_count in failed_auth_by_email.items():
        if failure_count >= 3:
            suspicious_activity.append(
                {
                    "email": email,
                    "reason": f"{failure_count} failed login attempts in recent logs.",
                }
            )
    for row in admin_attempt_rows:
        locked_until = (row["locked_until"] or "").strip()
        if locked_until and datetime.fromisoformat(locked_until) > now:
            suspicious_activity.append(
                {
                    "email": normalize_email(row["email"]),
                    "reason": f"Admin login lock active for IP {row['ip_address']}.",
                }
            )

    ip_activity: dict[str, dict[str, Any]] = {}
    for log in logs[:300]:
        ip_address = log["ip_address"] or "unknown"
        entry = ip_activity.setdefault(ip_address, {"ip_address": ip_address, "users": set(), "actions": 0})
        if log["email"]:
            entry["users"].add(log["email"])
        entry["actions"] += 1

    device_activity: dict[str, int] = {}
    for log in logs[:300]:
        device = (log["user_agent"] or "Unknown device")[:120]
        device_activity[device] = device_activity.get(device, 0) + 1

    return {
        "overview": {
            "kpis": {
                "total_users": total_users,
                "active_users_1h": len(active_1h),
                "active_users_24h": len(active_24h),
                "active_users_7d": len(active_7d),
                "new_users_7d": new_users_7d,
                "lectures_uploaded_today": lectures_today,
                "lectures_uploaded_week": lectures_week,
                "lectures_transcribed": transcriptions_total,
                "lectures_transcribed_week": transcriptions_week,
                "study_guides_generated": guide_count,
                "study_materials_generated": study_material_count,
                "tests_generated": test_count,
                "storage_used_bytes": storage_usage_breakdown["total_bytes"],
                "active_sessions": session_analytics["totals"]["active_sessions"],
                "avg_session_duration_seconds": session_analytics["totals"]["avg_session_duration_seconds"],
                "bounce_rate_percent": session_analytics["totals"]["bounce_rate_percent"],
                "avg_processing_time_ms": avg_processing_time,
                "error_rate_percent": error_rate,
                "system_load": len(processing_jobs),
                "transcription_queue": transcription_queue["in_queue"] + transcription_queue["processing"],
            },
            "charts": {
                "real_time_activity": real_time_activity,
                "daily_active_users": daily_activity,
                "user_sessions": session_analytics["timeline"],
                "feature_usage_breakdown": [
                    {"label": label, "count": count}
                    for label, count in sorted(feature_counts.items(), key=lambda item: item[1], reverse=True)
                ],
                "conversion_funnel": [
                    {"label": "Uploaded", "count": upload_count},
                    {"label": "Processed", "count": processed_count},
                    {"label": "Generated Output", "count": generated_count or guide_count},
                ],
                "wau": len(active_7d),
                "mau": len({log["email"] for log in logs if log["email"]}),
            },
        },
        "users": user_records,
        "sessions": session_analytics,
        "activity_logs": activity_logs,
        "content": {
            "items": content_items,
            "storage_insights": {
                "tracked_study_packs": len(history_items),
                "total_bytes": storage_usage_breakdown["total_bytes"],
                "breakdown": storage_usage_breakdown,
                "top_users": top_users_by_usage,
            },
        },
        "ai_generation": {
            "totals": {
                "study_guides": guide_count,
                "presentations": sum(
                    1
                    for item in history_items
                    if isinstance(item.get("presentationData"), dict)
                    and (
                        item["presentationData"].get("slides")
                        or item["presentationData"].get("presentation_slides")
                    )
                ),
                "podcasts": sum(
                    1
                    for item in history_items
                    if isinstance(item.get("podcastData"), dict)
                    and (
                        item["podcastData"].get("script")
                        or item["podcastData"].get("podcast_script")
                    )
                ),
            },
            "success_rate_percent": round(
                (
                    len([log for log in logs if log["action"].endswith(".completed") and log["status"] == "success"])
                    / max(1, len([log for log in logs if log["action"].endswith(".completed")]))
                ) * 100,
                1,
            ),
            "avg_generation_time_ms": avg_processing_time,
            "failed_jobs": failed_jobs,
        },
        "analytics": {
            "session_heatmap": session_heatmap,
            "retention": {
                "day_1": compute_retention_rate(logs, 1),
                "day_7": compute_retention_rate(logs, 7),
                "day_30": compute_retention_rate(logs, 30),
            },
            "most_used_tools": [
                {"label": label, "count": count}
                for label, count in sorted(feature_counts.items(), key=lambda item: item[1], reverse=True)[:8]
            ],
            "drop_off_points": [
                {"label": "Upload to Processed", "count": max(upload_count - processed_count, 0)},
                {"label": "Processed to Generated", "count": max(processed_count - (generated_count or guide_count), 0)},
            ],
            "performance": {
                "avg_response_time_ms": avg_processing_time,
                "actions_per_session": round(
                    len(logs) / max(1, session_analytics["totals"]["tracked_sessions_30d"]),
                    1,
                ),
            },
        },
        "system_health": {
            "state": system_state,
            "api_response_time_ms": avg_processing_time,
            "queue_length": len(processing_jobs),
            "active_sessions": session_analytics["totals"]["active_sessions"],
            "transcription_queue": transcription_queue,
            "active_jobs": [
                {
                    "job_id": job.get("job_id", ""),
                    "job_type": job.get("job_type", ""),
                    "status": job.get("status", ""),
                    "stage": job.get("stage", ""),
                    "progress": job.get("progress", 0),
                    "owner_email": job.get("owner_email", ""),
                }
                for job in processing_jobs[:20]
            ],
            "recent_failures": failed_jobs[:12],
        },
        "security": {
            "failed_logins": [
                {
                    "timestamp": log["created_at"],
                    "email": log["email"] or "unknown",
                    "ip_address": log["ip_address"],
                    "status": log["status"],
                    "reason": str(log["metadata"].get("reason") or "Failed login"),
                }
                for log in logs
                if log["action"].startswith("auth.") and log["status"] != "success"
            ][:40],
            "suspicious_activity": suspicious_activity[:20],
            "ip_tracking": [
                {
                    "ip_address": item["ip_address"],
                    "users": sorted(item["users"]),
                    "actions": item["actions"],
                }
                for item in sorted(ip_activity.values(), key=lambda value: value["actions"], reverse=True)[:20]
            ],
            "device_tracking": [
                {"device": device, "actions": count}
                for device, count in sorted(device_activity.items(), key=lambda item: item[1], reverse=True)[:20]
            ],
        },
        "billing": {
            "subscriptions": [],
            "usage": [],
            "revenue": [],
        },
        "settings": {
            "available_languages": sorted(set(SUPPORTED_OUTPUT_LANGUAGES.values())),
            "admin_email_count": len(get_admin_email_set()),
            "feature_flags": {
                "powerpoint_generation": Presentation is not None,
                "podcast_generation": True,
                "collaboration": True,
            },
        },
        "generated_at": now.isoformat(),
    }


def ensure_openai_key():
    if not os.getenv("OPENAI_API_KEY"):
        raise HTTPException(
            status_code=500,
            detail="OPENAI_API_KEY is not configured on the backend.",
        )


def is_text_upload(filename: str, content_type: str | None) -> bool:
    suffix = Path(filename or "").suffix.lower()
    return bool(
        (content_type and content_type.startswith("text/"))
        or suffix in {".txt", ".md", ".text"}
    )


def is_pdf_upload(filename: str, content_type: str | None) -> bool:
    suffix = Path(filename or "").suffix.lower()
    return suffix == ".pdf" or content_type == "application/pdf"


def is_pptx_upload(filename: str, content_type: str | None) -> bool:
    suffix = Path(filename or "").suffix.lower()
    return suffix == ".pptx" or content_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation"


def is_docx_upload(filename: str, content_type: str | None) -> bool:
    suffix = Path(filename or "").suffix.lower()
    return suffix == ".docx" or content_type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document"


def build_data_url(file_bytes: bytes, content_type: str | None, filename: str = "") -> str:
    mime_type = content_type or mimetypes.guess_type(filename)[0] or "application/octet-stream"
    encoded = base64.b64encode(file_bytes).decode("ascii")
    return f"data:{mime_type};base64,{encoded}"


def normalize_zip_member_path(base_member: str, target: str) -> str:
    combined_parts = list(PurePosixPath(base_member).parent.joinpath(target).parts)
    normalized_parts: list[str] = []
    for part in combined_parts:
        if part in {"", "."}:
            continue
        if part == "..":
            if normalized_parts:
                normalized_parts.pop()
            continue
        normalized_parts.append(part)
    return "/".join(normalized_parts)


def merge_text_blocks(blocks: list[str]) -> str:
    seen: set[str] = set()
    merged: list[str] = []
    for block in blocks:
        cleaned = compact_text(block)
        if not cleaned or cleaned in seen:
            continue
        seen.add(cleaned)
        merged.append(cleaned)
    return "\n\n".join(merged).strip()


def extract_text_from_embedded_image(image_bytes: bytes, filename: str) -> str:
    content_type = mimetypes.guess_type(filename)[0] or "image/png"
    image_data_url = build_data_url(image_bytes, content_type, filename)
    return extract_slide_text_from_image(image_data_url, filename)


def iter_pdf_page_images(page: Any) -> list[tuple[str, bytes]]:
    images = getattr(page, "images", None) or []
    extracted: list[tuple[str, bytes]] = []
    for index, image in enumerate(images, start=1):
        image_bytes = getattr(image, "data", None)
        image_name = getattr(image, "name", "") or f"page-image-{index}.png"
        if image_bytes:
            extracted.append((image_name, image_bytes))
    return extracted


def extract_slide_image_text_blocks(image_items: list[tuple[str, bytes]], *, limit: int = 2) -> list[str]:
    text_blocks: list[str] = []
    for file_name, image_bytes in sorted(image_items, key=lambda item: len(item[1]), reverse=True):
        if len(text_blocks) >= limit:
            break
        if len(image_bytes) < 2048:
            continue
        try:
            text = extract_text_from_embedded_image(image_bytes, file_name)
        except Exception as exc:
            logger.warning("Slide OCR fallback failed for %s: %s", file_name, exc)
            continue
        if text:
            text_blocks.append(text)
    return text_blocks


def build_reference_image_data_url(image_bytes: bytes, filename: str = "") -> str:
    content_type = mimetypes.guess_type(filename or "")[0] or "image/png"
    encoded = base64.b64encode(image_bytes).decode("ascii")
    return f"data:{content_type};base64,{encoded}"


def extract_reference_images_from_items(image_items: list[tuple[str, bytes]], *, limit: int = 6) -> list[str]:
    if limit <= 0:
        return []

    extracted: list[str] = []
    seen_hashes: set[str] = set()
    for file_name, image_bytes in sorted(image_items, key=lambda item: len(item[1]), reverse=True):
        if len(extracted) >= limit:
            break
        if len(image_bytes) < 4096:
            continue
        fingerprint = hashlib.sha256(image_bytes).hexdigest()
        if fingerprint in seen_hashes:
            continue
        seen_hashes.add(fingerprint)
        extracted.append(build_reference_image_data_url(image_bytes, file_name))
    return extracted


def extract_reference_images_from_pdf(file_bytes: bytes, *, limit: int = 6) -> list[str]:
    if PdfReader is None:
        return []

    reader = PdfReader(BytesIO(file_bytes))
    image_items: list[tuple[str, bytes]] = []
    for page in reader.pages:
        image_items.extend(iter_pdf_page_images(page))
        if len(image_items) >= limit * 3:
            break
    return extract_reference_images_from_items(image_items, limit=limit)


def extract_reference_images_from_pptx(file_bytes: bytes, *, limit: int = 6) -> list[str]:
    namespaces = {
        "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
        "rels": "http://schemas.openxmlformats.org/package/2006/relationships",
    }

    with zipfile.ZipFile(BytesIO(file_bytes)) as archive:
        slide_names = sorted(
            [name for name in archive.namelist() if re.fullmatch(r"ppt/slides/slide\d+\.xml", name)],
            key=lambda value: int(re.search(r"slide(\d+)\.xml", value).group(1)),
        )
        image_items: list[tuple[str, bytes]] = []
        for slide_name in slide_names:
            if len(image_items) >= limit * 3:
                break
            root = ET.fromstring(archive.read(slide_name))
            rels_name = normalize_zip_member_path(slide_name, f"_rels/{PurePosixPath(slide_name).name}.rels")
            if rels_name not in archive.namelist():
                continue
            rels_root = ET.fromstring(archive.read(rels_name))
            relationship_targets = {
                relation.attrib.get("Id", ""): normalize_zip_member_path(rels_name, relation.attrib.get("Target", ""))
                for relation in rels_root.findall("rels:Relationship", namespaces)
                if relation.attrib.get("Type", "").endswith("/image")
            }
            slide_image_targets: list[str] = []
            for blip in root.findall(".//a:blip", namespaces):
                target = relationship_targets.get(blip.attrib.get("{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed", ""))
                if target and target in archive.namelist() and target not in slide_image_targets:
                    slide_image_targets.append(target)
            image_items.extend((target, archive.read(target)) for target in slide_image_targets)
        return extract_reference_images_from_items(image_items, limit=limit)


def extract_slide_text_from_pdf(file_bytes: bytes) -> str:
    if PdfReader is None:
        raise HTTPException(
            status_code=500,
            detail="PDF slide support needs the pypdf package on the backend. Install backend requirements and try again.",
        )

    reader = PdfReader(BytesIO(file_bytes))
    pages: list[str] = []
    for index, page in enumerate(reader.pages, start=1):
        page_text = (page.extract_text() or "").strip()
        text_blocks = [page_text] if page_text else []
        if len(page_text) < 80:
            text_blocks.extend(extract_slide_image_text_blocks(iter_pdf_page_images(page), limit=2))
        merged_text = merge_text_blocks(text_blocks)
        if merged_text:
            pages.append(f"SLIDE PAGE {index}\n{merged_text}")
    return "\n\n".join(pages).strip()


def extract_slide_text_from_pptx(file_bytes: bytes) -> str:
    namespaces = {
        "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
        "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
        "rels": "http://schemas.openxmlformats.org/package/2006/relationships",
    }

    with zipfile.ZipFile(BytesIO(file_bytes)) as archive:
        slide_names = sorted(
            [name for name in archive.namelist() if re.fullmatch(r"ppt/slides/slide\d+\.xml", name)],
            key=lambda value: int(re.search(r"slide(\d+)\.xml", value).group(1)),
        )

        slide_text_parts: list[str] = []
        for slide_index, slide_name in enumerate(slide_names, start=1):
            xml_bytes = archive.read(slide_name)
            root = ET.fromstring(xml_bytes)
            text_runs = [node.text.strip() for node in root.findall(".//a:t", namespaces) if node.text and node.text.strip()]
            slide_blocks = ["\n".join(text_runs)] if text_runs else []

            if len("\n".join(text_runs).strip()) < 80:
                rels_name = normalize_zip_member_path(slide_name, f"_rels/{PurePosixPath(slide_name).name}.rels")
                image_items: list[tuple[str, bytes]] = []
                if rels_name in archive.namelist():
                    rels_root = ET.fromstring(archive.read(rels_name))
                    relationship_targets = {
                        relation.attrib.get("Id", ""): normalize_zip_member_path(rels_name, relation.attrib.get("Target", ""))
                        for relation in rels_root.findall("rels:Relationship", namespaces)
                        if relation.attrib.get("Type", "").endswith("/image")
                    }
                    slide_image_targets: list[str] = []
                    for blip in root.findall(".//a:blip", namespaces):
                        target = relationship_targets.get(blip.attrib.get("{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed", ""))
                        if target and target in archive.namelist() and target not in slide_image_targets:
                            slide_image_targets.append(target)
                    image_items = [(target, archive.read(target)) for target in slide_image_targets]

                slide_blocks.extend(extract_slide_image_text_blocks(image_items, limit=2))

            merged_text = merge_text_blocks(slide_blocks)
            if merged_text:
                slide_text_parts.append(f"SLIDE {slide_index}\n{merged_text}")

    return "\n\n".join(slide_text_parts).strip()


def extract_text_from_docx(file_bytes: bytes) -> str:
    namespaces = {"w": "http://schemas.openxmlformats.org/wordprocessingml/2006/main"}
    paragraphs: list[str] = []

    with zipfile.ZipFile(BytesIO(file_bytes)) as archive:
        if "word/document.xml" not in archive.namelist():
            return ""

        root = ET.fromstring(archive.read("word/document.xml"))
        for paragraph in root.findall(".//w:p", namespaces):
            text_runs = [node.text for node in paragraph.findall(".//w:t", namespaces) if node.text]
            paragraph_text = "".join(text_runs).strip()
            if paragraph_text:
                paragraphs.append(paragraph_text)

    return "\n\n".join(paragraphs).strip()


def extract_slide_text_from_image(image_data_url: str, file_name: str = "") -> str:
    response = client.with_options(timeout=VISION_REQUEST_TIMEOUT).chat.completions.create(
        model=VISION_MODEL,
        max_completion_tokens=1200,
        messages=[
            {
                "role": "system",
                "content": (
                    "You are extracting useful study content from a lecture slide image. "
                    "Return only the slide content in clean readable Markdown. Preserve headings, bullets, "
                    "key definitions, and formulas. Rewrite formulas in plain human-readable board style, "
                    "never in LaTeX."
                ),
            },
            {
                "role": "user",
                "content": [
                    {
                        "type": "text",
                        "text": (
                            f"Extract the important lecture slide content from this image. "
                            f"File name: {file_name or 'slide image'}. "
                            "Keep the result concise but complete enough to help build better study guides and quiz questions."
                        ),
                    },
                    {"type": "image_url", "image_url": {"url": image_data_url}},
                ],
            },
        ],
    )
    return (response.choices[0].message.content or "").strip()


def clamp_score(value: Any, max_score: int) -> int:
    try:
        numeric = int(round(float(str(value).strip())))
    except (TypeError, ValueError):
        return 0
    return max(0, min(max_score, numeric))


def parse_quiz_marking_response(content: str, max_score: int) -> dict[str, Any]:
    parsed_object = parse_json_object(content)
    mistakes = parsed_object.get("mistakes")
    if not isinstance(mistakes, list):
        mistakes = []
    return {
        "score": clamp_score(parsed_object.get("score"), max_score),
        "max_score": max_score,
        "extracted_answer": compact_text(parsed_object.get("extracted_answer")),
        "feedback": compact_text(parsed_object.get("feedback"), "The answer was reviewed."),
        "mistakes": [compact_text(item) for item in mistakes if compact_text(item)],
    }


def mark_quiz_answer_with_ai(
    question: str,
    expected_answer: str,
    student_answer: str = "",
    answer_image_data_urls: list[str] | None = None,
    answer_points: list[str] | None = None,
    max_score: int = 1,
) -> dict[str, Any]:
    point_lines = "\n".join(f"- {point}" for point in (answer_points or []) if compact_text(point))
    user_parts: list[dict] = [
        {
            "type": "text",
            "text": (
                "Grade this student answer against the expected answer.\n\n"
                f"Question: {question}\n"
                f"Maximum marks: {max_score}\n"
                f"Expected answer: {expected_answer}\n"
                f"Key marking points:\n{point_lines or '- Use the expected answer directly.'}\n"
                f"Typed student answer: {student_answer or 'No typed answer provided.'}\n\n"
                "Be fair if the meaning is correct even when wording differs. "
                "Award partial credit when some key points are correct. "
                "Return JSON only with these keys:\n"
                f'{{"score": 0 to {max_score}, "extracted_answer": "...", "feedback": "...", "mistakes": ["..."]}}'
            ),
        }
    ]

    for answer_image_data_url in answer_image_data_urls or []:
        if answer_image_data_url:
            user_parts.append({"type": "image_url", "image_url": {"url": answer_image_data_url}})

    response = client.with_options(timeout=VISION_REQUEST_TIMEOUT).chat.completions.create(
        model=VISION_MODEL,
        max_completion_tokens=500,
        messages=[
            {
                "role": "system",
                "content": (
                    "You mark student quiz answers carefully. "
                    "You may read typed answers, one answer photo, several answer photos, neat handwriting, messy handwriting, or any combination of them. "
                    "Return strict JSON only. "
                    "Never award more than the maximum marks provided."
                ),
            },
            {"role": "user", "content": user_parts},
        ],
    )
    content = (response.choices[0].message.content or "").strip()
    return parse_quiz_marking_response(content, max_score)


def grade_option_based_question(
    question_type: str,
    subparts: list[dict[str, Any]],
    student_selection: dict[str, Any],
) -> dict[str, Any]:
    normalized_selection = {
        compact_text(key).lower(): compact_text(value)
        for key, value in (student_selection or {}).items()
        if compact_text(key)
    }
    subpart_results: list[dict[str, Any]] = []
    incorrect_items: list[str] = []
    extracted_lines: list[str] = []
    total_score = 0
    max_score = 0

    for subpart in subparts or []:
        label = compact_text(subpart.get("label")).lower()
        marks = clamp_score(subpart.get("marks"), 100) or 1
        max_score += marks
        expected_answer = compact_text(subpart.get("answer"))
        selected_answer = normalized_selection.get(label, "")
        extracted_lines.append(f"{label}) {selected_answer or 'No answer'}")
        is_correct = bool(selected_answer) and selected_answer.lower() == expected_answer.lower()
        marks_awarded = marks if is_correct else 0
        total_score += marks_awarded
        if not is_correct:
            incorrect_items.append(label)
        explanation = compact_text(subpart.get("explanation"))
        subpart_results.append(
            {
                "label": label,
                "marks": marks,
                "marks_awarded": marks_awarded,
                "student_answer": selected_answer,
                "expected_answer": expected_answer,
                "is_correct": is_correct,
                "feedback": (
                    "Correct."
                    if is_correct
                    else compact_text(
                        f"{expected_answer} was correct. {explanation}".strip(),
                        f"{expected_answer} was correct.",
                    )
                ),
            }
        )

    if max_score <= 0:
        max_score = sum(int(subpart.get("marks") or 1) for subpart in subparts or []) or 1

    if incorrect_items:
        feedback = (
            f"You scored {total_score}/{max_score}. Review {', '.join(incorrect_items)} "
            f"and check the correct answers shown below."
        )
    else:
        feedback = f"Excellent. You scored {total_score}/{max_score} on this {question_type.replace('_', ' ')} question."

    return {
        "score": total_score,
        "max_score": max_score,
        "extracted_answer": "\n".join(extracted_lines),
        "feedback": feedback,
        "mistakes": [f"{label}) needs correction." for label in incorrect_items],
        "incorrect_items": incorrect_items,
        "subpart_results": subpart_results,
    }


@app.get("/health")
async def health_check():
    return {"status": "ok"}


@app.post("/auth/request-code")
async def request_login_code(payload: RequestCodeRequest, request: Request):
    started_at = utc_now()
    email = validate_email_address(payload.email)
    ensure_user_account_is_active(email)
    code = create_login_code(email)
    await asyncio.to_thread(send_verification_email, email, code)
    record_audit_log(
        action="auth.code.request",
        email=email,
        request=request,
        resource_type="auth",
        resource_name="email-code",
        duration_ms=int((utc_now() - started_at).total_seconds() * 1000),
    )
    return {"message": "Verification code sent.", "email": email}


@app.post("/auth/verify-code")
async def verify_login(payload: VerifyCodeRequest, request: Request):
    started_at = utc_now()
    email = validate_email_address(payload.email)
    code = payload.code.strip()
    if not code:
        raise HTTPException(status_code=400, detail="Verification code is required.")
    session_token = verify_login_code(email, code)
    record_audit_log(
        action="auth.code.verify",
        email=email,
        request=request,
        resource_type="auth",
        resource_name="email-code",
        duration_ms=int((utc_now() - started_at).total_seconds() * 1000),
    )
    return build_auth_response(email, session_token)


@app.post("/auth/email-password/request-code")
async def request_email_password_login_code(payload: EmailPasswordAuthRequest, request: Request):
    started_at = utc_now()
    email = validate_email_address(payload.email)
    code = await asyncio.to_thread(request_email_password_code, email, payload.password, payload.mode)
    await asyncio.to_thread(send_verification_email, email, code)
    record_audit_log(
        action="auth.email_password.code_request",
        email=email,
        request=request,
        resource_type="auth",
        resource_name=normalize_email_password_auth_mode(payload.mode),
        duration_ms=int((utc_now() - started_at).total_seconds() * 1000),
    )
    return {"message": "Verification code sent.", "email": email, "mode": normalize_email_password_auth_mode(payload.mode)}


@app.post("/auth/email-password/register/request-code")
async def request_email_password_registration_code_route(payload: RequestCodeRequest, request: Request):
    started_at = utc_now()
    email = validate_email_address(payload.email)
    code = await asyncio.to_thread(request_email_password_registration_code, email)
    await asyncio.to_thread(send_verification_email, email, code)
    record_audit_log(
        action="auth.email_password.register_code_request",
        email=email,
        request=request,
        resource_type="auth",
        resource_name="user-register",
        duration_ms=int((utc_now() - started_at).total_seconds() * 1000),
    )
    return {"message": "Verification code sent.", "email": email}


@app.post("/auth/email-password/register/verify-code")
async def verify_email_password_registration_code_route(payload: VerifyCodeRequest, request: Request):
    started_at = utc_now()
    email = validate_email_address(payload.email)
    code = payload.code.strip()
    if not code:
        raise HTTPException(status_code=400, detail="Verification code is required.")
    registration_token = await asyncio.to_thread(verify_email_password_registration_code, email, code)
    record_audit_log(
        action="auth.email_password.register_code_verify",
        email=email,
        request=request,
        resource_type="auth",
        resource_name="user-register",
        duration_ms=int((utc_now() - started_at).total_seconds() * 1000),
    )
    return {"message": "Email verified.", "email": email, "registration_token": registration_token}


@app.post("/auth/email-password/register/complete")
async def complete_email_password_registration_route(payload: EmailPasswordRegistrationCompleteRequest, request: Request):
    started_at = utc_now()
    email = validate_email_address(payload.email)
    session_token = await asyncio.to_thread(
        complete_email_password_registration,
        email,
        payload.registration_token,
        payload.password,
    )
    record_audit_log(
        action="auth.email_password.register_complete",
        email=email,
        request=request,
        resource_type="auth",
        resource_name="user-register",
        duration_ms=int((utc_now() - started_at).total_seconds() * 1000),
    )
    return build_auth_response(email, session_token)


@app.post("/auth/email-password/login")
async def login_with_email_password_route(payload: EmailPasswordAuthRequest, request: Request):
    started_at = utc_now()
    email = validate_email_address(payload.email)
    if normalize_email_password_auth_mode(payload.mode) != "login":
        raise HTTPException(status_code=400, detail="Direct sign-in is only available for login mode.")
    ip_address = get_client_ip(request)

    if is_admin_email(email):
        ensure_admin_login_allowed(email, ip_address)

    try:
        session_token = await asyncio.to_thread(login_with_email_password, email, payload.password)
    except HTTPException as exc:
        if is_admin_email(email):
            failure_state = record_admin_login_failure(email, ip_address)
            locked_until = failure_state.get("locked_until")
            failure_count = int(failure_state.get("failure_count") or 0)
            if locked_until:
                detail = (
                    "Admin login has been locked after 3 failed attempts. "
                    f"Try again after {locked_until.astimezone(timezone.utc).strftime('%Y-%m-%d %H:%M UTC')}."
                )
                status_code = 429
            else:
                remaining = max(ADMIN_LOGIN_MAX_ATTEMPTS - failure_count, 0)
                detail = f"Email or password is incorrect. {remaining} admin attempt{'s' if remaining != 1 else ''} remaining."
                status_code = 400
            record_audit_log(
                action="auth.email_password.login",
                status="failed",
                email=email,
                request=request,
                resource_type="auth",
                resource_name="admin-login",
                duration_ms=int((utc_now() - started_at).total_seconds() * 1000),
                metadata={"reason": detail},
            )
            raise HTTPException(status_code=status_code, detail=detail) from exc
        record_audit_log(
            action="auth.email_password.login",
            status="failed",
            email=email,
            request=request,
            resource_type="auth",
            resource_name="user-login",
            duration_ms=int((utc_now() - started_at).total_seconds() * 1000),
            metadata={"reason": exc.detail},
        )
        raise

    if is_admin_email(email):
        clear_admin_login_attempts(email, ip_address)

    record_audit_log(
        action="auth.email_password.login",
        email=email,
        request=request,
        resource_type="auth",
        resource_name="admin-login" if is_admin_email(email) else "user-login",
        duration_ms=int((utc_now() - started_at).total_seconds() * 1000),
    )
    return build_auth_response(email, session_token)


@app.post("/auth/email-password/register")
async def register_with_email_password_route(payload: EmailPasswordAuthRequest, request: Request):
    started_at = utc_now()
    email = validate_email_address(payload.email)
    if normalize_email_password_auth_mode(payload.mode) != "register":
        raise HTTPException(status_code=400, detail="Direct account creation is only available for register mode.")

    session_token = await asyncio.to_thread(register_with_email_password, email, payload.password)
    record_audit_log(
        action="auth.email_password.register",
        email=email,
        request=request,
        resource_type="auth",
        resource_name="user-register",
        duration_ms=int((utc_now() - started_at).total_seconds() * 1000),
    )
    return build_auth_response(email, session_token)


@app.post("/auth/email-password/verify-code")
async def verify_email_password_login(payload: EmailPasswordVerifyRequest, request: Request):
    started_at = utc_now()
    email = validate_email_address(payload.email)
    code = payload.code.strip()
    if not code:
        raise HTTPException(status_code=400, detail="Verification code is required.")
    session_token = await asyncio.to_thread(verify_email_password_auth, email, payload.password, code, payload.mode)
    record_audit_log(
        action="auth.email_password.verify_code",
        email=email,
        request=request,
        resource_type="auth",
        resource_name=normalize_email_password_auth_mode(payload.mode),
        duration_ms=int((utc_now() - started_at).total_seconds() * 1000),
    )
    return build_auth_response(email, session_token)


@app.post("/auth/google")
async def google_login(payload: GoogleAuthRequest, request: Request):
    started_at = utc_now()
    credential = payload.credential.strip()
    if not credential:
        raise HTTPException(status_code=400, detail="Google credential is required.")
    session_token, email = await asyncio.to_thread(create_session_from_google_credential, credential)
    record_audit_log(
        action="auth.google.login",
        email=email,
        request=request,
        resource_type="auth",
        resource_name="google",
        duration_ms=int((utc_now() - started_at).total_seconds() * 1000),
    )
    return build_auth_response(email, session_token)


@app.post("/auth/apple")
async def apple_login(payload: AppleAuthRequest, request: Request):
    started_at = utc_now()
    if not compact_text(payload.authorization_code) and not compact_text(payload.id_token):
        raise HTTPException(status_code=400, detail="Apple sign-in did not return a usable credential.")
    session_token, email = await asyncio.to_thread(create_session_from_apple_auth, payload)
    record_audit_log(
        action="auth.apple.login",
        email=email,
        request=request,
        resource_type="auth",
        resource_name="apple",
        duration_ms=int((utc_now() - started_at).total_seconds() * 1000),
    )
    return build_auth_response(email, session_token)


@app.get("/auth/me")
async def auth_me(request: Request, authorization: str | None = Header(None)):
    started_at = utc_now()
    token = get_authorization_token(authorization)
    context = get_session_context(token)
    if not context:
        raise HTTPException(status_code=401, detail="Your session is invalid or has expired.")

    refreshed_token = ""
    if should_refresh_session_token(token):
        refreshed_token = create_session(context["email"], session_mode=context["mode"])
    record_audit_log(
        action="auth.session.resume",
        email=context["email"],
        request=request,
        resource_type="auth",
        resource_name=context["mode"],
        duration_ms=int((utc_now() - started_at).total_seconds() * 1000),
        metadata={"token_refreshed": bool(refreshed_token)},
    )
    response = build_auth_response(context["email"], refreshed_token or token)
    response["token"] = refreshed_token
    return response


@app.post("/auth/select-mode")
async def select_auth_mode(
    payload: SessionModeRequest,
    request: Request,
    authorization: str | None = Header(None),
):
    started_at = utc_now()
    token = get_authorization_token(authorization)
    context = get_session_context(token)
    if not context:
        raise HTTPException(status_code=401, detail="Your session is invalid or has expired.")

    next_mode = normalize_session_mode(payload.mode, context["email"])
    session_token = create_session(context["email"], session_mode=next_mode)
    record_audit_log(
        action="auth.mode.select",
        email=context["email"],
        request=request,
        resource_type="auth",
        resource_name=next_mode,
        duration_ms=int((utc_now() - started_at).total_seconds() * 1000),
    )
    return build_auth_response(context["email"], session_token)


@app.post("/auth/logout")
async def logout(request: Request, authorization: str | None = Header(None)):
    token = get_authorization_token(authorization)
    context = get_session_context(token)
    revoke_session(token)
    if context:
        record_audit_log(
            action="auth.logout",
            email=context["email"],
            request=request,
            resource_type="auth",
            resource_name=context["mode"],
        )
    return {"message": "Logged out."}


@app.post("/support/contact")
async def submit_support_request(
    payload: SupportMessageRequest,
    request: Request,
    current_user: str = Depends(require_authenticated_user),
):
    started_at = utc_now()
    message = (payload.message or "").strip()
    page = compact_text(payload.page, "unknown-page")[:120]
    if not message:
        raise HTTPException(status_code=400, detail="Support message cannot be empty.")

    support_message = await asyncio.to_thread(
        create_or_load_support_message,
        current_user,
        message,
        page,
        payload.client_request_id,
    )
    saved_message = compact_text(support_message.get("message"), message)
    saved_page = compact_text(support_message.get("page"), page)[:120]
    delivery_status = compact_text(support_message.get("email_delivery_status"), "queued")
    email_error = compact_text(support_message.get("email_error"))
    if delivery_status != "sent":
        delivery_result = await deliver_support_message_email(
            support_message["id"],
            current_user,
            saved_message,
            saved_page,
        )
        delivery_status = compact_text(delivery_result.get("status"), delivery_status)
        email_error = compact_text(delivery_result.get("error"), email_error)

    response_message = (
        "Support message sent."
        if delivery_status == "sent"
        else "Support message was saved, but the support inbox email could not be sent right now."
    )
    record_audit_log(
        action="support.contact",
        status="success" if delivery_status == "sent" else "error",
        email=current_user,
        request=request,
        resource_type="support",
        resource_name=page or "unknown-page",
        duration_ms=int((utc_now() - started_at).total_seconds() * 1000),
        metadata={
            "support_message_id": support_message.get("id", ""),
            "client_request_id": support_message.get("client_request_id", ""),
            "delivery_status": delivery_status,
            "email_error": email_error,
        },
    )
    return {
        "message": response_message,
        "delivery_status": delivery_status,
        "support_message_id": support_message.get("id", ""),
        "email_error": email_error,
    }


@app.get("/history")
async def get_study_history(current_user: str = Depends(require_authenticated_user)):
    return {"items": get_history_items_for_user(current_user)}


@app.put("/history")
async def sync_study_history(
    payload: HistorySyncRequest,
    current_user: str = Depends(require_authenticated_user),
):
    items = payload.items if isinstance(payload.items, list) else []
    return {"items": replace_history_items_for_user(current_user, items)}


@app.get("/admin/dashboard")
async def get_admin_dashboard(
    request: Request,
    current_admin: str = Depends(require_admin_user),
):
    started_at = utc_now()
    snapshot = build_admin_dashboard_snapshot()
    record_audit_log(
        action="admin.dashboard.view",
        email=current_admin,
        request=request,
        resource_type="admin",
        resource_name="dashboard",
        duration_ms=int((utc_now() - started_at).total_seconds() * 1000),
    )
    return snapshot


@app.post("/admin/users/{target_email}/status")
async def update_admin_user_status(
    target_email: str,
    payload: AdminUserStatusRequest,
    request: Request,
    current_admin: str = Depends(require_admin_user),
):
    normalized_email = validate_email_address(target_email)
    next_status = normalize_user_account_status(payload.status)
    set_user_account_status(normalized_email, next_status, current_admin)
    if next_status == "suspended":
        revoke_all_sessions_for_user(normalized_email)
    record_audit_log(
        action="admin.user.status",
        email=current_admin,
        request=request,
        resource_type="admin",
        resource_name=normalized_email,
        metadata={"status": next_status},
    )
    return {"email": normalized_email, "status": next_status}


@app.post("/admin/users/{target_email}/force-logout")
async def force_logout_admin_user(
    target_email: str,
    request: Request,
    current_admin: str = Depends(require_admin_user),
):
    normalized_email = validate_email_address(target_email)
    revoke_all_sessions_for_user(normalized_email)
    record_audit_log(
        action="admin.user.force_logout",
        email=current_admin,
        request=request,
        resource_type="admin",
        resource_name=normalized_email,
    )
    return {"email": normalized_email, "message": "User sessions revoked."}


@app.get("/jobs/{job_id}")
async def get_job(job_id: str, current_user: str = Depends(require_authenticated_user)):
    job = jobs.get(job_id)
    if not job:
        raise HTTPException(status_code=404, detail="Job not found.")
    if job.get("owner_email") and job["owner_email"] != current_user:
        raise HTTPException(status_code=404, detail="Job not found.")
    return serialize_job(job)


@app.get("/jobs/{job_id}/podcast-audio/{segment_index}")
async def get_podcast_audio_segment(
    job_id: str,
    segment_index: int,
    current_user: str = Depends(require_authenticated_user),
):
    job = jobs.get(job_id)
    if not job:
        raise HTTPException(status_code=404, detail="Job not found.")
    if job.get("owner_email") and job["owner_email"] != current_user:
        raise HTTPException(status_code=404, detail="Job not found.")

    audio_files = job.get("_podcast_audio_files") or []
    if segment_index < 0 or segment_index >= len(audio_files):
        raise HTTPException(status_code=404, detail="Podcast audio segment not found.")

    file_path = Path(audio_files[segment_index])
    if not file_path.exists():
        raise HTTPException(status_code=404, detail="Podcast audio segment is no longer available.")

    return Response(content=file_path.read_bytes(), media_type="audio/mpeg")


@app.get("/jobs/{job_id}/podcast-download")
async def download_podcast_audio(
    job_id: str,
    current_user: str = Depends(require_authenticated_user),
):
    job = jobs.get(job_id)
    if not job:
        raise HTTPException(status_code=404, detail="Job not found.")
    if job.get("owner_email") and job["owner_email"] != current_user:
        raise HTTPException(status_code=404, detail="Job not found.")

    file_name = job.get("_podcast_download_file") or ""
    file_path = Path(file_name)
    if not file_path.exists():
        raise HTTPException(status_code=404, detail="The downloadable podcast audio is no longer available.")

    return Response(
        content=file_path.read_bytes(),
        media_type="audio/mpeg",
        headers={
            "Content-Disposition": f'attachment; filename="{sanitize_download_filename(job.get("podcast_title") or "lecture-podcast")}.mp3"',
        },
    )


@app.get("/jobs/{job_id}/presentation-download")
async def download_presentation_file(
    job_id: str,
    current_user: str = Depends(require_authenticated_user),
):
    job = jobs.get(job_id)
    if not job:
        raise HTTPException(status_code=404, detail="Job not found.")
    if job.get("owner_email") and job["owner_email"] != current_user:
        raise HTTPException(status_code=404, detail="Job not found.")

    file_name = job.get("_presentation_download_file") or ""
    file_path = Path(file_name)
    if not file_path.exists():
        raise HTTPException(status_code=404, detail="The downloadable PowerPoint file is no longer available.")

    return Response(
        content=file_path.read_bytes(),
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers={
            "Content-Disposition": f'attachment; filename="{sanitize_download_filename(job.get("presentation_title") or "lecture-presentation")}.pptx"',
        },
    )


async def save_upload_to_disk(file: UploadFile, job_id: str | None = None) -> Path:
    suffix = Path(file.filename or "audio.wav").suffix or ".wav"

    def _write_file() -> Path:
        total_written = 0
        update_points = {
            1 * 1024 * 1024: 2,
            5 * 1024 * 1024: 4,
            10 * 1024 * 1024: 6,
            20 * 1024 * 1024: 8,
        }
        with tempfile.NamedTemporaryFile(
            delete=False,
            suffix=suffix,
            prefix="lecture_",
            dir=UPLOAD_DIR,
        ) as temp_file:
            file.file.seek(0)
            while chunk := file.file.read(1024 * 1024):
                temp_file.write(chunk)
                total_written += len(chunk)
                if job_id:
                    for threshold, progress in update_points.items():
                        if total_written >= threshold:
                            update_job(
                                job_id,
                                status="processing",
                                stage="Uploading lecture file",
                                progress=progress,
                            )
            return Path(temp_file.name)

    return await asyncio.to_thread(_write_file)


def get_file_size(file_path: Path) -> int:
    return file_path.stat().st_size


def normalize_video_url(value: str) -> str:
    cleaned = (value or "").strip()
    parsed = urlparse(cleaned)
    if parsed.scheme not in {"http", "https"} or not parsed.netloc:
        raise ValueError("Paste a full video link that starts with http:// or https://.")
    return cleaned


def resolve_youtube_cookiefile() -> Path | None:
    if YOUTUBE_COOKIES_FILE:
        cookie_path = Path(YOUTUBE_COOKIES_FILE)
        if cookie_path.exists() and cookie_path.is_file():
            return cookie_path
        logger.warning("Configured YOUTUBE_COOKIES_FILE does not exist: %s", cookie_path)

    if not YOUTUBE_COOKIES_TXT:
        return None

    cookie_text = YOUTUBE_COOKIES_TXT.replace("\\n", "\n").strip()
    if not cookie_text:
        return None

    cookie_path = UPLOAD_DIR / "youtube_cookies.txt"
    try:
        cookie_path.write_text(f"{cookie_text}\n", encoding="utf-8")
    except OSError as exc:
        logger.warning("Could not write YouTube cookies file: %s", exc)
        return None
    return cookie_path


def has_youtube_cookie_source() -> bool:
    return bool(YOUTUBE_COOKIES_FILE or YOUTUBE_COOKIES_TXT)


def load_youtube_request_cookies() -> dict[str, str]:
    cookies = {
        "CONSENT": "YES+cb.20210328-17-p0.en+FX+700",
        "PREF": "hl=en&gl=US",
    }
    cookie_path = resolve_youtube_cookiefile()
    if not cookie_path:
        return cookies

    try:
        cookie_jar = MozillaCookieJar()
        cookie_jar.load(str(cookie_path), ignore_discard=True, ignore_expires=True)
    except Exception as exc:
        logger.warning("Could not load YouTube cookies from %s: %s", cookie_path, exc)
        return cookies

    for cookie in cookie_jar:
        if "youtube.com" in (cookie.domain or "").lower() or "google.com" in (cookie.domain or "").lower():
            cookies[cookie.name] = cookie.value
    return cookies


def build_youtube_request_headers() -> dict[str, str]:
    return {
        "User-Agent": YOUTUBE_USER_AGENT,
        "Accept-Language": "en-US,en;q=0.9",
        "Referer": "https://www.youtube.com/",
        "Origin": "https://www.youtube.com",
    }


def build_webshare_proxy_url() -> str:
    if not (YOUTUBE_WEBSHARE_PROXY_USERNAME and YOUTUBE_WEBSHARE_PROXY_PASSWORD):
        return ""

    host = YOUTUBE_WEBSHARE_PROXY_HOST or "p.webshare.io"
    port = YOUTUBE_WEBSHARE_PROXY_PORT or "80"
    username = quote(YOUTUBE_WEBSHARE_PROXY_USERNAME, safe="")
    password = quote(YOUTUBE_WEBSHARE_PROXY_PASSWORD, safe="")
    return f"http://{username}:{password}@{host}:{port}"


def has_youtube_proxy_source() -> bool:
    return bool(YOUTUBE_PROXY_HTTP_URL or YOUTUBE_PROXY_HTTPS_URL or build_webshare_proxy_url())


def get_youtube_proxy_url(*, prefer_https: bool = True) -> str:
    if prefer_https and YOUTUBE_PROXY_HTTPS_URL:
        return YOUTUBE_PROXY_HTTPS_URL
    if YOUTUBE_PROXY_HTTP_URL:
        return YOUTUBE_PROXY_HTTP_URL
    if not prefer_https and YOUTUBE_PROXY_HTTPS_URL:
        return YOUTUBE_PROXY_HTTPS_URL
    return build_webshare_proxy_url()


def build_youtube_request_proxies(*, use_proxy: bool = True) -> dict[str, str]:
    if not use_proxy:
        return {}

    proxies: dict[str, str] = {}
    fallback_proxy_url = build_webshare_proxy_url()
    if YOUTUBE_PROXY_HTTP_URL:
        proxies["http"] = YOUTUBE_PROXY_HTTP_URL
    elif fallback_proxy_url:
        proxies["http"] = fallback_proxy_url
    if YOUTUBE_PROXY_HTTPS_URL:
        proxies["https"] = YOUTUBE_PROXY_HTTPS_URL
    elif YOUTUBE_PROXY_HTTP_URL:
        proxies["https"] = YOUTUBE_PROXY_HTTP_URL
    elif fallback_proxy_url:
        proxies["https"] = fallback_proxy_url
    return proxies


def iter_youtube_proxy_attempts() -> list[bool]:
    return [True, False] if has_youtube_proxy_source() else [False]


def is_proxy_auth_error(exc: Exception) -> bool:
    lowered = str(exc or "").strip().lower()
    return (
        "response 407" in lowered
        or "proxy authentication required" in lowered
        or ("connect tunnel failed" in lowered and "407" in lowered)
    )


def create_youtube_requests_session(*, use_proxy: bool = True) -> Any:
    if curl_requests is not None:
        try:
            session = curl_requests.Session(impersonate="chrome")
        except Exception:
            session = curl_requests.Session()
    else:
        session = requests.Session()
    try:
        session.trust_env = False
    except Exception:
        pass
    session.headers.update(build_youtube_request_headers())
    session.cookies.update(load_youtube_request_cookies())
    proxy_map = build_youtube_request_proxies(use_proxy=use_proxy)
    if proxy_map:
        session.proxies.update(proxy_map)
    return session


def build_ytdlp_options(
    *,
    output_template: str | None = None,
    progress_hook: Any = None,
    skip_download: bool = False,
    use_cookiefile: bool = True,
    use_impersonation: bool = True,
    use_proxy: bool = True,
) -> dict[str, Any]:
    options: dict[str, Any] = {
        "noplaylist": True,
        "quiet": True,
        "no_warnings": True,
        "restrictfilenames": True,
        "geo_bypass": True,
        "retries": 2,
        "fragment_retries": 2,
        "extractor_retries": 2,
        "socket_timeout": 20,
        "http_headers": build_youtube_request_headers(),
        "extractor_args": {
            "youtube": {
                "player_client": list(YTDLP_YOUTUBE_PLAYER_CLIENTS),
                "skip": ["translated_subs"],
            }
        },
    }
    if use_impersonation and YTDLP_IMPERSONATE_TARGET:
        options["impersonate"] = YTDLP_IMPERSONATE_TARGET
    if output_template:
        options["outtmpl"] = output_template
    if progress_hook:
        options["progress_hooks"] = [progress_hook]
    if skip_download:
        options["skip_download"] = True
    cookie_path = resolve_youtube_cookiefile()
    if use_cookiefile and cookie_path:
        options["cookiefile"] = str(cookie_path)
    proxy_url = get_youtube_proxy_url() if use_proxy else ""
    if proxy_url:
        options["proxy"] = proxy_url
    elif not use_proxy:
        options["proxy"] = ""
    return options


def is_unavailable_impersonation_target_error(exc: Exception) -> bool:
    lowered = str(exc or "").strip().lower()
    return "impersonate target" in lowered and "not available" in lowered


def iter_ytdlp_option_variants(options: dict[str, Any]):
    yield options
    if options.get("impersonate"):
        fallback_options = dict(options)
        fallback_options.pop("impersonate", None)
        yield fallback_options


def extract_youtube_video_id(video_url: str) -> str | None:
    parsed = urlparse(video_url)
    host = (parsed.netloc or "").lower().split(":")[0]
    if host.startswith("www."):
        host = host[4:]

    path_parts = [part for part in parsed.path.split("/") if part]
    if host == "youtu.be":
        return path_parts[0] if path_parts else None

    if host not in {"youtube.com", "m.youtube.com", "music.youtube.com", "youtube-nocookie.com"}:
        return None

    if parsed.path == "/watch":
        video_id = parse_qs(parsed.query).get("v", [""])[0].strip()
        return video_id or None

    if path_parts and path_parts[0] in {"embed", "shorts", "live", "v"} and len(path_parts) > 1:
        return path_parts[1]

    return None


def captions_to_transcript_text(items: Any) -> str:
    lines: list[str] = []
    for item in items or []:
        text = item.get("text", "") if isinstance(item, dict) else getattr(item, "text", "")
        cleaned = re.sub(r"\s+", " ", str(text or "").replace("\xa0", " ")).strip()
        if cleaned:
            lines.append(cleaned)
    return "\n".join(lines).strip()


def build_youtube_transcript_api() -> Any | None:
    if YouTubeTranscriptApi is None:
        return None

    if (
        WebshareProxyConfig is not None
        and YOUTUBE_WEBSHARE_PROXY_USERNAME
        and YOUTUBE_WEBSHARE_PROXY_PASSWORD
    ):
        proxy_kwargs: dict[str, Any] = {
            "proxy_username": YOUTUBE_WEBSHARE_PROXY_USERNAME,
            "proxy_password": YOUTUBE_WEBSHARE_PROXY_PASSWORD,
        }
        if YOUTUBE_WEBSHARE_PROXY_LOCATIONS:
            proxy_kwargs["filter_ip_locations"] = list(YOUTUBE_WEBSHARE_PROXY_LOCATIONS)
        logger.info("Using Webshare proxy configuration for YouTube transcript requests")
        return YouTubeTranscriptApi(proxy_config=WebshareProxyConfig(**proxy_kwargs))

    if GenericProxyConfig is not None and (YOUTUBE_PROXY_HTTP_URL or YOUTUBE_PROXY_HTTPS_URL):
        proxy_kwargs = {}
        if YOUTUBE_PROXY_HTTP_URL:
            proxy_kwargs["http_url"] = YOUTUBE_PROXY_HTTP_URL
        if YOUTUBE_PROXY_HTTPS_URL:
            proxy_kwargs["https_url"] = YOUTUBE_PROXY_HTTPS_URL
        logger.info("Using generic proxy configuration for YouTube transcript requests")
        return YouTubeTranscriptApi(proxy_config=GenericProxyConfig(**proxy_kwargs))

    return YouTubeTranscriptApi()


def fetch_youtube_transcript(video_url: str, job_id: str) -> str | None:
    if YouTubeTranscriptApi is None:
        return None

    video_id = extract_youtube_video_id(video_url)
    if not video_id:
        return None

    update_job(job_id, status="processing", stage="Checking YouTube captions", progress=3)

    transcript_api = build_youtube_transcript_api()
    if transcript_api is None:
        return None

    preferred_languages = list(dict.fromkeys(["en", *YOUTUBE_LANGUAGE_PREFERENCES]))
    try:
        transcript_text = captions_to_transcript_text(
            transcript_api.fetch(video_id, languages=preferred_languages)
        )
        if transcript_text:
            update_job(job_id, status="processing", stage="YouTube captions found. Preparing transcript", progress=14)
            return transcript_text
    except Exception as exc:
        logger.info("Could not fetch YouTube transcript directly for %s: %s", video_url, exc)

    try:
        transcript_list = transcript_api.list(video_id)
    except Exception as exc:
        logger.info("Could not read YouTube transcript list for %s: %s", video_url, exc)
        return None

    candidates = []
    for finder_name in ("find_transcript", "find_manually_created_transcript", "find_generated_transcript"):
        finder = getattr(transcript_list, finder_name, None)
        if not callable(finder):
            continue
        try:
            candidates.append(finder(preferred_languages))
        except Exception:
            continue

    try:
        candidates.extend(list(transcript_list))
    except Exception as exc:
        logger.info("Could not iterate YouTube transcripts for %s: %s", video_url, exc)

    seen_candidates: set[tuple[str, bool]] = set()
    for transcript_candidate in candidates:
        candidate_key = (
            str(getattr(transcript_candidate, "language_code", "")),
            bool(getattr(transcript_candidate, "is_generated", False)),
        )
        if candidate_key in seen_candidates:
            continue
        seen_candidates.add(candidate_key)
        try:
            transcript_text = captions_to_transcript_text(transcript_candidate.fetch())
        except Exception as exc:
            logger.info("Could not fetch YouTube transcript candidate for %s: %s", video_url, exc)
            continue
        if transcript_text:
            update_job(job_id, status="processing", stage="YouTube captions found. Preparing transcript", progress=14)
            return transcript_text

    return None


def can_process_video_url(video_url: str) -> bool:
    return yt_dlp is not None or (YouTubeTranscriptApi is not None and extract_youtube_video_id(video_url) is not None)


def subtitle_body_to_text(raw_text: str) -> str:
    lines: list[str] = []
    for raw_line in (raw_text or "").splitlines():
        line = raw_line.strip()
        if (
            not line
            or line.upper() == "WEBVTT"
            or line.startswith("NOTE")
            or line.startswith("Kind:")
            or line.startswith("Language:")
            or "-->" in line
            or line.isdigit()
        ):
            continue
        line = html.unescape(re.sub(r"<[^>]+>", " ", line))
        line = re.sub(r"\s+", " ", line).strip()
        if line and (not lines or line != lines[-1]):
            lines.append(line)
    return "\n".join(lines).strip()


def read_subtitle_file(file_path: Path) -> str:
    try:
        raw_text = file_path.read_text(encoding="utf-8")
    except UnicodeDecodeError:
        raw_text = file_path.read_text(encoding="utf-8-sig", errors="ignore")
    return subtitle_body_to_text(raw_text)


def subtitle_file_priority(file_path: Path) -> tuple[int, int, str]:
    name = file_path.name.lower()
    language_score = 3
    if ".en." in name or ".en-" in name or name.endswith(".en.vtt") or name.endswith(".en.srt"):
        language_score = 0
    elif ".en-us." in name or ".en-gb." in name or ".en-orig." in name or ".orig." in name:
        language_score = 1
    elif ".a." in name:
        language_score = 2

    extension_score = 0 if file_path.suffix.lower() == ".vtt" else 1
    return (language_score, extension_score, name)


def cleanup_caption_files(prefix: str):
    for path in UPLOAD_DIR.glob(f"{prefix}*"):
        if path.is_file():
            try:
                path.unlink()
            except OSError:
                logger.warning("Could not delete subtitle artifact: %s", path)


def cleanup_download_artifacts(prefix: str):
    for path in UPLOAD_DIR.glob(f"{prefix}*"):
        if path.is_file():
            try:
                path.unlink()
            except OSError:
                logger.warning("Could not delete download artifact: %s", path)


def extract_json_object_from_text(source: str, marker: str) -> dict[str, Any] | None:
    marker_index = source.find(marker)
    if marker_index == -1:
        return None

    start_index = source.find("{", marker_index)
    if start_index == -1:
        return None

    depth = 0
    in_string = False
    escaped = False

    for index in range(start_index, len(source)):
        char = source[index]
        if in_string:
            if escaped:
                escaped = False
            elif char == "\\":
                escaped = True
            elif char == '"':
                in_string = False
            continue

        if char == '"':
            in_string = True
        elif char == "{":
            depth += 1
        elif char == "}":
            depth -= 1
            if depth == 0:
                try:
                    return json.loads(source[start_index:index + 1])
                except json.JSONDecodeError:
                    return None

    return None


def choose_youtube_caption_track(caption_tracks: list[dict[str, Any]]) -> dict[str, Any] | None:
    if not caption_tracks:
        return None

    def track_priority(track: dict[str, Any]) -> tuple[int, int, str]:
        language_code = str(track.get("languageCode", "")).lower()
        kind = str(track.get("kind", "")).lower()
        language_score = 3
        if language_code in {language.lower() for language in YOUTUBE_LANGUAGE_PREFERENCES}:
            language_score = 0
        elif language_code.startswith("en"):
            language_score = 1
        elif language_code.endswith("-orig") or language_code == "orig":
            language_score = 2

        kind_score = 0 if kind != "asr" else 1
        return (language_score, kind_score, language_code)

    return sorted(caption_tracks, key=track_priority)[0]


def fetch_youtube_watch_page_captions(video_url: str, job_id: str) -> str | None:
    video_id = extract_youtube_video_id(video_url)
    if not video_id:
        return None

    update_job(job_id, status="processing", stage="Checking YouTube watch-page captions", progress=7)
    watch_urls = [
        f"https://www.youtube.com/watch?v={video_id}&hl=en&bpctr=9999999999&has_verified=1",
        f"https://m.youtube.com/watch?v={video_id}&hl=en&bpctr=9999999999&has_verified=1",
        f"https://www.youtube.com/embed/{video_id}?hl=en",
        f"https://www.youtube-nocookie.com/embed/{video_id}?hl=en",
    ]

    for use_proxy in iter_youtube_proxy_attempts():
        if not use_proxy:
            update_job(job_id, status="processing", stage="Retrying watch-page captions without proxy", progress=8)
        with create_youtube_requests_session(use_proxy=use_proxy) as session:
            for watch_url in watch_urls:
                try:
                    response = session.get(watch_url, timeout=20)
                    response.raise_for_status()
                except Exception as exc:
                    logger.info(
                        "Could not fetch YouTube watch page for %s via %s (use_proxy=%s): %s",
                        video_url,
                        watch_url,
                        use_proxy,
                        exc,
                    )
                    continue

                player_response = None
                for marker in (
                    "var ytInitialPlayerResponse = ",
                    "ytInitialPlayerResponse = ",
                    'window["ytInitialPlayerResponse"] = ',
                    "window['ytInitialPlayerResponse'] = ",
                ):
                    player_response = extract_json_object_from_text(response.text, marker)
                    if player_response:
                        break

                caption_tracks = (
                    player_response.get("captions", {})
                    .get("playerCaptionsTracklistRenderer", {})
                    .get("captionTracks", [])
                    if isinstance(player_response, dict)
                    else []
                )
                caption_track = choose_youtube_caption_track(caption_tracks)
                base_url = caption_track.get("baseUrl", "") if isinstance(caption_track, dict) else ""
                if not base_url:
                    continue

                caption_url = base_url if "fmt=" in base_url else f"{base_url}&fmt=vtt"
                try:
                    caption_response = session.get(caption_url, timeout=20)
                    caption_response.raise_for_status()
                except Exception as exc:
                    logger.info(
                        "Could not fetch YouTube caption track for %s via %s (use_proxy=%s): %s",
                        video_url,
                        watch_url,
                        use_proxy,
                        exc,
                    )
                    continue

                transcript_text = subtitle_body_to_text(caption_response.text)
                if transcript_text:
                    update_job(job_id, status="processing", stage="YouTube watch-page captions found. Preparing transcript", progress=15)
                    return transcript_text
    return None


def timedtext_body_to_text(raw_text: str) -> str:
    cleaned = compact_text(raw_text)
    if not cleaned:
        return ""
    if cleaned.lstrip().startswith("WEBVTT"):
        return subtitle_body_to_text(cleaned)

    try:
        root = ET.fromstring(cleaned)
    except ET.ParseError:
        return ""

    lines: list[str] = []
    for node in root.findall(".//text"):
        text = html.unescape("".join(node.itertext()).replace("\xa0", " "))
        normalized = re.sub(r"\s+", " ", text).strip()
        if normalized and (not lines or normalized != lines[-1]):
            lines.append(normalized)
    return "\n".join(lines).strip()


def extract_youtube_timedtext_tracks(raw_text: str) -> list[dict[str, str]]:
    cleaned = compact_text(raw_text)
    if not cleaned:
        return []
    try:
        root = ET.fromstring(cleaned)
    except ET.ParseError:
        return []

    tracks: list[dict[str, str]] = []
    for node in root.findall(".//track"):
        language_code = compact_text(node.get("lang_code") or node.get("lang-code"))
        if not language_code:
            continue
        tracks.append(
            {
                "lang_code": language_code,
                "kind": compact_text(node.get("kind")),
                "name": compact_text(node.get("name")),
            }
        )
    return tracks


def choose_youtube_timedtext_track(tracks: list[dict[str, str]]) -> dict[str, str] | None:
    if not tracks:
        return None

    preferred_languages = {language.lower(): index for index, language in enumerate(YOUTUBE_LANGUAGE_PREFERENCES)}

    def track_priority(track: dict[str, str]) -> tuple[int, int, str]:
        language_code = compact_text(track.get("lang_code")).lower()
        kind = compact_text(track.get("kind")).lower()
        language_score = preferred_languages.get(language_code, 10)
        if language_score == 10 and language_code.startswith("en"):
            language_score = 3
        kind_score = 1 if kind == "asr" else 0
        return (language_score, kind_score, language_code)

    return sorted(tracks, key=track_priority)[0]


def build_youtube_timedtext_urls(video_id: str, track: dict[str, str] | None = None) -> list[str]:
    endpoints = (
        "https://www.youtube.com/api/timedtext",
        "https://video.google.com/timedtext",
    )
    candidate_params: list[dict[str, str]] = []

    if track:
        candidate_params.append(
            {
                "lang": compact_text(track.get("lang_code")),
                "kind": compact_text(track.get("kind")),
                "name": compact_text(track.get("name")),
            }
        )

    for language_code in dict.fromkeys(["en", *YOUTUBE_LANGUAGE_PREFERENCES]):
        candidate_params.append({"lang": language_code, "kind": "", "name": ""})
        candidate_params.append({"lang": language_code, "kind": "asr", "name": ""})

    urls: list[str] = []
    seen: set[str] = set()
    for endpoint in endpoints:
        for params in candidate_params:
            language_code = compact_text(params.get("lang"))
            if not language_code:
                continue
            pieces = [f"v={quote(video_id)}", f"lang={quote(language_code)}", "fmt=vtt"]
            if params.get("kind"):
                pieces.append(f"kind={quote(params['kind'])}")
            if params.get("name"):
                pieces.append(f"name={quote(params['name'])}")
            url = f"{endpoint}?{'&'.join(pieces)}"
            if url in seen:
                continue
            seen.add(url)
            urls.append(url)
    return urls


def fetch_youtube_timedtext_captions(video_url: str, job_id: str) -> str | None:
    video_id = extract_youtube_video_id(video_url)
    if not video_id:
        return None

    update_job(job_id, status="processing", stage="Checking direct YouTube caption endpoints", progress=10)
    track_list_urls = [
        f"https://www.youtube.com/api/timedtext?v={quote(video_id)}&type=list",
        f"https://video.google.com/timedtext?v={quote(video_id)}&type=list",
    ]

    for use_proxy in iter_youtube_proxy_attempts():
        if not use_proxy:
            update_job(job_id, status="processing", stage="Retrying direct YouTube captions without proxy", progress=11)
        with create_youtube_requests_session(use_proxy=use_proxy) as session:
            tracks: list[dict[str, str]] = []
            for track_list_url in track_list_urls:
                try:
                    response = session.get(track_list_url, timeout=20)
                    response.raise_for_status()
                except Exception as exc:
                    logger.info(
                        "Could not fetch YouTube timedtext track list for %s via %s (use_proxy=%s): %s",
                        video_url,
                        track_list_url,
                        use_proxy,
                        exc,
                    )
                    continue
                tracks = extract_youtube_timedtext_tracks(response.text)
                if tracks:
                    break

            selected_track = choose_youtube_timedtext_track(tracks)
            for caption_url in build_youtube_timedtext_urls(video_id, selected_track):
                try:
                    response = session.get(caption_url, timeout=20)
                    response.raise_for_status()
                except Exception as exc:
                    logger.info(
                        "Could not fetch YouTube timedtext captions for %s via %s (use_proxy=%s): %s",
                        video_url,
                        caption_url,
                        use_proxy,
                        exc,
                    )
                    continue

                transcript_text = timedtext_body_to_text(response.text)
                if transcript_text:
                    update_job(job_id, status="processing", stage="Direct YouTube captions found. Preparing transcript", progress=16)
                    return transcript_text
    return None


def choose_ytdlp_caption_entry(info: dict[str, Any]) -> dict[str, Any] | None:
    candidates: list[dict[str, Any]] = []
    preferred_languages = tuple(language.lower() for language in YOUTUBE_LANGUAGE_PREFERENCES)

    for source_name in ("subtitles", "automatic_captions"):
        language_map = info.get(source_name, {})
        if not isinstance(language_map, dict):
            continue
        for language_code, entries in language_map.items():
            if not isinstance(entries, list):
                continue
            for entry in entries:
                if not isinstance(entry, dict):
                    continue
                url = compact_text(entry.get("url"))
                if not url:
                    continue
                candidates.append(
                    {
                        "url": url,
                        "language_code": compact_text(language_code).lower(),
                        "ext": compact_text(entry.get("ext")).lower(),
                        "is_auto": source_name == "automatic_captions",
                    }
                )

    if not candidates:
        return None

    def caption_priority(entry: dict[str, Any]) -> tuple[int, int, int, str]:
        language_code = entry["language_code"]
        language_score = 4
        if language_code in preferred_languages:
            language_score = 0
        elif any(language_code.startswith(language) for language in preferred_languages):
            language_score = 1
        elif language_code.startswith("en"):
            language_score = 2
        elif language_code.endswith("-orig") or language_code == "orig":
            language_score = 3

        ext = entry["ext"]
        format_score = 3
        if ext == "vtt":
            format_score = 0
        elif ext in {"srv3", "srv2", "ttml", "xml"}:
            format_score = 1
        elif ext in {"json3", "srt"}:
            format_score = 2

        return (language_score, 1 if entry["is_auto"] else 0, format_score, language_code)

    return sorted(candidates, key=caption_priority)[0]


def fetch_ytdlp_caption_track(video_url: str, job_id: str) -> str | None:
    if yt_dlp is None:
        return None

    normalized_url = normalize_video_url(video_url)
    update_job(job_id, status="processing", stage="Checking yt-dlp caption metadata", progress=11)

    metadata_attempts = [True, False] if has_youtube_cookie_source() else [False]
    info: dict[str, Any] | None = None

    for attempt_index, use_cookiefile in enumerate(metadata_attempts, start=1):
        if attempt_index > 1:
            update_job(job_id, status="processing", stage="Retrying caption metadata without saved cookies", progress=12)

        for use_proxy in iter_youtube_proxy_attempts():
            base_options = build_ytdlp_options(
                skip_download=True,
                use_cookiefile=use_cookiefile,
                use_proxy=use_proxy,
            )
            base_options.update(
                {
                    "writesubtitles": False,
                    "writeautomaticsub": False,
                    "simulate": True,
                }
            )
            if not use_proxy:
                update_job(job_id, status="processing", stage="Retrying caption metadata without proxy", progress=12)

            try:
                last_variant_error: Exception | None = None
                for variant_index, options in enumerate(iter_ytdlp_option_variants(base_options), start=1):
                    if variant_index > 1:
                        update_job(job_id, status="processing", stage="Retrying caption metadata without browser impersonation", progress=12)
                    try:
                        with yt_dlp.YoutubeDL(options) as downloader:
                            info = downloader.extract_info(normalized_url, download=False)
                        last_variant_error = None
                        break
                    except Exception as variant_exc:
                        last_variant_error = variant_exc
                        if variant_index == 1 and is_unavailable_impersonation_target_error(variant_exc):
                            logger.info("Retrying caption metadata without impersonation for %s: %s", video_url, variant_exc)
                            continue
                        raise
                if last_variant_error is not None:
                    raise last_variant_error
                if info:
                    break
            except Exception as exc:
                logger.info(
                    "Could not inspect yt-dlp caption metadata for %s (use_cookiefile=%s, use_proxy=%s): %s",
                    video_url,
                    use_cookiefile,
                    use_proxy,
                    exc,
                )
                info = None
                continue

        if info:
            break

    if not isinstance(info, dict):
        return None

    caption_entry = choose_ytdlp_caption_entry(info)
    if not caption_entry:
        return None

    caption_url = caption_entry["url"]
    if "fmt=" not in caption_url:
        caption_url = f"{caption_url}&fmt=vtt"

    for use_proxy in iter_youtube_proxy_attempts():
        if not use_proxy:
            update_job(job_id, status="processing", stage="Retrying caption URL without proxy", progress=13)
        with create_youtube_requests_session(use_proxy=use_proxy) as session:
            try:
                response = session.get(caption_url, timeout=20)
                response.raise_for_status()
            except Exception as exc:
                logger.info("Could not fetch yt-dlp caption URL for %s (use_proxy=%s): %s", video_url, use_proxy, exc)
                continue

        transcript_text = subtitle_body_to_text(response.text)
        if transcript_text:
            update_job(job_id, status="processing", stage="Caption metadata found. Preparing transcript", progress=16)
            return transcript_text
    return None


def download_subtitles_from_video_url(video_url: str, job_id: str) -> str | None:
    if yt_dlp is None:
        return None

    normalized_url = normalize_video_url(video_url)
    subtitle_prefix = f"captions_{uuid4().hex}"
    output_template = str(UPLOAD_DIR / f"{subtitle_prefix}.%(ext)s")
    update_job(job_id, status="processing", stage="Checking downloadable captions", progress=9)
    subtitle_attempts = [True, False] if has_youtube_cookie_source() else [False]
    for attempt_index, use_cookiefile in enumerate(subtitle_attempts, start=1):
        if attempt_index > 1:
            update_job(job_id, status="processing", stage="Retrying caption download without saved cookies", progress=10)
        for use_proxy in iter_youtube_proxy_attempts():
            cleanup_caption_files(subtitle_prefix)
            base_options = build_ytdlp_options(
                output_template=output_template,
                skip_download=True,
                use_cookiefile=use_cookiefile,
                use_proxy=use_proxy,
            )
            base_options.update({
                "writesubtitles": True,
                "writeautomaticsub": True,
                "subtitleslangs": ["all"],
                "subtitlesformat": "vtt/best",
            })
            if not use_proxy:
                update_job(job_id, status="processing", stage="Retrying caption download without proxy", progress=10)
            try:
                last_variant_error: Exception | None = None
                for variant_index, options in enumerate(iter_ytdlp_option_variants(base_options), start=1):
                    if variant_index > 1:
                        update_job(
                            job_id,
                            status="processing",
                            stage="Retrying caption download without browser impersonation",
                            progress=10,
                        )
                    try:
                        with yt_dlp.YoutubeDL(options) as downloader:
                            downloader.extract_info(normalized_url, download=True)
                        last_variant_error = None
                        break
                    except Exception as variant_exc:
                        last_variant_error = variant_exc
                        if variant_index == 1 and is_unavailable_impersonation_target_error(variant_exc):
                            logger.info("Retrying subtitle download without impersonation for %s: %s", video_url, variant_exc)
                            continue
                        raise
                if last_variant_error is not None:
                    raise last_variant_error
                subtitle_files = sorted(
                    [
                        path
                        for path in UPLOAD_DIR.glob(f"{subtitle_prefix}*")
                        if path.is_file() and path.suffix.lower() in {".vtt", ".srt"}
                    ],
                    key=subtitle_file_priority,
                )
                transcript_text = ""
                for subtitle_file in subtitle_files:
                    transcript_text = read_subtitle_file(subtitle_file)
                    if transcript_text:
                        break
                cleanup_caption_files(subtitle_prefix)
                if transcript_text:
                    update_job(job_id, status="processing", stage="Downloadable captions found. Preparing transcript", progress=16)
                    return transcript_text
            except Exception as exc:
                logger.info(
                    "Could not download subtitles with yt-dlp for %s (use_cookiefile=%s, use_proxy=%s): %s",
                    video_url,
                    use_cookiefile,
                    use_proxy,
                    exc,
                )
                cleanup_caption_files(subtitle_prefix)
                continue

    return None


def download_audio_from_video_url(video_url: str, job_id: str) -> Path:
    if yt_dlp is None:
        raise RuntimeError("Video-link transcription needs yt-dlp on the backend server.")

    normalized_url = normalize_video_url(video_url)
    download_prefix = f"video_{uuid4().hex}"
    output_template = str(UPLOAD_DIR / f"{download_prefix}.%(ext)s")

    def progress_hook(data: dict[str, Any]):
        status = data.get("status")
        if status == "downloading":
            total_bytes = data.get("total_bytes") or data.get("total_bytes_estimate") or 0
            downloaded = data.get("downloaded_bytes") or 0
            if total_bytes:
                progress = min(18, 4 + int((downloaded / total_bytes) * 14))
            else:
                progress = 8
            update_job(job_id, status="processing", stage="Downloading audio from the video link", progress=progress)
        elif status == "finished":
            update_job(job_id, status="processing", stage="Video downloaded. Preparing audio", progress=18)

    update_job(job_id, status="processing", stage="Connecting to the video link", progress=3)
    candidate_paths: list[Path] = []
    last_error = ""
    download_attempts = [
        {
            "format": "bestaudio[ext=m4a]/bestaudio[acodec!=none]/bestaudio/best",
            "drop_extractor_args": False,
            "use_cookiefile": has_youtube_cookie_source(),
        },
        {
            "format": "best[acodec!=none]/best",
            "drop_extractor_args": True,
            "use_cookiefile": has_youtube_cookie_source(),
        },
    ]
    if has_youtube_cookie_source():
        download_attempts.extend([
            {
                "format": "bestaudio[ext=m4a]/bestaudio[acodec!=none]/bestaudio/best",
                "drop_extractor_args": False,
                "use_cookiefile": False,
            },
            {
                "format": "best[acodec!=none]/best",
                "drop_extractor_args": True,
                "use_cookiefile": False,
            },
        ])

    for attempt_index, attempt in enumerate(download_attempts, start=1):
        for use_proxy in iter_youtube_proxy_attempts():
            cleanup_download_artifacts(download_prefix)
            base_options = build_ytdlp_options(
                output_template=output_template,
                progress_hook=progress_hook,
                use_cookiefile=attempt["use_cookiefile"],
                use_proxy=use_proxy,
            )
            base_options["format"] = attempt["format"]
            if attempt["drop_extractor_args"]:
                base_options.pop("extractor_args", None)
                update_job(
                    job_id,
                    status="processing",
                    stage=(
                        "Retrying video download with a fallback format"
                        if attempt["use_cookiefile"]
                        else "Retrying video download without saved cookies"
                    ),
                    progress=4,
                )
            elif not attempt["use_cookiefile"]:
                update_job(
                    job_id,
                    status="processing",
                    stage="Retrying video download without saved cookies",
                    progress=4,
                )
            if not use_proxy:
                update_job(
                    job_id,
                    status="processing",
                    stage="Retrying video download without proxy",
                    progress=4,
                )

            try:
                last_variant_error: Exception | None = None
                for variant_index, options in enumerate(iter_ytdlp_option_variants(base_options), start=1):
                    if variant_index > 1:
                        update_job(
                            job_id,
                            status="processing",
                            stage="Retrying video download without browser impersonation",
                            progress=4,
                        )
                    try:
                        with yt_dlp.YoutubeDL(options) as downloader:
                            info = downloader.extract_info(normalized_url, download=True)
                            requested_downloads = info.get("requested_downloads") or []
                            candidate_paths = [
                                Path(item["filepath"])
                                for item in requested_downloads
                                if isinstance(item, dict) and item.get("filepath")
                            ]
                            prepared_path = Path(downloader.prepare_filename(info))
                            if prepared_path not in candidate_paths:
                                candidate_paths.append(prepared_path)
                        last_variant_error = None
                        break
                    except Exception as variant_exc:
                        last_variant_error = variant_exc
                        if variant_index == 1 and is_unavailable_impersonation_target_error(variant_exc):
                            logger.info("Retrying video download without impersonation for %s: %s", video_url, variant_exc)
                            continue
                        raise
                if last_variant_error is not None:
                    raise last_variant_error
                last_error = ""
                break
            except Exception as exc:
                last_error = str(exc).strip()
                logger.info(
                    "Video download attempt %s failed for %s with format %s (use_cookiefile=%s, use_proxy=%s): %s",
                    attempt_index,
                    video_url,
                    attempt["format"],
                    attempt["use_cookiefile"],
                    use_proxy,
                    exc,
                )
                candidate_paths = []
                continue

        if not last_error:
            break

    if last_error:
        cleanup_download_artifacts(download_prefix)
        raise RuntimeError(last_error)

    candidates = [
        path
        for path in candidate_paths
        if path.exists() and path.is_file() and path.suffix.lower() != ".part"
    ]
    if not candidates:
        candidates = sorted(
            [path for path in UPLOAD_DIR.glob(f"{download_prefix}*") if path.is_file() and path.suffix.lower() != ".part"],
            key=lambda path: path.stat().st_mtime,
            reverse=True,
        )

    if not candidates:
        raise RuntimeError("The video link was reachable, but no downloadable audio file was produced.")

    file_path = candidates[0]
    file_size = get_file_size(file_path)
    if file_size > MAX_FILE_SIZE_BYTES:
        try:
            file_path.unlink()
        except OSError:
            logger.warning("Could not delete oversized video download: %s", file_path)
        raise RuntimeError(
            f"The downloaded video audio is too large ({file_size / (1024 * 1024):.1f} MB). "
            f"Please use a shorter video or a link with a smaller audio track."
        )
    return file_path


def require_ffmpeg() -> str:
    ffmpeg_path = shutil.which("ffmpeg")
    if not ffmpeg_path:
        raise RuntimeError(
            "Large files need ffmpeg for audio compression and splitting, but ffmpeg is not installed or not available in PATH."
        )
    return ffmpeg_path


def require_ffprobe() -> str:
    ffprobe_path = shutil.which("ffprobe")
    if not ffprobe_path:
        raise RuntimeError(
            "ffprobe is required for large-file chunking, but it is not installed or not available in PATH."
        )
    return ffprobe_path


def format_job_error(exc: Exception, source_url: str = "") -> str:
    if isinstance(exc, APIStatusError):
        return (
            f"OpenAI request failed with status {exc.status_code}. "
            "This can happen because of temporary upstream issues or an unsupported file."
        )
    if isinstance(exc, TimeoutError):
        return "Processing timed out. Try again, compress the file, or use a shorter lecture segment."
    message = str(exc).strip()
    lowered = message.lower()
    is_youtube_source = bool(source_url and extract_youtube_video_id(source_url))
    if is_proxy_auth_error(exc):
        if is_youtube_source:
            return (
                "The configured YouTube proxy rejected authentication (HTTP 407). "
                "Recheck YOUTUBE_PROXY_HTTP_URL and YOUTUBE_PROXY_HTTPS_URL, or verify your "
                "YOUTUBE_WEBSHARE_PROXY_USERNAME and YOUTUBE_WEBSHARE_PROXY_PASSWORD values on Render."
            )
        return (
            "The configured outbound proxy rejected authentication (HTTP 407). "
            "Recheck the proxy URL or credentials configured on the backend."
        )
    if "impersonate target" in lowered and "not available" in lowered:
        return (
            "The backend downloader could not use its browser impersonation target. "
            "Try the link again after redeploying this fallback, or leave YTDLP_IMPERSONATE_TARGET empty and use public captions, backend cookies, or a proxy for blocked YouTube links."
        )
    if "sign in to confirm you're not a bot" in lowered or "cookies-from-browser" in lowered:
        if is_youtube_source:
            if has_youtube_cookie_source():
                if has_youtube_proxy_source():
                    return (
                        "This YouTube video still blocked direct server-side download after the hosted backend tried the saved YouTube cookies "
                        "and proxy settings. Try a different public YouTube link with captions, adjust the proxy route on Render, or upload "
                        "the lecture file directly."
                    )
                return (
                    "This YouTube video still blocked direct server-side download after the hosted backend tried the saved YouTube cookies. "
                    "Try a different public YouTube link, add a YouTube proxy on Render, or upload the lecture file directly."
                )
            if has_youtube_proxy_source():
                return (
                    "This YouTube video blocked direct server-side download from the hosted backend even after the configured proxy route was tried. "
                    "Public captions are still checked first, but this video may still need working YouTube cookies or a different public link."
                )
            return (
                "This YouTube video blocked direct server-side download from the hosted backend. "
                "Public captions are still checked first, but this video may need YOUTUBE_COOKIES_TXT or a YouTube proxy on Render."
            )
        return (
            "This video host blocked direct server-side download from the hosted backend. "
            "Try another public link, or upload the media file directly."
        )
    if "requested format is not available" in lowered or "no video formats found" in lowered:
        if is_youtube_source:
            return (
                "This YouTube link was reachable, but the hosted backend could not open a downloadable audio format for it. "
                "Try another public YouTube URL, or upload the lecture file directly."
            )
        return (
            "This video link was reachable, but the hosted backend could not open a downloadable audio format from that site. "
            "Try another public link, or upload the media file directly."
        )
    if "unsupported url" in lowered:
        return "That video site is not supported by the backend downloader yet. Try another public link or upload the file directly."
    return message or "An unexpected processing error occurred."


def get_media_duration_seconds(file_path: Path) -> float:
    ffprobe_path = require_ffprobe()
    file_size_mb = get_file_size(file_path) / (1024 * 1024)
    result = subprocess.run(
        [
            ffprobe_path,
            "-v",
            "error",
            "-show_entries",
            "format=duration",
            "-of",
            "default=noprint_wrappers=1:nokey=1",
            str(file_path),
        ],
        capture_output=True,
        text=True,
        check=False,
    )
    if result.returncode != 0:
        raise RuntimeError(f"ffprobe could not read media duration: {result.stderr.strip()}")
    try:
        if file_size_mb > 5:
            # Surface a visible stage while ffprobe inspects larger files.
            logger.info("Inspecting media duration for %.1f MB file", file_size_mb)
        return float(result.stdout.strip())
    except ValueError as exc:
        raise RuntimeError("ffprobe returned an invalid media duration.") from exc


def build_audio_chunks(file_path: Path, job_id: str) -> list[Path]:
    if get_file_size(file_path) <= OPENAI_AUDIO_LIMIT_BYTES:
        update_job(job_id, status="processing", stage="Audio ready for transcription", progress=15)
        return [file_path]

    ffmpeg_path = require_ffmpeg()
    chunk_dir = UPLOAD_DIR / f"{file_path.stem}_chunks"
    chunk_dir.mkdir(parents=True, exist_ok=True)
    update_job(job_id, status="processing", stage="Preparing lecture audio", progress=12)
    copied_audio = chunk_dir / "audio_track.m4a"
    copy_command = [
        ffmpeg_path,
        "-y",
        "-i",
        str(file_path),
        "-vn",
        "-c:a",
        "copy",
        str(copied_audio),
    ]

    copy_result = subprocess.run(
        copy_command,
        capture_output=True,
        text=True,
        check=False,
    )

    extracted_audio = copied_audio if copy_result.returncode == 0 and copied_audio.exists() else None

    if extracted_audio and get_file_size(extracted_audio) <= OPENAI_AUDIO_LIMIT_BYTES:
        update_job(job_id, status="processing", stage="Audio prepared successfully", progress=20)
        return [extracted_audio]

    source_for_encoding = extracted_audio or file_path
    profile_candidates: list[tuple[str, str, int]] = []
    for bitrate, sample_rate, duration in [
        (TRANSCRIPTION_AUDIO_BITRATE, TRANSCRIPTION_AUDIO_SAMPLE_RATE, CHUNK_DURATION_SECONDS),
        ("40k", "16000", max(180, CHUNK_DURATION_SECONDS - 60)),
        ("32k", "16000", max(120, CHUNK_DURATION_SECONDS - 120)),
        ("24k", "12000", max(90, CHUNK_DURATION_SECONDS - 180)),
    ]:
        candidate = (bitrate, sample_rate, duration)
        if candidate not in profile_candidates:
            profile_candidates.append(candidate)

    last_error = ""

    for index, (bitrate, sample_rate, chunk_duration) in enumerate(profile_candidates, start=1):
        profile_dir = chunk_dir / f"profile_{index}"
        profile_dir.mkdir(parents=True, exist_ok=True)
        compressed_audio = profile_dir / f"compressed_{index}.mp3"
        update_job(
            job_id,
            status="processing",
            stage=f"Compressing lecture audio (pass {index})",
            progress=min(20, 13 + index * 2),
        )
        compress_command = [
            ffmpeg_path,
            "-y",
            "-i",
            str(source_for_encoding),
            "-vn",
            "-ac",
            "1",
            "-ar",
            sample_rate,
            "-b:a",
            bitrate,
            str(compressed_audio),
        ]

        compress_result = subprocess.run(
            compress_command,
            capture_output=True,
            text=True,
            check=False,
        )
        if compress_result.returncode != 0:
            last_error = compress_result.stderr.strip()
            continue

        if compressed_audio.exists() and get_file_size(compressed_audio) <= OPENAI_AUDIO_LIMIT_BYTES:
            update_job(job_id, status="processing", stage="Audio prepared successfully", progress=20)
            return [compressed_audio]

        update_job(
            job_id,
            status="processing",
            stage=f"Splitting lecture into parts (pass {index})",
            progress=min(22, 16 + index * 2),
        )
        segment_pattern = profile_dir / "chunk_%03d.mp3"
        split_command = [
            ffmpeg_path,
            "-y",
            "-i",
            str(compressed_audio),
            "-f",
            "segment",
            "-segment_time",
            str(chunk_duration),
            "-reset_timestamps",
            "1",
            "-c",
            "copy",
            str(segment_pattern),
        ]

        split_result = subprocess.run(
            split_command,
            capture_output=True,
            text=True,
            check=False,
        )
        if split_result.returncode != 0:
            last_error = split_result.stderr.strip()
            continue

        chunk_files = sorted(profile_dir.glob("chunk_*.mp3"))
        if not chunk_files:
            last_error = "Audio splitting completed, but no chunks were produced."
            continue

        oversized_chunks = [chunk for chunk in chunk_files if get_file_size(chunk) > OPENAI_AUDIO_LIMIT_BYTES]
        if oversized_chunks:
            last_error = (
                "Generated chunks are still above the OpenAI upload limit after adaptive compression. "
                "Try a shorter lecture segment or a stronger server."
            )
            continue

        update_job(job_id, status="processing", stage=f"Prepared {len(chunk_files)} lecture parts", progress=22)
        return chunk_files

    raise RuntimeError(last_error or "Audio preparation failed for this file.")


async def transcribe_audio(file_path: Path, job_id: str) -> str:
    def _transcribe() -> str:
        transcript_parts: list[str] = []
        models_to_try = [TRANSCRIPTION_MODEL]
        if FALLBACK_TRANSCRIPTION_MODEL and FALLBACK_TRANSCRIPTION_MODEL != TRANSCRIPTION_MODEL:
            models_to_try.append(FALLBACK_TRANSCRIPTION_MODEL)

        update_job(job_id, status="processing", stage="Checking lecture audio", progress=8)
        chunk_files = build_audio_chunks(file_path, job_id)
        total_chunks = max(len(chunk_files), 1)

        try:
            for index, chunk_path in enumerate(chunk_files, start=1):
                start_progress = 25 + int(((index - 1) / total_chunks) * 65)
                end_progress = 25 + int((index / total_chunks) * 65)
                update_job(
                    job_id,
                    status="processing",
                    stage=f"Transcribing lecture part {index} of {total_chunks}",
                    progress=start_progress,
                )
                last_error: Exception | None = None

                for model_name in models_to_try:
                    for attempt in range(1, TRANSCRIPTION_RETRIES + 1):
                        try:
                            update_job(
                                job_id,
                                status="processing",
                                stage=f"Transcribing lecture part {index} of {total_chunks}",
                                progress=start_progress,
                            )
                            with chunk_path.open("rb") as audio_file:
                                response = client.with_options(
                                    timeout=TRANSCRIPTION_REQUEST_TIMEOUT
                                ).audio.transcriptions.create(
                                    model=model_name,
                                    file=audio_file,
                                )
                            transcript_parts.append(response.text.strip())
                            update_job(
                                job_id,
                                status="processing",
                                stage=f"Finished chunk {index} of {total_chunks}",
                                progress=end_progress,
                            )
                            last_error = None
                            break
                        except InternalServerError as exc:
                            last_error = exc
                            if attempt < TRANSCRIPTION_RETRIES:
                                continue
                        except APIStatusError as exc:
                            last_error = exc
                            break
                    if last_error is None:
                        break
                else:
                    if last_error:
                        raise last_error

            return "\n\n".join(part for part in transcript_parts if part)
        finally:
            chunk_dir = UPLOAD_DIR / f"{file_path.stem}_chunks"
            if chunk_dir.exists():
                try:
                    for temp_path in chunk_dir.iterdir():
                        if temp_path.exists():
                            temp_path.unlink()
                    chunk_dir.rmdir()
                except OSError:
                    logger.warning("Could not delete chunk directory: %s", chunk_dir)

    return await asyncio.wait_for(asyncio.to_thread(_transcribe), timeout=TRANSCRIPTION_JOB_TIMEOUT)


def build_fallback_study_guide(transcript: str) -> str:
    cleaned = " ".join(transcript.split())
    sentence_candidates = re.split(r"(?<=[.!?])\s+", cleaned)
    sentences = [s.strip() for s in sentence_candidates if s.strip()]
    keywords = re.findall(r"\b[A-Za-z][A-Za-z0-9-]{4,}\b", transcript)

    frequency: dict[str, int] = {}
    stopwords = {
        "about",
        "after",
        "because",
        "before",
        "could",
        "their",
        "there",
        "these",
        "those",
        "which",
        "where",
        "while",
        "would",
    }
    for word in keywords:
        token = word.lower()
        if token in stopwords:
            continue
        frequency[token] = frequency.get(token, 0) + 1

    key_terms = [term.replace("-", " ") for term, _ in sorted(frequency.items(), key=lambda item: (-item[1], item[0]))[:6]]
    summary_sentences = sentences[:3] or ["The lecture transcript was available, but the AI summary step failed."]
    concept_bullets = key_terms[:5] or ["Main concepts were not clearly detected from the transcript."]
    question_topics = concept_bullets or ["the main lecture idea"]
    question_prompts = [
        "Define",
        "Describe",
        "Summarize",
        "Explain how to use",
        "Compare and contrast",
        "Apply",
        "Analyse",
        "Evaluate",
        "Connect",
        "Design an approach for",
    ]
    title_seed = key_terms[:3]
    title = " / ".join(term.title() for term in title_seed) if title_seed else "Lecture Notes"

    return "\n".join(
        [
            "**LECTURE TITLE**",
            title,
            "",
            "**SHORT SUMMARY**",
            *[f"- {sentence}" for sentence in summary_sentences],
            "",
            "**KEY CONCEPTS**",
            *[f"- {concept}" for concept in concept_bullets],
            "",
            "**IMPORTANT DEFINITIONS**",
            *[f"- {concept}: Review this concept in the transcript and expand it with your class notes." for concept in concept_bullets[:4]],
            "",
            "**IMPORTANT FORMULAS**",
            "- Not clearly detected in the fallback summary. Check the transcript for equations or numeric steps.",
            "",
            "**WORKED EXAMPLES**",
            "- Example problem: Choose one important concept from the lecture and explain how you would solve or apply it step by step.",
            "- Step 1: Identify the concept, formula, or method being used.",
            "- Step 2: Write down the known values, assumptions, or definitions.",
            "- Step 3: Apply the method in a clear sequence and explain each step.",
            "- Step 4: State the final result and what it means.",
            "",
            "**STEP-BY-STEP EXPLANATIONS**",
            "- Read the transcript from top to bottom and group the lecture into introduction, core method, and examples.",
            "- Highlight repeated ideas and convert them into short revision bullets.",
            "- Rewrite any lecturer demonstrations into your own words for exam preparation.",
            "",
            "**ADVANTAGES AND DISADVANTAGES**",
            "Advantages:",
            "- Compact notes make the core idea faster to revise before a quiz or exam.",
            "- Key concepts, flashcards, and questions help students move from reading to active recall.",
            "",
            "Disadvantages / cautions:",
            "- A short summary can hide missing detail if you never return to the transcript or class notes.",
            "- Students may remember words without understanding how to apply them if they skip worked examples.",
            "",
            "**COMMON MISTAKES TO AVOID**",
            "- Do not revise only the headings without checking the lecturer's examples.",
            "- Do not leave formulas or definitions in a half-understood state before test day.",
            "- Do not assume recognition means mastery; answer a question from memory to check yourself.",
            "",
            "**QUICK REVISION PLAN**",
            "- First 5 minutes: scan the summary and list the main ideas from memory.",
            "- Next 10 minutes: review definitions, formulas, and one worked example.",
            "- Next 10 minutes: answer two or three practice questions without notes.",
            "- Final 5 minutes: turn weak points into flashcards for the next study block.",
            "",
            "**VISUAL AIDS**",
            "| Study Tool | Simple Layout |",
            "| --- | --- |",
            "| Flow | Topic -> Method -> Result |",
            "| Relationship | Input -> Process -> Output |",
            "",
            "ASCII sketch:",
            "Start -> Key concept -> Example -> Exam point",
            "",
            "**REAL-WORLD EXAMPLES**",
            "- Link each key concept to one practical situation from your course or field.",
            "",
            "**FLASHCARDS**",
            *[
                (
                    f"Q: What is {item}?\n\n"
                    f"A: It is one of the main ideas highlighted in the lecture and should be revised from the transcript."
                )
                for item in concept_bullets[:3]
            ],
            "",
            "**EXAM TIPS**",
            "- Start with the summary, then revise the transcript for missing detail.",
            "- Turn repeated terms into flashcards and test yourself after each revision block.",
        ]
    )


def make_formulas_human_readable(text: str) -> str:
    superscript_map = str.maketrans({
        "0": "⁰",
        "1": "¹",
        "2": "²",
        "3": "³",
        "4": "⁴",
        "5": "⁵",
        "6": "⁶",
        "7": "⁷",
        "8": "⁸",
        "9": "⁹",
        "+": "⁺",
        "-": "⁻",
        "=": "⁼",
        "(": "⁽",
        ")": "⁾",
        "a": "ᵃ",
        "b": "ᵇ",
        "c": "ᶜ",
        "d": "ᵈ",
        "e": "ᵉ",
        "f": "ᶠ",
        "g": "ᵍ",
        "h": "ʰ",
        "i": "ⁱ",
        "j": "ʲ",
        "k": "ᵏ",
        "l": "ˡ",
        "m": "ᵐ",
        "n": "ⁿ",
        "o": "ᵒ",
        "p": "ᵖ",
        "r": "ʳ",
        "s": "ˢ",
        "t": "ᵗ",
        "u": "ᵘ",
        "v": "ᵛ",
        "w": "ʷ",
        "x": "ˣ",
        "y": "ʸ",
        "z": "ᶻ",
        "A": "ᴬ",
        "B": "ᴮ",
        "D": "ᴰ",
        "E": "ᴱ",
        "G": "ᴳ",
        "H": "ᴴ",
        "I": "ᴵ",
        "J": "ᴶ",
        "K": "ᴷ",
        "L": "ᴸ",
        "M": "ᴹ",
        "N": "ᴺ",
        "O": "ᴼ",
        "P": "ᴾ",
        "R": "ᴿ",
        "T": "ᵀ",
        "U": "ᵁ",
        "V": "ⱽ",
        "W": "ᵂ",
        " ": " ",
    })
    subscript_map = str.maketrans({
        "0": "₀",
        "1": "₁",
        "2": "₂",
        "3": "₃",
        "4": "₄",
        "5": "₅",
        "6": "₆",
        "7": "₇",
        "8": "₈",
        "9": "₉",
        "+": "₊",
        "-": "₋",
        "=": "₌",
        "(": "₍",
        ")": "₎",
        "a": "ₐ",
        "e": "ₑ",
        "h": "ₕ",
        "i": "ᵢ",
        "j": "ⱼ",
        "k": "ₖ",
        "l": "ₗ",
        "m": "ₘ",
        "n": "ₙ",
        "o": "ₒ",
        "p": "ₚ",
        "r": "ᵣ",
        "s": "ₛ",
        "t": "ₜ",
        "u": "ᵤ",
        "v": "ᵥ",
        "x": "ₓ",
        " ": " ",
    })

    def to_superscript(value: str) -> str:
        return value.translate(superscript_map)

    def to_subscript(value: str) -> str:
        return value.translate(subscript_map)

    replacements = {
        r"\mathcal{L}": "L",
        r"\Gamma": "Gamma",
        r"\infty": "infinity",
        r"\geq": ">=",
        r"\leq": "<=",
        r"\quad": " ",
        r"\cdot": "*",
        r"\,": " ",
        r"\left": "",
        r"\right": "",
        r"\{": "{",
        r"\}": "}",
        r"\[": "",
        r"\]": "",
        "$$": "",
        "$": "",
    }

    cleaned = text
    for old, new in replacements.items():
        cleaned = cleaned.replace(old, new)

    cleaned = re.sub(r"\\text\{([^}]*)\}", r"\1", cleaned)
    cleaned = re.sub(r"\\frac\{([^}]*)\}\{([^}]*)\}", r"(\1)/(\2)", cleaned)
    cleaned = re.sub(r"\\int_0\^infinity", "integral from 0 to infinity", cleaned)
    cleaned = re.sub(r"\\int_a\^infinity", "integral from a to infinity", cleaned)
    cleaned = re.sub(r"\\int", "integral", cleaned)
    cleaned = re.sub(r"\bpi\b", "π", cleaned, flags=re.IGNORECASE)
    cleaned = re.sub(r"\btheta\b", "θ", cleaned, flags=re.IGNORECASE)
    cleaned = re.sub(r"\bomega\b", "ω", cleaned, flags=re.IGNORECASE)
    cleaned = re.sub(r"\blambda\b", "λ", cleaned, flags=re.IGNORECASE)
    cleaned = cleaned.replace(">=", "≥").replace("<=", "≤").replace("!=", "≠")
    cleaned = cleaned.replace("integral from 0 to infinity", f"∫{to_subscript('0')}∞")
    cleaned = cleaned.replace("integral from a to infinity", f"∫{to_subscript('a')}∞")
    cleaned = re.sub(
        r"\^(\(([^)]+)\)|\{([^}]+)\}|([A-Za-z0-9+\- ]+))",
        lambda match: to_superscript(match.group(2) or match.group(3) or match.group(4) or ""),
        cleaned,
    )
    cleaned = re.sub(
        r"_(\(([^)]+)\)|\{([^}]+)\}|([A-Za-z0-9+\- ]+))",
        lambda match: to_subscript(match.group(2) or match.group(3) or match.group(4) or ""),
        cleaned,
    )
    cleaned = re.sub(r"\n([A-Z][A-Z \-]+)\n", lambda m: f"\n**{m.group(1).strip()}**\n\n", cleaned)
    cleaned = re.sub(r"\n{3,}", "\n\n", cleaned)
    return cleaned.strip()


def normalize_guide_heading(value: str) -> str:
    cleaned = compact_text(value).lower()
    cleaned = re.sub(r"^[#*\-\s]+", "", cleaned)
    cleaned = re.sub(r"[*:]+$", "", cleaned)
    cleaned = re.sub(r"\s+", " ", cleaned).strip()
    return cleaned


def canonical_guide_heading(value: str) -> str:
    normalized = normalize_guide_heading(value)
    return GUIDE_SECTION_ALIASES.get(normalized, normalized.upper())


def parse_study_guide_sections(markdown: str) -> list[dict[str, str]]:
    text = (markdown or "").replace("\r\n", "\n").strip()
    if not text:
        return []

    sections: list[dict[str, str]] = []
    intro_lines: list[str] = []
    current_heading = ""
    current_lines: list[str] = []
    title_text = ""

    def flush_current() -> None:
        nonlocal current_heading, current_lines
        content = "\n".join(current_lines).strip()
        if current_heading and content:
            sections.append(
                {
                    "heading": canonical_guide_heading(current_heading),
                    "raw_heading": compact_text(current_heading),
                    "content": content,
                }
            )
        current_heading = ""
        current_lines = []

    for raw_line in text.splitlines():
        line = raw_line.rstrip()
        stripped = line.strip()
        if not stripped:
            if current_heading:
                current_lines.append("")
            elif intro_lines:
                intro_lines.append("")
            continue

        markdown_heading_match = re.match(r"^(#{1,6})\s+(.+?)\s*$", stripped)
        if markdown_heading_match:
            heading_level = len(markdown_heading_match.group(1))
            heading_text = markdown_heading_match.group(2).strip().strip("*").strip()
            if heading_level == 1 and heading_text:
                if current_heading:
                    flush_current()
                if not title_text:
                    title_text = heading_text
                    continue
            flush_current()
            current_heading = heading_text
            continue

        bold_heading_match = re.match(r"^\*\*(.+?)\*\*\s*:?\s*(.*)$", stripped)
        if bold_heading_match:
            heading_text = bold_heading_match.group(1).strip()
            inline_content = bold_heading_match.group(2).strip()
            if normalize_guide_heading(heading_text) in GUIDE_SECTION_ALIASES:
                flush_current()
                current_heading = heading_text
                if inline_content:
                    current_lines.append(inline_content)
                continue

        if current_heading:
            current_lines.append(line)
        else:
            intro_lines.append(line)

    flush_current()

    if title_text:
        sections.insert(
            0,
            {
                "heading": "LECTURE TITLE",
                "raw_heading": title_text,
                "content": title_text,
            },
        )

    intro_text = "\n".join(intro_lines).strip()
    has_title = any(item["heading"] == "LECTURE TITLE" for item in sections)
    has_summary = any(item["heading"] == "SHORT SUMMARY" for item in sections)

    if intro_text:
        if not has_title:
            first_line = next((line.strip() for line in intro_text.splitlines() if line.strip()), "")
            if first_line:
                sections.insert(
                    0,
                    {
                        "heading": "LECTURE TITLE",
                        "raw_heading": first_line,
                        "content": first_line,
                    },
                )
                remainder = intro_text[len(first_line):].strip()
                if remainder and not has_summary:
                    sections.insert(
                        1,
                        {
                            "heading": "SHORT SUMMARY",
                            "raw_heading": "SHORT SUMMARY",
                            "content": remainder,
                        },
                    )
        elif not has_summary:
            insert_at = 1 if sections and sections[0]["heading"] == "LECTURE TITLE" else 0
            sections.insert(
                insert_at,
                {
                    "heading": "SHORT SUMMARY",
                    "raw_heading": "SHORT SUMMARY",
                    "content": intro_text,
                },
            )

    return sections


def extract_section(markdown: str, heading: str) -> str:
    target_heading = canonical_guide_heading(heading)
    for section in parse_study_guide_sections(markdown):
        if section["heading"] == target_heading:
            return section["content"].strip()
    return ""


def parse_flashcards(section_text: str) -> list[dict[str, str]]:
    lines = [line.strip() for line in section_text.splitlines() if line.strip()]
    cards: list[dict[str, str]] = []
    current_question = ""

    for line in lines:
        if line.startswith("- "):
            line = line[2:].strip()
        if line.startswith("Q:"):
            current_question = line[2:].strip()
        elif line.startswith("A:") and current_question:
            cards.append({"question": current_question, "answer": line[2:].strip()})
            current_question = ""

    return cards


def parse_quiz_questions(section_text: str) -> list[dict[str, str]]:
    pattern = re.compile(
        r"(?:^|\n)\s*(\d+)\.\s*Question:\s*(.*?)(?:\n\s*)+Answer:\s*(.*?)(?=\n\s*\d+\.\s*Question:|\Z)",
        re.DOTALL,
    )
    questions: list[dict[str, str]] = []
    for match in pattern.finditer(section_text):
        questions.append(
            {
                "number": match.group(1).strip(),
                "question": match.group(2).strip(),
                "answer": match.group(3).strip(),
            }
        )
    return questions


def replace_section(markdown: str, heading: str, new_body: str) -> str:
    pattern = re.compile(
        rf"((?:\*\*{re.escape(heading)}\*\*|##\s+{re.escape(heading)}))\s*(.*?)(?=\n(?:\*\*[A-Z][A-Z \-&]+\*\*|##\s+\S)|\Z)",
        re.DOTALL | re.IGNORECASE,
    )

    def _replace(match: re.Match[str]) -> str:
        heading_text = match.group(1)
        body = new_body.strip()
        return f"{heading_text}\n\n{body}\n\n"

    return pattern.sub(_replace, markdown, count=1)


def append_section_if_missing(markdown: str, heading: str, body: str) -> str:
    if extract_section(markdown, heading):
        return markdown.strip()
    cleaned = markdown.strip()
    separator = "\n\n" if cleaned else ""
    return f"{cleaned}{separator}**{heading}**\n\n{body.strip()}".strip()


def render_study_guide_sections(sections: list[dict[str, str]]) -> str:
    rendered: list[str] = []
    for section in sections:
        heading = compact_text(section.get("heading") or section.get("raw_heading"))
        content = (section.get("content") or "").strip()
        if not heading:
            continue
        if heading == "LECTURE TITLE":
            title = compact_text(content, compact_text(section.get("raw_heading"), "Study Guide"))
            rendered.append(f"# {title}")
            continue
        display_heading = compact_text(section.get("raw_heading"), heading)
        rendered.append(f"## {display_heading}")
        if content:
            rendered.append("")
            rendered.append(content)
        rendered.append("")
    return "\n".join(rendered).strip()


def upsert_guide_section(markdown: str, heading: str, new_body: str, after_heading: str = "") -> str:
    target_heading = canonical_guide_heading(heading)
    insert_after_heading = canonical_guide_heading(after_heading) if after_heading else ""
    sections = parse_study_guide_sections(markdown)
    if not sections:
        return append_section_if_missing(markdown, target_heading, new_body)

    updated_sections: list[dict[str, str]] = []
    replaced = False
    for section in sections:
        if section["heading"] == target_heading:
            updated_sections.append(
                {
                    "heading": target_heading,
                    "raw_heading": target_heading,
                    "content": new_body.strip(),
                }
            )
            replaced = True
            continue
        updated_sections.append(dict(section))

    if not replaced:
        insert_at = len(updated_sections)
        if insert_after_heading:
            for index, section in enumerate(updated_sections):
                if section["heading"] == insert_after_heading:
                    insert_at = index + 1
                    break
        updated_sections.insert(
            insert_at,
            {
                "heading": target_heading,
                "raw_heading": target_heading,
                "content": new_body.strip(),
            },
        )

    return render_study_guide_sections(updated_sections)


def extract_bullet_points(section_text: str) -> list[str]:
    items: list[str] = []
    for line in (section_text or "").splitlines():
        cleaned = line.strip()
        if not cleaned:
            continue
        if cleaned.startswith("- "):
            items.append(cleaned[2:].strip())
        elif re.match(r"^\d+\.\s+", cleaned):
            items.append(re.sub(r"^\d+\.\s+", "", cleaned).strip())
    return items


def build_missing_step_by_step_section(summary: str) -> str:
    worked_points = extract_bullet_points(extract_section(summary, "WORKED EXAMPLES"))
    concept_points = extract_bullet_points(extract_section(summary, "KEY CONCEPTS"))
    definition_points = extract_bullet_points(extract_section(summary, "IMPORTANT DEFINITIONS"))
    summary_points = extract_bullet_points(extract_section(summary, "SHORT SUMMARY"))
    formula_lines = [
        line.strip()
        for line in extract_section(summary, "IMPORTANT FORMULAS").splitlines()
        if compact_text(line) and not line.strip().startswith("|")
    ]

    explicit_steps: list[str] = []
    for point in worked_points:
        cleaned_point = compact_text(re.sub(r"^step\s*\d+\s*[:.\-]?\s*", "", point, flags=re.IGNORECASE))
        if cleaned_point:
            explicit_steps.append(cleaned_point)

    if explicit_steps:
        formatted_steps: list[str] = []
        for index, item in enumerate(explicit_steps[:4], start=1):
            ending = item if re.search(r"[.!?]$", item) else f"{item}."
            formatted_steps.append(f"- Step {index}: {ending}")
        return "\n".join(formatted_steps)

    topic_anchor = concept_points[0] if concept_points else (summary_points[0] if summary_points else "the main lecture topic")
    rule_anchor = formula_lines[0] if formula_lines else (definition_points[0] if definition_points else "the key rule or definition")
    ending_anchor = concept_points[1] if len(concept_points) > 1 else topic_anchor

    return "\n".join(
        [
            f"- Step 1: Identify the exact question, process, or idea and connect it to {topic_anchor}.",
            f"- Step 2: Recall the main rule, definition, formula, or principle that controls it: {rule_anchor}.",
            "- Step 3: Follow the lecture method carefully and explain each transition instead of jumping straight to the final result.",
            f"- Step 4: Check the final answer, interpretation, or conclusion and link it back to {ending_anchor}.",
        ]
    )


def add_student_support_sections(summary: str) -> str:
    cleaned = (summary or "").strip()
    cleaned = re.sub(r"\n{3,}", "\n\n", cleaned)
    if not compact_text(extract_section(cleaned, "STEP-BY-STEP EXPLANATIONS")):
        cleaned = upsert_guide_section(
            cleaned,
            "STEP-BY-STEP EXPLANATIONS",
            build_missing_step_by_step_section(cleaned),
            after_heading="WORKED EXAMPLES",
        )
    return cleaned.strip()


def build_study_image_queries(
    summary: str,
    transcript: str,
    lecture_notes: str,
    lecture_slides: str,
) -> list[str]:
    visual_aids = extract_section(summary, "VISUAL AIDS").lower()
    structured_sources = "\n\n".join(part for part in [lecture_notes.strip(), lecture_slides.strip()] if part).lower()
    concrete_markers = (
        "recognise",
        "recognize",
        "visible",
        "photo",
        "real-world appearance",
        "subtype",
        "machine",
        "instrument",
        "valve",
        "organ",
        "structure",
        "component",
        "equipment",
    )
    should_fetch_images = any(marker in visual_aids for marker in concrete_markers) or any(
        marker in structured_sources for marker in concrete_markers
    )
    if not should_fetch_images:
        return []

    context_parts = [
        trimmed_context_block("STUDY GUIDE", summary, 5000),
        trimmed_context_block("LECTURE NOTES", lecture_notes, 2500),
        trimmed_context_block("LECTURE SLIDES", lecture_slides, 2500),
    ]
    combined_context = "\n\n".join(part for part in context_parts if part).strip()
    if not combined_context:
        return []

    try:
        response = client.with_options(timeout=STUDY_IMAGE_QUERY_TIMEOUT).chat.completions.create(
            model=ASSET_GENERATION_MODEL,
            max_completion_tokens=350,
            messages=[
                {
                    "role": "system",
                    "content": (
                        "You pick real-photo search queries for a university study app. "
                        "Return strict JSON only in this shape: {\"queries\": [\"...\"]}.\n\n"
                        "Rules:\n"
                        "- Return 0 to 6 short search queries.\n"
                        "- Only return queries for concrete things that students should literally see, such as organs, machines, lab tools, landmarks, species, hardware, or physical processes.\n"
                        "- If the topic is mostly abstract, symbolic, theoretical, or mathematical, return an empty array.\n"
                        "- Prefer specific nouns over vague phrases.\n"
                        "- If the guide compares concrete subtypes, return one query per important subtype instead of one generic umbrella term.\n"
                        "- Do not return diagram, illustration, formula, or stock-photo style phrases."
                    ),
                },
                {"role": "user", "content": combined_context},
            ],
        )
        parsed = parse_json_object(response.choices[0].message.content or "")
        return [compact_text(item) for item in parse_json_list(parsed.get("queries")) if compact_text(item)][:MAX_STUDY_IMAGES]
    except Exception as exc:
        logger.warning("Study image query generation failed: %s", exc)
        return []


def is_real_photo_candidate_title(title: str) -> bool:
    normalized = compact_text(title).lower()
    blocked_markers = (
        "diagram",
        "logo",
        "icon",
        "equation",
        "formula",
        "graph",
        "chart",
        "scheme",
        "vector",
        "flag",
        "map",
        "coat of arms",
    )
    return normalized and not any(marker in normalized for marker in blocked_markers)


def humanize_commons_title(title: str) -> str:
    cleaned = re.sub(r"^File:", "", compact_text(title))
    cleaned = re.sub(r"\.[A-Za-z0-9]{2,5}$", "", cleaned)
    cleaned = cleaned.replace("_", " ")
    return cleaned.strip() or "Reference photo"


def search_wikimedia_photos(query: str, limit: int = 1) -> list[dict[str, str]]:
    response = requests.get(
        WIKIMEDIA_API_URL,
        params={
            "action": "query",
            "format": "json",
            "generator": "search",
            "gsrnamespace": 6,
            "gsrsearch": query,
            "gsrlimit": max(4, limit * 5),
            "prop": "imageinfo",
            "iiprop": "url|mime",
            "iiurlwidth": 1400,
        },
        timeout=STUDY_IMAGE_SEARCH_TIMEOUT,
    )
    response.raise_for_status()
    payload = response.json()

    pages = list((payload.get("query", {}) or {}).get("pages", {}).values())
    pages.sort(key=lambda item: int(item.get("index", 999999)))

    results: list[dict[str, str]] = []
    for page in pages:
        title = compact_text(page.get("title"))
        if not is_real_photo_candidate_title(title):
            continue
        image_info = ((page.get("imageinfo") or [{}])[0]) if isinstance(page.get("imageinfo"), list) else {}
        mime_type = compact_text(image_info.get("mime")).lower()
        image_url = compact_text(image_info.get("thumburl") or image_info.get("url"))
        if mime_type not in {"image/jpeg", "image/png", "image/webp"} or not image_url:
            continue
        results.append(
            {
                "query": compact_text(query),
                "title": humanize_commons_title(title),
                "image_url": image_url,
                "source_url": compact_text(f"https://commons.wikimedia.org/wiki/{title.replace(' ', '_')}"),
            }
        )
        if len(results) >= limit:
            break
    return results


async def generate_study_images(
    summary: str,
    transcript: str,
    lecture_notes: str,
    lecture_slides: str,
    job_id: str,
) -> list[dict[str, str]]:
    queries = await asyncio.to_thread(
        build_study_image_queries,
        summary,
        transcript,
        lecture_notes,
        lecture_slides,
    )
    if not queries:
        return []

    update_job(job_id, status="processing", stage="Finding real study photos", progress=92)
    seen_urls: set[str] = set()
    images: list[dict[str, str]] = []

    for query in queries:
        try:
            results = await asyncio.to_thread(search_wikimedia_photos, query, 1)
        except Exception as exc:
            logger.warning("Study photo search failed for query '%s': %s", query, exc)
            continue

        for result in results:
            image_url = compact_text(result.get("image_url"))
            if not image_url or image_url in seen_urls:
                continue
            seen_urls.add(image_url)
            images.append(result)
            if len(images) >= MAX_STUDY_IMAGES:
                return images

    return images


def parse_visual_analysis_items(content: str) -> list[dict[str, str]]:
    parsed = parse_json_object(content)
    items = parsed.get("items") if isinstance(parsed.get("items"), list) else parsed if isinstance(parsed, list) else []
    normalized: list[dict[str, str]] = []
    for item in items:
        if not isinstance(item, dict):
            continue
        normalized.append(
            {
                "title": compact_text(item.get("title"), "Lecture visual"),
                "visual_type": compact_text(item.get("visual_type"), "diagram"),
                "matched_section": compact_text(item.get("matched_section"), "Key concept"),
                "key_highlight": compact_text(item.get("key_highlight"), "Useful visual from the uploaded lecture material."),
                "diagram_label": compact_text(item.get("diagram_label"), compact_text(item.get("title"), "Lecture visual")),
            }
        )
    return normalized


async def analyze_reference_images_for_study_guide(
    reference_images: list[str],
    summary: str,
    lecture_notes: str,
    lecture_slides: str,
    output_language: str,
) -> list[dict[str, str]]:
    valid_images = [compact_text(item) for item in reference_images if compact_text(item)][:6]
    if not valid_images:
        return []

    context_parts = [
        trimmed_context_block("STUDY GUIDE SUMMARY", summary, 2400),
        trimmed_context_block("LECTURER NOTES", lecture_notes, 1800),
        trimmed_context_block("LECTURE SLIDES", lecture_slides, 1800),
    ]
    context_text = "\n\n".join(part for part in context_parts if part).strip()

    user_parts: list[dict[str, Any]] = [
        {
            "type": "text",
            "text": (
                "Review the uploaded lecture visuals and return JSON only.\n\n"
                f"Academic context:\n{context_text or 'No extra text context was supplied.'}\n\n"
                f"Write all labels in {output_language}.\n"
                "Match each visual to the most relevant study-guide section or concept.\n"
                "Detect whether the visual is a table, chart, diagram, process flow, equation snapshot, photo, or mixed slide visual.\n"
                "Give each visual a short clear classroom label and one highlight explaining what a student should notice.\n"
                "Return strict JSON in this shape only:\n"
                '{"items":[{"title":"...","visual_type":"...","matched_section":"...","key_highlight":"...","diagram_label":"..."}]}'
            ),
        }
    ]
    for image in valid_images:
        user_parts.append({"type": "image_url", "image_url": {"url": image}})

    def _analyze() -> list[dict[str, str]]:
        response = client.with_options(timeout=VISION_REQUEST_TIMEOUT).chat.completions.create(
            model=VISION_MODEL,
            max_completion_tokens=900,
            messages=[
                {
                    "role": "system",
                    "content": (
                        "You are labeling lecture visuals for a university study guide. "
                        "Be concise, accurate, and practical. "
                        "Return strict JSON only."
                    ),
                },
                {"role": "user", "content": user_parts},
            ],
        )
        return parse_visual_analysis_items(response.choices[0].message.content or "")

    try:
        return await asyncio.to_thread(_analyze)
    except Exception as exc:
        logger.warning("Reference image analysis failed for study guide: %s", exc)
        fallback_items: list[dict[str, str]] = []
        for index, _ in enumerate(valid_images, start=1):
            fallback_items.append(
                {
                    "title": f"Lecture visual {index}",
                    "visual_type": "diagram",
                    "matched_section": "Key concept",
                    "key_highlight": "Review this uploaded lecture visual alongside the matching section in the study guide.",
                    "diagram_label": f"Visual {index}",
                }
            )
        return fallback_items


def build_study_guide_visual_notes(visual_items: list[dict[str, str]]) -> str:
    if not visual_items:
        return ""
    lines = ["REFERENCE VISUALS"]
    for index, item in enumerate(visual_items, start=1):
        lines.append(
            (
                f"{index}. Title: {compact_text(item.get('title'), f'Lecture visual {index}')}\n"
                f"   Type: {compact_text(item.get('visual_type'), 'diagram')}\n"
                f"   Matched section: {compact_text(item.get('matched_section'), 'Key concept')}\n"
                f"   Diagram label: {compact_text(item.get('diagram_label'), f'Visual {index}')}\n"
                f"   Key highlight: {compact_text(item.get('key_highlight'), 'Review this visual carefully.')}"
            )
        )
    return "\n".join(lines).strip()


def build_reference_image_catalog(visual_items: list[dict[str, str]]) -> str:
    if not visual_items:
        return ""
    lines = ["REFERENCE IMAGE CATALOG"]
    for index, item in enumerate(visual_items):
        lines.append(
            (
                f"{index}. Title: {compact_text(item.get('title'), f'Lecture visual {index + 1}')}\n"
                f"   Type: {compact_text(item.get('visual_type'), 'diagram')}\n"
                f"   Best match: {compact_text(item.get('matched_section'), 'Key concept')}\n"
                f"   Classroom label: {compact_text(item.get('diagram_label'), f'Visual {index + 1}')}\n"
                f"   Key detail: {compact_text(item.get('key_highlight'), 'Review the important visual clue before using this image on a slide.')}"
            )
        )
    return "\n".join(lines).strip()


def build_uploaded_study_visuals(
    reference_images: list[str],
    visual_items: list[dict[str, str]],
) -> list[dict[str, str]]:
    uploaded_visuals: list[dict[str, str]] = []
    for index, image_url in enumerate(reference_images[: len(visual_items) or len(reference_images)], start=0):
        if not compact_text(image_url):
            continue
        visual = visual_items[index] if index < len(visual_items) else {}
        uploaded_visuals.append(
            {
                "query": compact_text(visual.get("matched_section"), "Uploaded lecture visual"),
                "title": compact_text(visual.get("title"), f"Lecture visual {index + 1}"),
                "image_url": image_url,
                "source_url": image_url,
                "source_type": "uploaded",
                "visual_type": compact_text(visual.get("visual_type"), "diagram"),
                "matched_section": compact_text(visual.get("matched_section"), "Key concept"),
                "key_highlight": compact_text(visual.get("key_highlight"), "Useful lecture visual for this section."),
                "diagram_label": compact_text(visual.get("diagram_label"), f"Visual {index + 1}"),
            }
        )
    return uploaded_visuals[:MAX_STUDY_IMAGES]


def merge_study_image_results(*image_groups: list[dict[str, str]]) -> list[dict[str, str]]:
    merged: list[dict[str, str]] = []
    seen_urls: set[str] = set()
    for group in image_groups:
        for item in group or []:
            image_url = compact_text(item.get("image_url"))
            if not image_url or image_url in seen_urls:
                continue
            seen_urls.add(image_url)
            merged.append(item)
            if len(merged) >= MAX_STUDY_IMAGES:
                return merged
    return merged


def tidy_study_guide_layout(summary: str) -> str:
    cleaned = (summary or "").replace("\r\n", "\n").strip()
    quiz_section = extract_section(cleaned, "PRACTICE QUESTIONS AND ANSWERS")
    flashcard_section = extract_section(cleaned, "FLASHCARDS")

    quiz_questions = parse_quiz_questions(quiz_section)
    if quiz_questions:
        formatted_quiz = "\n\n".join(
            f"{item['number']}. Question: {item['question']}\n\nAnswer: {item['answer']}"
            for item in quiz_questions
        )
        cleaned = replace_section(cleaned, "PRACTICE QUESTIONS AND ANSWERS", formatted_quiz)

    flashcards = parse_flashcards(flashcard_section)
    if flashcards:
        formatted_flashcards = "\n\n".join(
            f"Q: {item['question']}\n\nA: {item['answer']}"
            for item in flashcards
        )
        cleaned = replace_section(cleaned, "FLASHCARDS", formatted_flashcards)

    cleaned = re.sub(r"\n{3,}", "\n\n", cleaned)
    return cleaned.strip()


def extract_study_assets(summary: str) -> dict:
    formula_section = extract_section(summary, "IMPORTANT FORMULAS")
    flashcard_section = extract_section(summary, "FLASHCARDS")
    quiz_section = extract_section(summary, "PRACTICE QUESTIONS AND ANSWERS")

    return {
        "formula": formula_section or "No formula section was detected in the notes.",
        "worked_example": build_worked_example_asset(summary),
        "flashcards": parse_flashcards(flashcard_section),
        "quiz_questions": parse_quiz_questions(quiz_section),
    }


def build_worked_example_asset(summary: str, generated_worked_example: str = "") -> str:
    example_section = compact_text(extract_section(summary, "WORKED EXAMPLES"))
    step_section = compact_text(extract_section(summary, "STEP-BY-STEP EXPLANATIONS"))
    generated_section = compact_text(generated_worked_example)
    sections: list[str] = []

    if example_section:
        sections.append(f"**WORKED EXAMPLES**\n\n{example_section}")

    if step_section:
        sections.append(f"**STEP-BY-STEP EXPLANATIONS**\n\n{step_section}")

    if sections:
        if generated_section and not step_section:
            sections.append(f"**STEP-BY-STEP WALKTHROUGH**\n\n{generated_section}")
        return "\n\n".join(section.strip() for section in sections if section.strip()).strip()

    return generated_section or "No worked example section was detected in the notes."


def clamp_podcast_speaker_count(value: int) -> int:
    return 3 if int(value or 2) >= 3 else 2


def clamp_podcast_target_minutes(value: int) -> int:
    return max(6, min(18, int(value or 10)))


def build_podcast_speaker_profiles(speaker_count: int) -> list[dict[str, str]]:
    profiles = [
        {
            "key": "speaker_1",
            "name": "Njabulo",
            "role": "the calm explainer who keeps the lesson academically solid",
            "voice": "ash",
            "voice_style": "Sound warm, clear, confident, and distinctly like a young man hosting a serious revision podcast.",
        },
        {
            "key": "speaker_2",
            "name": "Olwethu",
            "role": "the curious challenger who adds light jokes, asks obvious student questions, and pushes for simpler wording",
            "voice": "coral",
            "voice_style": "Sound lively, sharp, conversational, and distinctly like a young woman with light humor.",
        },
        {
            "key": "speaker_3",
            "name": "Melusi",
            "role": "the exam coach who keeps connecting the topic to worked examples, likely test traps, and revision advice",
            "voice": "echo",
            "voice_style": "Sound focused, practical, grounded, and distinctly like a young man guiding exam revision.",
        },
    ]
    return profiles[: clamp_podcast_speaker_count(speaker_count)]


TEACHER_GUIDE_SECTION_HEADINGS = [
    "LECTURE TITLE",
    "SHORT SUMMARY",
    "KEY CONCEPTS",
    "IMPORTANT DEFINITIONS",
    "IMPORTANT FORMULAS",
    "WORKED EXAMPLES",
    "STEP-BY-STEP EXPLANATIONS",
    "ADVANTAGES AND DISADVANTAGES",
    "COMMON MISTAKES TO AVOID",
    "QUICK REVISION PLAN",
    "VISUAL AIDS",
    "REAL-WORLD EXAMPLES",
    "EXAM TIPS",
]
TEACHER_TARGET_MINUTES = 22
TEACHER_MINIMUM_MINUTES = 20.0
TEACHER_TARGET_WORDS = 3200
TEACHER_SEGMENT_LIMIT = 42
TEACHER_REQUIRED_SECTION_COUNTS = {
    "SHORT SUMMARY": 1,
    "KEY CONCEPTS": 1,
    "IMPORTANT DEFINITIONS": 1,
    "IMPORTANT FORMULAS": 1,
    "WORKED EXAMPLES": 2,
    "STEP-BY-STEP EXPLANATIONS": 2,
}
TEACHER_SOURCE_HEADING_HINTS = {
    "SHORT SUMMARY": ("summary", "overview", "introduction", "objective", "topic"),
    "KEY CONCEPTS": ("concept", "principle", "main idea", "key point", "important"),
    "IMPORTANT DEFINITIONS": ("definition", "defined", "means", "refers to", "is called"),
    "IMPORTANT FORMULAS": ("formula", "equation", "law", "rule", "="),
    "WORKED EXAMPLES": ("worked example", "example", "sample problem", "illustration", "calculate"),
    "STEP-BY-STEP EXPLANATIONS": ("step", "steps", "procedure", "method", "process", "how to"),
    "COMMON MISTAKES TO AVOID": ("mistake", "avoid", "common error", "trap"),
}
TEACHER_CONTEXT_PRIORITY_TERMS = tuple(
    dict.fromkeys(term for terms in TEACHER_SOURCE_HEADING_HINTS.values() for term in terms)
)
TEACHER_SIMILARITY_STOPWORDS = {
    "a",
    "an",
    "and",
    "are",
    "as",
    "at",
    "be",
    "because",
    "before",
    "but",
    "by",
    "can",
    "class",
    "could",
    "do",
    "does",
    "each",
    "for",
    "from",
    "had",
    "has",
    "have",
    "here",
    "how",
    "into",
    "its",
    "just",
    "let",
    "lets",
    "lesson",
    "like",
    "may",
    "might",
    "more",
    "most",
    "much",
    "next",
    "now",
    "our",
    "out",
    "over",
    "really",
    "same",
    "should",
    "some",
    "still",
    "student",
    "students",
    "teacher",
    "than",
    "that",
    "the",
    "their",
    "them",
    "then",
    "there",
    "these",
    "they",
    "this",
    "those",
    "through",
    "topic",
    "very",
    "was",
    "were",
    "what",
    "when",
    "where",
    "which",
    "while",
    "with",
    "would",
    "you",
    "your",
}
TEACHER_SECTION_TARGET_SEGMENTS = {
    "LECTURE TITLE": 1,
    "SHORT SUMMARY": 2,
    "KEY CONCEPTS": 2,
    "IMPORTANT DEFINITIONS": 2,
    "IMPORTANT FORMULAS": 3,
    "WORKED EXAMPLES": 5,
    "STEP-BY-STEP EXPLANATIONS": 5,
    "ADVANTAGES AND DISADVANTAGES": 1,
    "COMMON MISTAKES TO AVOID": 2,
    "QUICK REVISION PLAN": 1,
    "VISUAL AIDS": 1,
    "REAL-WORLD EXAMPLES": 2,
    "EXAM TIPS": 2,
}


def teacher_heading_rank(heading: str) -> int:
    cleaned = compact_text(heading)
    try:
        return TEACHER_GUIDE_SECTION_HEADINGS.index(cleaned)
    except ValueError:
        return len(TEACHER_GUIDE_SECTION_HEADINGS)


def split_teacher_source_passages(text: str) -> list[str]:
    cleaned = compact_text(text)
    if not cleaned:
        return []

    paragraphs = [compact_text(item) for item in re.split(r"\n\s*\n+", cleaned) if compact_text(item)]
    if len(paragraphs) >= 4:
        return paragraphs

    lines = [line.strip() for line in cleaned.splitlines() if line.strip()]
    if not lines:
        return paragraphs or [cleaned]

    grouped: list[str] = []
    current: list[str] = []
    for line in lines:
        current.append(line)
        joined = " ".join(current).strip()
        if len(joined) >= 320 or (line.endswith((".", ":", ";", "?", "!")) and len(joined) >= 160):
            grouped.append(joined)
            current = []

    if current:
        grouped.append(" ".join(current).strip())

    return grouped or paragraphs or [cleaned]


def score_teacher_source_passage(passage: str, heading: str) -> int:
    lowered = (passage or "").lower()
    if not lowered:
        return 0

    score = 0
    for keyword in TEACHER_SOURCE_HEADING_HINTS.get(heading, ()):
        if keyword in lowered:
            score += 2 if " " in keyword else 1

    if heading == "IMPORTANT FORMULAS":
        score += lowered.count("=")
        score += lowered.count("therefore")
    elif heading == "WORKED EXAMPLES":
        score += len(re.findall(r"\bexample\b", lowered))
        score += len(re.findall(r"\bsolve\b|\bcalculate\b|\bfind\b", lowered))
    elif heading == "STEP-BY-STEP EXPLANATIONS":
        score += len(re.findall(r"\bstep\s*\d+\b", lowered)) * 2
        score += len(re.findall(r"\bfirst\b|\bthen\b|\bnext\b|\bfinally\b", lowered))
    elif heading == "IMPORTANT DEFINITIONS":
        score += len(re.findall(r"\bdefined as\b|\bmeans\b|\brefers to\b", lowered))

    return score


def extract_teacher_source_excerpt(source_text: str, heading: str, limit: int = 900) -> str:
    passages = split_teacher_source_passages(source_text)
    if not passages:
        return ""

    matches: list[tuple[int, int, str]] = []
    for index, passage in enumerate(passages):
        score = score_teacher_source_passage(passage, heading)
        if score > 0:
            matches.append((score, index, compact_text(passage)))

    if not matches:
        if heading == "SHORT SUMMARY":
            return compact_text(passages[0][:limit])
        if heading == "KEY CONCEPTS":
            return compact_text("\n\n".join(passages[:2])[:limit])
        return ""

    selected: list[tuple[int, str]] = []
    used_indexes: set[int] = set()
    budget = 0
    for score, index, passage in sorted(matches, key=lambda item: (-item[0], item[1])):
        if not passage or index in used_indexes:
            continue
        addition = len(passage) + (2 if selected else 0)
        if selected and budget + addition > limit:
            continue
        selected.append((index, passage))
        used_indexes.add(index)
        budget += addition
        if budget >= limit or len(selected) >= 3:
            break

    if not selected:
        top_passage = matches[0][2]
        return compact_text(top_passage[:limit])

    selected.sort(key=lambda item: item[0])
    excerpt = "\n\n".join(passage for _, passage in selected).strip()
    return compact_text(excerpt[:limit])


def extract_teacher_priority_excerpt(
    source_text: str,
    limit: int,
    excluded_passages: list[str] | None = None,
) -> str:
    passages = split_teacher_source_passages(source_text)
    if not passages:
        return ""

    excluded = [compact_text(item) for item in (excluded_passages or []) if compact_text(item)]
    candidates: list[tuple[int, int, str]] = []
    for index, passage in enumerate(passages):
        normalized = compact_text(passage)
        if not normalized:
            continue
        if any(normalized in blocked or blocked in normalized for blocked in excluded):
            continue
        lowered = normalized.lower()
        score = sum(1 for term in TEACHER_CONTEXT_PRIORITY_TERMS if term in lowered)
        score += len(re.findall(r"\bstep\s*\d+\b", lowered)) * 2
        if "=" in normalized:
            score += 2
        if score > 0:
            candidates.append((score, index, normalized))

    if not candidates:
        return ""

    selected: list[tuple[int, str]] = []
    budget = 0
    for score, index, passage in sorted(candidates, key=lambda item: (-item[0], item[1])):
        addition = len(passage) + (2 if selected else 0)
        if selected and budget + addition > limit:
            continue
        selected.append((index, passage))
        budget += addition
        if budget >= limit or len(selected) >= 4:
            break

    selected.sort(key=lambda item: item[0])
    excerpt = "\n\n".join(passage for _, passage in selected).strip()
    return compact_text(excerpt[:limit])


def build_teacher_context_block(label: str, value: str, limit: int) -> str:
    cleaned = compact_text(value)
    if not cleaned:
        return ""
    if len(cleaned) <= limit:
        return f"{label}\n{cleaned}"

    head_budget = min(max(limit // 3, 1200), 2200)
    tail_budget = min(max(limit // 6, 700), 1400)
    priority_budget = max(650, limit - head_budget - tail_budget - 260)
    if head_budget + tail_budget + priority_budget > limit:
        priority_budget = max(500, limit - head_budget - tail_budget - 120)

    opening_excerpt = cleaned[:head_budget].rstrip()
    later_excerpt = cleaned[-tail_budget:].lstrip()
    priority_excerpt = extract_teacher_priority_excerpt(
        cleaned,
        priority_budget,
        excluded_passages=[opening_excerpt, later_excerpt],
    )

    parts = [f"OPENING NOTES\n{opening_excerpt}"]
    if priority_excerpt:
        parts.append(f"PRIORITY EXCERPTS\n{priority_excerpt}")
    parts.append(f"LATER NOTES\n{later_excerpt}")
    parts.append(
        "NOTE: This source was condensed to preserve the opening lesson flow together with later definitions, formulas, worked examples, and method steps."
    )
    return f"{label}\n" + "\n\n".join(part for part in parts if part)


def build_teacher_section_outline(
    summary: str,
    lecture_notes: str = "",
    lecture_slides: str = "",
    past_question_papers: str = "",
    transcript: str = "",
) -> list[dict[str, str]]:
    sections: list[dict[str, str]] = []
    seen_headings: set[str] = set()
    for heading in TEACHER_GUIDE_SECTION_HEADINGS:
        content = compact_text(extract_section(summary, heading))
        if not content or heading in seen_headings:
            continue
        seen_headings.add(heading)
        sections.append({"section_heading": heading, "content": content})

    source_texts = [lecture_notes, lecture_slides, past_question_papers, transcript]
    for heading in TEACHER_REQUIRED_SECTION_COUNTS:
        if heading in seen_headings:
            continue
        for source_text in source_texts:
            content = extract_teacher_source_excerpt(source_text, heading)
            if not content:
                continue
            seen_headings.add(heading)
            sections.append({"section_heading": heading, "content": content})
            break
    return sorted(sections, key=lambda item: teacher_heading_rank(item["section_heading"]))


def normalize_teacher_similarity_text(text: str) -> str:
    normalized = compact_text(text).lower()
    if not normalized:
        return ""
    normalized = normalized.replace("let us", "lets")
    normalized = re.sub(r"[^a-z0-9\s]", " ", normalized)
    normalized = re.sub(r"\s+", " ", normalized).strip()
    return normalized


def teacher_content_tokens(text: str) -> set[str]:
    normalized = normalize_teacher_similarity_text(text)
    return {
        token
        for token in normalized.split()
        if len(token) > 2 and not token.isdigit() and token not in TEACHER_SIMILARITY_STOPWORDS
    }


def teacher_text_similarity(left: str, right: str) -> float:
    left_normalized = normalize_teacher_similarity_text(left)
    right_normalized = normalize_teacher_similarity_text(right)
    if not left_normalized or not right_normalized:
        return 0.0
    if left_normalized == right_normalized:
        return 1.0
    if len(left_normalized) >= 80 and (left_normalized in right_normalized or right_normalized in left_normalized):
        return 0.96

    left_tokens = teacher_content_tokens(left_normalized)
    right_tokens = teacher_content_tokens(right_normalized)
    if not left_tokens or not right_tokens:
        return 0.0

    shared = len(left_tokens & right_tokens)
    if not shared:
        return 0.0

    jaccard = shared / max(1, len(left_tokens | right_tokens))
    overlap = shared / max(1, min(len(left_tokens), len(right_tokens)))
    return max(jaccard, overlap)


def clean_teacher_segment_text(text: str) -> str:
    cleaned = compact_text(text)
    if not cleaned:
        return ""

    sentences = [compact_text(item) for item in re.split(r"(?<=[.!?])\s+", cleaned) if compact_text(item)]
    if len(sentences) <= 1:
        return cleaned

    filtered_sentences: list[str] = []
    for sentence in sentences:
        if not filtered_sentences:
            filtered_sentences.append(sentence)
            continue
        if any(teacher_text_similarity(sentence, existing) >= 0.93 for existing in filtered_sentences[-3:]):
            continue
        filtered_sentences.append(sentence)

    return compact_text(" ".join(filtered_sentences), cleaned)


def is_teacher_segment_redundant(
    section_heading: str,
    text: str,
    existing_segments: list[dict[str, Any]],
) -> bool:
    cleaned_heading = compact_text(section_heading)
    cleaned_text = clean_teacher_segment_text(text)
    if not cleaned_text:
        return True

    for existing in existing_segments:
        existing_heading = compact_text(existing.get("section_heading"))
        existing_text = compact_text(existing.get("text"))
        if not existing_text:
            continue
        similarity = teacher_text_similarity(cleaned_text, existing_text)
        threshold = 0.78 if existing_heading == cleaned_heading else 0.9
        if similarity >= threshold:
            return True
    return False


def dedupe_teacher_segments(
    segments: list[dict[str, Any]],
    allowed_headings: list[str] | None = None,
    *,
    preserve_required_sections: bool = False,
) -> list[dict[str, Any]]:
    if not segments:
        return []

    allowed = set(allowed_headings or [])
    required_counts = TEACHER_REQUIRED_SECTION_COUNTS if preserve_required_sections else {}
    filtered_segments: list[dict[str, Any]] = []
    heading_counts: dict[str, int] = {}

    for raw_segment in segments:
        heading = compact_text(raw_segment.get("section_heading"))
        if allowed and heading not in allowed:
            continue
        prompt = compact_text(raw_segment.get("prompt"))
        text = clean_teacher_segment_text(raw_segment.get("text"))
        if not text:
            continue

        is_redundant = is_teacher_segment_redundant(heading, text, filtered_segments)
        if is_redundant and heading_counts.get(heading, 0) < required_counts.get(heading, 0):
            is_redundant = False
        if is_redundant:
            continue

        filtered_segments.append(
            {
                **raw_segment,
                "index": len(filtered_segments) + 1,
                "section_heading": heading,
                "prompt": prompt,
                "text": text,
                "estimated_minutes": estimate_spoken_minutes(text),
            }
        )
        heading_counts[heading] = heading_counts.get(heading, 0) + 1
        if len(filtered_segments) >= TEACHER_SEGMENT_LIMIT:
            break

    return filtered_segments[:TEACHER_SEGMENT_LIMIT]


def normalize_teacher_segments(
    raw_segments: Any,
    fallback_segments: list[dict[str, Any]],
    allowed_headings: list[str],
) -> list[dict[str, Any]]:
    allowed = set(allowed_headings)
    source_segments = raw_segments if isinstance(raw_segments, list) else []
    normalized_segments: list[dict[str, Any]] = []

    for raw_segment in source_segments:
        if not isinstance(raw_segment, dict):
            continue
        section_heading = compact_text(raw_segment.get("section_heading"))
        if section_heading not in allowed:
            continue
        prompt = compact_text(raw_segment.get("prompt"))
        text = clean_teacher_segment_text(raw_segment.get("text"))
        if not text:
            continue
        for chunk_index, chunk in enumerate(split_podcast_text(text, max_chars=920, max_words=145), start=1):
            normalized_segments.append(
                {
                    "index": len(normalized_segments) + 1,
                    "section_heading": section_heading,
                    "prompt": prompt if chunk_index == 1 else "",
                    "text": chunk,
                    "estimated_minutes": estimate_spoken_minutes(chunk),
                }
            )
            if len(normalized_segments) >= TEACHER_SEGMENT_LIMIT:
                return normalized_segments[:TEACHER_SEGMENT_LIMIT]

    if normalized_segments:
        return normalized_segments
    return fallback_segments[:TEACHER_SEGMENT_LIMIT]


def build_teacher_lesson_fallback(
    summary: str,
    transcript: str,
    lecture_notes: str = "",
    lecture_slides: str = "",
    past_question_papers: str = "",
) -> dict[str, Any]:
    outline = build_teacher_section_outline(
        summary,
        lecture_notes,
        lecture_slides,
        past_question_papers,
        transcript,
    )
    title_lines = [line.strip() for line in extract_section(summary, "LECTURE TITLE").splitlines() if line.strip()]
    topic = title_lines[0] if title_lines else "This Lecture Topic"
    short_summary = compact_text(extract_section(summary, "SHORT SUMMARY"))
    fallback_source_text = "\n\n".join(
        part.strip()
        for part in [transcript, lecture_notes, lecture_slides, past_question_papers]
        if compact_text(part)
    )
    transcript_chunks = split_podcast_text(fallback_source_text, max_chars=430, max_words=70)
    segments: list[dict[str, Any]] = []
    transcript_index = 0

    section_openers = {
        "LECTURE TITLE": "Let us anchor the topic first so every later detail has a clear home.",
        "SHORT SUMMARY": "Let us start with the big picture before we dive into the detail.",
        "KEY CONCEPTS": "These are the core ideas holding the rest of the topic together.",
        "IMPORTANT DEFINITIONS": "This is one of those sections where precise wording protects marks.",
        "IMPORTANT FORMULAS": "Let us make the formula intuitive before we try to use it.",
        "WORKED EXAMPLES": "This is where the method becomes visible step by step.",
        "STEP-BY-STEP EXPLANATIONS": "Let us slow the process down and make every step earn its place.",
        "COMMON MISTAKES TO AVOID": "This is where we catch the traps before they catch your marks.",
        "REAL-WORLD EXAMPLES": "This is where the topic becomes easier to picture in real situations.",
        "EXAM TIPS": "Now let us translate understanding into marks.",
    }

    section_prompts = {
        "LECTURE TITLE": "Topic purpose",
        "SHORT SUMMARY": "Big-picture meaning",
        "KEY CONCEPTS": "Core idea links",
        "IMPORTANT DEFINITIONS": "Definition precision",
        "IMPORTANT FORMULAS": "Formula intuition",
        "WORKED EXAMPLES": "Next method step",
        "STEP-BY-STEP EXPLANATIONS": "Hidden step check",
        "COMMON MISTAKES TO AVOID": "Trap to avoid",
        "REAL-WORLD EXAMPLES": "Real-world anchor",
        "EXAM TIPS": "Marks strategy",
    }

    section_questions = {
        "LECTURE TITLE": "Before we even calculate anything, what kind of problem is this topic really helping you solve?",
        "SHORT SUMMARY": "What is the main message the lecturer wants you to keep in mind here?",
        "KEY CONCEPTS": "If one of these concepts disappeared, what would stop making sense first?",
        "IMPORTANT DEFINITIONS": "Which word in this definition would you not want to misread in a test?",
        "IMPORTANT FORMULAS": "What changes in the result when one term in this formula changes?",
        "WORKED EXAMPLES": "Before I reveal the next move, what do you think the next step should be?",
        "STEP-BY-STEP EXPLANATIONS": "Can you see which step students most often rush past too quickly?",
        "COMMON MISTAKES TO AVOID": "Which mistake feels easiest to make when you are under pressure?",
        "REAL-WORLD EXAMPLES": "Where would you actually see this outside the classroom?",
        "EXAM TIPS": "If this appeared for marks tomorrow, what would you write first?",
    }

    section_focus = {
        "LECTURE TITLE": "Connect the topic name to the kind of reasoning, calculation, comparison, or argument it usually requires.",
        "SHORT SUMMARY": "Stay with the high-level meaning first so the later details have somewhere sensible to live.",
        "KEY CONCEPTS": "Notice how the ideas connect instead of memorizing them as isolated pieces.",
        "IMPORTANT DEFINITIONS": "Definitions become easier when you spot which words are doing the precise technical work.",
        "IMPORTANT FORMULAS": "The goal is to know when the formula applies, what each term is doing, and what assumption is hiding inside it.",
        "WORKED EXAMPLES": "This is where we slow down properly because method matters more than pretending the final answer appeared by itself.",
        "STEP-BY-STEP EXPLANATIONS": "Think of this as the careful road map that stops you from skipping hidden steps.",
        "ADVANTAGES AND DISADVANTAGES": "This helps you judge when the idea is useful and where its limitations begin.",
        "COMMON MISTAKES TO AVOID": "If you can spot the trap early, you save both marks and stress.",
        "QUICK REVISION PLAN": "A short plan is useful because focused revision beats rushed revision.",
        "VISUAL AIDS": "Picture the structure, motion, or pattern so the concept feels visible rather than abstract only.",
        "REAL-WORLD EXAMPLES": "Real examples make the topic easier to recognize, remember, and explain.",
        "EXAM TIPS": "This is where understanding gets converted into marks under time pressure.",
    }

    def next_transcript_chunk() -> str:
        nonlocal transcript_index
        if transcript_index >= len(transcript_chunks):
            return ""
        chunk = transcript_chunks[transcript_index]
        transcript_index += 1
        return chunk

    for section in outline:
        heading = section["section_heading"]
        section_content = short_summary if heading == "SHORT SUMMARY" and short_summary else section["content"]
        content_chunks = split_podcast_text(section_content, max_chars=360, max_words=58)
        if not content_chunks:
            content_chunks = [compact_text(section_content[:900])]
        content_chunks = [chunk for chunk in content_chunks if chunk]
        if not content_chunks:
            continue
        target_segments = TEACHER_SECTION_TARGET_SEGMENTS.get(heading, 2)
        for segment_number in range(target_segments):
            primary_detail = content_chunks[segment_number % len(content_chunks)]
            transcript_detail = ""
            if heading in {"WORKED EXAMPLES", "STEP-BY-STEP EXPLANATIONS", "IMPORTANT FORMULAS", "KEY CONCEPTS"}:
                transcript_detail = next_transcript_chunk()
            elif segment_number == target_segments - 1:
                transcript_detail = next_transcript_chunk()

            spoken_text = (
                f"{section_openers.get(heading, 'Let us unpack this carefully.')} "
                f"{section_focus.get(heading, 'Keep the idea simple first, then build the detail.')} "
                f"{primary_detail} "
            )
            if heading in {"WORKED EXAMPLES", "STEP-BY-STEP EXPLANATIONS"}:
                spoken_text += (
                    "Let us slow this down carefully: identify what the question is asking, choose the rule or principle, "
                    "check why that choice fits, and only then move to the next line. "
                )
            if heading == "IMPORTANT FORMULAS":
                spoken_text += (
                    "As you revise, keep asking what each term is doing, what changes when one quantity increases, "
                    "and what assumption the formula quietly depends on. "
                )
            if transcript_detail:
                spoken_text += f"{transcript_detail} "
            spoken_text += (
                f"{section_questions.get(heading, 'What do you think is the key idea here?')} "
                "If the answer still feels fuzzy, that is fine. We are building understanding one clear step at a time."
            )

            for chunk_index, chunk in enumerate(split_podcast_text(spoken_text, max_chars=920, max_words=145), start=1):
                prompt = section_prompts.get(heading, "What do you think is the key idea here?") if chunk_index == 1 and segment_number == 0 else ""
                segments.append(
                    {
                        "index": len(segments) + 1,
                        "section_heading": heading,
                        "prompt": prompt,
                        "text": chunk,
                        "estimated_minutes": estimate_spoken_minutes(chunk),
                    }
                )
                if len(segments) >= TEACHER_SEGMENT_LIMIT:
                    break
            if len(segments) >= TEACHER_SEGMENT_LIMIT:
                break
        if len(segments) >= TEACHER_SEGMENT_LIMIT:
            break

    if estimate_podcast_total_minutes(segments) < TEACHER_MINIMUM_MINUTES:
        recap_sections = outline or [{"section_heading": "SHORT SUMMARY", "content": short_summary or compact_text(transcript[:700])}]
        recap_round = 0
        while estimate_podcast_total_minutes(segments) < TEACHER_MINIMUM_MINUTES and len(segments) < TEACHER_SEGMENT_LIMIT:
            section = recap_sections[recap_round % len(recap_sections)]
            heading = section["section_heading"]
            recap_content = compact_text(section["content"])
            transcript_detail = next_transcript_chunk()
            spoken_text = (
                f"Let us circle back to {heading.lower()} because understanding usually becomes stronger on the second pass. "
                f"{recap_content[:360]} "
            )
            if heading in {"WORKED EXAMPLES", "STEP-BY-STEP EXPLANATIONS"}:
                spoken_text += (
                    "Say the next move to yourself before I say it, because that habit is what turns an example into a skill. "
                )
            if transcript_detail:
                spoken_text += f"{transcript_detail} "
            spoken_text += "Ask yourself what changed in your understanding compared with a few minutes ago."
            for chunk in split_podcast_text(spoken_text, max_chars=920, max_words=145):
                segments.append(
                    {
                        "index": len(segments) + 1,
                        "section_heading": heading,
                        "prompt": "",
                        "text": chunk,
                        "estimated_minutes": estimate_spoken_minutes(chunk),
                    }
                )
                if estimate_podcast_total_minutes(segments) >= TEACHER_MINIMUM_MINUTES or len(segments) >= TEACHER_SEGMENT_LIMIT:
                    break
            recap_round += 1

    if not segments:
        fallback_points = transcript_chunks[:12] or [compact_text(transcript[:700], "Let us revise the most important lecture ideas clearly and calmly.")]
        for point in fallback_points:
            spoken_text = (
                "Let us work through this like a real revision class. "
                f"{point} "
                "Pause for a second, ask yourself what should happen next, and then check whether your reasoning still fits the topic."
            )
            segments.append(
                {
                    "index": len(segments) + 1,
                    "section_heading": "SHORT SUMMARY",
                    "prompt": "Next reasoning step",
                    "text": spoken_text,
                    "estimated_minutes": estimate_spoken_minutes(spoken_text),
                }
            )
            if len(segments) >= TEACHER_SEGMENT_LIMIT:
                break

    return {
        "title": f"{topic} Teacher Lesson",
        "overview": (
            f"A tutor-style walkthrough of {topic} designed for understanding, retention, and exam performance, "
            f"running for about {TEACHER_TARGET_MINUTES} minutes or more while following the guide section by section."
        ),
        "segments": dedupe_teacher_segments(
            segments[:TEACHER_SEGMENT_LIMIT],
            [item["section_heading"] for item in outline] or ["SHORT SUMMARY"],
        ),
    }


def sort_teacher_segments_by_heading_order(
    segments: list[dict[str, Any]],
    allowed_headings: list[str],
) -> list[dict[str, Any]]:
    heading_rank = {
        heading: teacher_heading_rank(heading)
        for heading in allowed_headings
    }
    ordered = sorted(
        enumerate(segments),
        key=lambda item: (
            heading_rank.get(
                compact_text(item[1].get("section_heading")),
                teacher_heading_rank(compact_text(item[1].get("section_heading"))),
            ),
            item[0],
        ),
    )
    return [
        {
            **segment,
            "index": index,
        }
        for index, (_, segment) in enumerate(ordered[:TEACHER_SEGMENT_LIMIT], start=1)
    ]


def ensure_teacher_section_coverage(
    segments: list[dict[str, Any]],
    fallback_segments: list[dict[str, Any]],
    allowed_headings: list[str],
) -> list[dict[str, Any]]:
    if not segments:
        return fallback_segments[:TEACHER_SEGMENT_LIMIT]

    allowed_set = set(allowed_headings)
    heading_counts: dict[str, int] = {}
    for segment in segments:
        heading = compact_text(segment.get("section_heading"))
        if heading:
            heading_counts[heading] = heading_counts.get(heading, 0) + 1

    needed_headings = {
        heading: count
        for heading, count in TEACHER_REQUIRED_SECTION_COUNTS.items()
        if heading in allowed_set and heading_counts.get(heading, 0) < count
    }
    if not needed_headings:
        return sort_teacher_segments_by_heading_order(segments, allowed_headings)

    extended_segments = list(segments)

    for fallback_segment in fallback_segments:
        heading = compact_text(fallback_segment.get("section_heading"))
        if heading not in needed_headings:
            continue
        cleaned_text = clean_teacher_segment_text(fallback_segment.get("text"))
        if not cleaned_text:
            continue
        if is_teacher_segment_redundant(heading, cleaned_text, extended_segments):
            continue
        extended_segments.append(
            {
                **fallback_segment,
                "index": len(extended_segments) + 1,
                "text": cleaned_text,
                "estimated_minutes": estimate_spoken_minutes(cleaned_text),
            }
        )
        heading_counts[heading] = heading_counts.get(heading, 0) + 1
        if heading_counts[heading] >= needed_headings[heading]:
            del needed_headings[heading]
        if not needed_headings or len(extended_segments) >= TEACHER_SEGMENT_LIMIT:
            break

    return sort_teacher_segments_by_heading_order(extended_segments[:TEACHER_SEGMENT_LIMIT], allowed_headings)


def extend_teacher_segments_to_target(
    segments: list[dict[str, Any]],
    fallback_segments: list[dict[str, Any]],
) -> list[dict[str, Any]]:
    if estimate_podcast_total_minutes(segments) >= TEACHER_MINIMUM_MINUTES:
        return segments[:TEACHER_SEGMENT_LIMIT]

    extended_segments = list(segments)

    for fallback_segment in fallback_segments:
        heading = compact_text(fallback_segment.get("section_heading"))
        cleaned_text = clean_teacher_segment_text(fallback_segment.get("text"))
        if not cleaned_text:
            continue
        if is_teacher_segment_redundant(heading, cleaned_text, extended_segments):
            continue
        extended_segments.append(
            {
                **fallback_segment,
                "index": len(extended_segments) + 1,
                "text": cleaned_text,
                "estimated_minutes": estimate_spoken_minutes(cleaned_text),
            }
        )
        if estimate_podcast_total_minutes(extended_segments) >= TEACHER_MINIMUM_MINUTES or len(extended_segments) >= TEACHER_SEGMENT_LIMIT:
            break

    return extended_segments[:TEACHER_SEGMENT_LIMIT]


def estimate_spoken_minutes(text: str) -> float:
    words = len(re.findall(r"\b\w+\b", text or ""))
    if not words:
        return 0.0
    return round(words / 140, 1)


def estimate_podcast_total_minutes(segments: list[dict[str, Any]]) -> float:
    return round(sum(float(segment.get("estimated_minutes", 0) or 0) for segment in segments), 1)


def split_podcast_text(text: str, *, max_chars: int = 700, max_words: int = 110) -> list[str]:
    cleaned = re.sub(r"\s+", " ", compact_text(text)).strip()
    if not cleaned:
        return []

    sentences = [item.strip() for item in re.split(r"(?<=[.!?])\s+", cleaned) if item.strip()]
    if not sentences:
        return [cleaned[:max_chars].strip()]

    chunks: list[str] = []
    current: list[str] = []
    current_chars = 0
    current_words = 0

    for sentence in sentences:
        sentence_words = len(re.findall(r"\b\w+\b", sentence))
        sentence_chars = len(sentence) + (1 if current else 0)
        if current and (current_chars + sentence_chars > max_chars or current_words + sentence_words > max_words):
            chunks.append(" ".join(current).strip())
            current = [sentence]
            current_chars = len(sentence)
            current_words = sentence_words
            continue
        current.append(sentence)
        current_chars += sentence_chars
        current_words += sentence_words

    if current:
        chunks.append(" ".join(current).strip())

    return [chunk for chunk in chunks if chunk]


def build_podcast_fallback(
    summary: str,
    transcript: str,
    speaker_profiles: list[dict[str, str]],
    target_minutes: int,
) -> dict[str, Any]:
    title_lines = [line.strip() for line in extract_section(summary, "LECTURE TITLE").splitlines() if line.strip()]
    topic = title_lines[0] if title_lines else "This Lecture Topic"
    short_summary = compact_text(extract_section(summary, "SHORT SUMMARY"))
    concept_points = extract_bullet_points(extract_section(summary, "KEY CONCEPTS"))[:6]
    example_points = [line for line in extract_section(summary, "WORKED EXAMPLES").splitlines() if compact_text(line)][:4]
    mistake_points = extract_bullet_points(extract_section(summary, "COMMON MISTAKES TO AVOID"))[:3]
    revision_points = extract_bullet_points(extract_section(summary, "QUICK REVISION PLAN"))[:3]

    if not short_summary:
        short_summary = compact_text(transcript[:700], f"A revision debate about {topic}.")

    transcript_chunks = split_podcast_text(transcript, max_chars=520, max_words=85)
    speaking_points = [short_summary]
    speaking_points.extend(concept_points)
    speaking_points.extend(example_points)
    speaking_points.extend(mistake_points)
    speaking_points.extend(revision_points)
    speaking_points.extend(transcript_chunks[:18])
    speaking_points = [compact_text(item) for item in speaking_points if compact_text(item)]
    minimum_minutes = max(5.5, clamp_podcast_target_minutes(target_minutes) * 0.92)

    normalized_segments: list[dict[str, Any]] = []
    for index, point in enumerate(speaking_points, start=1):
        profile = speaker_profiles[(index - 1) % len(speaker_profiles)]
        text = point
        if index == 1:
            text = f"Today we're breaking down {topic}. {point}"
        elif index == 2:
            text = f"Let us slow that down for a student who is hearing it for the first time. {point}"
        elif "example" in point.lower():
            text = f"Here is the kind of example that makes the topic stick. {point}"
        elif "mistake" in point.lower() or "avoid" in point.lower():
            text = f"This is where students usually lose marks. {point}"
        elif "minute" in point.lower() or "plan" in point.lower():
            text = f"If exam week is close, this is the revision move to make. {point}"

        for chunk in split_podcast_text(text):
            normalized_segments.append(
                {
                    "index": len(normalized_segments) + 1,
                    "speaker_key": profile["key"],
                    "speaker_name": profile["name"],
                    "speaker_role": profile["role"],
                    "voice": profile["voice"],
                    "voice_style": profile["voice_style"],
                    "text": chunk,
                    "estimated_minutes": estimate_spoken_minutes(chunk),
                }
            )
        if estimate_podcast_total_minutes(normalized_segments) >= minimum_minutes and len(normalized_segments) >= 10:
            break

    overview = (
        f"A fallback revision debate about {topic}, mixing clear explanations, example-driven discussion, "
        "and exam-focused reminders."
    )
    return {
        "title": f"{topic} Debate Podcast",
        "overview": overview,
        "segments": normalized_segments,
    }


def extend_podcast_segments_to_target(
    segments: list[dict[str, Any]],
    fallback_segments: list[dict[str, Any]],
    target_minutes: int,
) -> list[dict[str, Any]]:
    minimum_minutes = max(5.5, clamp_podcast_target_minutes(target_minutes) * 0.92)
    if estimate_podcast_total_minutes(segments) >= minimum_minutes:
        return segments

    seen_texts = {compact_text(segment.get("text")) for segment in segments if compact_text(segment.get("text"))}
    extended_segments = list(segments)
    for fallback_segment in fallback_segments:
        cleaned_text = compact_text(fallback_segment.get("text"))
        if not cleaned_text or cleaned_text in seen_texts:
            continue
        seen_texts.add(cleaned_text)
        extended_segments.append(
            {
                **fallback_segment,
                "index": len(extended_segments) + 1,
            }
        )
        if estimate_podcast_total_minutes(extended_segments) >= minimum_minutes:
            break
    return extended_segments


def normalize_podcast_segments(
    raw_segments: Any,
    speaker_profiles: list[dict[str, str]],
) -> list[dict[str, Any]]:
    if not isinstance(raw_segments, list):
        return []

    profile_by_key = {profile["key"]: profile for profile in speaker_profiles}
    profile_keys = list(profile_by_key)
    normalized_segments: list[dict[str, Any]] = []
    last_key = ""

    for raw_segment in raw_segments:
        if not isinstance(raw_segment, dict):
            continue

        speaker_key = compact_text(raw_segment.get("speaker_key")).lower()
        if speaker_key not in profile_by_key:
            speaker_label = compact_text(raw_segment.get("speaker")).lower()
            matching_profile = next(
                (
                    profile
                    for profile in speaker_profiles
                    if speaker_label in {profile["key"].lower(), profile["name"].lower()}
                ),
                None,
            )
            if matching_profile:
                speaker_key = matching_profile["key"]

        if speaker_key not in profile_by_key:
            speaker_key = last_key or profile_keys[len(normalized_segments) % len(profile_keys)]

        text = compact_text(raw_segment.get("text"))
        if not text:
            continue

        profile = profile_by_key[speaker_key]
        for chunk in split_podcast_text(text):
            normalized_segments.append(
                {
                    "index": len(normalized_segments) + 1,
                    "speaker_key": profile["key"],
                    "speaker_name": profile["name"],
                    "speaker_role": profile["role"],
                    "voice": profile["voice"],
                    "voice_style": profile["voice_style"],
                    "text": chunk,
                    "estimated_minutes": estimate_spoken_minutes(chunk),
                }
            )
        last_key = speaker_key

    return normalized_segments[:24]


def build_podcast_script_markdown(
    title: str,
    overview: str,
    segments: list[dict[str, Any]],
) -> str:
    blocks = [
        "**PODCAST TITLE**",
        "",
        compact_text(title, "Lecture Debate Podcast"),
        "",
        "**PODCAST OVERVIEW**",
        "",
        compact_text(overview, "A multi-speaker revision debate."),
        "",
        "**DEBATE SCRIPT**",
        "",
    ]

    for segment in segments:
        blocks.append(f"{segment['index']}. {segment['speaker_name']}: {segment['text']}")
        blocks.append("")

    return "\n".join(blocks).strip()


def build_downloadable_podcast_file(output_dir: Path, audio_files: list[str]) -> str:
    combined_path = output_dir / "full-podcast.mp3"
    with combined_path.open("wb") as combined_file:
        for file_name in audio_files:
            combined_file.write(Path(file_name).read_bytes())
    return str(combined_path)


def synthesize_podcast_audio_segments(job_id: str, segments: list[dict[str, Any]]) -> tuple[list[str], str]:
    output_dir = PODCAST_OUTPUT_DIR / job_id
    if output_dir.exists():
        shutil.rmtree(output_dir, ignore_errors=True)
    output_dir.mkdir(parents=True, exist_ok=True)

    total_segments = max(1, len(segments))
    audio_files: list[str] = []
    for index, segment in enumerate(segments, start=1):
        update_job(
            job_id,
            status="processing",
            stage=f"Recording podcast voice {index} of {total_segments}",
            progress=min(98, 42 + int((index / total_segments) * 54)),
        )
        response = client.with_options(timeout=PODCAST_TTS_TIMEOUT).audio.speech.create(
            model=PODCAST_TTS_MODEL,
            voice=segment["voice"],
            input=segment["text"],
            instructions=(
                f"You are {segment['speaker_name']}. {segment['speaker_role']}. "
                f"{segment.get('voice_style', 'Sound natural, clear, and human.')} "
                "Keep the delivery natural, human, and podcast-like."
            ),
            response_format="mp3",
        )
        audio_bytes = response.content if response.content else response.read()
        file_path = output_dir / f"{index:03d}-{segment['speaker_key']}.mp3"
        file_path.write_bytes(audio_bytes)
        audio_files.append(str(file_path))

    download_file = build_downloadable_podcast_file(output_dir, audio_files)
    return audio_files, download_file


async def generate_podcast_package(
    summary: str,
    transcript: str,
    lecture_notes: str,
    lecture_slides: str,
    past_question_papers: str,
    speaker_count: int,
    target_minutes: int,
    job_id: str,
    output_language: str,
) -> dict[str, Any]:
    normalized_speaker_count = clamp_podcast_speaker_count(speaker_count)
    normalized_target_minutes = clamp_podcast_target_minutes(target_minutes)
    speaker_profiles = build_podcast_speaker_profiles(normalized_speaker_count)
    target_turns = 12 if normalized_speaker_count == 2 else 15
    target_words = normalized_target_minutes * 140

    source_blocks = [
        trimmed_context_block("STUDY GUIDE SUMMARY", summary, MAX_PODCAST_CONTEXT_CHARS // 2),
        trimmed_context_block("LECTURE NOTES", lecture_notes, MAX_PODCAST_CONTEXT_CHARS // 4),
        trimmed_context_block("LECTURE SLIDES", lecture_slides, MAX_PODCAST_CONTEXT_CHARS // 4),
        trimmed_context_block("PAST QUESTION PAPERS", past_question_papers, MAX_PODCAST_CONTEXT_CHARS // 4),
        trimmed_context_block("LECTURE TRANSCRIPT", transcript, MAX_PODCAST_CONTEXT_CHARS // 2),
        "SPEAKER PROFILES\n" + json.dumps(speaker_profiles, ensure_ascii=False, indent=2),
    ]
    combined_source = "\n\n".join(block for block in source_blocks if block)

    def _generate_podcast_script(revision_note: str = "") -> dict[str, Any]:
        response = client.with_options(timeout=PODCAST_REQUEST_TIMEOUT).chat.completions.create(
            model=PODCAST_SCRIPT_MODEL,
            max_completion_tokens=min(MAX_COMPLETION_TOKENS, 5000),
            messages=[
                {
                    "role": "system",
                    "content": (
                        "You create educational podcast debates for students. "
                        "Return strict JSON only with these keys: title, overview, estimated_minutes, segments.\n\n"
                        "Rules:\n"
                        "- `segments` must be an array of objects with `speaker_key` and `text` only.\n"
                        "- Use exactly the provided speaker keys.\n"
                        f"- Write a lively debate for {normalized_speaker_count} speakers.\n"
                        f"- Target about {normalized_target_minutes} minutes and roughly {target_words} spoken words in total.\n"
                        f"- Do not return a script shorter than about {max(6, normalized_target_minutes - 1)} minutes.\n"
                        f"- Aim for about {target_turns} turns before any automatic text splitting.\n"
                        "- Make it feel like a real spoken discussion, not lecture notes being read aloud.\n"
                        "- Mix humor with seriousness, but keep the academic content accurate.\n"
                        "- Use simple jokes or playful challenges, never nonsense or disrespectful humor.\n"
                        "- Keep each turn focused on understanding, application, or exam reasoning.\n"
                        "- Mention worked examples from the lecture when they are present in the source material.\n"
                        "- If an example is inferred rather than clearly from the lecture, frame it as a fresh example instead of claiming the lecturer said it.\n"
                        "- Include clarifying analogies, common mistakes, and exam-useful examples.\n"
                        "- Do not use stage directions, bullet points, markdown, or narration tags like [music].\n"
                        "- Do not use emojis.\n"
                        "- Keep every spoken turn between about 35 and 95 words.\n"
                        f"- Write the whole podcast in {output_language}."
                    ),
                },
                {
                    "role": "user",
                    "content": (
                        "Build the podcast debate from this lecture material.\n\n"
                        + (f"{revision_note.strip()}\n\n" if revision_note.strip() else "")
                        + combined_source
                    ),
                },
            ],
        )
        return parse_json_object(response.choices[0].message.content or "")

    update_job(job_id, status="processing", stage="Writing podcast debate", progress=18)
    fallback_package = build_podcast_fallback(summary, transcript, speaker_profiles, normalized_target_minutes)

    try:
        generated_package = await asyncio.to_thread(_generate_podcast_script)
    except Exception as exc:
        logger.warning("Podcast script generation failed, using fallback podcast script: %s", exc)
        generated_package = {}

    normalized_segments = normalize_podcast_segments(generated_package.get("segments"), speaker_profiles)
    estimated_minutes = estimate_podcast_total_minutes(normalized_segments)
    if normalized_segments and estimated_minutes < max(5.5, normalized_target_minutes * 0.9):
        update_job(job_id, status="processing", stage="Extending podcast to the requested length", progress=28)
        try:
            generated_package = await asyncio.to_thread(
                _generate_podcast_script,
                (
                    f"The first draft landed at about {estimated_minutes:.1f} minutes, which is too short. "
                    f"Rewrite the whole podcast so it lands near {normalized_target_minutes} minutes, "
                    "with more explanation, worked examples, exam traps, and recap moments."
                ),
            )
            normalized_segments = normalize_podcast_segments(generated_package.get("segments"), speaker_profiles)
        except Exception as exc:
            logger.warning("Podcast rewrite for target length failed: %s", exc)

    if not normalized_segments:
        normalized_segments = fallback_package["segments"]
    normalized_segments = extend_podcast_segments_to_target(
        normalized_segments,
        fallback_package["segments"],
        normalized_target_minutes,
    )

    title = compact_text(generated_package.get("title"), fallback_package["title"])
    overview = compact_text(generated_package.get("overview"), fallback_package["overview"])
    script_markdown = build_podcast_script_markdown(title, overview, normalized_segments)

    update_job(job_id, status="processing", stage="Preparing speaker voices", progress=40)
    audio_files, download_file = await asyncio.to_thread(synthesize_podcast_audio_segments, job_id, normalized_segments)

    return {
        "podcast_title": title,
        "podcast_overview": overview,
        "podcast_script": script_markdown,
        "podcast_segments": normalized_segments,
        "_podcast_audio_files": audio_files,
        "_podcast_download_file": download_file,
    }


async def generate_teacher_lesson_package(
    summary: str,
    transcript: str,
    lecture_notes: str,
    lecture_slides: str,
    past_question_papers: str,
    job_id: str,
    output_language: str,
) -> dict[str, Any]:
    outline = build_teacher_section_outline(
        summary,
        lecture_notes,
        lecture_slides,
        past_question_papers,
        transcript,
    )
    fallback_package = build_teacher_lesson_fallback(
        summary,
        transcript,
        lecture_notes,
        lecture_slides,
        past_question_papers,
    )
    allowed_headings = [item["section_heading"] for item in outline] or ["SHORT SUMMARY"]
    outline_block = "\n".join(
        f"- {item['section_heading']}: {compact_text(item['content'])[:260]}"
        for item in outline
    ) or "- SHORT SUMMARY: Explain the main lecture idea clearly."

    context_blocks = [
        build_teacher_context_block("STUDY GUIDE SUMMARY", summary, MAX_STUDY_GUIDE_INPUT_CHARS),
        build_teacher_context_block("LECTURE NOTES", lecture_notes, MAX_STUDY_GUIDE_INPUT_CHARS // 2),
        build_teacher_context_block("LECTURE SLIDES", lecture_slides, MAX_STUDY_GUIDE_INPUT_CHARS // 2),
        build_teacher_context_block("PAST QUESTION PAPERS", past_question_papers, MAX_STUDY_GUIDE_INPUT_CHARS // 2),
        build_teacher_context_block("LECTURE TRANSCRIPT", transcript, MAX_TRANSCRIPT_STUDY_GUIDE_INPUT_CHARS // 2),
        f"AVAILABLE GUIDE HEADINGS\n{outline_block}",
    ]
    combined_source = "\n\n".join(block for block in context_blocks if block)

    def _generate_teacher_script(revision_note: str = "") -> dict[str, Any]:
        response = client.with_options(timeout=TEACHER_REQUEST_TIMEOUT).chat.completions.create(
            model=TEACHER_SCRIPT_MODEL,
            max_completion_tokens=min(MAX_COMPLETION_TOKENS, 6800),
            messages=[
                {
                    "role": "system",
                    "content": (
                        "You are an elite AI teaching engine built for university-level learning. "
                        "Return strict JSON only with the keys title, overview, and segments.\n\n"
                        "Final goal:\n"
                        "- Teach for understanding, retention, engagement, exam performance, and long-term memory.\n"
                        "- Feel like a top private tutor: highly intelligent but easy to understand, conversational, motivating, adaptive to learner difficulty, and visually synchronized with the study guide.\n\n"
                        "JSON rules:\n"
                        "- `segments` must be an array of objects with `section_heading`, `prompt`, and `text` only.\n"
                        "- `section_heading` must be one of the provided guide headings.\n"
                        "- `prompt` must be a short natural cue of about 4 to 10 words naming the live teaching focus for the UI, such as a definition, process step, comparison, misconception, or exam clue.\n"
                        "- `overview` must be 1 to 2 sentences describing the lesson style and coverage.\n\n"
                        "Teaching behavior:\n"
                        "- Explain concepts step by step.\n"
                        "- Introduce the idea simply, explain what it means, explain why it matters, give an intuitive example, connect it to previous concepts, and reinforce the key takeaway.\n"
                        "- Teach like speaking to a real student, not like reading notes line by line.\n"
                        "- Use calm, clear, encouraging language.\n"
                        "- Focus on why concepts matter, where students get confused, and how the idea appears in exams.\n"
                        "- Use chunking, active recall, analogy, real-world examples, comparisons, and memory anchors when they genuinely help.\n"
                        "- Ask occasional reflective questions naturally.\n"
                        "- Reinforce important ideas with a new angle instead of repeating the same wording.\n"
                        "- If the concept is difficult, slow the pacing, use shorter sentences, and add more examples.\n"
                        "- If the concept is easier, be concise and focus on high-yield details.\n"
                        "- Spend the biggest share of time on WORKED EXAMPLES and STEP-BY-STEP EXPLANATIONS when those sections exist.\n"
                        "- Keep IMPORTANT FORMULAS conceptual by explaining what each term is doing, when the formula applies, and what assumption it depends on.\n"
                        "- Occasionally point out common exam questions or common confusions when it fits.\n\n"
                        "Sync with the guide:\n"
                        "- Follow the guide section by section.\n"
                        "- Make each segment clearly about the current section only.\n"
                        "- Use `prompt` to name the exact teaching focus so the interface can sync the lesson to the guide.\n"
                        "- When discussing definitions, processes, comparisons, theories, assumptions, criticisms, formulas, or key terms, make that focus obvious in the `prompt` and spoken explanation.\n\n"
                        "Length and output rules:\n"
                        f"- The full lesson must land at roughly {TEACHER_MINIMUM_MINUTES:.0f} to {TEACHER_TARGET_MINUTES + 2} minutes and about {TEACHER_TARGET_WORDS} spoken words in total.\n"
                        "- Aim for about 24 to 36 teaching segments before any automatic chunk splitting.\n"
                        "- Avoid repeating the same definition, example, or explanation unless you are adding a clearly new step, misconception fix, or application.\n"
                        "- Do not skip the opening lesson flow, top notes, or early definitions when the source starts with foundational explanations.\n"
                        "- Do not create spoken segments for FLASHCARDS or PRACTICE QUESTIONS AND ANSWERS.\n"
                        "- Keep each segment between about 95 and 160 spoken words.\n"
                        "- Do not use bullet points, markdown, stage directions, or sound-effect text.\n"
                        "- Do not output fake dialogue or multiple speakers.\n"
                        f"- Write everything in {output_language}."
                    ),
                },
                {
                    "role": "user",
                    "content": (
                        "Build a teacher-mode lesson that follows the guide section by section.\n"
                        "The student should feel personally guided by a premium AI tutor who teaches for true understanding, not just information delivery.\n"
                        "Do not explain flashcards. Do not explain practice questions and answers.\n"
                        "Spend extra time making worked examples, definitions, method steps, and exam-relevant misunderstandings very clear.\n"
                        "Do not restate the same point in later segments unless the later segment adds a new angle, deeper reason, or new application.\n"
                        "Cover the beginning of the notes properly before moving deeper into the lesson.\n"
                        "Think like an expert university tutor creating revision teaching before exams.\n\n"
                        + (f"{revision_note.strip()}\n\n" if revision_note.strip() else "")
                        + f"Allowed section headings: {', '.join(allowed_headings)}\n\n"
                        + combined_source
                    ),
                },
            ],
        )
        return parse_json_object(response.choices[0].message.content or "")

    update_job(job_id, status="processing", stage="Planning teacher lesson", progress=16)
    try:
        generated_package = await asyncio.to_thread(_generate_teacher_script)
    except Exception as exc:
        logger.warning("Teacher lesson generation failed, using fallback script: %s", exc)
        generated_package = {}

    normalized_segments = normalize_teacher_segments(
        generated_package.get("segments"),
        fallback_package["segments"],
        allowed_headings,
    )
    normalized_segments = dedupe_teacher_segments(
        normalized_segments,
        allowed_headings,
        preserve_required_sections=True,
    )
    normalized_segments = ensure_teacher_section_coverage(
        normalized_segments,
        fallback_package["segments"],
        allowed_headings,
    )
    estimated_minutes = estimate_podcast_total_minutes(normalized_segments)
    if normalized_segments and estimated_minutes < TEACHER_MINIMUM_MINUTES:
        update_job(job_id, status="processing", stage="Extending teacher lesson depth", progress=24)
        try:
            generated_package = await asyncio.to_thread(
                _generate_teacher_script,
                (
                    f"The first teacher draft landed at about {estimated_minutes:.1f} minutes, which is too short. "
                    f"Rewrite the whole lesson so it lands around {TEACHER_MINIMUM_MINUTES:.0f} to {TEACHER_TARGET_MINUTES + 2} minutes, "
                    "spend much more time on worked examples, definitions, and step-by-step reasoning, "
                    "do not skip the beginning of the lesson or the top lecturer notes, "
                    "and keep the explanation warm, natural, and non-repetitive."
                ),
            )
            normalized_segments = normalize_teacher_segments(
                generated_package.get("segments"),
                fallback_package["segments"],
                allowed_headings,
            )
            normalized_segments = dedupe_teacher_segments(
                normalized_segments,
                allowed_headings,
                preserve_required_sections=True,
            )
            normalized_segments = ensure_teacher_section_coverage(
                normalized_segments,
                fallback_package["segments"],
                allowed_headings,
            )
        except Exception as exc:
            logger.warning("Teacher lesson rewrite for target length failed: %s", exc)

    if not normalized_segments:
        normalized_segments = fallback_package["segments"]
    normalized_segments = extend_teacher_segments_to_target(normalized_segments, fallback_package["segments"])
    normalized_segments = ensure_teacher_section_coverage(
        normalized_segments,
        fallback_package["segments"],
        allowed_headings,
    )
    normalized_segments = dedupe_teacher_segments(
        normalized_segments,
        allowed_headings,
        preserve_required_sections=True,
    )
    if estimate_podcast_total_minutes(normalized_segments) < TEACHER_MINIMUM_MINUTES:
        normalized_segments = extend_teacher_segments_to_target(normalized_segments, fallback_package["segments"])
        normalized_segments = ensure_teacher_section_coverage(
            normalized_segments,
            fallback_package["segments"],
            allowed_headings,
        )
        normalized_segments = dedupe_teacher_segments(
            normalized_segments,
            allowed_headings,
            preserve_required_sections=True,
        )
    title = compact_text(generated_package.get("title"), fallback_package["title"])
    overview = compact_text(generated_package.get("overview"), fallback_package["overview"])

    return {
        "teacher_title": title,
        "teacher_overview": overview,
        "teacher_segments": normalized_segments,
    }


PRESENTATION_THEMES: dict[str, dict[str, str]] = {
    "emerald-scholar": {
        "id": "emerald-scholar",
        "style_family": "emerald-scholar",
        "name": "Emerald Scholar",
        "background": "061912",
        "surface": "0F2B20",
        "surface_alt": "15392A",
        "accent": "4ADE80",
        "accent_soft": "D9F99D",
        "text": "F8FAFC",
        "muted": "CFE8DA",
        "dark_text": "052E16",
    },
    "sunset-classroom": {
        "id": "sunset-classroom",
        "style_family": "sunset-classroom",
        "name": "Sunset Classroom",
        "background": "FFF8F1",
        "surface": "FFF1E1",
        "surface_alt": "F6D2B6",
        "accent": "EA580C",
        "accent_soft": "FDBA74",
        "text": "7C2D12",
        "muted": "9A3412",
        "dark_text": "7C2D12",
    },
    "midnight-grid": {
        "id": "midnight-grid",
        "style_family": "midnight-grid",
        "name": "Midnight Grid",
        "background": "020617",
        "surface": "0F172A",
        "surface_alt": "172554",
        "accent": "38BDF8",
        "accent_soft": "A5F3FC",
        "text": "E2E8F0",
        "muted": "BFDBFE",
        "dark_text": "082F49",
    },
    "aurora-waves": {
        "id": "aurora-waves",
        "style_family": "midnight-grid",
        "name": "Aurora Waves",
        "background": "04111F",
        "surface": "0B1730",
        "surface_alt": "1D2E67",
        "accent": "60A5FA",
        "accent_soft": "BFDBFE",
        "text": "EFF6FF",
        "muted": "DBEAFE",
        "dark_text": "1E3A8A",
    },
    "glass-cube": {
        "id": "glass-cube",
        "style_family": "sunset-classroom",
        "name": "Glass Cube",
        "background": "F0FDFA",
        "surface": "DDF8F3",
        "surface_alt": "BDEEE5",
        "accent": "0F766E",
        "accent_soft": "99F6E4",
        "text": "134E4A",
        "muted": "115E59",
        "dark_text": "134E4A",
    },
    "celebration-night": {
        "id": "celebration-night",
        "style_family": "midnight-grid",
        "name": "Celebration Night",
        "background": "111827",
        "surface": "1F2937",
        "surface_alt": "312E81",
        "accent": "F59E0B",
        "accent_soft": "FDE68A",
        "text": "F9FAFB",
        "muted": "E0E7FF",
        "dark_text": "78350F",
    },
    "amber-lux": {
        "id": "amber-lux",
        "style_family": "midnight-grid",
        "name": "Amber Lux",
        "background": "09090B",
        "surface": "18181B",
        "surface_alt": "3F3F46",
        "accent": "F59E0B",
        "accent_soft": "FCD34D",
        "text": "FAFAFA",
        "muted": "FDE68A",
        "dark_text": "78350F",
    },
    "editorial-sage": {
        "id": "editorial-sage",
        "style_family": "emerald-scholar",
        "name": "Editorial Sage",
        "background": "F7F7F2",
        "surface": "E9F2E4",
        "surface_alt": "D7E7CF",
        "accent": "2F6B4F",
        "accent_soft": "CBE7D6",
        "text": "1F3A2A",
        "muted": "466A55",
        "dark_text": "1F3A2A",
    },
    "clinical-blue": {
        "id": "clinical-blue",
        "style_family": "sunset-classroom",
        "name": "Clinical Blue",
        "background": "F4FBFF",
        "surface": "E7F4FB",
        "surface_alt": "D4ECF8",
        "accent": "2563EB",
        "accent_soft": "BFDBFE",
        "text": "1E3A8A",
        "muted": "3B82F6",
        "dark_text": "1E3A8A",
    },
    "festival-pop": {
        "id": "festival-pop",
        "style_family": "sunset-classroom",
        "name": "Festival Pop",
        "background": "FFF7FB",
        "surface": "FFE4F1",
        "surface_alt": "FED7E2",
        "accent": "DB2777",
        "accent_soft": "F9A8D4",
        "text": "9D174D",
        "muted": "BE185D",
        "dark_text": "9D174D",
    },
    "summit-minimal": {
        "id": "summit-minimal",
        "style_family": "sunset-classroom",
        "name": "Summit Minimal",
        "background": "FFFDEA",
        "surface": "F5F3FF",
        "surface_alt": "DBEAFE",
        "accent": "4F46E5",
        "accent_soft": "C7D2FE",
        "text": "3730A3",
        "muted": "4338CA",
        "dark_text": "3730A3",
    },
}


def ensure_presentation_support():
    if Presentation is None or RGBColor is None or MSO_SHAPE is None or PP_ALIGN is None or Inches is None or Pt is None:
        raise HTTPException(
            status_code=500,
            detail="PowerPoint generation is not configured on the server yet. Install python-pptx and redeploy.",
        )


def normalize_presentation_design_id(value: str) -> str:
    design_id = compact_text(value, "emerald-scholar").lower()
    return design_id if design_id in PRESENTATION_THEMES else "emerald-scholar"


def infer_presentation_visual_type(title: str, bullets: list[str]) -> str:
    combined = " ".join([title, *bullets]).lower()
    if any(marker in combined for marker in ["table", "tabulate", "rows", "columns", "summary table"]):
        return "table"
    if any(marker in combined for marker in ["graph", "trend", "increase", "decrease", "curve", "plot"]):
        return "graph"
    if any(marker in combined for marker in ["chart", "distribution", "breakdown", "percentage", "share"]):
        return "chart"
    if any(marker in combined for marker in ["component", "components", "input", "output", "module", "subsystem", "architecture", "framework", "pins", "signal path"]):
        return "components"
    if any(marker in combined for marker in ["photo", "photograph", "microscope", "specimen", "machine", "device", "organ", "appearance", "recognise"]):
        return "photo"
    if any(marker in combined for marker in ["compare", "difference", "advantage", "disadvantage", "versus"]):
        return "comparison"
    if any(marker in combined for marker in ["step", "process", "procedure", "method", "workflow"]):
        return "flow"
    if any(marker in combined for marker in ["timeline", "sequence", "revision plan", "history"]):
        return "timeline"
    if any(marker in combined for marker in ["cycle", "loop", "repeat"]):
        return "cycle"
    if any(marker in combined for marker in ["formula", "equation", "rule"]):
        return "formula"
    return "cluster"


PRESENTATION_EXCLUDED_PATTERNS = (
    "common mistake",
    "common mistakes",
    "exam tip",
    "exam tips",
    "revision plan",
    "revision focus",
    "quick revision",
    "study tip",
    "test strategy",
    "exam strategy",
    "exam trap",
    "exam traps",
    "before the test",
    "how to pass",
)

PRESENTATION_CLOSING_PATTERNS = (
    "thank you",
    "thanks",
    "q&a",
)

PRESENTATION_VISUAL_TYPES = {
    "flow",
    "comparison",
    "timeline",
    "cycle",
    "formula",
    "cluster",
    "components",
    "table",
    "chart",
    "graph",
    "photo",
    "closing",
}


def contains_presentation_excluded_text(value: str) -> bool:
    normalized = compact_text(value).lower()
    if not normalized:
        return False
    return any(pattern in normalized for pattern in PRESENTATION_EXCLUDED_PATTERNS)


def is_presentation_closing_text(value: str) -> bool:
    normalized = compact_text(value).lower()
    if not normalized:
        return False
    return any(pattern in normalized for pattern in PRESENTATION_CLOSING_PATTERNS)


def sanitize_presentation_text_items(values: list[str]) -> list[str]:
    return [
        compact_text(value)
        for value in values
        if compact_text(value)
        and not contains_presentation_excluded_text(compact_text(value))
        and not is_presentation_closing_text(compact_text(value))
    ]


def normalize_visual_items(raw_items: Any, fallback_items: list[str]) -> list[str]:
    items: list[str] = []
    if isinstance(raw_items, list):
        items = [compact_text(item) for item in raw_items if compact_text(item)]
    elif isinstance(raw_items, str):
        items = [compact_text(item) for item in re.split(r"\n+|;\s*", raw_items) if compact_text(item)]
    if items:
        return items[:4]
    return [compact_text(item) for item in fallback_items if compact_text(item)][:4]


def normalize_presentation_slides(raw_slides: Any) -> list[dict[str, Any]]:
    if not isinstance(raw_slides, list):
        return []

    normalized: list[dict[str, Any]] = []
    for raw_slide in raw_slides:
        if not isinstance(raw_slide, dict):
            continue
        title = compact_text(raw_slide.get("title"))
        if not title or contains_presentation_excluded_text(title) or is_presentation_closing_text(title):
            continue
        bullets = sanitize_presentation_text_items([
            compact_text(item)
            for item in (raw_slide.get("bullets") or [])
            if compact_text(item)
        ])
        if not title or len(bullets) < 2:
            continue
        visual_title = compact_text(raw_slide.get("visual_title"), compact_text(raw_slide.get("note"), "Visual summary"))
        if contains_presentation_excluded_text(visual_title) or is_presentation_closing_text(visual_title):
            visual_title = "Visual summary"
        visual_items = sanitize_presentation_text_items(normalize_visual_items(raw_slide.get("visual_items"), bullets[:4]))
        visual_type = compact_text(raw_slide.get("visual_type"), infer_presentation_visual_type(title, bullets)).lower()
        if visual_type not in PRESENTATION_VISUAL_TYPES:
            visual_type = infer_presentation_visual_type(title, bullets)
        flow_note = compact_text(
            raw_slide.get("flow_note"),
            f"This slide moves the deck from {title.lower()} into the next teachable idea.",
        )
        if contains_presentation_excluded_text(flow_note) or is_presentation_closing_text(flow_note):
            flow_note = f"This slide teaches {title.lower()} through a clearer classroom explanation."
        reference_image_index = -1
        try:
            reference_image_index = max(-1, int(raw_slide.get("reference_image_index")))
        except (TypeError, ValueError):
            if visual_type == "photo":
                reference_image_index = 0
        if visual_type == "photo" and reference_image_index < 0:
            visual_type = infer_presentation_visual_type(title, bullets)
            if visual_type == "photo":
                visual_type = "components"
        normalized.append(
            {
                "title": title,
                "bullets": bullets[:5],
                "visual_title": visual_title or "Visual summary",
                "visual_items": (visual_items or bullets[:3])[:4],
                "visual_type": visual_type,
                "flow_note": flow_note,
                "reference_image_index": reference_image_index,
            }
        )
        if len(normalized) >= 8:
            break
    return normalized


def build_presentation_fallback(summary: str, transcript: str) -> dict[str, Any]:
    title_lines = [line.strip() for line in extract_section(summary, "LECTURE TITLE").splitlines() if line.strip()]
    topic = title_lines[0] if title_lines else "Lecture Presentation"
    short_summary = compact_text(extract_section(summary, "SHORT SUMMARY"), compact_text(transcript[:240], "Lecture overview"))
    concept_points = extract_bullet_points(extract_section(summary, "KEY CONCEPTS"))[:8]
    definition_points = extract_bullet_points(extract_section(summary, "IMPORTANT DEFINITIONS"))[:6]
    formula_points = [line.strip() for line in extract_section(summary, "IMPORTANT FORMULAS").splitlines() if compact_text(line)][:4]
    example_points = [line.strip() for line in extract_section(summary, "WORKED EXAMPLES").splitlines() if compact_text(line)][:4]
    component_points = definition_points[:4] or concept_points[2:6]
    relationship_points = concept_points[3:8] or definition_points[1:5] or example_points[:4]

    slides = [
        {
            "title": "Overview",
            "bullets": [
                short_summary or "Introduce the lecture topic in one short overview.",
                concept_points[0] if len(concept_points) > 0 else "Explain the main direction of the lesson before going deeper.",
                concept_points[1] if len(concept_points) > 1 else "",
            ],
            "visual_title": "Lecture flow",
            "visual_type": "flow",
            "visual_items": ["Overview", "Core concepts", "Applied example"],
            "flow_note": "Open the deck with the topic, the route through the lesson, and what students should expect next.",
        },
        {
            "title": "Core Concepts",
            "bullets": concept_points[:5] or definition_points[:5] or [short_summary, "Explain the first major idea clearly."],
            "visual_title": "Concept cluster",
            "visual_type": "cluster",
            "visual_items": concept_points[:4],
            "flow_note": "Move from the overview into the key building blocks students must understand before details or examples.",
        },
        {
            "title": "Definitions and Terms",
            "bullets": definition_points[:5] or concept_points[3:8] or [short_summary, "Define the most important terminology clearly."],
            "visual_title": "Key terms",
            "visual_type": "comparison",
            "visual_items": definition_points[:4] or concept_points[:4],
            "flow_note": "Clarify vocabulary and distinctions so the rest of the deck uses the same language consistently.",
        },
        {
            "title": "System Components",
            "bullets": component_points[:4] or ["Show the main parts or inputs of the topic.", "Explain how the parts connect inside the overall system."],
            "visual_title": "Component diagram",
            "visual_type": "components",
            "visual_items": component_points[:4] or ["Core unit", "Input", "Output", "Control path"],
            "flow_note": "Map the important parts, inputs, or layers before moving into rules or calculations.",
        },
        {
            "title": "Formulas and Rules",
            "bullets": formula_points[:4] or ["Highlight the key formulas or rules from the lesson.", "Explain what each part represents."],
            "visual_title": "Formula sheet",
            "visual_type": "formula",
            "visual_items": formula_points[:4],
            "flow_note": "Introduce the rules, equations, or structured logic that support the worked reasoning in later slides.",
        },
        {
            "title": "Worked Example",
            "bullets": example_points[:4] or ["Walk through one representative example from the lecture.", "Point out the order of steps clearly."],
            "visual_title": "Method steps",
            "visual_type": "flow",
            "visual_items": example_points[:4],
            "flow_note": "Turn the theory into action with one clear example students can follow step by step.",
        },
        {
            "title": "Key Relationships",
            "bullets": relationship_points[:4] or ["Summarize how the main ideas connect to each other.", "Highlight the relationship between the concepts and the worked example."],
            "visual_title": "Compare or connect ideas",
            "visual_type": "comparison",
            "visual_items": relationship_points[:4] or component_points[:4],
            "flow_note": "Close the teaching sequence by connecting the main concepts into one final explanatory view.",
        },
    ]

    for slide in slides:
        bullets = [compact_text(item) for item in slide.get("bullets", []) if compact_text(item)]
        while len(bullets) < 2:
            bullets.append("Connect this point back to the lecture and explain why it matters.")
        slide["bullets"] = bullets[:5]

    normalized_slides = normalize_presentation_slides(slides)
    return {
        "title": topic,
        "subtitle": short_summary or "Lecture summary prepared for classroom presentation.",
        "slides": normalized_slides[:7],
    }


def build_presentation_closing_slide_content(title: str) -> dict[str, Any]:
    return {
        "title": "THANK YOU",
        "bullets": [],
        "visual_title": "",
        "visual_items": [],
        "visual_type": "closing",
        "flow_note": "",
        "reference_image_index": -1,
    }


def presentation_slides_to_text(title: str, subtitle: str, slides: list[dict[str, Any]]) -> str:
    blocks = [
        "POWERPOINT TITLE",
        title,
        "",
        "SUBTITLE",
        subtitle,
        "",
    ]
    for index, slide in enumerate(slides, start=1):
        blocks.append(f"SLIDE {index}: {slide['title']}")
        blocks.extend(f"- {bullet}" for bullet in slide.get("bullets", []))
        if slide.get("visual_title"):
            blocks.append(f"Visual panel: {slide['visual_title']} ({slide.get('visual_type', 'cluster')})")
        for item in slide.get("visual_items", []):
            blocks.append(f"  * {item}")
        if slide.get("flow_note"):
            blocks.append(f"Flow note: {slide['flow_note']}")
        blocks.append("")
    return "\n".join(blocks).strip()


def rgb_from_hex(value: str) -> RGBColor:
    return RGBColor.from_string(value.upper())


def get_presentation_style_family(theme: dict[str, str]) -> str:
    return compact_text(theme.get("style_family"), compact_text(theme.get("id"), "midnight-grid"))


def get_presentation_slide_layout(presentation: Any, *preferred_names: str, fallback_index: int = 6):
    layouts = list(getattr(presentation, "slide_layouts", []) or [])
    if not layouts:
        raise RuntimeError("The PowerPoint template has no slide layouts available.")

    lowered_names = [name.strip().lower() for name in preferred_names if name.strip()]
    for preferred_name in lowered_names:
        for layout in layouts:
            layout_name = compact_text(getattr(layout, "name", "")).lower()
            if preferred_name in layout_name:
                return layout

    safe_index = min(max(fallback_index, 0), len(layouts) - 1)
    return layouts[safe_index]


def clear_presentation_slides(presentation: Any):
    slide_id_list = getattr(presentation.slides, "_sldIdLst", None)
    if slide_id_list is None:
        return

    # Keep the template's masters/layouts, but rebuild the visible deck so the
    # downloaded PowerPoint matches the generated slide sequence shown in the app.
    for slide_id in list(slide_id_list):
        rel_id = slide_id.rId
        presentation.part.drop_rel(rel_id)
        slide_id_list.remove(slide_id)


def style_shape(shape: Any, fill_hex: str, *, transparency: float = 0.0):
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb_from_hex(fill_hex)
    if transparency:
        shape.fill.transparency = transparency
    shape.line.color.rgb = rgb_from_hex(fill_hex)
    shape.line.transparency = 1


def add_textbox(
    slide: Any,
    *,
    left: float,
    top: float,
    width: float,
    height: float,
    text: str,
    font_size: int,
    color_hex: str,
    bold: bool = False,
    align: Any = None,
    font_name: str = "Aptos",
):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    if MSO_ANCHOR is not None:
        frame.vertical_anchor = MSO_ANCHOR.TOP
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align or PP_ALIGN.LEFT
    run = paragraph.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = rgb_from_hex(color_hex)
    return box


def add_bullet_list(
    slide: Any,
    *,
    left: float,
    top: float,
    width: float,
    height: float,
    bullets: list[str],
    color_hex: str,
    font_name: str = "Aptos",
):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    if MSO_ANCHOR is not None:
        frame.vertical_anchor = MSO_ANCHOR.TOP

    for index, bullet in enumerate(bullets):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.alignment = PP_ALIGN.LEFT
        paragraph.level = 0
        paragraph.space_after = Pt(8)
        run = paragraph.add_run()
        run.text = f"• {bullet}"
        run.font.name = font_name
        run.font.size = Pt(19)
        run.font.color.rgb = rgb_from_hex(color_hex)
    return box


def add_visual_card(
    slide: Any,
    *,
    left: float,
    top: float,
    width: float,
    height: float,
    fill_hex: str,
    text: str,
    text_hex: str,
    font_size: int = 13,
    bold: bool = False,
):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    style_shape(card, fill_hex)
    add_textbox(
        slide,
        left=left + 0.14,
        top=top + 0.12,
        width=max(width - 0.28, 0.2),
        height=max(height - 0.24, 0.2),
        text=text,
        font_size=font_size,
        color_hex=text_hex,
        bold=bold,
    )
    return card


def load_presentation_reference_image(source: str) -> BytesIO | None:
    resolved_source = compact_text(source)
    if not resolved_source:
        return None

    if resolved_source.startswith("data:image/"):
        try:
            _, encoded = resolved_source.split(",", 1)
            return BytesIO(base64.b64decode(encoded))
        except (ValueError, binascii.Error):
            return None

    if resolved_source.startswith("http://") or resolved_source.startswith("https://"):
        try:
            response = requests.get(resolved_source, timeout=20)
            response.raise_for_status()
            return BytesIO(response.content)
        except Exception:
            return None
    return None


def resolve_reference_image_source(reference_images: list[str], slide_content: dict[str, Any]) -> str:
    if not reference_images:
        return ""

    try:
        requested_index = int(slide_content.get("reference_image_index"))
    except (TypeError, ValueError):
        requested_index = -1

    valid_sources = [compact_text(item) for item in reference_images if compact_text(item)]
    if not valid_sources:
        return ""
    if requested_index < 0:
        return ""
    if requested_index >= len(valid_sources):
        requested_index = 0
    return valid_sources[requested_index]


def draw_presentation_visual_panel(
    slide: Any,
    slide_content: dict[str, Any],
    theme: dict[str, str],
    reference_images: list[str],
):
    panel_left = 8.72
    panel_top = 1.95
    panel_width = 3.65
    panel_height = 4.73

    panel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(panel_left), Inches(panel_top), Inches(panel_width), Inches(panel_height))
    style_shape(panel, theme["accent"])
    add_textbox(
        slide,
        left=panel_left + 0.3,
        top=panel_top + 0.26,
        width=3.0,
        height=0.3,
        text=slide_content.get("visual_title") or "Visual frame",
        font_size=14,
        color_hex=theme["dark_text"],
        bold=True,
        font_name="Aptos Display",
    )

    style_family = get_presentation_style_family(theme)
    uses_dark_text = style_family == "sunset-classroom"
    primary_fill = theme["surface"] if not uses_dark_text else theme["surface_alt"]
    secondary_fill = theme["surface_alt"] if not uses_dark_text else theme["surface"]
    primary_text = theme["text"] if not uses_dark_text else theme["dark_text"]
    soft_fill = theme["accent_soft"] if not uses_dark_text else theme["surface"]

    visual_type = compact_text(slide_content.get("visual_type"), "cluster").lower()
    visual_items = [compact_text(item) for item in slide_content.get("visual_items", []) if compact_text(item)]
    if not visual_items:
        visual_items = [compact_text(item) for item in slide_content.get("bullets", []) if compact_text(item)][:4]
    reference_source = resolve_reference_image_source(reference_images, slide_content)
    reference_stream = load_presentation_reference_image(reference_source)

    def draw_chart_bars(items: list[str], *, line_mode: bool = False):
        axis = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(panel_left + 0.36), Inches(panel_top + 0.96), Inches(0.05), Inches(2.9))
        style_shape(axis, theme["surface"])
        base = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(panel_left + 0.36), Inches(panel_top + 3.81), Inches(2.78), Inches(0.05))
        style_shape(base, theme["surface"])
        for index, item in enumerate(items[:4] or ["Point A", "Point B", "Point C"]):
            normalized_item = compact_text(item, f"Point {index + 1}")
            bar_height = min(2.1, 0.8 + (0.22 * len(normalized_item.split())))
            x = panel_left + 0.62 + (index * 0.62)
            y = panel_top + 3.78 - bar_height
            if line_mode:
                marker = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(0.26), Inches(0.26))
                style_shape(marker, theme["accent_soft"])
                if index > 0:
                    connector = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x - 0.42), Inches(y + 0.12), Inches(0.44), Inches(0.04))
                    style_shape(connector, theme["surface"])
            else:
                bar = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(0.34), Inches(bar_height))
                style_shape(bar, theme["surface"])
            add_textbox(
                slide,
                left=x - 0.08,
                top=panel_top + 3.96,
                width=0.54,
                height=0.42,
                text=normalized_item[:18],
                font_size=8,
                color_hex=primary_text,
                align=PP_ALIGN.CENTER,
            )

    if visual_type == "flow":
        for index, item in enumerate(visual_items[:4]):
            y = panel_top + 0.74 + (index * 0.88)
            add_visual_card(
                slide,
                left=panel_left + 0.32,
                top=y,
                width=2.82,
                height=0.54,
                fill_hex=primary_fill,
                text=item,
                text_hex=primary_text,
                font_size=12,
                bold=True,
            )
            add_visual_card(
                slide,
                left=panel_left + 2.98,
                top=y,
                width=0.38,
                height=0.38,
                fill_hex=soft_fill,
                text=str(index + 1),
                text_hex=theme["dark_text"],
                font_size=11,
                bold=True,
            )
    elif visual_type == "comparison":
        left_items = visual_items[::2][:2]
        right_items = visual_items[1::2][:2] or visual_items[:2]
        add_visual_card(
            slide,
            left=panel_left + 0.24,
            top=panel_top + 0.84,
            width=1.42,
            height=2.76,
            fill_hex=primary_fill,
            text="\n\n".join(left_items) or "Point A",
            text_hex=primary_text,
            font_size=12,
            bold=True,
        )
        add_visual_card(
            slide,
            left=panel_left + 1.99,
            top=panel_top + 0.84,
            width=1.42,
            height=2.76,
            fill_hex=secondary_fill,
            text="\n\n".join(right_items) or "Point B",
            text_hex=primary_text,
            font_size=12,
            bold=True,
        )
    elif visual_type == "components":
        core_label = visual_items[0] if visual_items else "Core unit"
        center_card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(panel_left + 1.0), Inches(panel_top + 1.76), Inches(1.62), Inches(0.82))
        style_shape(center_card, primary_fill)
        add_textbox(
            slide,
            left=panel_left + 1.12,
            top=panel_top + 2.02,
            width=1.36,
            height=0.22,
            text=core_label,
            font_size=12,
            color_hex=primary_text,
            bold=True,
            align=PP_ALIGN.CENTER,
        )
        component_positions = [
            (panel_left + 0.2, panel_top + 0.88),
            (panel_left + 2.14, panel_top + 0.88),
            (panel_left + 0.2, panel_top + 3.0),
            (panel_left + 2.14, panel_top + 3.0),
        ]
        connector_positions = [
            (panel_left + 1.01, panel_top + 1.28, 0.24, 0.52),
            (panel_left + 2.42, panel_top + 1.28, 0.24, 0.52),
            (panel_left + 1.01, panel_top + 2.58, 0.24, 0.52),
            (panel_left + 2.42, panel_top + 2.58, 0.24, 0.52),
        ]
        for (x, y), item, (line_x, line_y, line_width, line_height) in zip(
            component_positions,
            (visual_items[1:5] or slide_content.get("bullets", [])[:4] or ["Input", "Output", "Stage", "Feedback"]),
            connector_positions,
            strict=False,
        ):
            connector = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(line_x), Inches(line_y), Inches(line_width), Inches(line_height))
            style_shape(connector, theme["surface"])
            add_visual_card(
                slide,
                left=x,
                top=y,
                width=1.28,
                height=0.64,
                fill_hex=secondary_fill,
                text=item,
                text_hex=primary_text,
                font_size=10,
                bold=True,
            )
    elif visual_type == "timeline":
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(panel_left + 0.46), Inches(panel_top + 0.86), Inches(0.08), Inches(3.2))
        style_shape(line, primary_fill)
        for index, item in enumerate(visual_items[:4]):
            y = panel_top + 0.78 + (index * 0.82)
            dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(panel_left + 0.25), Inches(y), Inches(0.42), Inches(0.42))
            style_shape(dot, secondary_fill)
            add_visual_card(
                slide,
                left=panel_left + 0.82,
                top=y - 0.04,
                width=2.35,
                height=0.5,
                fill_hex=primary_fill,
                text=item,
                text_hex=primary_text,
                font_size=11,
                bold=True,
            )
    elif visual_type == "cycle":
        positions = [(panel_left + 1.14, panel_top + 0.92), (panel_left + 0.32, panel_top + 2.08), (panel_left + 1.95, panel_top + 2.08)]
        for (x, y), item in zip(positions, visual_items[:3] or ["Phase 1", "Phase 2", "Phase 3"], strict=False):
            circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(1.05), Inches(1.05))
            style_shape(circle, primary_fill)
            add_textbox(
                slide,
                left=x + 0.12,
                top=y + 0.22,
                width=0.8,
                height=0.48,
                text=item,
                font_size=11,
                color_hex=primary_text,
                bold=True,
                align=PP_ALIGN.CENTER,
            )
    elif visual_type == "formula":
        for index, item in enumerate(visual_items[:4]):
            add_visual_card(
                slide,
                left=panel_left + 0.28,
                top=panel_top + 0.84 + (index * 0.76),
                width=3.02,
                height=0.5,
                fill_hex=primary_fill,
                text=item,
                text_hex=primary_text,
                font_size=12,
                bold=True,
            )
    elif visual_type == "table":
        header = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(panel_left + 0.24), Inches(panel_top + 0.82), Inches(3.02), Inches(0.46))
        style_shape(header, primary_fill)
        add_textbox(
            slide,
            left=panel_left + 0.38,
            top=panel_top + 0.94,
            width=1.12,
            height=0.18,
            text="Point",
            font_size=10,
            color_hex=primary_text,
            bold=True,
        )
        add_textbox(
            slide,
            left=panel_left + 1.74,
            top=panel_top + 0.94,
            width=1.25,
            height=0.18,
            text="Detail",
            font_size=10,
            color_hex=primary_text,
            bold=True,
        )
        for index, item in enumerate(visual_items[:4] or slide_content.get("bullets", [])[:4]):
            label, _, detail = compact_text(item, "Key point").partition(":")
            row_top = panel_top + 1.38 + (index * 0.68)
            left_cell = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(panel_left + 0.24), Inches(row_top), Inches(1.18), Inches(0.5))
            right_cell = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(panel_left + 1.56), Inches(row_top), Inches(1.7), Inches(0.5))
            style_shape(left_cell, secondary_fill)
            style_shape(right_cell, primary_fill)
            add_textbox(
                slide,
                left=panel_left + 0.34,
                top=row_top + 0.12,
                width=0.98,
                height=0.22,
                text=label[:22],
                font_size=10,
                color_hex=primary_text,
                bold=True,
            )
            add_textbox(
                slide,
                left=panel_left + 1.68,
                top=row_top + 0.12,
                width=1.44,
                height=0.22,
                text=(detail or label)[:38],
                font_size=9,
                color_hex=primary_text,
            )
    elif visual_type in {"chart", "graph"}:
        draw_chart_bars(visual_items, line_mode=visual_type == "graph")
        summary_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(panel_left + 0.38), Inches(panel_top + 4.04), Inches(2.74), Inches(0.42))
        style_shape(summary_box, primary_fill)
        add_textbox(
            slide,
            left=panel_left + 0.54,
            top=panel_top + 4.14,
            width=2.42,
            height=0.2,
            text=compact_text(slide_content.get("flow_note"), "Trend summary"),
            font_size=8,
            color_hex=primary_text,
        )
    elif visual_type == "photo":
        if reference_stream is not None:
            try:
                slide.shapes.add_picture(
                    reference_stream,
                    Inches(panel_left + 0.26),
                    Inches(panel_top + 0.82),
                    width=Inches(3.06),
                    height=Inches(2.36),
                )
            except Exception:
                reference_stream = None
        if reference_stream is None:
            add_visual_card(
                slide,
                left=panel_left + 0.28,
                top=panel_top + 0.84,
                width=3.02,
                height=1.38,
                fill_hex=primary_fill,
                text=visual_items[0] if visual_items else "Reference image",
                text_hex=primary_text,
                font_size=13,
                bold=True,
            )
        caption_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(panel_left + 0.28), Inches(panel_top + 3.36), Inches(3.02), Inches(0.96))
        style_shape(caption_box, primary_fill)
        add_textbox(
            slide,
            left=panel_left + 0.46,
            top=panel_top + 3.54,
            width=2.64,
            height=0.58,
            text="\n".join((visual_items[:2] or slide_content.get("bullets", [])[:2]))[:160],
            font_size=10,
            color_hex=primary_text,
        )
    else:
        add_visual_card(
            slide,
            left=panel_left + 0.95,
            top=panel_top + 1.28,
            width=1.78,
            height=0.7,
            fill_hex=primary_fill,
            text=visual_items[0] if visual_items else "Core idea",
            text_hex=primary_text,
            font_size=12,
            bold=True,
        )
        orbit_positions = [(panel_left + 0.24, panel_top + 2.42), (panel_left + 1.95, panel_top + 2.42), (panel_left + 1.08, panel_top + 3.38)]
        for (x, y), item in zip(orbit_positions, visual_items[1:4] or slide_content.get("bullets", [])[:3], strict=False):
            add_visual_card(
                slide,
                left=x,
                top=y,
                width=1.46,
                height=0.58,
                fill_hex=secondary_fill,
                text=item,
                text_hex=primary_text,
                font_size=11,
                bold=True,
            )


def decorate_presentation_slide(slide: Any, theme: dict[str, str], slide_index: int, *, is_title: bool = False):
    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    style_shape(background, theme["background"])

    style_family = get_presentation_style_family(theme)
    if style_family == "emerald-scholar":
        accent_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.45), Inches(0.45), Inches(0.24), Inches(6.55))
        style_shape(accent_bar, theme["accent"])
        orb = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(10.7), Inches(-0.5), Inches(3.2), Inches(3.2))
        style_shape(orb, theme["accent"], transparency=0.72)
        ribbon = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.9), Inches(6.55), Inches(3.7), Inches(0.5))
        style_shape(ribbon, theme["surface_alt"])
    elif style_family == "sunset-classroom":
        panel = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(4.1), Inches(7.5))
        style_shape(panel, theme["surface_alt"], transparency=0.18)
        orb = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(10.8), Inches(0.35), Inches(2.0), Inches(2.0))
        style_shape(orb, theme["accent_soft"], transparency=0.12)
        footer = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(6.55), Inches(4.6), Inches(0.48))
        style_shape(footer, theme["accent"])
    else:
        stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.45))
        style_shape(stripe, theme["accent"])
        side = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(11.85), Inches(0.45), Inches(1.0), Inches(6.5))
        style_shape(side, theme["surface_alt"])
        for offset in range(4):
            dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.8 + (offset * 0.38)), Inches(6.7), Inches(0.16), Inches(0.16))
            style_shape(dot, theme["accent_soft"])

    if not is_title:
        banner = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.85), Inches(0.58), Inches(2.8), Inches(0.48))
        style_shape(banner, theme["surface"])
        add_textbox(
            slide,
            left=1.05,
            top=0.68,
            width=2.3,
            height=0.2,
            text=f"Slide {slide_index}",
            font_size=11,
            color_hex=theme["accent_soft"] if style_family != "sunset-classroom" else "FFF7ED",
            bold=True,
        )
    add_textbox(
        slide,
        left=12.0,
        top=7.02,
        width=0.72,
        height=0.14,
        text="mabaso",
        font_size=9,
        color_hex=theme["muted"],
        bold=True,
        align=PP_ALIGN.RIGHT,
    )


def add_presentation_title_slide(
    presentation: Any,
    title: str,
    subtitle: str,
    theme: dict[str, str],
    *,
    use_custom_template: bool = False,
):
    slide = presentation.slides.add_slide(get_presentation_slide_layout(presentation, "blank", fallback_index=6))
    if not use_custom_template:
        decorate_presentation_slide(slide, theme, 0, is_title=True)
    else:
        title_panel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.82), Inches(0.94), Inches(9.22), Inches(3.18))
        style_shape(title_panel, "0F172A", transparency=0.16)
    title_color = "FFFFFF" if use_custom_template else theme["text"]
    subtitle_color = "E2E8F0" if use_custom_template else theme["muted"]
    add_textbox(
        slide,
        left=0.95,
        top=1.24,
        width=8.8,
        height=1.55,
        text=title,
        font_size=32 if use_custom_template else 30,
        color_hex=title_color,
        bold=True,
        font_name="Aptos Display",
    )
    add_textbox(
        slide,
        left=0.98,
        top=2.92,
        width=8.2,
        height=1.2,
        text=subtitle,
        font_size=18,
        color_hex=subtitle_color,
    )


def add_presentation_content_slide(
    presentation: Any,
    slide_index: int,
    slide_content: dict[str, Any],
    theme: dict[str, str],
    reference_images: list[str],
    *,
    use_custom_template: bool = False,
):
    slide = presentation.slides.add_slide(get_presentation_slide_layout(presentation, "blank", fallback_index=6))
    if not use_custom_template:
        decorate_presentation_slide(slide, theme, slide_index)
    style_family = get_presentation_style_family(theme)
    uses_dark_text = style_family == "sunset-classroom" and not use_custom_template
    title_fill = "0F172A" if use_custom_template else theme["surface"]
    body_fill = "111827" if use_custom_template else (theme["surface_alt"] if uses_dark_text else theme["surface"])
    title_color = "FFFFFF" if use_custom_template else (theme["text"] if not uses_dark_text else theme["dark_text"])
    body_text_color = "F8FAFC" if use_custom_template else (theme["text"] if not uses_dark_text else theme["dark_text"])

    title_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.82), Inches(1.18), Inches(7.7), Inches(0.95))
    style_shape(title_box, title_fill, transparency=0.16 if use_custom_template else 0.0)
    add_textbox(
        slide,
        left=1.05,
        top=1.42,
        width=7.2,
        height=0.3,
        text=slide_content["title"],
        font_size=23,
        color_hex=title_color,
        bold=True,
        font_name="Aptos Display",
    )

    body_panel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.85), Inches(2.35), Inches(7.5), Inches(4.35))
    style_shape(body_panel, body_fill, transparency=0.14 if use_custom_template else 0.0)
    add_bullet_list(
        slide,
        left=1.14,
        top=2.67,
        width=6.85,
        height=3.7,
        bullets=slide_content.get("bullets", []),
        color_hex=body_text_color,
    )
    draw_presentation_visual_panel(slide, slide_content, theme, reference_images)

    if not use_custom_template:
        footer = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.72), Inches(6.06), Inches(3.65), Inches(0.62))
        style_shape(footer, theme["surface_alt"] if not uses_dark_text else theme["surface"])
        add_textbox(
            slide,
            left=8.98,
            top=6.24,
            width=3.05,
            height=0.22,
            text=theme["name"],
            font_size=11,
            color_hex=theme["muted"] if not uses_dark_text else theme["dark_text"],
            bold=True,
        )


def add_presentation_closing_slide(
    presentation: Any,
    slide_content: dict[str, Any],
    theme: dict[str, str],
    *,
    use_custom_template: bool = False,
):
    slide = presentation.slides.add_slide(get_presentation_slide_layout(presentation, "blank", fallback_index=6))
    if not use_custom_template:
        decorate_presentation_slide(slide, theme, len(presentation.slides))
    else:
        closing_panel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3.05), Inches(2.35), Inches(7.2), Inches(1.95))
        style_shape(closing_panel, "0F172A", transparency=0.18)
    closing_title = compact_text(slide_content.get("title"), "THANK YOU")
    add_textbox(
        slide,
        left=1.1,
        top=2.58,
        width=11.1,
        height=1.1,
        text=closing_title,
        font_size=36,
        color_hex="FFFFFF" if use_custom_template else theme["text"],
        bold=True,
        font_name="Aptos Display",
        align=PP_ALIGN.CENTER,
    )


def build_presentation_file(
    job_id: str,
    *,
    title: str,
    subtitle: str,
    slides: list[dict[str, Any]],
    design_id: str,
    reference_images: list[str],
    template_file_bytes: bytes | None = None,
) -> str:
    ensure_presentation_support()
    output_dir = PRESENTATION_OUTPUT_DIR / job_id
    if output_dir.exists():
        shutil.rmtree(output_dir, ignore_errors=True)
    output_dir.mkdir(parents=True, exist_ok=True)

    use_custom_template = bool(template_file_bytes)
    if use_custom_template:
        try:
            presentation = Presentation(BytesIO(template_file_bytes))
        except Exception as exc:
            raise RuntimeError("Could not read the uploaded PowerPoint template. Upload a valid .pptx file.") from exc
        clear_presentation_slides(presentation)
    else:
        presentation = Presentation()
        presentation.slide_width = Inches(13.333)
        presentation.slide_height = Inches(7.5)
    theme = PRESENTATION_THEMES[normalize_presentation_design_id(design_id)]

    add_presentation_title_slide(presentation, title, subtitle, theme, use_custom_template=use_custom_template)
    for slide_index, slide_content in enumerate(slides, start=1):
        if compact_text(slide_content.get("visual_type")).lower() == "closing":
            add_presentation_closing_slide(presentation, slide_content, theme, use_custom_template=use_custom_template)
        else:
            add_presentation_content_slide(
                presentation,
                slide_index,
                slide_content,
                theme,
                reference_images,
                use_custom_template=use_custom_template,
            )

    file_path = output_dir / "lecture-presentation.pptx"
    presentation.save(str(file_path))
    return str(file_path)


async def generate_presentation_package(
    summary: str,
    transcript: str,
    lecture_notes: str,
    lecture_slides: str,
    past_question_papers: str,
    design_id: str,
    job_id: str,
    output_language: str,
    reference_images: list[str],
    template_file_bytes: bytes | None = None,
    template_file_name: str = "",
) -> dict[str, Any]:
    normalized_design_id = normalize_presentation_design_id(design_id)
    fallback_package = build_presentation_fallback(summary, transcript)
    reference_catalog_items = await analyze_reference_images_for_study_guide(
        reference_images,
        summary,
        lecture_notes,
        lecture_slides,
        output_language,
    )
    reference_catalog = build_reference_image_catalog(reference_catalog_items)
    context_blocks = [
        trimmed_context_block("STUDY GUIDE SUMMARY", summary, MAX_STUDY_GUIDE_INPUT_CHARS),
        trimmed_context_block("LECTURE NOTES", lecture_notes, MAX_STUDY_GUIDE_INPUT_CHARS // 2),
        trimmed_context_block("LECTURE SLIDES", lecture_slides, MAX_STUDY_GUIDE_INPUT_CHARS // 2),
        trimmed_context_block("PAST QUESTION PAPERS", past_question_papers, MAX_STUDY_GUIDE_INPUT_CHARS // 2),
        trimmed_context_block("LECTURE TRANSCRIPT", transcript, MAX_TRANSCRIPT_STUDY_GUIDE_INPUT_CHARS // 2),
    ]
    if reference_catalog:
        context_blocks.append(reference_catalog)
    combined_source = "\n\n".join(block for block in context_blocks if block)

    def _generate_plan() -> dict[str, Any]:
        response = client.with_options(timeout=PRESENTATION_REQUEST_TIMEOUT).chat.completions.create(
            model=PRESENTATION_MODEL,
            max_completion_tokens=min(MAX_COMPLETION_TOKENS, 4200),
            messages=[
                {
                    "role": "system",
                    "content": (
                        "You create concise academic PowerPoint structures for students. "
                        "Return strict JSON only with these keys: title, subtitle, slides.\n\n"
                        "Rules:\n"
                        "- `title` must be the lecture topic only.\n"
                        "- `subtitle` must be one short description sentence for the topic slide.\n"
                        "- `slides` must be an array of 5 to 7 content slides.\n"
                        "- Each slide object must contain `title`, `bullets`, `visual_title`, `visual_type`, `visual_items`, `flow_note`, and `reference_image_index`.\n"
                        "- `bullets` must contain 3 to 5 short bullet strings with useful teaching value, not copied fragments.\n"
                        "- `visual_items` must contain 2 to 4 short labels for a diagram panel on the slide.\n"
                        "- `visual_type` must be one of: flow, comparison, timeline, cycle, formula, cluster, components, table, chart, graph, photo.\n"
                        "- `flow_note` must explain how the slide advances the lesson flow in one sentence.\n"
                        f"- `reference_image_index` must be -1 when no uploaded image clearly matches, otherwise an integer from 0 to {max(len(reference_images) - 1, 0)}.\n"
                        "- This is a formal PowerPoint deck, not speaker notes. Do not write what the presenter should say.\n"
                        "- Keep bullet lines concise, readable, and presentation-ready.\n"
                        f"- Write the slide text in {output_language}.\n"
                        "- When the lecture material covers multiple topics, keep those topics separate instead of mixing them into one slide.\n"
                        "- Cover overview, core concepts, definitions or structure, formulas or rules when present, and worked examples or applications.\n"
                        "- Do not create slides about common mistakes, exam tips, revision plans, study advice, exam traps, or test strategy.\n"
                        "- Do not use phrases like 'common mistakes', 'exam tips', 'revision focus', 'exam traps', or 'before the test' in titles, bullets, flow notes, or visual panels.\n"
                        "- Use tables, graphs, charts, component diagrams, process flows, and photos when they help students understand the material quickly.\n"
                        "- Only use `photo` when the uploaded lecture images clearly match the topic on that slide.\n"
                        "- If the reference image catalog does not clearly match a slide, use a non-photo visual instead of forcing an image.\n"
                        "- When no real photo fits, create a designed diagram panel using components, graphs, charts, formulas, or process visuals instead of a generic placeholder.\n"
                        "- Do not create a thank-you slide or closing slide in JSON. The app adds the final closing slide automatically.\n"
                        "- Make the slide sequence feel like a high-grade classroom deck: clear concepts first, then structure, then worked explanation or interpretation.\n"
                        "- Use lecture language when it is reliable, but rewrite it cleanly for slides.\n"
                        "- Do not return markdown, numbering, or commentary outside JSON."
                    ),
                },
                {
                    "role": "user",
                    "content": (
                        (
                            f"REFERENCE IMAGES AVAILABLE: {len(reference_images)}\n"
                            "Match uploaded lecture images only when the catalog and slide topic clearly fit.\n"
                            "If there is no clear fit, set `reference_image_index` to -1.\n\n"
                        )
                        if reference_images
                        else ""
                    ) + (
                        combined_source or presentation_slides_to_text(
                            fallback_package["title"],
                            fallback_package["subtitle"],
                            fallback_package["slides"],
                        )
                    ),
                },
            ],
        )
        return parse_json_object(response.choices[0].message.content or "")

    update_job(job_id, status="processing", stage="Designing PowerPoint storyboard", progress=18)
    try:
        generated_package = await asyncio.to_thread(_generate_plan)
    except Exception as exc:
        logger.warning("Presentation plan generation failed, using fallback outline: %s", exc)
        generated_package = {}

    normalized_slides = normalize_presentation_slides(generated_package.get("slides"))
    if not normalized_slides:
        normalized_slides = fallback_package["slides"]

    title = compact_text(generated_package.get("title"), fallback_package["title"])
    subtitle = compact_text(generated_package.get("subtitle"), fallback_package["subtitle"])
    closing_slide = build_presentation_closing_slide_content(title)
    slides_with_closing = [*normalized_slides, closing_slide]

    update_job(
        job_id,
        status="processing",
        stage="Applying your PowerPoint template and building slides" if template_file_bytes else "Building PowerPoint file",
        progress=56,
    )
    download_file = await asyncio.to_thread(
        build_presentation_file,
        job_id,
        title=title,
        subtitle=subtitle,
        slides=slides_with_closing,
        design_id=normalized_design_id,
        reference_images=reference_images,
        template_file_bytes=template_file_bytes,
    )

    return {
        "presentation_title": title,
        "presentation_subtitle": subtitle,
        "presentation_design_id": normalized_design_id,
        "presentation_template_name": compact_text(template_file_name),
        "presentation_slides": slides_with_closing,
        "_presentation_download_file": download_file,
    }


def determine_quiz_total_marks(
    summary: str,
    transcript: str,
    lecture_notes: str,
    lecture_slides: str,
    past_question_papers: str,
) -> int:
    combined_size = sum(
        len((value or "").strip())
        for value in [summary, transcript, lecture_notes, lecture_slides, past_question_papers]
    )
    return 50 if combined_size >= 24000 else 40


def build_group_subparts(labels: list[str]) -> list[dict[str, Any]]:
    return [{"label": label, "marks": 1} for label in labels]


def build_quiz_blueprint(total_marks: int) -> list[dict[str, Any]]:
    if total_marks >= 50:
        return [
            {"number": "1", "type": "short_answer", "marks": 2},
            {"number": "2", "type": "multiple_choice_group", "marks": 4, "subparts": build_group_subparts(["a", "b", "c", "d"])},
            {"number": "3", "type": "true_false_group", "marks": 5, "subparts": build_group_subparts(["a", "b", "c", "d", "e"])},
            {"number": "4", "type": "short_answer", "marks": 4},
            {"number": "5", "type": "short_answer", "marks": 5},
            {"number": "6", "type": "multiple_choice_group", "marks": 4, "subparts": build_group_subparts(["a", "b", "c", "d"])},
            {"number": "7", "type": "true_false_group", "marks": 5, "subparts": build_group_subparts(["a", "b", "c", "d", "e"])},
            {"number": "8", "type": "short_answer", "marks": 6},
            {"number": "9", "type": "short_answer", "marks": 7},
            {"number": "10", "type": "short_answer", "marks": 8},
        ]
    return [
        {"number": "1", "type": "short_answer", "marks": 2},
        {"number": "2", "type": "multiple_choice_group", "marks": 4, "subparts": build_group_subparts(["a", "b", "c", "d"])},
        {"number": "3", "type": "true_false_group", "marks": 5, "subparts": build_group_subparts(["a", "b", "c", "d", "e"])},
        {"number": "4", "type": "short_answer", "marks": 3},
        {"number": "5", "type": "short_answer", "marks": 4},
        {"number": "6", "type": "multiple_choice_group", "marks": 4, "subparts": build_group_subparts(["a", "b", "c", "d"])},
        {"number": "7", "type": "true_false_group", "marks": 5, "subparts": build_group_subparts(["a", "b", "c", "d", "e"])},
        {"number": "8", "type": "short_answer", "marks": 4},
        {"number": "9", "type": "short_answer", "marks": 4},
        {"number": "10", "type": "short_answer", "marks": 5},
    ]


def normalize_answer_points(value: Any, fallback_answer: str, marks: int) -> list[str]:
    items: list[str] = []
    if isinstance(value, list):
        for entry in value:
            cleaned = compact_text(entry)
            if cleaned:
                items.append(cleaned)
    elif isinstance(value, str):
        lines = [re.sub(r"^[-*\d.\s]+", "", part).strip() for part in value.splitlines()]
        items.extend(line for line in lines if line)

    if not items and fallback_answer:
        items = [segment.strip() for segment in re.split(r"\n+|;\s*", fallback_answer) if segment.strip()]

    minimum_points = min(max(marks, 2), 6)
    return items[: max(minimum_points, 1)]


def normalize_flashcards(raw_cards: Any, fallback_cards: list[dict[str, str]]) -> list[dict[str, str]]:
    cards: list[dict[str, str]] = []
    if isinstance(raw_cards, list):
        for entry in raw_cards:
            if not isinstance(entry, dict):
                continue
            question = compact_text(entry.get("question"))
            answer = compact_text(entry.get("answer"))
            if question and answer:
                cards.append({"question": question, "answer": answer})
    if cards:
        return cards[:12]
    return fallback_cards[:12]


def make_short_answer_question(
    number: str,
    marks: int,
    question_text: str,
    answer_text: str,
    answer_points: list[str],
) -> dict[str, Any]:
    safe_answer = answer_text or "\n".join(answer_points) or "No answer supplied."
    return {
        "number": number,
        "type": "short_answer",
        "marks": marks,
        "question": question_text,
        "answer": safe_answer,
        "answer_points": answer_points or [safe_answer],
    }


def normalize_options(raw_options: Any) -> list[str]:
    if not isinstance(raw_options, list):
        return []
    options = [compact_text(option) for option in raw_options if compact_text(option)]
    seen: set[str] = set()
    unique_options: list[str] = []
    for option in options:
        if option in seen:
            continue
        seen.add(option)
        unique_options.append(option)
    return unique_options


def build_structured_quiz_fallback(simple_questions: list[dict[str, str]], blueprint: list[dict[str, Any]]) -> list[dict[str, Any]]:
    fallback_questions: list[dict[str, Any]] = []
    for index, blueprint_item in enumerate(blueprint):
        source = simple_questions[index] if index < len(simple_questions) else {}
        question_text = compact_text(source.get("question"), f"Explain the key idea in part {blueprint_item['number']}.")
        answer_text = compact_text(source.get("answer"), "No answer supplied.")
        answer_points = normalize_answer_points([], answer_text, int(blueprint_item["marks"]))
        fallback_questions.append(
            make_short_answer_question(
                str(blueprint_item["number"]),
                int(blueprint_item["marks"]),
                question_text,
                answer_text,
                answer_points,
            )
        )
    return fallback_questions


def normalize_generated_quiz_questions(
    raw_questions: Any,
    blueprint: list[dict[str, Any]],
    fallback_questions: list[dict[str, str]],
) -> list[dict[str, Any]]:
    source_questions = raw_questions if isinstance(raw_questions, list) else []
    normalized_questions: list[dict[str, Any]] = []

    for index, blueprint_item in enumerate(blueprint):
        raw_question = source_questions[index] if index < len(source_questions) and isinstance(source_questions[index], dict) else {}
        fallback_source = fallback_questions[index] if index < len(fallback_questions) else {}
        number = str(blueprint_item["number"])
        marks = int(blueprint_item["marks"])
        question_text = compact_text(raw_question.get("question"), compact_text(fallback_source.get("question"), f"Question {number}"))
        answer_text = compact_text(raw_question.get("answer"), compact_text(fallback_source.get("answer"), "No answer supplied."))
        question_type = compact_text(raw_question.get("type"), str(blueprint_item["type"])).lower()

        if question_type == "short_answer":
            normalized_questions.append(
                make_short_answer_question(
                    number,
                    marks,
                    question_text,
                    answer_text,
                    normalize_answer_points(raw_question.get("answer_points"), answer_text, marks),
                )
            )
            continue

        if question_type not in {"multiple_choice_group", "true_false_group"}:
            normalized_questions.append(
                make_short_answer_question(
                    number,
                    marks,
                    question_text,
                    answer_text,
                    normalize_answer_points(raw_question.get("answer_points"), answer_text, marks),
                )
            )
            continue

        raw_subparts = raw_question.get("subparts") if isinstance(raw_question.get("subparts"), list) else []
        subpart_lookup = {
            compact_text(item.get("label")).lower(): item
            for item in raw_subparts
            if isinstance(item, dict) and compact_text(item.get("label"))
        }
        subparts: list[dict[str, Any]] = []
        valid_group = True

        for sub_blueprint in blueprint_item.get("subparts", []):
            label = str(sub_blueprint["label"]).lower()
            raw_subpart = subpart_lookup.get(label)
            if not raw_subpart:
                valid_group = False
                break

            subpart_question = compact_text(raw_subpart.get("question"))
            if not subpart_question:
                valid_group = False
                break

            if question_type == "true_false_group":
                options = ["True", "False"]
            else:
                options = normalize_options(raw_subpart.get("options"))
                if len(options) < 4:
                    valid_group = False
                    break
                options = options[:4]

            answer = compact_text(raw_subpart.get("answer"))
            resolved_answer = next((option for option in options if option.lower() == answer.lower()), "")
            if not resolved_answer:
                valid_group = False
                break

            subparts.append(
                {
                    "label": label,
                    "marks": int(sub_blueprint["marks"]),
                    "question": subpart_question,
                    "options": options,
                    "answer": resolved_answer,
                    "explanation": compact_text(raw_subpart.get("explanation"), "Review the matching concept in the notes."),
                }
            )

        if not valid_group:
            normalized_questions.append(
                make_short_answer_question(
                    number,
                    marks,
                    question_text,
                    answer_text,
                    normalize_answer_points(raw_question.get("answer_points"), answer_text, marks),
                )
            )
            continue

        normalized_questions.append(
            {
                "number": number,
                "type": question_type,
                "marks": marks,
                "question": question_text or ("Choose the correct answer for each subpart." if question_type == "multiple_choice_group" else "State whether each statement is true or false."),
                "answer": "\n".join(f"{item['label']}) {item['answer']}" for item in subparts),
                "subparts": subparts,
            }
        )

    return normalized_questions


async def generate_structured_study_assets(
    summary: str,
    transcript: str,
    lecture_notes: str,
    lecture_slides: str,
    past_question_papers: str,
    job_id: str,
    output_language: str,
) -> dict[str, Any]:
    fallback_assets = extract_study_assets(summary)
    source_blocks = [
        trimmed_context_block("STUDY GUIDE SUMMARY", summary, MAX_STUDY_GUIDE_INPUT_CHARS),
        trimmed_context_block("LECTURER NOTES", lecture_notes, MAX_STUDY_GUIDE_INPUT_CHARS // 2),
        trimmed_context_block("LECTURE SLIDES", lecture_slides, MAX_STUDY_GUIDE_INPUT_CHARS // 2),
        trimmed_context_block("PAST QUESTION PAPERS", past_question_papers, MAX_STUDY_GUIDE_INPUT_CHARS // 2),
        trimmed_context_block("LECTURE TRANSCRIPT", transcript, MAX_TRANSCRIPT_STUDY_GUIDE_INPUT_CHARS // 2),
    ]
    combined_source = "\n\n".join(block for block in source_blocks if block)

    def _generate_assets() -> dict[str, Any]:
        response = client.with_options(timeout=STUDY_GUIDE_REQUEST_TIMEOUT).chat.completions.create(
            model=ASSET_GENERATION_MODEL,
            max_completion_tokens=min(MAX_COMPLETION_TOKENS, 3400),
            messages=[
                {
                    "role": "system",
                    "content": (
                        "You build structured study assets for a university revision app. "
                        "Return only valid JSON with the keys formula, worked_example, and flashcards.\n\n"
                        "Rules:\n"
                        "- Do not mention how a student should feel.\n"
                        "- Use plain readable formulas, never LaTeX.\n"
                        "- `formula` should be a compact markdown study sheet or a short note when no formula is relevant.\n"
                        "- `worked_example` should explain every example that appears in the study guide WORKED EXAMPLES section.\n"
                        "- For each worked example, write the method step by step in markdown, using labels like Step 1, Step 2, and so on.\n"
                        "- Use the STEP-BY-STEP EXPLANATIONS section to expand the reasoning, not to replace any example from the guide.\n"
                        "- `flashcards` should contain 10 to 12 items, each with `question` and `answer`.\n"
                        "- If past question papers are provided, use them only as reference for topic coverage, phrasing style, and likely mark patterns. Do not copy them verbatim.\n"
                        f"- Write every returned field in {output_language}.\n"
                        "- Return JSON only, with no markdown code fence."
                    ),
                },
                {"role": "user", "content": combined_source},
            ],
        )
        return parse_json_object(response.choices[0].message.content or "")

    try:
        update_job(job_id, status="processing", stage="Building flashcards and worked examples", progress=82)
        generated_assets = await asyncio.to_thread(_generate_assets)
    except Exception as exc:
        logger.warning("Structured asset generation failed, using extracted fallback assets: %s", exc)
        generated_assets = {}

    return {
        "formula": compact_text(generated_assets.get("formula"), fallback_assets["formula"]),
        "worked_example": build_worked_example_asset(
            summary,
            compact_text(generated_assets.get("worked_example"), fallback_assets["worked_example"]),
        ),
        "flashcards": normalize_flashcards(generated_assets.get("flashcards"), fallback_assets["flashcards"]),
    }


def build_quiz_generation_source_text(
    summary: str,
    transcript: str,
    lecture_notes: str,
    lecture_slides: str,
    past_question_papers: str,
) -> tuple[dict[str, Any], list[dict[str, Any]], str]:
    fallback_assets = extract_study_assets(summary)
    total_marks = determine_quiz_total_marks(
        summary,
        transcript,
        lecture_notes,
        lecture_slides,
        past_question_papers,
    )
    blueprint = build_quiz_blueprint(total_marks)
    source_blocks = [
        trimmed_context_block("STUDY GUIDE SUMMARY", summary, MAX_STUDY_GUIDE_INPUT_CHARS),
        trimmed_context_block("LECTURER NOTES", lecture_notes, MAX_STUDY_GUIDE_INPUT_CHARS // 2),
        trimmed_context_block("LECTURE SLIDES", lecture_slides, MAX_STUDY_GUIDE_INPUT_CHARS // 2),
        trimmed_context_block("PAST QUESTION PAPERS", past_question_papers, MAX_STUDY_GUIDE_INPUT_CHARS // 2),
        trimmed_context_block("LECTURE TRANSCRIPT", transcript, MAX_TRANSCRIPT_STUDY_GUIDE_INPUT_CHARS // 2),
        f"QUIZ BLUEPRINT\n{json.dumps(blueprint, ensure_ascii=False, indent=2)}",
    ]
    return fallback_assets, blueprint, "\n\n".join(block for block in source_blocks if block)


async def generate_quiz_questions(
    summary: str,
    transcript: str,
    lecture_notes: str,
    lecture_slides: str,
    past_question_papers: str,
    job_id: str,
    output_language: str,
) -> list[dict[str, Any]]:
    fallback_assets, blueprint, combined_source = build_quiz_generation_source_text(
        summary,
        transcript,
        lecture_notes,
        lecture_slides,
        past_question_papers,
    )
    total_marks = sum(int(item.get("marks", 0)) for item in blueprint)

    def _generate_quiz() -> dict[str, Any]:
        response = client.with_options(timeout=STUDY_GUIDE_REQUEST_TIMEOUT).chat.completions.create(
            model=ASSET_GENERATION_MODEL,
            max_completion_tokens=min(MAX_COMPLETION_TOKENS, 3200),
            messages=[
                {
                    "role": "system",
                    "content": (
                        "You build structured university tests for a revision app. "
                        "Return only valid JSON with the key quiz_questions.\n\n"
                        "Rules:\n"
                        "- Keep the test rising in difficulty from the first question to the last question, but do not label difficulty levels.\n"
                        "- `quiz_questions` must follow the blueprint exactly for question number, question type, and marks.\n"
                        "- If past question papers are provided, use them only as reference for topic coverage, phrasing style, and likely mark patterns. Do not copy them verbatim.\n"
                        "- Short-answer questions need `question`, `answer`, and `answer_points` for partial-credit marking.\n"
                        "- Multiple-choice group questions need `question` and `subparts`. Each subpart needs `label`, `question`, `marks`, `options`, `answer`, and `explanation`.\n"
                        "- True/false group questions need `question` and `subparts`. Each subpart needs `label`, `question`, `marks`, `options`, `answer`, and `explanation`.\n"
                        "- For true/false questions, the options must be exactly [\"True\", \"False\"].\n"
                        "- Keep each option-based subpart worth 1 mark.\n"
                        f"- The full test must add up to {total_marks} marks.\n"
                        f"- Write every returned field in {output_language}.\n"
                        "- Return JSON only, with no markdown code fence."
                    ),
                },
                {"role": "user", "content": combined_source},
            ],
        )
        return parse_json_object(response.choices[0].message.content or "")

    try:
        update_job(job_id, status="processing", stage="Planning test structure", progress=34)
        generated_assets = await asyncio.to_thread(_generate_quiz)
    except Exception as exc:
        logger.warning("Structured quiz generation failed, using fallback quiz: %s", exc)
        generated_assets = {}

    quiz_questions = normalize_generated_quiz_questions(
        generated_assets.get("quiz_questions"),
        blueprint,
        fallback_assets["quiz_questions"],
    )
    if quiz_questions:
        return quiz_questions
    return build_structured_quiz_fallback(fallback_assets["quiz_questions"], blueprint)


async def generate_study_guide(
    transcript: str,
    lecture_notes: str,
    lecture_slides: str,
    past_question_papers: str,
    job_id: str,
    output_language: str,
) -> tuple[str, bool]:
    trimmed_transcript = transcript[:MAX_TRANSCRIPT_STUDY_GUIDE_INPUT_CHARS].strip()
    trimmed_notes = lecture_notes[:MAX_STUDY_GUIDE_INPUT_CHARS].strip()
    trimmed_slides = lecture_slides[:MAX_STUDY_GUIDE_INPUT_CHARS].strip()
    trimmed_past_papers = past_question_papers[:MAX_STUDY_GUIDE_INPUT_CHARS].strip()

    if len(transcript) > MAX_TRANSCRIPT_STUDY_GUIDE_INPUT_CHARS:
        trimmed_transcript += (
            "\n\nNOTE: The transcript was shortened for faster study-guide generation. "
            "Focus on the most important themes, formulas, definitions, and examples."
        )

    if len(lecture_notes) > MAX_STUDY_GUIDE_INPUT_CHARS:
        trimmed_notes += (
            "\n\nNOTE: The lecture notes were shortened. Prioritize the clearest formulas, worked examples, and lecturer definitions."
        )

    if len(lecture_slides) > MAX_STUDY_GUIDE_INPUT_CHARS:
        trimmed_slides += (
            "\n\nNOTE: The lecture slides were shortened. Prioritize slide headings, formulas, worked examples, and assessment clues."
        )

    if len(past_question_papers) > MAX_STUDY_GUIDE_INPUT_CHARS:
        trimmed_past_papers += (
            "\n\nNOTE: The past question papers were shortened. Prioritize recurring topics, command words, and mark allocations."
        )

    user_content_parts = []
    if trimmed_notes:
        user_content_parts.append(f"LECTURER NOTES\n{trimmed_notes}")
    if trimmed_slides:
        user_content_parts.append(f"LECTURE SLIDES\n{trimmed_slides}")
    if trimmed_past_papers:
        user_content_parts.append(f"PAST QUESTION PAPERS\n{trimmed_past_papers}")
    if trimmed_transcript:
        user_content_parts.append(f"LECTURE TRANSCRIPT\n{trimmed_transcript}")
    user_content_parts.append(
        (
            "OUTPUT INSTRUCTIONS\n"
            f"- Write the full study pack in {output_language}.\n"
            "- Detect when the sources cover multiple distinct topics, chapters, or subtopics.\n"
            "- Keep each topic separate with clear headings and do not mix unrelated notes, formulas, or examples.\n"
            "- If one uploaded source belongs to a different topic, isolate it under its own topic heading instead of blending it into another topic.\n"
            "- Label diagrams, tables, charts, and process visuals clearly inside the guide."
        )
    )
    combined_user_content = "\n\n".join(user_content_parts)

    def _generate() -> str:
        update_job(job_id, status="processing", stage="Preparing study material for notes", progress=20)
        response = client.with_options(timeout=STUDY_GUIDE_REQUEST_TIMEOUT).chat.completions.create(
            model=STUDY_GUIDE_MODEL,
            max_completion_tokens=MAX_COMPLETION_TOKENS,
            messages=[
                {"role": "system", "content": STUDY_GUIDE_PROMPT},
                {"role": "user", "content": combined_user_content},
            ],
        )
        return (response.choices[0].message.content or "").strip()

    try:
        update_job(job_id, status="processing", stage="Generating study guide", progress=55)
        summary = await asyncio.to_thread(_generate)
        cleaned_summary = tidy_study_guide_layout(make_formulas_human_readable(summary))
        return add_student_support_sections(cleaned_summary), False
    except Exception as exc:
        logger.warning("Primary study guide generation failed, using fallback summary: %s", exc)
        update_job(job_id, status="processing", stage="Generating fallback study guide", progress=75)
        fallback_source = "\n\n".join(
            part
            for part in [lecture_notes.strip(), lecture_slides.strip(), past_question_papers.strip(), transcript.strip()]
            if part
        )
        cleaned_summary = tidy_study_guide_layout(make_formulas_human_readable(build_fallback_study_guide(fallback_source)))
        return add_student_support_sections(cleaned_summary), True


async def run_transcription_job(job_id: str, file_path: Path):
    try:
        update_job(job_id, status="processing", stage="Lecture received", progress=1)
        transcript = await transcribe_audio(file_path, job_id)
        if not transcript:
            raise RuntimeError("Transcription returned no text.")

        update_job(
            job_id,
            status="completed",
            stage="Transcription completed",
            progress=100,
            transcript=transcript,
        )
        job = jobs.get(job_id, {})
        started_at = parse_history_datetime(job.get("created_at"), utc_now())
        record_audit_log(
            action="transcription.completed",
            email=job.get("owner_email", ""),
            resource_type="transcription",
            resource_name=file_path.name,
            duration_ms=int((utc_now() - started_at).total_seconds() * 1000),
            metadata={"job_id": job_id},
        )
    except Exception as exc:
        logger.exception("Transcription job failed")
        update_job(
            job_id,
            status="failed",
            stage="Transcription failed",
            progress=100,
            error=format_job_error(exc),
        )
        job = jobs.get(job_id, {})
        started_at = parse_history_datetime(job.get("created_at"), utc_now())
        record_audit_log(
            action="transcription.completed",
            status="failed",
            email=job.get("owner_email", ""),
            resource_type="transcription",
            resource_name=file_path.name,
            duration_ms=int((utc_now() - started_at).total_seconds() * 1000),
            metadata={"job_id": job_id, "error": format_job_error(exc)},
        )
    finally:
        if file_path.exists():
            try:
                file_path.unlink()
            except OSError:
                logger.warning("Could not delete temp file: %s", file_path)


async def run_video_transcription_job(job_id: str, video_url: str):
    file_path: Path | None = None
    try:
        update_job(job_id, status="processing", stage="Video link received", progress=1)
        if extract_youtube_video_id(video_url):
            transcript = await asyncio.to_thread(fetch_youtube_transcript, video_url, job_id)
            if transcript:
                update_job(
                    job_id,
                    status="completed",
                    stage="Video transcription completed",
                    progress=100,
                    transcript=transcript,
                )
                job = jobs.get(job_id, {})
                started_at = parse_history_datetime(job.get("created_at"), utc_now())
                record_audit_log(
                    action="transcription.completed",
                    email=job.get("owner_email", ""),
                    resource_type="transcription",
                    resource_name=video_url,
                    duration_ms=int((utc_now() - started_at).total_seconds() * 1000),
                    metadata={"job_id": job_id, "source": "youtube_transcript"},
                )
                return

            transcript = await asyncio.to_thread(fetch_youtube_watch_page_captions, video_url, job_id)
            if transcript:
                update_job(
                    job_id,
                    status="completed",
                    stage="Video transcription completed",
                    progress=100,
                    transcript=transcript,
                )
                job = jobs.get(job_id, {})
                started_at = parse_history_datetime(job.get("created_at"), utc_now())
                record_audit_log(
                    action="transcription.completed",
                    email=job.get("owner_email", ""),
                    resource_type="transcription",
                    resource_name=video_url,
                    duration_ms=int((utc_now() - started_at).total_seconds() * 1000),
                    metadata={"job_id": job_id, "source": "watch_page"},
                )
                return

            transcript = await asyncio.to_thread(fetch_youtube_timedtext_captions, video_url, job_id)
            if transcript:
                update_job(
                    job_id,
                    status="completed",
                    stage="Video transcription completed",
                    progress=100,
                    transcript=transcript,
                )
                job = jobs.get(job_id, {})
                started_at = parse_history_datetime(job.get("created_at"), utc_now())
                record_audit_log(
                    action="transcription.completed",
                    email=job.get("owner_email", ""),
                    resource_type="transcription",
                    resource_name=video_url,
                    duration_ms=int((utc_now() - started_at).total_seconds() * 1000),
                    metadata={"job_id": job_id, "source": "timedtext"},
                )
                return

            transcript = await asyncio.to_thread(fetch_ytdlp_caption_track, video_url, job_id)
            if transcript:
                update_job(
                    job_id,
                    status="completed",
                    stage="Video transcription completed",
                    progress=100,
                    transcript=transcript,
                )
                job = jobs.get(job_id, {})
                started_at = parse_history_datetime(job.get("created_at"), utc_now())
                record_audit_log(
                    action="transcription.completed",
                    email=job.get("owner_email", ""),
                    resource_type="transcription",
                    resource_name=video_url,
                    duration_ms=int((utc_now() - started_at).total_seconds() * 1000),
                    metadata={"job_id": job_id, "source": "caption_track"},
                )
                return

        transcript = await asyncio.to_thread(download_subtitles_from_video_url, video_url, job_id)
        if transcript:
            update_job(
                job_id,
                status="completed",
                stage="Video transcription completed",
                progress=100,
                transcript=transcript,
            )
            job = jobs.get(job_id, {})
            started_at = parse_history_datetime(job.get("created_at"), utc_now())
            record_audit_log(
                action="transcription.completed",
                email=job.get("owner_email", ""),
                resource_type="transcription",
                resource_name=video_url,
                duration_ms=int((utc_now() - started_at).total_seconds() * 1000),
                metadata={"job_id": job_id, "source": "subtitle_download"},
            )
            return

        if yt_dlp is None:
            raise RuntimeError(
                "This video does not expose downloadable captions, and direct video downloading is not installed on the backend yet."
            )

        file_path = await asyncio.wait_for(
            asyncio.to_thread(download_audio_from_video_url, video_url, job_id),
            timeout=VIDEO_DOWNLOAD_TIMEOUT,
        )
        transcript = await transcribe_audio(file_path, job_id)
        if not transcript:
            raise RuntimeError("The video link was processed, but no transcript text was returned.")

        update_job(
            job_id,
            status="completed",
            stage="Video transcription completed",
            progress=100,
            transcript=transcript,
        )
        job = jobs.get(job_id, {})
        started_at = parse_history_datetime(job.get("created_at"), utc_now())
        record_audit_log(
            action="transcription.completed",
            email=job.get("owner_email", ""),
            resource_type="transcription",
            resource_name=video_url,
            duration_ms=int((utc_now() - started_at).total_seconds() * 1000),
            metadata={"job_id": job_id, "source": "audio_download"},
        )
    except Exception as exc:
        logger.exception("Video transcription job failed")
        update_job(
            job_id,
            status="failed",
            stage="Video transcription failed",
            progress=100,
            error=format_job_error(exc, video_url),
        )
        job = jobs.get(job_id, {})
        started_at = parse_history_datetime(job.get("created_at"), utc_now())
        record_audit_log(
            action="transcription.completed",
            status="failed",
            email=job.get("owner_email", ""),
            resource_type="transcription",
            resource_name=video_url,
            duration_ms=int((utc_now() - started_at).total_seconds() * 1000),
            metadata={"job_id": job_id, "error": format_job_error(exc, video_url)},
        )
    finally:
        if file_path and file_path.exists():
            try:
                file_path.unlink()
            except OSError:
                logger.warning("Could not delete downloaded video audio: %s", file_path)


async def run_summary_job(
    job_id: str,
    transcript: str,
    lecture_notes: str,
    lecture_slides: str,
    past_question_papers: str,
    output_language: str,
    reference_images: list[str],
):
    try:
        update_job(job_id, status="processing", stage="Starting study guide generation", progress=10)
        summary, used_fallback = await generate_study_guide(
            transcript,
            lecture_notes,
            lecture_slides,
            past_question_papers,
            job_id,
            output_language,
        )
        visual_analysis = await analyze_reference_images_for_study_guide(
            reference_images,
            summary,
            lecture_notes,
            lecture_slides,
            output_language,
        )
        uploaded_visuals = build_uploaded_study_visuals(reference_images, visual_analysis)
        assets = await generate_structured_study_assets(
            summary,
            transcript,
            lecture_notes,
            lecture_slides,
            past_question_papers,
            job_id,
            output_language,
        )
        generated_study_images = await generate_study_images(
            summary,
            transcript,
            lecture_notes,
            lecture_slides,
            job_id,
        )
        study_images = merge_study_image_results(uploaded_visuals, generated_study_images)
        update_job(
            job_id,
            status="completed",
            stage="Study guide ready",
            progress=100,
            summary=summary,
            formula=assets["formula"],
            worked_example=assets["worked_example"],
            flashcards=assets["flashcards"],
            quiz_questions=[],
            study_images=study_images,
            used_fallback=used_fallback,
        )
        job = jobs.get(job_id, {})
        started_at = parse_history_datetime(job.get("created_at"), utc_now())
        record_audit_log(
            action="study_guide.completed",
            email=job.get("owner_email", ""),
            resource_type="study_guide",
            resource_name=output_language,
            duration_ms=int((utc_now() - started_at).total_seconds() * 1000),
            metadata={"job_id": job_id, "used_fallback": used_fallback},
        )
    except Exception as exc:
        logger.exception("Study guide job failed")
        update_job(
            job_id,
            status="failed",
            stage="Study guide failed",
            progress=100,
            error=format_job_error(exc),
        )
        job = jobs.get(job_id, {})
        started_at = parse_history_datetime(job.get("created_at"), utc_now())
        record_audit_log(
            action="study_guide.completed",
            status="failed",
            email=job.get("owner_email", ""),
            resource_type="study_guide",
            resource_name=output_language,
            duration_ms=int((utc_now() - started_at).total_seconds() * 1000),
            metadata={"job_id": job_id, "error": format_job_error(exc)},
        )


async def run_teacher_lesson_job(
    job_id: str,
    summary: str,
    transcript: str,
    lecture_notes: str,
    lecture_slides: str,
    past_question_papers: str,
    output_language: str,
):
    try:
        update_job(job_id, status="processing", stage="Starting teacher lesson", progress=8)
        teacher_package = await generate_teacher_lesson_package(
            summary,
            transcript,
            lecture_notes,
            lecture_slides,
            past_question_papers,
            job_id,
            output_language,
        )
        update_job(
            job_id,
            status="completed",
            stage="Teacher lesson ready",
            progress=100,
            **teacher_package,
        )
        job = jobs.get(job_id, {})
        started_at = parse_history_datetime(job.get("created_at"), utc_now())
        record_audit_log(
            action="teacher_lesson.completed",
            email=job.get("owner_email", ""),
            resource_type="teacher_lesson",
            resource_name=output_language,
            duration_ms=int((utc_now() - started_at).total_seconds() * 1000),
            metadata={"job_id": job_id},
        )
    except Exception as exc:
        logger.exception("Teacher lesson job failed")
        update_job(
            job_id,
            status="failed",
            stage="Teacher lesson failed",
            progress=100,
            error=format_job_error(exc),
        )
        job = jobs.get(job_id, {})
        started_at = parse_history_datetime(job.get("created_at"), utc_now())
        record_audit_log(
            action="teacher_lesson.completed",
            status="failed",
            email=job.get("owner_email", ""),
            resource_type="teacher_lesson",
            resource_name=output_language,
            duration_ms=int((utc_now() - started_at).total_seconds() * 1000),
            metadata={"job_id": job_id, "error": format_job_error(exc)},
        )


async def run_quiz_job(
    job_id: str,
    summary: str,
    transcript: str,
    lecture_notes: str,
    lecture_slides: str,
    past_question_papers: str,
    output_language: str,
):
    try:
        update_job(job_id, status="processing", stage="Starting test generation", progress=8)
        quiz_questions = await generate_quiz_questions(
            summary,
            transcript,
            lecture_notes,
            lecture_slides,
            past_question_papers,
            job_id,
            output_language,
        )
        update_job(
            job_id,
            status="completed",
            stage="Test ready",
            progress=100,
            quiz_questions=quiz_questions,
        )
        job = jobs.get(job_id, {})
        started_at = parse_history_datetime(job.get("created_at"), utc_now())
        record_audit_log(
            action="quiz.completed",
            email=job.get("owner_email", ""),
            resource_type="quiz",
            resource_name=output_language,
            duration_ms=int((utc_now() - started_at).total_seconds() * 1000),
            metadata={"job_id": job_id, "question_count": len(quiz_questions)},
        )
    except Exception as exc:
        logger.exception("Quiz generation failed")
        update_job(
            job_id,
            status="failed",
            stage="Test generation failed",
            progress=100,
            error=format_job_error(exc),
        )
        job = jobs.get(job_id, {})
        started_at = parse_history_datetime(job.get("created_at"), utc_now())
        record_audit_log(
            action="quiz.completed",
            status="failed",
            email=job.get("owner_email", ""),
            resource_type="quiz",
            resource_name=output_language,
            duration_ms=int((utc_now() - started_at).total_seconds() * 1000),
            metadata={"job_id": job_id, "error": format_job_error(exc)},
        )


async def run_podcast_job(
    job_id: str,
    summary: str,
    transcript: str,
    lecture_notes: str,
    lecture_slides: str,
    past_question_papers: str,
    speaker_count: int,
    target_minutes: int,
    output_language: str,
):
    try:
        update_job(job_id, status="processing", stage="Starting podcast generation", progress=8)
        podcast_package = await generate_podcast_package(
            summary,
            transcript,
            lecture_notes,
            lecture_slides,
            past_question_papers,
            speaker_count,
            target_minutes,
            job_id,
            output_language,
        )
        update_job(
            job_id,
            status="completed",
            stage="Podcast ready",
            progress=100,
            **podcast_package,
        )
        job = jobs.get(job_id, {})
        started_at = parse_history_datetime(job.get("created_at"), utc_now())
        record_audit_log(
            action="podcast.completed",
            email=job.get("owner_email", ""),
            resource_type="podcast",
            resource_name=output_language,
            duration_ms=int((utc_now() - started_at).total_seconds() * 1000),
            metadata={"job_id": job_id},
        )
    except Exception as exc:
        logger.exception("Podcast generation failed")
        update_job(
            job_id,
            status="failed",
            stage="Podcast generation failed",
            progress=100,
            error=format_job_error(exc),
        )
        job = jobs.get(job_id, {})
        started_at = parse_history_datetime(job.get("created_at"), utc_now())
        record_audit_log(
            action="podcast.completed",
            status="failed",
            email=job.get("owner_email", ""),
            resource_type="podcast",
            resource_name=output_language,
            duration_ms=int((utc_now() - started_at).total_seconds() * 1000),
            metadata={"job_id": job_id, "error": format_job_error(exc)},
        )


async def run_presentation_job(
    job_id: str,
    summary: str,
    transcript: str,
    lecture_notes: str,
    lecture_slides: str,
    past_question_papers: str,
    design_id: str,
    output_language: str,
    reference_images: list[str],
    template_file_bytes: bytes | None = None,
    template_file_name: str = "",
):
    try:
        update_job(job_id, status="processing", stage="Starting PowerPoint generation", progress=8)
        presentation_package = await generate_presentation_package(
            summary,
            transcript,
            lecture_notes,
            lecture_slides,
            past_question_papers,
            design_id,
            job_id,
            output_language,
            reference_images,
            template_file_bytes,
            template_file_name,
        )
        update_job(
            job_id,
            status="completed",
            stage="PowerPoint ready",
            progress=100,
            **presentation_package,
        )
        job = jobs.get(job_id, {})
        started_at = parse_history_datetime(job.get("created_at"), utc_now())
        record_audit_log(
            action="presentation.completed",
            email=job.get("owner_email", ""),
            resource_type="presentation",
            resource_name=design_id,
            duration_ms=int((utc_now() - started_at).total_seconds() * 1000),
            metadata={"job_id": job_id, "language": output_language},
        )
    except Exception as exc:
        logger.exception("PowerPoint generation failed")
        update_job(
            job_id,
            status="failed",
            stage="PowerPoint generation failed",
            progress=100,
            error=format_job_error(exc),
        )
        job = jobs.get(job_id, {})
        started_at = parse_history_datetime(job.get("created_at"), utc_now())
        record_audit_log(
            action="presentation.completed",
            status="failed",
            email=job.get("owner_email", ""),
            resource_type="presentation",
            resource_name=design_id,
            duration_ms=int((utc_now() - started_at).total_seconds() * 1000),
            metadata={"job_id": job_id, "language": output_language, "error": format_job_error(exc)},
        )


@app.post("/upload-audio/")
async def upload_audio(
    request: Request,
    file: UploadFile = File(...),
    current_user: str = Depends(require_authenticated_user),
):
    started_at = utc_now()
    if not file.filename:
        raise HTTPException(status_code=400, detail="No file selected.")

    job_id = create_job("transcription", owner_email=current_user)
    update_job(job_id, status="processing", stage="Starting upload", progress=1)

    try:
        ensure_openai_key()
        file_path = await save_upload_to_disk(file, job_id)
        file_size = get_file_size(file_path)
        if file_size > MAX_FILE_SIZE_BYTES:
            if file_path.exists():
                file_path.unlink()
            jobs.pop(job_id, None)
            raise HTTPException(
                status_code=413,
                detail=(
                    f"File is too large ({file_size / (1024 * 1024):.1f} MB). "
                    f"Please upload a file smaller than {MAX_FILE_SIZE_BYTES / (1024 * 1024):.0f} MB."
                ),
            )

        asyncio.create_task(run_transcription_job(job_id, file_path))
        record_audit_log(
            action="lecture.upload.request",
            email=current_user,
            request=request,
            resource_type="transcription",
            resource_name=file.filename,
            duration_ms=int((utc_now() - started_at).total_seconds() * 1000),
            metadata={"job_id": job_id},
        )
        return {"job_id": job_id}
    finally:
        await file.close()


@app.post("/transcribe-video-url/")
async def transcribe_video_url(
    payload: VideoUrlTranscriptionRequest,
    request: Request,
    current_user: str = Depends(require_authenticated_user),
):
    started_at = utc_now()
    try:
        video_url = normalize_video_url(payload.video_url)
    except ValueError as exc:
        raise HTTPException(status_code=400, detail=str(exc)) from exc

    ensure_openai_key()
    if not can_process_video_url(video_url):
        raise HTTPException(status_code=500, detail="Video-link transcription is not configured on the backend yet.")

    job_id = create_job("video_transcription", owner_email=current_user)
    update_job(job_id, status="processing", stage="Preparing video link", progress=1)
    asyncio.create_task(run_video_transcription_job(job_id, video_url))
    record_audit_log(
        action="lecture.video_link.request",
        email=current_user,
        request=request,
        resource_type="transcription",
        resource_name=video_url,
        duration_ms=int((utc_now() - started_at).total_seconds() * 1000),
        metadata={"job_id": job_id},
    )
    return {"job_id": job_id}


@app.post("/extract-slide-text/")
async def extract_slide_text(
    request: Request,
    file: UploadFile = File(...),
    current_user: str = Depends(require_authenticated_user),
):
    started_at = utc_now()
    if not file.filename:
        raise HTTPException(status_code=400, detail="No study source file selected.")

    ensure_openai_key()
    content_type = file.content_type or mimetypes.guess_type(file.filename)[0] or ""
    reference_images: list[str] = []

    try:
        file_bytes = await file.read()
        if not file_bytes:
            raise HTTPException(status_code=400, detail="The selected slide file is empty.")

        size_limit = MAX_IMAGE_UPLOAD_BYTES if content_type.startswith("image/") else MAX_SLIDE_UPLOAD_BYTES
        if len(file_bytes) > size_limit:
            raise HTTPException(
                status_code=413,
                detail=(
                    f"Study source file is too large ({len(file_bytes) / (1024 * 1024):.1f} MB). "
                    f"Please keep it below {size_limit / (1024 * 1024):.0f} MB."
                ),
            )

        if is_text_upload(file.filename, content_type):
            text = file_bytes.decode("utf-8", errors="ignore").strip()
            if not text:
                raise HTTPException(status_code=400, detail="The slide text file does not contain readable text.")
            return {"text": text, "image_urls": []}

        if is_pdf_upload(file.filename, content_type):
            text = await asyncio.to_thread(extract_slide_text_from_pdf, file_bytes)
            reference_images = await asyncio.to_thread(extract_reference_images_from_pdf, file_bytes, limit=6)
            if not text:
                raise HTTPException(status_code=422, detail="MABASO could not extract readable text from that PDF.")
            return {"text": text, "image_urls": reference_images}

        if is_pptx_upload(file.filename, content_type):
            try:
                text = await asyncio.to_thread(extract_slide_text_from_pptx, file_bytes)
                reference_images = await asyncio.to_thread(extract_reference_images_from_pptx, file_bytes, limit=6)
            except zipfile.BadZipFile as exc:
                raise HTTPException(status_code=400, detail="That PowerPoint file could not be opened. Use a .pptx file.") from exc

            if not text:
                raise HTTPException(status_code=422, detail="MABASO could not extract readable text from that PowerPoint file.")
            return {"text": text, "image_urls": reference_images}

        if is_docx_upload(file.filename, content_type):
            try:
                text = await asyncio.to_thread(extract_text_from_docx, file_bytes)
            except zipfile.BadZipFile as exc:
                raise HTTPException(status_code=400, detail="That Word document could not be opened. Use a .docx file.") from exc

            if not text:
                raise HTTPException(status_code=422, detail="MABASO could not extract readable text from that Word document.")
            return {"text": text, "image_urls": []}

        if not content_type.startswith("image/"):
            raise HTTPException(
                status_code=400,
                detail="Upload a slide image, .pdf, .pptx, .docx, or a text-based slide file such as .txt or .md.",
            )

        image_data_url = build_data_url(file_bytes, content_type, file.filename)
        text = await asyncio.to_thread(extract_slide_text_from_image, image_data_url, file.filename)
        if not text:
            raise HTTPException(status_code=422, detail="MABASO could not read text from that slide image.")
        record_audit_log(
            action="study_source.extract",
            email=current_user,
            request=request,
            resource_type="study_source",
            resource_name=file.filename,
            duration_ms=int((utc_now() - started_at).total_seconds() * 1000),
        )
        return {"text": text, "image_urls": [image_data_url]}
    finally:
        await file.close()


@app.post("/generate-study-guide/")
async def create_study_guide(
    payload: StudyGuideRequest,
    request: Request,
    current_user: str = Depends(require_authenticated_user),
):
    started_at = utc_now()
    transcript = payload.transcript.strip()
    lecture_notes = payload.lecture_notes.strip()
    lecture_slides = payload.lecture_slides.strip()
    past_question_papers = payload.past_question_papers.strip()
    output_language = normalize_output_language(payload.language)
    reference_images = [
        compact_text(item)
        for item in (getattr(payload, "reference_images", []) or [])
        if compact_text(item)
    ][:6]
    if not any([transcript, lecture_notes, lecture_slides, past_question_papers]):
        raise HTTPException(
            status_code=400,
            detail="Upload a transcript, notes, slides, or past question paper before generating a study guide.",
        )

    ensure_openai_key()
    job_id = create_job("study_guide", owner_email=current_user)
    update_job(job_id, _output_language=output_language)
    asyncio.create_task(
        run_summary_job(
            job_id,
            transcript,
            lecture_notes,
            lecture_slides,
            past_question_papers,
            output_language,
            reference_images,
        )
    )
    record_audit_log(
        action="study_guide.request",
        email=current_user,
        request=request,
        resource_type="study_guide",
        resource_name=output_language,
        duration_ms=int((utc_now() - started_at).total_seconds() * 1000),
        metadata={"job_id": job_id, "language": output_language, "reference_images": len(reference_images)},
    )
    return {"job_id": job_id}


@app.post("/generate-quiz/")
async def create_quiz(
    payload: QuizGenerationRequest,
    request: Request,
    current_user: str = Depends(require_authenticated_user),
):
    started_at = utc_now()
    transcript = payload.transcript.strip()
    summary = payload.summary.strip()
    lecture_notes = payload.lecture_notes.strip()
    lecture_slides = payload.lecture_slides.strip()
    past_question_papers = payload.past_question_papers.strip()
    output_language = normalize_output_language(payload.language)

    if not any([summary, transcript, lecture_notes, lecture_slides, past_question_papers]):
        raise HTTPException(
            status_code=400,
            detail="Generate a study guide or add lecture material before creating the test.",
        )

    ensure_openai_key()
    job_id = create_job("quiz", owner_email=current_user)
    update_job(job_id, _output_language=output_language)
    asyncio.create_task(
        run_quiz_job(
            job_id,
            summary,
            transcript,
            lecture_notes,
            lecture_slides,
            past_question_papers,
            output_language,
        )
    )
    record_audit_log(
        action="quiz.request",
        email=current_user,
        request=request,
        resource_type="quiz",
        resource_name=output_language,
        duration_ms=int((utc_now() - started_at).total_seconds() * 1000),
        metadata={"job_id": job_id, "language": output_language},
    )
    return {"job_id": job_id}


@app.post("/generate-teacher-lesson/")
async def create_teacher_lesson(
    payload: TeacherLessonRequest,
    request: Request,
    current_user: str = Depends(require_authenticated_user),
):
    started_at = utc_now()
    transcript = payload.transcript.strip()
    summary = payload.summary.strip()
    lecture_notes = payload.lecture_notes.strip()
    lecture_slides = payload.lecture_slides.strip()
    past_question_papers = payload.past_question_papers.strip()
    output_language = normalize_output_language(payload.language)

    if not any([summary, transcript, lecture_notes, lecture_slides, past_question_papers]):
        raise HTTPException(
            status_code=400,
            detail="Generate a study guide or add lecture material before opening teacher mode.",
        )

    ensure_openai_key()
    job_id = create_job("teacher_lesson", owner_email=current_user)
    update_job(job_id, _output_language=output_language)
    asyncio.create_task(
        run_teacher_lesson_job(
            job_id,
            summary,
            transcript,
            lecture_notes,
            lecture_slides,
            past_question_papers,
            output_language,
        )
    )
    record_audit_log(
        action="teacher_lesson.request",
        email=current_user,
        request=request,
        resource_type="teacher_lesson",
        resource_name=output_language,
        duration_ms=int((utc_now() - started_at).total_seconds() * 1000),
        metadata={"job_id": job_id, "language": output_language},
    )
    return {"job_id": job_id}


@app.post("/generate-podcast/")
async def create_podcast(
    payload: PodcastGenerationRequest,
    request: Request,
    current_user: str = Depends(require_authenticated_user),
):
    started_at = utc_now()
    transcript = payload.transcript.strip()
    summary = payload.summary.strip()
    lecture_notes = payload.lecture_notes.strip()
    lecture_slides = payload.lecture_slides.strip()
    past_question_papers = payload.past_question_papers.strip()
    output_language = normalize_output_language(payload.language)

    if not any([summary, transcript, lecture_notes, lecture_slides, past_question_papers]):
        raise HTTPException(
            status_code=400,
            detail="Generate a study guide or add lecture material before creating the podcast debate.",
        )

    ensure_openai_key()
    speaker_count = clamp_podcast_speaker_count(payload.speaker_count)
    target_minutes = clamp_podcast_target_minutes(payload.target_minutes)
    job_id = create_job("podcast", owner_email=current_user)
    update_job(job_id, _output_language=output_language)
    asyncio.create_task(
        run_podcast_job(
            job_id,
            summary,
            transcript,
            lecture_notes,
            lecture_slides,
            past_question_papers,
            speaker_count,
            target_minutes,
            output_language,
        )
    )
    record_audit_log(
        action="podcast.request",
        email=current_user,
        request=request,
        resource_type="podcast",
        resource_name=output_language,
        duration_ms=int((utc_now() - started_at).total_seconds() * 1000),
        metadata={"job_id": job_id, "language": output_language},
    )
    return {"job_id": job_id}


@app.post("/generate-presentation/")
async def create_presentation(
    request: Request,
    current_user: str = Depends(require_authenticated_user),
):
    started_at = utc_now()
    template_file_bytes: bytes | None = None
    template_file_name = ""
    content_type = (request.headers.get("content-type") or "").lower()
    if "multipart/form-data" in content_type:
        form = await request.form()
        template_upload = form.get("template_file")
        if isinstance(template_upload, UploadFile) and compact_text(template_upload.filename):
            if not is_pptx_upload(template_upload.filename, template_upload.content_type):
                raise HTTPException(status_code=400, detail="Upload a PowerPoint template in .pptx format.")
            template_file_bytes = await template_upload.read()
            if not template_file_bytes:
                raise HTTPException(status_code=400, detail="The selected PowerPoint template is empty.")
            template_file_name = compact_text(template_upload.filename)
        payload = PresentationGenerationRequest(
            transcript=compact_text(form.get("transcript")),
            summary=compact_text(form.get("summary")),
            lecture_notes=compact_text(form.get("lecture_notes")),
            lecture_slides=compact_text(form.get("lecture_slides")),
            past_question_papers=compact_text(form.get("past_question_papers")),
            design_id=compact_text(form.get("design_id"), "emerald-scholar"),
            language=compact_text(form.get("language"), "English"),
            reference_images=[
                compact_text(item)
                for item in form.getlist("reference_images")
                if compact_text(item)
            ][:6],
        )
    else:
        payload = PresentationGenerationRequest(**(await request.json()))
    transcript = payload.transcript.strip()
    summary = payload.summary.strip()
    lecture_notes = payload.lecture_notes.strip()
    lecture_slides = payload.lecture_slides.strip()
    past_question_papers = payload.past_question_papers.strip()
    output_language = normalize_output_language(payload.language)
    reference_images = [compact_text(item) for item in (getattr(payload, "reference_images", []) or []) if compact_text(item)][:6]

    if not any([summary, transcript, lecture_notes, lecture_slides, past_question_papers]):
        raise HTTPException(
            status_code=400,
            detail="Generate a study guide or add lecture material before creating the PowerPoint presentation.",
        )

    ensure_openai_key()
    ensure_presentation_support()
    design_id = normalize_presentation_design_id(payload.design_id)
    job_id = create_job("presentation", owner_email=current_user)
    update_job(
        job_id,
        _output_language=output_language,
        presentation_template_name=template_file_name,
    )
    asyncio.create_task(
        run_presentation_job(
            job_id,
            summary,
            transcript,
            lecture_notes,
            lecture_slides,
            past_question_papers,
            design_id,
            output_language,
            reference_images,
            template_file_bytes,
            template_file_name,
        )
    )
    record_audit_log(
        action="presentation.request",
        email=current_user,
        request=request,
        resource_type="presentation",
        resource_name=design_id,
        duration_ms=int((utc_now() - started_at).total_seconds() * 1000),
        metadata={
            "job_id": job_id,
            "design_id": design_id,
            "language": output_language,
            "reference_images": len(reference_images),
            "custom_template": bool(template_file_name),
            "template_name": template_file_name,
        },
    )
    return {"job_id": job_id}


@app.post("/mark-quiz-answer/")
async def mark_quiz_answer(
    question: str = Form(...),
    expected_answer: str = Form(...),
    student_answer: str = Form(""),
    question_type: str = Form("short_answer"),
    max_score: int = Form(1),
    answer_points_json: str = Form("[]"),
    subparts_json: str = Form("[]"),
    student_selection_json: str = Form("{}"),
    answer_image: UploadFile | None = File(None),
    answer_images: list[UploadFile] | None = File(None),
    current_user: str = Depends(require_authenticated_user),
):
    resolved_question_type = compact_text(question_type, "short_answer").lower()
    resolved_max_score = max(1, clamp_score(max_score, 100))

    if resolved_question_type in {"multiple_choice_group", "true_false_group"}:
        subparts = [item for item in parse_json_list(subparts_json) if isinstance(item, dict)]
        student_selection = parse_json_object(student_selection_json)
        return grade_option_based_question(resolved_question_type, subparts, student_selection)

    ensure_openai_key()
    answer_points = [compact_text(item) for item in parse_json_list(answer_points_json) if compact_text(item)]
    uploaded_images = [file for file in (answer_images or []) if file is not None]
    if answer_image is not None:
        uploaded_images.insert(0, answer_image)

    if not student_answer.strip() and not uploaded_images:
        return {
            "score": 0,
            "max_score": resolved_max_score,
            "feedback": "No answer was submitted yet.",
            "extracted_answer": "",
            "mistakes": [],
        }

    image_data_urls: list[str] = []

    try:
        for index, uploaded_image in enumerate(uploaded_images, start=1):
            content_type = uploaded_image.content_type or mimetypes.guess_type(uploaded_image.filename or "")[0] or ""
            if not content_type.startswith("image/"):
                raise HTTPException(status_code=400, detail="Please upload an image for answer-photo marking.")

            image_bytes = await uploaded_image.read()
            if not image_bytes:
                raise HTTPException(status_code=400, detail="The uploaded answer image is empty.")
            if len(image_bytes) > MAX_IMAGE_UPLOAD_BYTES:
                raise HTTPException(
                    status_code=413,
                    detail=(
                        f"Answer image is too large ({len(image_bytes) / (1024 * 1024):.1f} MB). "
                        f"Please keep it below {MAX_IMAGE_UPLOAD_BYTES / (1024 * 1024):.0f} MB."
                    ),
                )

            image_data_urls.append(
                build_data_url(
                    image_bytes,
                    content_type,
                    uploaded_image.filename or f"answer-image-{index}",
                )
            )

        result = await asyncio.to_thread(
            mark_quiz_answer_with_ai,
            question.strip(),
            expected_answer.strip(),
            student_answer.strip(),
            image_data_urls,
            answer_points,
            resolved_max_score,
        )
        return result
    finally:
        for uploaded_image in uploaded_images:
            await uploaded_image.close()


@app.post("/ask-study-assistant/")
async def ask_study_assistant(
    payload: StudyChatRequest,
    current_user: str = Depends(require_authenticated_user),
):
    ensure_openai_key()
    if not payload.question.strip():
        raise HTTPException(status_code=400, detail="A question is required.")

    def _ask() -> str:
        use_vision = bool(payload.reference_images)
        response = client.with_options(timeout=VISION_REQUEST_TIMEOUT if use_vision else 45).chat.completions.create(
            model=VISION_MODEL if use_vision else STUDY_CHAT_MODEL,
            max_completion_tokens=1200,
            messages=build_chat_messages(payload),
        )
        return (response.choices[0].message.content or "").strip()

    answer = await asyncio.to_thread(_ask)
    answer_with_follow_up = ensure_study_chat_follow_up(answer, payload.question)
    return {"answer": make_formulas_human_readable(answer_with_follow_up)}


@app.get("/collaboration/rooms")
async def list_collaboration_rooms(current_user: str = Depends(require_authenticated_user)):
    with get_db_connection() as connection:
        rows = connection.execute(
            """
            SELECT DISTINCT r.*
            FROM collaboration_rooms r
            LEFT JOIN collaboration_room_members m
                ON m.room_id = r.id
            WHERE r.owner_email = ? OR m.email = ?
            ORDER BY r.updated_at DESC
            """,
            (current_user, current_user),
        ).fetchall()

    return {"rooms": [serialize_collaboration_room_card(row) for row in rows]}


@app.post("/collaboration/rooms")
async def create_collaboration_room(
    payload: CollaborationRoomCreateRequest,
    current_user: str = Depends(require_authenticated_user),
):
    title = payload.title.strip()
    if not title:
        raise HTTPException(status_code=400, detail="Room title is required.")

    transcript = payload.transcript.strip()
    summary = payload.summary.strip()
    lecture_notes = payload.lecture_notes.strip()
    lecture_slides = payload.lecture_slides.strip()
    if not any([transcript, summary, lecture_notes, lecture_slides]):
        raise HTTPException(status_code=400, detail="Create the collaboration room from a lecture that already has content.")

    room_id = uuid4().hex
    now_iso = utc_now().isoformat()
    invited_emails = normalize_invited_emails(payload.invited_emails, current_user)

    with get_db_connection() as connection:
        connection.execute(
            """
            INSERT INTO collaboration_rooms (
                id, owner_email, title, transcript, summary, formula, example,
                lecture_notes, lecture_slides, shared_notes, flashcards_json,
                quiz_questions_json, active_tab, test_visibility, created_at, updated_at
            )
            VALUES (?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?)
            """,
            (
                room_id,
                current_user,
                title,
                transcript,
                summary,
                payload.formula.strip(),
                payload.example.strip(),
                lecture_notes,
                lecture_slides,
                payload.shared_notes.strip(),
                dump_json(payload.flashcards or []),
                dump_json(payload.quiz_questions or []),
                sanitize_collaboration_tab(payload.active_tab),
                normalize_test_visibility(payload.test_visibility),
                now_iso,
                now_iso,
            ),
        )
        connection.execute(
            """
            INSERT INTO collaboration_room_members (room_id, email, role, created_at)
            VALUES (?, ?, ?, ?)
            """,
            (room_id, current_user, "owner", now_iso),
        )
        for email in invited_emails:
            connection.execute(
                """
                INSERT OR REPLACE INTO collaboration_room_members (room_id, email, role, created_at)
                VALUES (?, ?, ?, ?)
                """,
                (room_id, email, "member", now_iso),
            )

    room = get_accessible_collaboration_room(room_id, current_user)
    return {"room": serialize_collaboration_room(room, current_user)}


@app.get("/collaboration/rooms/{room_id}")
async def get_collaboration_room(room_id: str, current_user: str = Depends(require_authenticated_user)):
    room = get_accessible_collaboration_room(room_id, current_user)
    return {"room": serialize_collaboration_room(room, current_user)}


@app.post("/collaboration/rooms/{room_id}/messages")
async def send_collaboration_message(
    room_id: str,
    payload: CollaborationMessageRequest,
    current_user: str = Depends(require_authenticated_user),
):
    room = get_accessible_collaboration_room(room_id, current_user)
    content = payload.content.strip()
    if not content:
        raise HTTPException(status_code=400, detail="Type a collaboration message first.")

    now_iso = utc_now().isoformat()
    with get_db_connection() as connection:
        connection.execute(
            """
            INSERT INTO collaboration_room_messages (id, room_id, author_email, content, created_at)
            VALUES (?, ?, ?, ?, ?)
            """,
            (uuid4().hex, room["id"], current_user, content, now_iso),
        )
        connection.execute(
            "UPDATE collaboration_rooms SET updated_at = ? WHERE id = ?",
            (now_iso, room["id"]),
        )

    updated_room = get_accessible_collaboration_room(room_id, current_user)
    return {"room": serialize_collaboration_room(updated_room, current_user)}


@app.post("/collaboration/rooms/{room_id}/notes")
async def save_collaboration_notes(
    room_id: str,
    payload: CollaborationSharedNotesRequest,
    current_user: str = Depends(require_authenticated_user),
):
    room = get_accessible_collaboration_room(room_id, current_user)
    now_iso = utc_now().isoformat()
    with get_db_connection() as connection:
        connection.execute(
            """
            UPDATE collaboration_rooms
            SET shared_notes = ?, updated_at = ?
            WHERE id = ?
            """,
            (payload.shared_notes.strip(), now_iso, room["id"]),
        )

    updated_room = get_accessible_collaboration_room(room_id, current_user)
    return {"room": serialize_collaboration_room(updated_room, current_user)}


@app.post("/collaboration/rooms/{room_id}/active-tab")
async def update_collaboration_active_tab(
    room_id: str,
    payload: CollaborationActiveTabRequest,
    current_user: str = Depends(require_authenticated_user),
):
    room = get_accessible_collaboration_room(room_id, current_user)
    now_iso = utc_now().isoformat()
    with get_db_connection() as connection:
        connection.execute(
            """
            UPDATE collaboration_rooms
            SET active_tab = ?, updated_at = ?
            WHERE id = ?
            """,
            (sanitize_collaboration_tab(payload.active_tab), now_iso, room["id"]),
        )

    updated_room = get_accessible_collaboration_room(room_id, current_user)
    return {"room": serialize_collaboration_room(updated_room, current_user)}


@app.post("/collaboration/rooms/{room_id}/test-visibility")
async def update_collaboration_test_visibility(
    room_id: str,
    payload: CollaborationTestVisibilityRequest,
    current_user: str = Depends(require_authenticated_user),
):
    room = get_accessible_collaboration_room(room_id, current_user)
    if room["owner_email"] != current_user:
        raise HTTPException(status_code=403, detail="Only the room owner can change the test visibility.")

    now_iso = utc_now().isoformat()
    with get_db_connection() as connection:
        connection.execute(
            """
            UPDATE collaboration_rooms
            SET test_visibility = ?, updated_at = ?
            WHERE id = ?
            """,
            (normalize_test_visibility(payload.test_visibility), now_iso, room["id"]),
        )

    updated_room = get_accessible_collaboration_room(room_id, current_user)
    return {"room": serialize_collaboration_room(updated_room, current_user)}


@app.post("/collaboration/rooms/{room_id}/quiz-answers")
async def save_collaboration_quiz_answer(
    room_id: str,
    payload: CollaborationQuizAnswerRequest,
    current_user: str = Depends(require_authenticated_user),
):
    room = get_accessible_collaboration_room(room_id, current_user)
    question_number = payload.question_number.strip()
    if not question_number:
        raise HTTPException(status_code=400, detail="Question number is required.")

    now_iso = utc_now().isoformat()
    answer_text = payload.answer_text.strip()

    with get_db_connection() as connection:
        if answer_text:
            connection.execute(
                """
                INSERT INTO collaboration_room_answers (room_id, question_number, author_email, answer_text, updated_at)
                VALUES (?, ?, ?, ?, ?)
                ON CONFLICT(room_id, question_number, author_email) DO UPDATE SET
                    answer_text = excluded.answer_text,
                    updated_at = excluded.updated_at
                """,
                (room["id"], question_number, current_user, answer_text, now_iso),
            )
        else:
            connection.execute(
                """
                DELETE FROM collaboration_room_answers
                WHERE room_id = ? AND question_number = ? AND author_email = ?
                """,
                (room["id"], question_number, current_user),
            )
        connection.execute(
            "UPDATE collaboration_rooms SET updated_at = ? WHERE id = ?",
            (now_iso, room["id"]),
        )

    updated_room = get_accessible_collaboration_room(room_id, current_user)
    return {"room": serialize_collaboration_room(updated_room, current_user)}


@app.post("/export-study-pack-pdf/")
async def export_study_pack_pdf(
    payload: PdfExportRequest,
    current_user: str = Depends(require_authenticated_user),
):
    title = payload.title.strip() or "MABASO Study Pack"
    pdf_bytes = await asyncio.to_thread(build_pdf_document, title, payload.sections)
    safe_name = sanitize_download_filename(title)
    return Response(
        content=pdf_bytes,
        media_type="application/pdf",
        headers={"Content-Disposition": f'attachment; filename="{safe_name}.pdf"'},
    )
